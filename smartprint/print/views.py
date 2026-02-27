# -*- coding: utf-8 -*-
from decimal import Decimal, InvalidOperation
from django.shortcuts import render, redirect
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.conf import settings
from django.contrib.auth import login
from django.contrib.auth.models import User
from django.contrib import messages
from django.contrib.admin.views.decorators import staff_member_required
from django.contrib.auth.decorators import login_required
from django.db import connection
import boto3
import datetime
import json
import requests
import uuid
import random
import re
import time
import jwt  # For local token decoding fallback
from django.contrib.auth.hashers import make_password, check_password
from django.utils import timezone
from django.urls import reverse
import os
import base64
from django.core.mail import send_mail
from django.template.loader import render_to_string
from django.utils.html import strip_tags
import traceback
from PIL import Image, ImageDraw
import io
from django.views.decorators.http import require_POST, require_http_methods, require_GET
import threading
import schedule
import subprocess
import tempfile
import shutil
from urllib.parse import urlparse, urlunparse, unquote
from channels.layers import get_channel_layer
from asgiref.sync import async_to_sync
from .models import VendorLocationSession

# Lock dictionary to handle concurrent job completions
_job_completion_locks = {}
_job_completion_lock = threading.Lock()  # Lock for the locks dictionary itself

KNOWN_WORKER_ENDPOINTS = {
    'add-contact',
    'add-vendor-register',
    'add-vendor-pricing',
    'add-vendor-tokens',
    'get-all-vendors',
    'get-vendor-pricing',
    'get-vendor-by-email',
    'get-vendor-by-id',
    'add-vendor-print-job',
    'add-user-print-job',
    'add-user-notification',
    'add-vendor-notification',
    'add-user-points',
    'get-user-total-points',
    'upsert-vendor-service',
    'get-vendor-service',
    'add-user-signup',
    'get-user-signup',
}


def get_ist_timestamp() -> str:
    """
    Return the current timestamp in ISO format using the project's default
    timezone (configured as Asia/Kolkata in settings).
    """
    return timezone.localtime(timezone.now()).isoformat()


def get_worker_base_url():
    """
    Normalize WORKER_API_URL so we can append different endpoints reliably.
    Removes the last path segment only when it matches a known Worker endpoint.
    """
    api_url = getattr(settings, 'WORKER_API_URL', '').strip()
    if not api_url:
        return ''

    parsed = urlparse(api_url)
    path = (parsed.path or '').rstrip('/')
    if path:
        segments = [segment for segment in path.split('/') if segment]
        if segments and segments[-1] in KNOWN_WORKER_ENDPOINTS:
            segments = segments[:-1]
        normalized_path = '/' + '/'.join(segments) if segments else ''
    else:
        normalized_path = ''

    base_url = urlunparse((parsed.scheme, parsed.netloc, normalized_path, '', '', '')).rstrip('/')
    return base_url or ''


def build_worker_endpoint(path):
    """
    Build a fully-qualified Worker endpoint for the provided relative path.
    """
    base_url = get_worker_base_url()
    if not base_url:
        return ''

    if not path.startswith('/'):
        path = '/' + path

    return f"{base_url}{path}"


def post_to_worker(path, payload=None, timeout=10):
    """
    Helper to send JSON POST requests to the Worker API.
    Returns (endpoint, response) so callers can log details if needed.
    """
    api_key = getattr(settings, 'WORKER_API_KEY', '')
    endpoint = build_worker_endpoint(path)

    if not endpoint or not api_key:
        raise RuntimeError("Worker API not configured")

    headers = {
        'Content-Type': 'application/json',
        'x-api-key': api_key
    }

    response = requests.post(
        endpoint,
        json=payload or {},
        headers=headers,
        timeout=timeout
    )

    return endpoint, response


def _normalize_r2_key(raw_value):
    """
    Convert an R2 key or URL-like value into a clean object key.
    """
    if not raw_value:
        return ''

    value = str(raw_value).strip()
    if not value:
        return ''

    try:
        if '://' in value:
            parsed = urlparse(value)
            key = unquote((parsed.path or '').lstrip('/'))
        else:
            key = unquote(value.lstrip('/'))
    except Exception:
        key = str(raw_value).strip().lstrip('/')

    bucket = (settings.R2_BUCKET or '').strip('/')
    if bucket and key.startswith(bucket + '/'):
        key = key[len(bucket) + 1:]

    return key.lstrip('/')


def create_or_update_admin_user_in_d1(
    username,
    password_hash,
    email=None,
    first_name=None,
    last_name=None,
    is_superuser=False,
    is_staff=True,
    is_active=True,
    permissions=None,
):
    """
    Create or update an admin user in D1 database.

    NOTE: This helper is only responsible for admin user management and must
    not depend on print‑job metadata. A previous version accidentally
    referenced a non‑existent ``metadata`` variable, which could raise a
    NameError during execution.
    """
    try:
        api_url = getattr(settings, "WORKER_API_URL", "").strip()
        api_key = getattr(settings, "WORKER_API_KEY", "")

        if not api_url or not api_key:
            raise RuntimeError("Worker API not configured")

        # Build endpoint
        base_url = api_url.rstrip("/")
        if "/add-contact" in base_url:
            endpoint = base_url.replace("/add-contact", "/create-admin-user")
        elif "/add-vendor-register" in base_url:
            endpoint = base_url.replace("/add-vendor-register", "/create-admin-user")
        else:
            endpoint = base_url + "/create-admin-user"

        payload = {
            "username": username,
            "password_hash": password_hash,
            "email": email,
            "first_name": first_name,
            "last_name": last_name,
            "is_superuser": is_superuser,
            "is_staff": is_staff,
            "is_active": is_active,
            "permissions": permissions,
        }

        response = requests.post(
            endpoint,
            json=payload,
            headers={
                "Content-Type": "application/json",
                "x-api-key": api_key,
            },
            timeout=10,
        )

        if response.status_code == 200:
            data = response.json()
            return data.get("success", False)
        return False

    except Exception as e:
        print(f"Error creating/updating admin user in D1: {e}")
        return False


def get_admin_user_from_d1(username):
    """
    Get admin user from D1 database by username
    """
    try:
        api_url = getattr(settings, 'WORKER_API_URL', '').strip()
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            raise RuntimeError("Worker API not configured")
        
        # Build endpoint
        base_url = api_url.rstrip('/')
        if '/add-contact' in base_url:
            endpoint = base_url.replace('/add-contact', '/get-admin-user')
        elif '/add-vendor-register' in base_url:
            endpoint = base_url.replace('/add-vendor-register', '/get-admin-user')
        else:
            endpoint = base_url + '/get-admin-user'
        
        response = requests.post(
            endpoint,
            json={'username': username},
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        if response.status_code == 200:
            data = response.json()
            if data.get('success'):
                return data.get('user')
        return None
        
    except Exception as e:
        print(f"Error getting admin user from D1: {e}")
        return None


def get_all_admin_users_from_d1():
    """
    Get all admin users from D1 database
    """
    try:
        api_url = getattr(settings, 'WORKER_API_URL', '').strip()
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            raise RuntimeError("Worker API not configured")
        
        # Build endpoint
        base_url = api_url.rstrip('/')
        if '/add-contact' in base_url:
            endpoint = base_url.replace('/add-contact', '/get-all-admin-users')
        elif '/add-vendor-register' in base_url:
            endpoint = base_url.replace('/add-vendor-register', '/get-all-admin-users')
        else:
            endpoint = base_url + '/get-all-admin-users'
        
        response = requests.post(
            endpoint,
            json={},
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        if response.status_code == 200:
            data = response.json()
            if data.get('success'):
                return data.get('users', [])
        return []
        
    except Exception as e:
        print(f"Error getting all admin users from D1: {e}")
        return []


def _to_decimal(value, precision=2):
    """
    Safely convert inputs to floats with optional rounding.
    Handles strings, decimals, and None without raising.
    """
    try:
        if value is None or value == '':
            return 0.0
        number = float(value)
        return round(number, precision) if precision is not None else number
    except (TypeError, ValueError):
        return 0.0


def firebase_messaging_sw(request):
    """
    Serve Firebase messaging service worker from root with dynamic config
    Service workers must be served from root for proper scope
    """
    from django.http import HttpResponse
    from django.template import Template, Context
    
    # Template for service worker with dynamic Firebase config
    sw_template = """// Import Firebase scripts for service worker
importScripts('https://www.gstatic.com/firebasejs/10.7.1/firebase-app-compat.js');
importScripts('https://www.gstatic.com/firebasejs/10.7.1/firebase-messaging-compat.js');

// Initialize Firebase in the service worker with config from Django settings
const firebaseConfig = {
  apiKey: "{{ apiKey }}",
  authDomain: "{{ authDomain }}",
  projectId: "{{ projectId }}",
  storageBucket: "{{ storageBucket }}",
  messagingSenderId: "{{ messagingSenderId }}",
  appId: "{{ appId }}"
};

firebase.initializeApp(firebaseConfig);

// Retrieve an instance of Firebase Messaging
const messaging = firebase.messaging();

// Handle background messages
messaging.onBackgroundMessage(function(payload) {
  console.debug('[firebase-messaging-sw.js] Received background message', payload);
  
  const notificationTitle = payload.notification?.title || payload.data?.title || 'PrintMax Notification';
  // Default to PrintMax colored logo (shows like Ola/Rapido style icons)
  const defaultIcon = '/static/images/printmax-color-512.png';
  const notificationIcon = payload.notification?.icon || payload.data?.icon || defaultIcon;
  const notificationBadge = payload.data?.badge || payload.notification?.badge || defaultIcon;
  
  const notificationOptions = {
    body: payload.notification?.body || payload.data?.message || 'You have a new notification',
    icon: notificationIcon,
    badge: notificationBadge,
    tag: payload.data?.notification_id || 'printmax-notification',
    requireInteraction: true,
    data: payload.data || {}
  };

  // Add click action if provided
  if (payload.data?.click_action || payload.fcmOptions?.link) {
    notificationOptions.data.url = payload.data?.click_action || payload.fcmOptions?.link;
  }

  return self.registration.showNotification(notificationTitle, notificationOptions);
});

// Handle notification clicks
self.addEventListener('notificationclick', function(event) {
  console.debug('[firebase-messaging-sw.js] Notification click received.');
  
  event.notification.close();

  // Open the app URL if provided
  const urlToOpen = event.notification.data?.url || '/userdashboard/';
  
  event.waitUntil(
    clients.matchAll({ type: 'window', includeUncontrolled: true }).then(function(clientList) {
      // Check if there's already a window/tab open with the target URL
      for (let i = 0; i < clientList.length; i++) {
        const client = clientList[i];
        if (client.url === urlToOpen && 'focus' in client) {
          return client.focus();
        }
      }
      // If no window is open, open a new one
      if (clients.openWindow) {
        return clients.openWindow(urlToOpen);
      }
    })
  );
});
"""
    
    # Get Firebase config from settings
    firebase_config = getattr(settings, 'FIREBASE_CONFIG', {})
    
    # Render template with config
    template = Template(sw_template)
    context = Context({
        'apiKey': firebase_config.get('apiKey', ''),
        'authDomain': firebase_config.get('authDomain', ''),
        'projectId': firebase_config.get('projectId', ''),
        'storageBucket': firebase_config.get('storageBucket', ''),
        'messagingSenderId': firebase_config.get('messagingSenderId', ''),
        'appId': firebase_config.get('appId', ''),
    })
    
    content = template.render(context)
    return HttpResponse(content, content_type='application/javascript')


def terms(request):
    """
    Terms & Conditions page view
    """
    return render(request, 'terms.html')

def privacy(request):
    """
    Privacy Policy page view
    """
    return render(request, 'privacy.html')

# Token management for sequential token generation (100-200)
_used_tokens = set()

# Simple in-memory cache for vendordashboard
_vendor_dashboard_cache = {}
_cache_timestamp = {}
CACHE_DURATION = 30  # seconds

def clear_vendor_cache(vendor_email, vendor_id=None):
    """Clear cache for a specific vendor"""
    cache_key = f"vendor_dashboard_{vendor_email}_{vendor_id}" if vendor_id else f"vendor_dashboard_{vendor_email}"
    if cache_key in _vendor_dashboard_cache:
        del _vendor_dashboard_cache[cache_key]
    if cache_key in _cache_timestamp:
        del _cache_timestamp[cache_key]
    print(f"🗑️ Cleared cache for vendor: {vendor_email}")
_current_token = 100

# Fields we no longer want to persist inside R2 object metadata
_R2_METADATA_BLOCKLIST = {
    'color',
    'copies',
    'lamination',
    'orientation',
    'pagerange',
    'pagesize',
    'size',
    'payment_id',
    'pricing_details',
    'pricingdetails',
    'printer_name',
    'rendered_status',
    'service_name',
    'service_type',
    'specificpages',
    'spiralbinding',
    'total_price',
    'user',
}

# Hard upper bound for any single metadata value stored in R2 headers.
# This prevents "Request max total header size exceeded" errors when
# verbose text (e.g. feedback, long page lists) is present.
_R2_METADATA_MAX_VALUE_LENGTH = 200


def sanitize_r2_metadata(metadata):
    """
    Return a copy of metadata without verbose per-job fields that are now stored in D1.
    Keeps values stringified for R2 compatibility.
    """
    if not metadata:
        return {}

    cleaned = {}
    for key, value in metadata.items():
        if key is None:
            continue

        lower_key = str(key).lower()
        if lower_key in _R2_METADATA_BLOCKLIST:
            # These fields are now stored in D1 and don't need to be
            # duplicated in R2 object metadata.
            continue

        # Always coerce to string for R2 compatibility and hard‑limit length
        str_key = str(key)
        if value is None:
            str_value = ''
        else:
            str_value = str(value)
            if len(str_value) > _R2_METADATA_MAX_VALUE_LENGTH:
                # Truncate excessively long values to keep total header size small
                str_value = str_value[:_R2_METADATA_MAX_VALUE_LENGTH]

        cleaned[str_key] = str_value

    return cleaned

def get_minimal_r2_metadata(filename):
    """
    Return minimal R2 metadata containing only filename.
    Used for userdashboard and photoprint modals where all details are stored in D1.
    """
    return {
        'filename': str(filename) if filename else ''
    }

def get_next_sequential_token():
    """
    Generate token based on total pending requests + random digit.
    Format: (total_pending + 1) * 100 + random_digit
    """
    try:
        # Count total pending requests across all users and vendors
        total_pending = 0
        
        # Count user pending jobs
        try:
            s3 = boto3.client('s3',
                              aws_access_key_id=settings.R2_ACCESS_KEY,
                              aws_secret_access_key=settings.R2_SECRET_KEY,
                              endpoint_url=settings.R2_ENDPOINT,
                              region_name='auto')
            
            user_jobs_data = s3.get_object(
                Bucket=settings.R2_BUCKET,
                Key='users/jobs_data.json'
            )
            user_jobs = json.loads(user_jobs_data['Body'].read().decode('utf-8'))
            total_pending += len([job for job in user_jobs if job.get('status', '').lower() in ['pending', 'processing', 'in_progress']])
        except:
            pass
            
        # Count vendor pending jobs
        try:
            response = s3.list_objects_v2(
                Bucket=settings.R2_BUCKET,
                Prefix='vendor_print_jobs/'
            )
            
            for obj in response.get('Contents', []):
                if obj['Key'].endswith('jobs_data.json'):
                    try:
                        vendor_jobs_data = s3.get_object(
                            Bucket=settings.R2_BUCKET,
                            Key=obj['Key']
                        )
                        vendor_jobs = json.loads(vendor_jobs_data['Body'].read().decode('utf-8'))
                        total_pending += len([job for job in vendor_jobs if job.get('status', '').lower() in ['pending', 'processing', 'in_progress']])
                    except:
                        continue
        except:
            pass
        
        # Generate token: (total_pending + 1) + 100 + random_digit
        random_digit = random.randint(0, 9)
        token = (total_pending + 1) + 100 + random_digit
        
        return str(token)
        
    except Exception as e:
        print(f"Error generating token: {e}")
        # Fallback to random token
        return str(random.randint(100, 200))

def contact_view(request):
    """
    Render Contact page and handle basic form POST.
    On success, show a Django success message and stay on the page.
    """
    if request.method == 'POST':
        name = request.POST.get('name', '').strip()
        email = request.POST.get('email', '').strip()
        subject = request.POST.get('subject', '').strip()
        message_text = request.POST.get('message', '').strip()

        if not name or not email or not subject or not message_text:
            messages.error(request, 'Please fill out all fields before submitting.')
        else:
            messages.success(request, 'Thanks for contacting us ! Your message has been sent successfully.')

        return render(request, 'contact.html', {
            'prefill': {
                'name': name,
                'email': email,
                'subject': subject,
                'message': message_text,
            }
        })

    return render(request, 'contact.html')

@csrf_exempt
def save_contact_details(request):
    """
    Save contact form details into Cloudflare D1 via the Worker API.
    Accepts JSON or form data. Returns JSON.
    """
    try:
        if request.method != 'POST':
            return JsonResponse({'success': False, 'error': 'Invalid method'}, status=405)

        # Accept both JSON and form posts
        if request.content_type and 'application/json' in request.content_type:
            try:
                payload = json.loads(request.body or '{}')
            except Exception:
                payload = {}
        else:
            payload = {
                'name': request.POST.get('name', ''),
                'email': request.POST.get('email', ''),
                'subject': request.POST.get('subject', ''),
                'message': request.POST.get('message', ''),
            }

        # Basic validation on Django side
        name = (payload.get('name') or '').strip()
        email = (payload.get('email') or '').strip()
        subject = (payload.get('subject') or '').strip()
        message = (payload.get('message') or '').strip()

        if not name or not email or not subject or not message:
            return JsonResponse({'success': False, 'error': 'All fields are required'}, status=400)

        # Forward to Worker API (D1 database only - no R2 storage)
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')

        if not api_url or not api_key:
            return JsonResponse({'success': False, 'error': 'Server misconfigured: missing Worker API settings'}, status=500)
        
        # Construct the Worker API endpoint for adding contact
        # Remove any existing endpoint paths and add /add-contact
        base_url = api_url.rstrip('/')
        # Remove common endpoint paths if present
        for endpoint in ['/add-contact', '/add-vendor-register', '/add-vendor-pricing', '/get-all-vendors']:
            if base_url.endswith(endpoint):
                base_url = base_url[:-len(endpoint)]
        worker_endpoint = base_url.rstrip('/') + '/add-contact'
        
        # Debug: Log API key info (first/last 4 chars only for security)
        api_key_preview = f"{api_key[:4]}...{api_key[-4:]}" if len(api_key) > 8 else "***"
        print(f"Sending request to Worker API: {worker_endpoint}")
        print(f"API Key (preview): {api_key_preview} (length: {len(api_key)})")

        # Prepare payload for Worker
        worker_payload = {
            'name': name,
            'email': email,
            'subject': subject,
            'message': message,
        }

        try:
            resp = requests.post(
                worker_endpoint,
                json=worker_payload,
                headers={
                    'Content-Type': 'application/json',
                    'x-api-key': api_key
                },
                timeout=10
            )
            
            # Log response for debugging
            print(f"Worker API Response Status: {resp.status_code}")
            print(f"Worker API URL: {worker_endpoint}")
            
            try:
                data = resp.json()
                print(f"Worker API Response Data: {data}")
            except Exception as json_err:
                response_text = resp.text[:500]  # First 500 chars
                print(f"Failed to parse JSON response: {json_err}")
                print(f"Response text: {response_text}")
                data = {'success': False, 'error': f'Invalid response from API ({resp.status_code}): {response_text}'}

            if resp.status_code == 200 and data.get('success'):
                return JsonResponse({'success': True})

            error_msg = data.get('error', f'API error (status {resp.status_code})')
            print(f"Worker API Error: {error_msg}")
            return JsonResponse({'success': False, 'error': error_msg}, status=resp.status_code or 500)
            
        except requests.exceptions.RequestException as req_err:
            error_msg = f'Failed to connect to Worker API: {str(req_err)}'
            print(f"Request Exception: {error_msg}")
            return JsonResponse({'success': False, 'error': error_msg}, status=500)

    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"Error saving contact details: {str(e)}")
        print(f"Traceback: {error_trace}")
        return JsonResponse({'success': False, 'error': f'Server error: {str(e)}'}, status=500)

def assign_token_from_vendor_pool(vendor_email):
    """
    Assign a token for the vendor using the D1 Vendor_tokens table via Worker API.
    Returns the reserved token number or None if allocation failed.
    """
    if not vendor_email:
        return None

    try:
        endpoint, resp = post_to_worker('/assign-vendor-token', {'vendor_email': vendor_email})
        if resp.status_code != 200:
            print(f"⚠️ Worker assign token failed ({resp.status_code}) via {endpoint}: {resp.text[:300]}")
            return None

        payload = resp.json()
        if payload.get('success'):
            token_number = payload.get('token_number')
            print(f"✅ Assigned token {token_number} from D1 for vendor {vendor_email}")
            return token_number

        print(f"⚠️ Worker assign token responded with error: {payload}")
        return None
    except Exception as exc:
        print(f"❌ Error assigning token from D1 for {vendor_email}: {exc}")
        return None
# ─────────────────────────────────────────────────────────────
# BASIC PAGE VIEWS
# ─────────────────────────────────────────────────────────────


def home(request):
    # ✅ Auto-redirect authenticated users to dashboard ONLY if they have valid Google session
    if request.user.is_authenticated:
        # Check if user has valid Google session (google_user_id in session)
        google_user_id = request.session.get('google_user_id')
        user_email = request.session.get('user_email') or (request.user.email if request.user.is_authenticated else None)
        
        if google_user_id and user_email:
            # Verify user exists in D1 User_signup_details
            user_details = get_user_details_from_d1(user_email)
            if user_details:
                print(f"✅ Valid Google session found for {user_email}, redirecting to dashboard")
                return redirect('userdashboard')
            else:
                print(f"⚠️ User {user_email} authenticated but not found in D1, clearing session")
                from django.contrib.auth import logout
                logout(request)
        else:
            print(f"⚠️ User authenticated but no valid Google session found, clearing session")
            from django.contrib.auth import logout
            logout(request)
    
    return render(request, 'home.html')


def vendor_appointment_page(request):
    """Render appointment booking page for a vendor (expects ?email=)."""
    email = request.GET.get('email', '')
    return render(request, 'vendor_appointment.html', { 'vendor_email': email })


@require_http_methods(["GET"]) 
def vendor_appointment_get_availability(request):
    """Return current month availability with only available slots per date.
    Query: month=YYYY-MM (optional); defaults to current month.
    """
    try:
        month = request.GET.get('month')
        if not month:
            now = datetime.datetime.now()
            month = f"{now.year}-{str(now.month).zfill(2)}"

        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        year, mon = month.split('-')
        month_name = datetime.datetime(int(year), int(mon), 1).strftime('%B %Y')
        folder = "Printmax Support availability Calendar"
        legacy_folder = "Printmax Support avialablity Calendar"
        key = f"{folder}/{month_name}.json"
        data = { 'dates': {} }
        try:
            resp = s3.get_object(Bucket=settings.R2_BUCKET, Key=key)
            data = json.loads(resp['Body'].read().decode('utf-8'))
        except Exception:
            try:
                legacy_key = f"{legacy_folder}/{month_name}.json"
                resp = s3.get_object(Bucket=settings.R2_BUCKET, Key=legacy_key)
                data = json.loads(resp['Body'].read().decode('utf-8'))
            except Exception:
                data = { 'dates': {} }

        # Filter to only available slots for each date
        pruned = {}
        for d, slots in (data.get('dates') or {}).items():
            avail = [s for s, st in (slots or {}).items() if st == 'available']
            if avail:
                pruned[d] = avail

        # Build day-of-week metadata for client calendar header
        return JsonResponse({ 'success': True, 'month': month, 'dates': pruned })
    except Exception as e:
        return JsonResponse({ 'success': False, 'error': str(e) }, status=500)


@csrf_exempt
@require_http_methods(["POST"]) 
def vendor_appointment_book(request):
    """Book an appointment for a vendor: save appointment.json and mark slot busy.
    Body: { vendor_email, date: 'YYYY-MM-DD', slot: 'HH:MM-HH:MM' }
    """
    try:
        body = json.loads(request.body or '{}')
        vendor_email = (body.get('vendor_email') or '').strip()
        date_str = (body.get('date') or '').strip()
        slot = (body.get('slot') or '').strip()
        if not vendor_email or not date_str or not slot:
            return JsonResponse({ 'success': False, 'error': 'vendor_email, date, slot are required' }, status=400)

        # Save appointment.json into vendor folder
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        appointment = {
            'vendor_email': vendor_email,
            'appointment_date': date_str,
            'appointment_slot': slot,
            'booked_at': datetime.datetime.utcnow().isoformat() + 'Z'
        }
        app_key = f"vendor_register_details/{sanitize_email(vendor_email)}/appointment.json"
        s3.put_object(Bucket=settings.R2_BUCKET, Key=app_key, Body=json.dumps(appointment, indent=2), ContentType='application/json')

        # Mark the selected slot as busy (not available) in month file
        year, mon, _ = date_str.split('-')
        month = f"{year}-{mon}"
        month_name = datetime.datetime(int(year), int(mon), 1).strftime('%B %Y')
        folder = "Printmax Support availability Calendar"
        key = f"{folder}/{month_name}.json"
        data = { 'dates': {} }
        try:
            resp = s3.get_object(Bucket=settings.R2_BUCKET, Key=key)
            data = json.loads(resp['Body'].read().decode('utf-8'))
        except Exception:
            data = { 'dates': {} }

        # Ensure structure and flip the slot to not available
        dates = data.get('dates') or {}
        if date_str not in dates:
            dates[date_str] = {}
        if slot not in dates[date_str]:
            dates[date_str][slot] = 'not available'
        else:
            dates[date_str][slot] = 'not available'
        data['dates'] = dates

        s3.put_object(Bucket=settings.R2_BUCKET, Key=key, Body=json.dumps(data, indent=2), ContentType='application/json')

        return JsonResponse({ 'success': True, 'appointment': appointment })
    except Exception as e:
        return JsonResponse({ 'success': False, 'error': str(e) }, status=500)


def get_vendor_details_by_email(email):
    """
    Fetch vendor details from D1 via Worker. Falls back to legacy R2 storage if needed.
    Returns data_source so callers can trust only D1 metadata for authentication.
    """
    if not email:
        return None

    vendor_data = None
    data_source = None

    try:
        endpoint, resp = post_to_worker('/get-vendor-by-email', {'email': email})
        if resp.status_code == 200:
            payload = resp.json()
            if payload.get('success'):
                vendor_data = payload.get('vendor')
                data_source = 'd1'
            else:
                print(f"⚠️ Worker returned error for {email}: {payload}")
        elif resp.status_code != 404:
            print(f"⚠️ Worker API error ({resp.status_code}) for {endpoint}: {resp.text[:300]}")
    except Exception as worker_error:
        print(f"⚠️ Worker fetch failed for {email}: {worker_error}")

    if vendor_data:
        profile_image_url = vendor_data.get('profile_image_url', '')
        try:
            s3 = boto3.client(
                's3',
                aws_access_key_id=settings.R2_ACCESS_KEY,
                aws_secret_access_key=settings.R2_SECRET_KEY,
                endpoint_url=settings.R2_ENDPOINT,
                region_name='auto'
            )

            if profile_image_url:
                try:
                    test_response = requests.head(profile_image_url, timeout=5)
                    if test_response.status_code != 200:
                        profile_image_url = refresh_profile_image_url(s3, email, vendor_data)
                except Exception as e:
                    print(f"⚠️ Profile image head request failed: {e}, refreshing...")
                    profile_image_url = refresh_profile_image_url(s3, email, vendor_data)
            else:
                profile_image_url = find_and_generate_profile_image_url(s3, email)
        except Exception as image_error:
            print(f"⚠️ Unable to refresh profile image from R2: {image_error}")

        vendor_data['profile_image_url'] = profile_image_url
        return {
            'vendor_id': vendor_data.get('vendor_id', ''),
            'vendor_name': vendor_data.get('vendor_name', ''),
            'vendor_email': vendor_data.get('email') or vendor_data.get('vendor_email', ''),
            'phone_number': vendor_data.get('phone_number', ''),
            'shop_address': vendor_data.get('shop_address', ''),
            'city': vendor_data.get('city', ''),
            'pincode': vendor_data.get('pincode', ''),
            'vendor_token': vendor_data.get('vendor_token', ''),
            'profile_image_url': vendor_data.get('profile_image_url', ''),
            'data_source': data_source or 'd1'
        }

    legacy_details = _get_vendor_details_from_r2_storage(email)
    if legacy_details:
        legacy_details['data_source'] = 'r2'
    return legacy_details


def _get_vendor_details_from_r2_storage(email):
    """
    Legacy helper that fetches vendor details from R2. Kept as fallback only.
    """
    s3 = boto3.client('s3',
        aws_access_key_id=settings.R2_ACCESS_KEY,
        aws_secret_access_key=settings.R2_SECRET_KEY,
        endpoint_url=settings.R2_ENDPOINT,
        region_name='auto'
    )
    try:
        reg_key = f'vendor_register_details/{sanitize_email(email)}/registration_details.json'
        response = s3.get_object(Bucket=settings.R2_BUCKET, Key=reg_key)
        vendor_data = json.loads(response['Body'].read().decode('utf-8'))
        
        print(f"🔍 DEBUG - Vendor data for {email}:")
        print(f"   - vendor_name: {vendor_data.get('vendor_name', '')}")
        print(f"   - profile_image_url: {vendor_data.get('profile_image_url', '')}")
        
        profile_image_url = vendor_data.get('profile_image_url', '')
        if profile_image_url:
            try:
                test_response = requests.head(profile_image_url, timeout=5)
                if test_response.status_code != 200:
                    print(f"⚠️ Profile image URL not accessible, attempting to refresh...")
                    profile_image_url = refresh_profile_image_url(s3, email, vendor_data)
            except Exception as e:
                print(f"⚠️ Profile image URL test failed: {e}, attempting to refresh...")
                profile_image_url = refresh_profile_image_url(s3, email, vendor_data)
        else:
            profile_image_url = find_and_generate_profile_image_url(s3, email)
        
        return {
            'vendor_id': vendor_data.get('vendor_id', ''),
            'vendor_name': vendor_data.get('vendor_name', ''),
            'vendor_email': vendor_data.get('vendor_email', ''),
            'phone_number': vendor_data.get('phone_number', ''),
            'shop_address': vendor_data.get('shop_address', ''),
            'city': vendor_data.get('city', ''),
            'pincode': vendor_data.get('pincode', ''),
            'vendor_token': vendor_data.get('vendor_token', ''),
            'profile_image_url': profile_image_url,
            'data_source': 'r2',
        }
    except Exception as e:
        print(f"Error fetching vendor details for {email} from R2: {str(e)}")
        return None


def refresh_profile_image_url(s3, email, vendor_data):
    """Refresh profile image URL by generating a new presigned URL"""
    try:
        # Look for existing profile image in the profile folder
        profile_prefix = f"vendor_register_details/{sanitize_email(email)}/profile/"
        
        # List objects in the profile folder
        response = s3.list_objects_v2(
            Bucket=settings.R2_BUCKET,
            Prefix=profile_prefix
        )
        
        if 'Contents' in response and response['Contents']:
            # Get the most recent profile image
            latest_image = max(response['Contents'], key=lambda x: x['LastModified'])
            image_key = latest_image['Key']
            
            # Generate a new presigned URL
            new_url = s3.generate_presigned_url(
                'get_object',
                Params={'Bucket': settings.R2_BUCKET, 'Key': image_key},
                ExpiresIn=3600 * 24 * 7  # 7 days
            )
            
            print(f"✅ Refreshed profile image URL: {new_url}")
            
            # Update the vendor data with the new URL
            vendor_data['profile_image_url'] = new_url
            
            # Save the updated data back to R2
            reg_key = f'vendor_register_details/{sanitize_email(email)}/registration_details.json'
            s3.put_object(
                Bucket=settings.R2_BUCKET,
                Key=reg_key,
                Body=json.dumps(vendor_data, indent=2),
                ContentType='application/json'
            )
            
            return new_url
        else:
            print(f"❌ No profile image found for {email}")
            return None
            
    except Exception as e:
        print(f"❌ Error refreshing profile image URL: {str(e)}")
        return None


def find_and_generate_profile_image_url(s3, email):
    """Find existing profile image and generate URL"""
    try:
        # Look for existing profile image in the profile folder
        profile_prefix = f"vendor_register_details/{sanitize_email(email)}/profile/"
        
        # List objects in the profile folder
        response = s3.list_objects_v2(
            Bucket=settings.R2_BUCKET,
            Prefix=profile_prefix
        )
        
        if 'Contents' in response and response['Contents']:
            # Get the most recent profile image
            latest_image = max(response['Contents'], key=lambda x: x['LastModified'])
            image_key = latest_image['Key']
            
            # Generate a presigned URL
            new_url = s3.generate_presigned_url(
                'get_object',
                Params={'Bucket': settings.R2_BUCKET, 'Key': image_key},
                ExpiresIn=3600 * 24 * 7  # 7 days
            )
            
            print(f"✅ Generated new profile image URL: {new_url}")
            return new_url
        else:
            print(f"❌ No profile image found for {email}")
            return None
            
    except Exception as e:
        print(f"❌ Error finding profile image: {str(e)}")
        return None


def vendordashboard(request):
    try:
        # Get vendor details from session
        vendor_details = None
        vendor_email = request.session.get('vendor_email')
        vendor_id = request.session.get('vendor_id')  # Get vendor_id directly from session
        vendor_status = request.session.get('vendor_status', 'pending').strip().lower()
        
        print(f"🔍 Session data - Email: {vendor_email}, Vendor ID: {vendor_id}, Status: {vendor_status}")
        
        if vendor_email and vendor_id:
            vendor_details = get_vendor_details_by_email(vendor_email)
            if vendor_details:
                vendor_details_source = vendor_details.get('data_source')
                details_vendor_id = vendor_details.get('vendor_id')
                # Only let D1 data update the session vendor_id
                if details_vendor_id and vendor_details_source == 'd1' and details_vendor_id != vendor_id:
                    vendor_id = details_vendor_id
                    request.session['vendor_id'] = vendor_id  # Update session with correct vendor_id
                elif vendor_details_source != 'd1':
                    print("⚠️ Vendor details loaded from legacy R2 store; keeping existing session vendor_id")
                print(f"🔍 Loading dashboard for vendor: {vendor_details.get('vendor_name')} (ID: {vendor_id})")
            else:
                print("❌ Could not fetch vendor details from R2")
        elif vendor_email and not vendor_id:
            # Try to get vendor_id from vendor details
            vendor_details = get_vendor_details_by_email(vendor_email)
            if vendor_details:
                if vendor_details.get('data_source') == 'd1' and vendor_details.get('vendor_id'):
                    vendor_id = vendor_details.get('vendor_id')
                    request.session['vendor_id'] = vendor_id
                    print(f"🔍 Retrieved vendor ID from D1 details: {vendor_id}")
                else:
                    print("⚠️ Vendor details available only from R2; skipping session update until D1 sync completes")
        
        if not vendor_id:
            print("❌ No vendor ID found in session or vendor details")
            return render(request, 'vendordashboard.html', {
                'manual_print_jobs': [],
                'print_requests': [],
                'completed_jobs': [],
                'vendor_details': vendor_details,
                'vendor_details_error': 'Vendor not authenticated. Please login again.',
                'total_jobs': 0,
                'manual_print_count': 0,
                'print_requests_count': 0,
                'completed_jobs_count': 0,
                'vendor_status': vendor_status,
            })

        # Check for service update needed flag (set during login)
        show_service_update_modal = request.session.pop('service_update_needed', False)

        # Check vendor status - only show jobs if status is 'verified'
        if vendor_status != 'verified':
            print(f"⚠️ Vendor status is '{vendor_status}', not verified. Hiding print jobs.")
            return render(request, 'vendordashboard.html', {
                'manual_print_jobs': [],
                'print_requests': [],
                'completed_jobs': [],
                'vendor_details': vendor_details,
                'vendor_status': vendor_status,
                'vendor_status_message': 'Your verification is in process. Once verified within 24hrs, you can start receiving jobs too.',
                'total_jobs': 0,
                'manual_print_count': 0,
                'print_requests_count': 0,
                'completed_jobs_count': 0,
                'daily_earnings': 0,
            })

        # Fetch vendor-specific jobs strictly from D1 database (pending + completed for KPIs)
        pending_jobs = get_vendor_jobs_from_d1(vendor_id=vendor_id, vendor_email=vendor_email, job_status='NO') or []
        completed_jobs_raw = get_vendor_jobs_from_d1(vendor_id=vendor_id, vendor_email=vendor_email, job_status='YES') or []
        files = pending_jobs + completed_jobs_raw
        
        # Calculate Daily Earnings (Total Price - Platform Profit) for today's completed jobs
        daily_earnings = 0
        try:
            today_iso = datetime.date.today().isoformat()
            # If using local time, ensure alignment with completion_time format (typically ISO with maybe Z or offset)
            # completion_time example: "2023-10-27T10:00:00"
            
            for job in completed_jobs_raw:
                c_time = job.get('completion_time')
                if c_time:
                    # simplistic check: does it start with YYYY-MM-DD?
                    if c_time.startswith(today_iso):
                        total = float(job.get('total_price', 0) or 0)
                        platform = float(job.get('platform_profit', 0) or 0)
                        
                        # Fallback for platform_profit if not top-level, check pricing_details
                        if platform == 0:
                             pd = job.get('pricing_details')
                             if pd and isinstance(pd, dict):
                                 platform = float(pd.get('platform_profit', 0) or 0)

                        daily_earnings += (total - platform)
                        
        except Exception as e:
            print(f"⚠️ Error calculating daily earnings: {e}")
            daily_earnings = 0

        if not files:
            print(f"ℹ️ No jobs returned from D1 for vendor {vendor_id}")
        
        # Define service types for categorization
        manual_services = [
            'digital_print', 'project_binding', 'gloss_printing', 'jumbo_printing', 'golden_embossing'
        ]
        regular_services = [
            'regular_print', 'passport_print', 'photo_print', 'regular print', 'passport_photo'
        ]
        
        manual_print_jobs = []
        print_requests = []
        completed_jobs = []
        
        for job in files:
            # Default rendered_status
            job['rendered_status'] = job.get('rendered_status') or job.get('metadata', {}).get('rendered_status', 'NO')
            # Get job status - prioritize job_completed over job_completed_status
            job_completed_raw = job.get('job_completed')
            if not job_completed_raw:
                job_completed_raw = job.get('job_completed_status')
            if job_completed_raw is None or str(job_completed_raw).strip() == '':
                job_completed_raw = 'NO'
            job_completed = str(job_completed_raw).strip().upper()
            
            # Normalize job_completed in the job object to ensure frontend receives consistent format
            job['job_completed'] = job_completed
            
            service_type = (job.get('service_type') or '').strip().lower()
            vendor_status = job.get('vendor_status', 'not sended').lower()
            is_hidden = job.get('is_hidden', 'false').lower() == 'true'
            
            # Ensure service_type is preserved in job object for frontend categorization
            # Frontend needs the original service_type to categorize into passport, golden, etc. sections
            job['service_type'] = job.get('service_type', '').strip() or service_type
            
            print(f"🔍 Job: {job.get('filename', 'unknown')} - job_completed: {job_completed}, service_type: {service_type}, vendor_status: {vendor_status}")
            
            # Skip cancelled jobs - they should not be displayed
            if job_completed == 'CANCELLED':
                continue
            
            if job_completed == 'YES':
                # Completed jobs - regardless of service type
                if not is_hidden:
                    completed_jobs.append(job)
            elif job_completed == 'NO':
                # Pending jobs - categorize by service type
                if service_type in manual_services:
                    # Manual print jobs (digital_print, project_binding, etc.)
                    manual_print_jobs.append(job)
                elif service_type in regular_services or service_type == '' or service_type == 'regular print':
                    # Regular print jobs - show both pending and accepted
                    print_requests.append(job)
                else:
                    # Default to print requests for unknown service types
                    print_requests.append(job)

        print(f"📊 Job categorization - Manual: {len(manual_print_jobs)}, Requests: {len(print_requests)}, Completed: {len(completed_jobs)}")

        context = {
            'manual_print_jobs': manual_print_jobs,
            'print_requests': print_requests,
            'completed_jobs': completed_jobs,
            'vendor_details': vendor_details,
            'vendor_status': vendor_status,
            'total_jobs': len(manual_print_jobs) + len(print_requests) + len(completed_jobs),
            'manual_print_count': len(manual_print_jobs),
            'print_requests_count': len(print_requests),
            'completed_jobs_count': len(completed_jobs),
            # Provide a flat list for initial render JS to consume without extra fetch
            'initial_jobs': manual_print_jobs + print_requests + completed_jobs,
             'daily_earnings': daily_earnings,
            'show_service_update_modal': show_service_update_modal,
        }
        
        # Cache the context for faster subsequent loads
        cache_key = f"vendor_dashboard_{vendor_email}_{vendor_id}"
        _vendor_dashboard_cache[cache_key] = context
        _cache_timestamp[cache_key] = time.time()
        
        return render(request, 'vendordashboard.html', context)
        
    except Exception as e:
        print(f"Error loading vendor dashboard data: {str(e)}")
        vendor_status = request.session.get('vendor_status', 'pending').strip().lower()
        return render(request, 'vendordashboard.html', {
            'manual_print_jobs': [],
            'print_requests': [],
            'completed_jobs': [],
            'vendor_details': None,
            'vendor_details_error': 'Dashboard error. Please try again later.',
            'vendor_status': vendor_status,
            'total_jobs': 0,
            'manual_print_count': 0,
            'print_requests_count': 0,
            'completed_jobs_count': 0,
        })



def save_user_signup_to_d1(payload):
    """
    Persist user signup/login metadata to the D1 User_signup_details table via Worker.
    """
    try:
        endpoint, resp = post_to_worker('/add-user-signup', payload)
        if resp.status_code != 200:
            print(f"⚠️ Failed to persist signup details to D1 ({resp.status_code}) via {endpoint}: {resp.text[:300]}")
        else:
            response_json = resp.json()
            if not response_json.get('success'):
                print(f"⚠️ Worker add-user-signup responded with error: {response_json}")
    except Exception as exc:
        print(f"❌ Error saving signup details to D1: {exc}")


def get_user_details_from_d1(user_email):
    """
    Fetch user signup details from D1 via Worker (User_signup_details table)
    """
    if not user_email:
        return None

    try:
        endpoint, resp = post_to_worker('/get-user-signup', {'email': user_email})
        if resp.status_code != 200:
            print(f"⚠️ Worker get-user-signup failed ({resp.status_code}) via {endpoint}: {resp.text[:300]}")
            return None

        payload = resp.json()
        if not payload.get('success'):
            print(f"⚠️ Worker get-user-signup responded with error: {payload}")
            return None

        user_data = payload.get('user_signup') or {}
        return {
            'name': user_data.get('name', ''),
            'email': user_data.get('email', ''),
            'profile_picture': user_data.get('picture', ''),
            'given_name': user_data.get('given_name', ''),
            'family_name': user_data.get('family_name', ''),
            'locale': user_data.get('locale', ''),
            'email_verified': bool(user_data.get('email_verified'))
        }
    except Exception as exc:
        print(f"❌ Error fetching user signup details from D1: {exc}")
        return None


def get_user_jobs_from_r2(user_email):
    """
    Get all jobs uploaded by a specific user from R2 storage - OPTIMIZED for speed
    """
    s3 = boto3.client('s3',
                      aws_access_key_id=settings.R2_ACCESS_KEY,
                      aws_secret_access_key=settings.R2_SECRET_KEY,
                      endpoint_url=settings.R2_ENDPOINT,
                      region_name='auto')

    try:
        # List all files in the user's folder
        user_prefix = f"users/{user_email}/"
        objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=user_prefix)
        user_jobs = []

        # ULTRA-FAST: Process all jobs in parallel using ThreadPoolExecutor
        import concurrent.futures
        
        def process_single_job(obj):
            try:
                key = obj["Key"]
                filename = key.split("/")[-1]

                # Skip if it's just the folder itself
                if filename == "":
                    return None

                # Generate presigned URL for preview
                url = s3.generate_presigned_url(
                    ClientMethod='get_object',
                    Params={
                        'Bucket': settings.R2_BUCKET,
                        'Key': key
                    },
                    ExpiresIn=3600
                )

                # Get object metadata
                head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                metadata = head_response.get('Metadata', {})

                # Determine file type and icon
                file_extension = filename.split('.')[-1].lower() if '.' in filename else ''
                file_type = get_file_type(file_extension)

                # Calculate estimated pages if not in metadata
                pages = metadata.get('pages', estimate_pages_from_size(obj.get('Size', 0), file_extension))

                # Get vendor coordinates for tracking
                vendor_id = metadata.get('vendor', 'firozshop')
                vendor_lat = None
                vendor_lng = None
                vendor_email = None
                
                # Try to get vendor coordinates from vendor_register_details
                try:
                    vendor_email = get_vendor_email_by_vendor_id(vendor_id)
                    if vendor_email:
                        # Get vendor coordinates from registration details
                        vendor_coords = get_vendor_coordinates_from_email(vendor_email)
                        if vendor_coords:
                            vendor_lat = vendor_coords.get('latitude')
                            vendor_lng = vendor_coords.get('longitude')
                except Exception as e:
                    print(f"Error getting vendor coordinates for {vendor_id}: {str(e)}")

                # Check job completion and vendor status for filtering
                job_completed = metadata.get('job_completed', 'NO').upper()
                vendor_status = metadata.get('vendor_status', 'not sended').lower()
                
                # Only include jobs that are not completed and vendor status is not 'sended'
                if job_completed == 'YES' or vendor_status == 'sended':
                    return None

                # Build job info
                job_info = {
                    "filename": filename,
                    "preview_url": url,
                    "file_type": file_type,
                    "file_extension": file_extension,
                    "size": format_file_size(obj.get('Size', 0)),
                    "pages": pages,
                    "status": metadata.get('status', 'pending').title(),
                    "uploaded_at": obj["LastModified"].strftime("%Y-%m-%d %H:%M"),
                    "priority": metadata.get('priority', 'Medium'),
                    "copies": metadata.get('copies', '1'),
                    "color": metadata.get('color', 'Black and White'),
                    "orientation": metadata.get('orientation', 'portrait'),
                    "pageRange": metadata.get('pagerange', 'all'),
                    "specificPages": metadata.get('specificpages', ''),
                    "pageSize": metadata.get('pagesize', 'A4'),
                    "spiralBinding": metadata.get('spiralbinding', 'No'),
                    "lamination": metadata.get('lamination', 'No'),
                    "job_completed": job_completed,
                    "vendor_status": vendor_status,
                    "timestamp": metadata.get('timestamp', obj["LastModified"].isoformat()),
                    "vendor": metadata.get('vendor', 'firozshop'),
                    "vendor_lat": vendor_lat,
                    "vendor_lng": vendor_lng,
                    "vendor_email": vendor_email,
                    "service_type": metadata.get('service_type', ''),
                    "job_id": metadata.get('job_id', ''),
                    "token": metadata.get('token', ''),
                    "feedback": metadata.get('feedback', ''),
                    "quality": metadata.get('quality', ''),
                    "thickness": metadata.get('thickness', ''),
                    "service_name": metadata.get('service_name', '')
                }
                
                # Parse pricing details from metadata if available (simplified for speed)
                pricing_details_str = metadata.get('pricing_details')
                if pricing_details_str:
                    try:
                        # Try to parse as compact JSON format first
                        try:
                            compact_pricing = json.loads(pricing_details_str)
                            if isinstance(compact_pricing, dict) and 'total' in compact_pricing:
                                # Parse compact JSON format
                                total_price = compact_pricing.get('total', 0)
                                price_per_page = compact_pricing.get('per_page', 0)
                                page_count = compact_pricing.get('pages', 0)
                                num_copies = compact_pricing.get('copies', 0)
                                pricing_key = compact_pricing.get('key', '')
                                quality_upgrade = compact_pricing.get('quality', 0)
                                
                                # Reconstruct full pricing details object
                                pricing_details = {
                                    'total_price': total_price,
                                    'pricing_breakdown': {
                                        'price_per_page': price_per_page,
                                        'page_count': page_count,
                                        'num_copies': num_copies,
                                        'pricing_key_used': pricing_key,
                                        'base_price': price_per_page,
                                        'quality_upgrade': quality_upgrade
                                    },
                                    'calculation_timestamp': metadata.get('timestamp', ''),
                                    'vendor_email': None
                                }
                                job_info['pricing_details'] = pricing_details
                                # print(f"💰 Pricing details loaded for {filename}: Rs{total_price}")  # Removed for speed
                            else:
                                # Try to parse as old full JSON format
                                pricing_details = json.loads(pricing_details_str)
                                job_info['pricing_details'] = pricing_details
                                # print(f"💰 Pricing details loaded for {filename}: Rs{pricing_details.get('total_price', 'N/A')}")  # Removed for speed
                        except json.JSONDecodeError:
                            # Try to parse as old pipe-separated format
                            if '|' in pricing_details_str:
                                parts = pricing_details_str.split('|')
                                if len(parts) >= 5:
                                    total_price = float(parts[0].replace('Rs', ''))
                                    price_per_page = float(parts[1])
                                    page_count = int(parts[2])
                                    num_copies = int(parts[3])
                                    pricing_key = parts[4]
                                    
                                    # Reconstruct pricing details object
                                    pricing_details = {
                                        'total_price': total_price,
                                        'pricing_breakdown': {
                                            'price_per_page': price_per_page,
                                            'page_count': page_count,
                                            'num_copies': num_copies,
                                            'pricing_key_used': pricing_key,
                                            'base_price': price_per_page,
                                            'quality_upgrade': 0
                                        },
                                        'calculation_timestamp': metadata.get('timestamp', ''),
                                        'vendor_email': None
                                    }
                                    job_info['pricing_details'] = pricing_details
                                    # print(f"💰 Pricing details loaded for {filename}: Rs{total_price}")  # Removed for speed
                                else:
                                    # print(f"⚠️ Invalid pricing format for {filename}")
                                    job_info['pricing_details'] = None
                            else:
                                # print(f"⚠️ Unknown pricing format for {filename}")
                                job_info['pricing_details'] = None
                    except (json.JSONDecodeError, ValueError) as e:
                        # print(f"⚠️ Error parsing pricing details for {filename}: {e}")
                        job_info['pricing_details'] = None
                else:
                    job_info['pricing_details'] = None

                return job_info
            except Exception as e:
                # print(f"Error processing job {obj.get('Key', 'unknown')}: {str(e)}")
                return None

        # Process all jobs in parallel with maximum workers
        with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
            future_to_obj = {executor.submit(process_single_job, obj): obj for obj in objects.get("Contents", [])}
            
            for future in concurrent.futures.as_completed(future_to_obj):
                result = future.result()
                if result is not None:
                    user_jobs.append(result)

        # Sort by upload date (newest first)
        user_jobs.sort(key=lambda x: x.get('timestamp', ''), reverse=True)
        
        print(f"⚡ ULTRA-FAST: Loaded {len(user_jobs)} jobs in parallel for {user_email}")
        return user_jobs

    except Exception as e:
        print(f"Error getting user jobs from R2: {str(e)}")
        return []


def userdashboard(request):
    # Check if user is authenticated
    print(f"🔍 Userdashboard access - User authenticated: {request.user.is_authenticated}")
    print(f"🔍 User: {request.user}")
    print(f"🔍 Session keys: {list(request.session.keys())}")
    
    if not request.user.is_authenticated:
        print("❌ User not authenticated, redirecting to login")
        return redirect('/login/')
    
    # ✅ STRICT: Verify user has valid Google session
    google_user_id = request.session.get('google_user_id')
    user_email = request.session.get('user_email') or request.user.email
    
    if not google_user_id:
        print(f"❌ No Google session found for user {user_email}, redirecting to login")
        from django.contrib.auth import logout
        logout(request)
        return redirect('/login/')
    
    # ✅ STRICT: Verify user exists in D1 User_signup_details (must have signed up via Google)
    user_details = get_user_details_from_d1(user_email)
    if not user_details:
        print(f"❌ User {user_email} not found in D1 User_signup_details, redirecting to login")
        from django.contrib.auth import logout
        logout(request)
        return redirect('/login/')
    
    print(f"✅ Valid Google session and D1 record confirmed for {user_email}")

    try:
        # ULTRA-FAST: Avoid synchronous job loading; fetch details only
        user_details = get_user_details_from_d1(request.user.email) or {}
        preloaded_jobs = get_user_jobs_from_d1(request.user.email)
        if not isinstance(preloaded_jobs, list):
            preloaded_jobs = []
        pending_job_list = [
            job for job in preloaded_jobs
            if str(job.get('job_completed', 'NO')).upper() != 'YES'
        ]
        total_jobs = len(preloaded_jobs)
        pending_jobs = len(pending_job_list)
        completed_jobs = total_jobs - pending_jobs
        current_month = datetime.datetime.now().strftime("%Y-%m")
        current_month_jobs = sum(
            1 for job in preloaded_jobs
            if str(job.get('uploaded_at') or job.get('timestamp', '')).startswith(current_month)
        )
        total_earnings = current_month_jobs * 50

        session_name = (request.session.get('user_name') or '').strip()
        session_picture = (request.session.get('user_picture') or '').strip()

        if session_picture and not user_details.get('profile_picture'):
            user_details['profile_picture'] = session_picture
        if not user_details.get('name') and session_name:
            user_details['name'] = session_name
        if not user_details.get('email'):
            user_details['email'] = request.user.email

        first_name_candidates = [
            user_details.get('given_name'),
            request.user.first_name,
            session_name.split()[0] if session_name else '',
            (request.user.email.split('@')[0] if request.user.email else ''),
            request.user.username
        ]
        nav_first_name = next((name for name in first_name_candidates if name), '').strip()
        if not user_details.get('given_name') and nav_first_name:
            user_details['given_name'] = nav_first_name

        nav_profile_image = user_details.get('profile_picture') or session_picture

        # Serialize preloaded jobs for instant client rendering
        try:
            user_jobs_json = json.dumps(preloaded_jobs, default=str)
        except Exception as serialize_error:
            print(f"⚠️ Error serializing preloaded jobs: {serialize_error}")
            user_jobs_json = '[]'
        try:
            user_points = get_total_user_points(request.user.email)
        except Exception:
            user_points = 0

        # Get Firebase config from settings
        firebase_config = getattr(settings, 'FIREBASE_CONFIG', {})
        firebase_vapid_key = getattr(settings, 'FIREBASE_VAPID_PUBLIC_KEY', '')
        
        context = {
            'user': request.user,
            'user_details': user_details,
            'user_nav_first_name': nav_first_name,
            'user_nav_profile_image': nav_profile_image,
            'firebase_uid': request.session.get('firebase_uid'),
            'auth_method': request.session.get('auth_method', 'unknown'),
            'user_jobs': preloaded_jobs,
            'user_jobs_json': user_jobs_json,
            'total_jobs': total_jobs,
            'pending_jobs': pending_jobs,
            'completed_jobs': completed_jobs,
            'total_earnings': total_earnings,
            'user_points': user_points,
            'GOOGLE_CLIENT_ID': settings.GOOGLE_CLIENT_ID,
            'GOOGLE_DEVELOPER_KEY': settings.GOOGLE_DEVELOPER_KEY,
            # Note: GOOGLE_MAPS_API is NOT passed to frontend - it's server-side only for Distance Matrix API
            # Firebase FCM Configuration
            'FIREBASE_CONFIG': firebase_config,
            'FIREBASE_VAPID_KEY': firebase_vapid_key,
        }
        return render(request, 'userdashboard.html', context)
        
    except Exception as e:
        print(f"Error loading user dashboard: {str(e)}")
        session_name = (request.session.get('user_name') or '').strip()
        nav_first_name = (
            request.user.first_name or
            (session_name.split()[0] if session_name else '') or
            (request.user.email.split('@')[0] if request.user.email else request.user.username)
        )
        nav_profile_image = (request.session.get('user_picture') or '').strip()
        fallback_details = {
            'name': session_name or request.user.get_full_name() or request.user.username,
            'email': request.user.email,
            'profile_picture': nav_profile_image,
            'given_name': nav_first_name,
        }
        # Return minimal context on error for faster fallback
        return render(request, 'userdashboard.html', {
            'user': request.user,
            'user_details': fallback_details,
            'user_nav_first_name': nav_first_name,
            'user_nav_profile_image': nav_profile_image,
            'firebase_uid': request.session.get('firebase_uid'),
            'auth_method': request.session.get('auth_method', 'unknown'),
            'user_jobs': [],
            'user_jobs_json': json.dumps([]),
            'total_jobs': 0,
            'pending_jobs': 0,
            'completed_jobs': 0,
            'total_earnings': 0,
            'user_points': 0
        })


def get_user_jobs_from_d1(user_email):
    """
    Fetch user print jobs from D1 database via Worker API ONLY
    Returns list of jobs with R2 paths from database
    STRICTLY NO R2 FALLBACK - Only returns data from User_print_jobs table
    """
    try:
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            print("⚠️ Worker API not configured for D1 database")
            return []
        
        # Construct Worker API endpoint
        worker_endpoint = api_url.rstrip('/') + '/get-user-print-jobs'
        
        # Prepare payload
        worker_payload = {
            'user_email': user_email
        }
        
        # Make request to Worker API - use x-api-key header (not Authorization Bearer)
        resp = requests.post(
            worker_endpoint,
            json=worker_payload,
            headers={
                'x-api-key': api_key,
                'Content-Type': 'application/json'
            },
            timeout=10
        )
        
        if resp.status_code == 200:
            data = resp.json()
            if data.get('success') and data.get('data'):
                jobs = data['data']
                # Process jobs to generate presigned URLs from R2 paths stored in database
                s3 = boto3.client('s3',
                                aws_access_key_id=settings.R2_ACCESS_KEY,
                                aws_secret_access_key=settings.R2_SECRET_KEY,
                                endpoint_url=settings.R2_ENDPOINT,
                                region_name='auto')
                
                processed_jobs = []
                vendor_email_cache = {}
                vendor_coords_cache = {}

                def _normalize_identifier(value):
                    if value is None:
                        return ''
                    return str(value).strip().lower()
                for job in jobs:
                    # Get storage folder and R2 path from database
                    storage_folder = job.get('storage_folder') or 'user_print_jobs'
                    job['storage_folder'] = storage_folder
                    job.setdefault('user', job.get('user_email', ''))
                    job.setdefault('vendor', job.get('vendor_id', ''))

                    # Get R2 path from database (r2_path field) - ONLY use R2 for file URL
                    r2_path = job.get('r2_path')
                    if not r2_path and job.get('filename'):
                        r2_path = f"{storage_folder.rstrip('/')}/{job.get('user_email', '')}/{job.get('filename')}"
                        job['r2_path'] = r2_path

                    if r2_path:
                        try:
                            # Generate presigned URL from R2 path - ONLY for file rendering
                            url = s3.generate_presigned_url(
                                ClientMethod='get_object',
                                Params={
                                    'Bucket': settings.R2_BUCKET,
                                    'Key': r2_path
                                },
                                ExpiresIn=3600
                            )
                            job['url'] = url
                            job['preview_url'] = url
                            job['download_url'] = url
                        except Exception as e:
                            print(f"⚠️ Error generating presigned URL for {r2_path}: {e}")
                            job['url'] = ''
                            job['preview_url'] = ''
                            job['download_url'] = ''
                    else:
                        job['url'] = job['preview_url'] = job['download_url'] = ''
                    
                    # Ensure preview/download URLs exist even if R2 path missing
                    job.setdefault('preview_url', job.get('url', ''))
                    job.setdefault('download_url', job.get('url', ''))
                    
                    # Map ALL metadata from D1 database (not from R2)
                    # All fields come from D1, only file URL comes from R2
                    job_completed_value = job.get('job_completed', job.get('job_completed_status', 'NO')) or 'NO'
                    job['job_completed'] = job_completed_value
                    job.setdefault('service_type', job.get('service_type', ''))
                    job.setdefault('pages', str(job.get('page_count', job.get('pages', '0'))))
                    job.setdefault('rendered_status', job.get('rendered_status', 'NO'))
                    job.setdefault('vendor_status', job.get('vendor_status', 'not sended'))
                    job.setdefault('token', job.get('token', ''))
                    job.setdefault('job_id', job.get('job_id', ''))
                    copies_value = job.get('copies', job.get('num_copies'))
                    job['copies'] = str(copies_value) if copies_value not in [None, ''] else '1'
                    job.setdefault('color', job.get('color', ''))
                    job.setdefault('orientation', job.get('orientation', ''))
                    job.setdefault('pageSize', job.get('pageSize', ''))
                    job.setdefault('timestamp', job.get('timestamp', ''))
                    job.setdefault('uploaded_at', job.get('timestamp', ''))
                    job.setdefault('completion_time', job.get('completion_time', ''))
                    job.setdefault('total_price', job.get('total_price', 0))
                    job.setdefault('final_amount', job.get('final_amount', 0))
                    job.setdefault('user_email', job.get('user_email', ''))
                    vendor_identifier = job.get('vendor_id') or job.get('vendor')
                    vendor_email_value = (job.get('vendor_email') or '').strip()
                    if not vendor_email_value and vendor_identifier:
                        cache_key = _normalize_identifier(vendor_identifier)
                        if cache_key not in vendor_email_cache:
                            try:
                                resolved_email = get_vendor_email_by_vendor_id(vendor_identifier) or ''
                            except Exception as vendor_lookup_error:
                                print(f"⚠️ Error resolving vendor email for {vendor_identifier}: {vendor_lookup_error}")
                                resolved_email = ''
                            vendor_email_cache[cache_key] = resolved_email
                        vendor_email_value = vendor_email_cache.get(cache_key, '')
                    job['vendor_email'] = vendor_email_value or job.get('vendor_email', '')
                    job.setdefault('vendor_id', job.get('vendor_id', ''))
                    job.setdefault('filename', job.get('filename', ''))

                    # Parse pricing_details if present
                    pricing_details_raw = job.get('pricing_details')
                    parsed_pricing = None
                    if isinstance(pricing_details_raw, str) and pricing_details_raw.strip():
                        try:
                            parsed_pricing = json.loads(pricing_details_raw)
                        except json.JSONDecodeError:
                            parsed_pricing = None
                    elif isinstance(pricing_details_raw, dict):
                        parsed_pricing = pricing_details_raw
                    if parsed_pricing is not None:
                        job['pricing_details'] = parsed_pricing

                    # Page range and B&W/Color range values for invoice (Mix type)
                    page_range_raw = (job.get('page_range') or job.get('pageRange') or '').strip()
                    job['page_range'] = page_range_raw
                    bw_range_display = (job.get('bw_page_range_value') or job.get('bwPageRangeValue') or '').strip()
                    color_range_display = (job.get('color_page_range_value') or job.get('colorPageRangeValue') or '').strip()
                    if not bw_range_display and not color_range_display and page_range_raw:
                        import re
                        bw_m = re.search(r'BW:\s*([^|]+)', page_range_raw, re.I)
                        color_m = re.search(r'Color:\s*(.+)$', page_range_raw, re.I)
                        if bw_m:
                            bw_range_display = bw_m.group(1).strip()
                        if color_m:
                            color_range_display = color_m.group(1).strip()
                    job['bw_page_range_value'] = bw_range_display or None
                    job['color_page_range_value'] = color_range_display or None

                    # Create metadata structure from D1 fields (not R2)
                    job['metadata'] = {
                        'status': job.get('status', 'pending'),
                        'job_completed': job_completed_value,
                        'copies': job.get('copies', '1'),
                        'color': job.get('color', ''),
                        'orientation': job.get('orientation', ''),
                        'page_size': job.get('pageSize', ''),
                        'pages': str(job.get('page_count', job.get('pages', '0'))),
                        'timestamp': job.get('timestamp', ''),
                        'vendor': job.get('vendor_id', ''),
                        'user': job.get('user_email', ''),
                        'service_type': job.get('service_type', ''),
                        'job_id': job.get('job_id', ''),
                        'token': job.get('token', ''),
                        'vendor_id': job.get('vendor_id', ''),
                        'rendered_status': job.get('rendered_status', 'NO'),
                        'vendor_status': job.get('vendor_status', 'not sended'),
                        'total_price': job.get('total_price', 0),
                        'final_amount': job.get('final_amount', 0),
                        'pageRange': page_range_raw,
                        'bwPageRangeValue': bw_range_display or '',
                        'colorPageRangeValue': color_range_display or '',
                    }
                    
                    job['metadata']['pricing_details'] = parsed_pricing

                    vendor_lat_value = job.get('vendor_lat')
                    vendor_lng_value = job.get('vendor_lng')
                    if (not vendor_lat_value or not vendor_lng_value) and vendor_email_value:
                        coords_key = _normalize_identifier(vendor_email_value)
                        if coords_key not in vendor_coords_cache:
                            coords = get_vendor_coordinates_from_email(vendor_email_value) or {}
                            vendor_coords_cache[coords_key] = coords
                        coords = vendor_coords_cache.get(coords_key) or {}
                        if coords:
                            vendor_lat_value = vendor_lat_value or coords.get('latitude')
                            vendor_lng_value = vendor_lng_value or coords.get('longitude')
                            if coords.get('vendor_name'):
                                job.setdefault('vendor_name', coords.get('vendor_name'))
                            if coords.get('shop_address'):
                                job.setdefault('vendor_shop_address', coords.get('shop_address'))
                            if coords.get('city'):
                                job.setdefault('vendor_city', coords.get('city'))

                    job['vendor_lat'] = str(vendor_lat_value) if vendor_lat_value not in [None, ''] else ''
                    job['vendor_lng'] = str(vendor_lng_value) if vendor_lng_value not in [None, ''] else ''

                    vendor_label = (
                        job.get('vendor_name') or
                        vendor_email_value or
                        job.get('vendor') or
                        vendor_identifier or ''
                    )
                    job['vendor_label'] = vendor_label
                    
                    processed_jobs.append(job)
                
                if processed_jobs:
                    processed_jobs.sort(
                        key=lambda job: (
                            job.get('timestamp')
                            or job.get('uploaded_at')
                            or ''
                        ),
                        reverse=True,
                    )
                    return processed_jobs
                print(f"ℹ️ D1 returned no jobs for {user_email}")
                return []
            else:
                print(f"⚠️ Worker API returned error: {data.get('error', 'Unknown error')}")
                return []
        else:
            print(f"⚠️ Worker API request failed with status {resp.status_code}")
            return []
            
    except Exception as e:
        print(f"⚠️ Error fetching user jobs from D1: {str(e)}")
        return []


def userdashboard_data(request):
    """
    AJAX endpoint to load user dashboard data quickly from D1 database
    """
    if not request.user.is_authenticated:
        return JsonResponse({'success': False, 'error': 'Not authenticated'}, status=401)
    
    try:
        # Load user data quickly
        user_details = get_user_details_from_d1(request.user.email)
        # Fetch jobs from D1 database instead of R2
        user_jobs = get_user_jobs_from_d1(request.user.email)
        
        # Calculate statistics
        total_jobs = len(user_jobs)
        pending_jobs = 0
        completed_jobs = 0
        current_month_jobs = 0
        current_month = datetime.datetime.now().strftime("%Y-%m")
        
        for job in user_jobs:
            job_completed = job.get('job_completed', 'NO')
            if job_completed == 'NO':
                pending_jobs += 1
            elif job_completed == 'YES':
                completed_jobs += 1
            
            uploaded_at = job.get('uploaded_at', job.get('timestamp', ''))
            if uploaded_at and uploaded_at.startswith(current_month):
                current_month_jobs += 1

        total_earnings = current_month_jobs * 50

        # Points start at 0 and accrue only via vendor cancellations
        try:
            user_points = get_total_user_points(request.user.email)
        except Exception:
            user_points = 0

        return JsonResponse({
            'success': True,
            'user_details': user_details,
            'user_jobs': user_jobs,
            'total_jobs': total_jobs,
            'pending_jobs': pending_jobs,
            'completed_jobs': completed_jobs,
            'total_earnings': total_earnings,
            'user_points': user_points
        })
        
    except Exception as e:
        print(f"Error loading user dashboard data: {str(e)}")
        return JsonResponse({
            'success': False,
            'error': str(e),
            'user_details': None,
            'user_jobs': [],
            'total_jobs': 0,
            'pending_jobs': 0,
            'completed_jobs': 0,
            'total_earnings': 0,
            'user_points': 0
        })


# ─────────────────────────────────────────────────────────────
# FILE LISTING FROM R2
# ─────────────────────────────────────────────────────────────

def get_vendor_printer_configuration(vendor_email):
    """
    Get printer configuration for a vendor from their pricing data
    """
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        # Get vendor pricing file
        pricing_key = f'vendor_register_details/{sanitize_email(vendor_email)}/pricing.json'
        response = s3.get_object(Bucket=settings.R2_BUCKET, Key=pricing_key)
        pricing_data = json.loads(response['Body'].read().decode('utf-8'))
        
        # Return printer configuration
        return pricing_data.get('printer_configuration', {})
        
    except Exception as e:
        print(f"❌ Error getting printer configuration for {vendor_email}: {str(e)}")
        return {}

def get_printer_name_for_service(printer_config, service_type):
    """
    Get the appropriate printer name for a given service type
    """
    if not printer_config:
        return "No Printer Configured"
    
    service_type = service_type.lower().strip()
    
    # Map service types to printer fields
    printer_mapping = {
        'digital_print': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'project_binding': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'gloss_printing': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'jumbo_printing': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'regular_print': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'passport_print': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3'],
        'photo_print': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3'],
        'passport_photo': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3']
    }
    
    # Get the appropriate printer fields for this service type
    printer_fields = printer_mapping.get(service_type, ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'])
    
    # Find the first non-NA printer
    for field in printer_fields:
        printer_name = printer_config.get(field, 'NA')
        if printer_name and printer_name.strip() and printer_name != 'NA':
            return printer_name
    
    return "No Printer Configured"

# ─────────────────────────────────────────────────────────────
# Printer assignment (alternating) logic for vendor dashboard/jobs
# ─────────────────────────────────────────────────────────────

# In-memory assignment state per vendor. Keys:
#  - 'small_last' for jobs with pages < 20
#  - 'large_last' for jobs with pages >= 20
_VENDOR_ASSIGNMENT_STATE = {}

def _get_candidate_printers_for_service(printer_config, service_type):
    """Return an ordered list of configured printer names (max 3) for a service."""
    if not printer_config:
        return []
    service_type = (service_type or '').lower().strip()
    mapping = {
        'digital_print': ['digital_printer_1', 'digital_printer_2', 'a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'project_binding': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'gloss_printing': ['gloss_printer_1', 'gloss_printer_2', 'a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'jumbo_printing': ['jumbo_printer_1', 'jumbo_printer_2', 'a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'regular_print': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'regular print': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'passport_print': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3'],
        'photo_print': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3'],
        'passport_photo': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3']
    }
    fields = mapping.get(service_type, ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'])
    candidates = []
    for field in fields:
        name = (printer_config.get(field) or '').strip()
        if name and name != 'NA' and name not in candidates:
            candidates.append(name)
    return candidates

def _assign_printer_alternating(vendor_id, printer_config, service_type, pages_str):
    """Assign a printer according to rules:
    - < 20 pages: strictly alternate between first two candidates
    - >= 20 pages: alternate between first two candidates, ensuring no consecutive large jobs on same printer
    - If only one candidate exists: always that one
    """
    try:
        pages = int(str(pages_str or '0'))
    except Exception:
        pages = 0

    candidates = _get_candidate_printers_for_service(printer_config, service_type)
    if not candidates:
        return "No Printer Configured"
    if len(candidates) == 1:
        return candidates[0]

    # Round-robin across all available candidates (2 or 3+)
    # Maintain separate indices for small and large jobs to avoid clustering
    state = _VENDOR_ASSIGNMENT_STATE.setdefault(str(vendor_id), {
        'small_idx': -1, 'large_idx': -1, 'small_last': None, 'large_last': None
    })

    # Backward-compatibility: if only *_last present, derive index once
    if state.get('small_idx', -1) == -1 and state.get('small_last') in candidates:
        try:
            state['small_idx'] = candidates.index(state['small_last'])
        except ValueError:
            state['small_idx'] = -1
    if state.get('large_idx', -1) == -1 and state.get('large_last') in candidates:
        try:
            state['large_idx'] = candidates.index(state['large_last'])
        except ValueError:
            state['large_idx'] = -1

    if pages < 20:
        state['small_idx'] = (state.get('small_idx', -1) + 1) % len(candidates)
        assigned = candidates[state['small_idx']]
        state['small_last'] = assigned
        return assigned
    else:
        state['large_idx'] = (state.get('large_idx', -1) + 1) % len(candidates)
        assigned = candidates[state['large_idx']]
        state['large_last'] = assigned
        return assigned

# ─────────────────────────────────────────────────────────────
# Printer assignment based on lowest count (load balancing)
# ─────────────────────────────────────────────────────────────
def _select_printer_by_lowest_count(printer_config, printer_counts, service_type):
    """Pick configured printer with the smallest count; ties by config order.
    Returns empty string if none configured.
    """
    candidates_fields = {
        'digital_print': ['digital_printer_1', 'digital_printer_2', 'a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'project_binding': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'gloss_printing': ['gloss_printer_1', 'gloss_printer_2', 'a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'jumbo_printing': ['jumbo_printer_1', 'jumbo_printer_2', 'a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'regular_print': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'regular print': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        'passport_print': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3'],
        'photo_print': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3'],
        'passport_photo': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3']
    }
    st = (service_type or '').lower().strip()
    fields = candidates_fields.get(st, ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'])

    chosen_name = ''
    chosen_count = None
    for field in fields:
        name = (printer_config or {}).get(field, '') or ''
        if not name or name == 'NA':
            continue
        count_key = f"{field}_count"
        try:
            cnt = int((printer_counts or {}).get(count_key, 0))
        except Exception:
            cnt = 0
        if chosen_count is None or cnt < chosen_count:
            chosen_count = cnt
            chosen_name = name
    return chosen_name


def assign_printer_and_increment_count(vendor_email, service_type):
    """Assign printer by lowest count and persist increment in pricing.json."""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        key = f"vendor_register_details/{sanitize_email(vendor_email)}/pricing.json"
        resp = s3.get_object(Bucket=settings.R2_BUCKET, Key=key)
        pricing_data = json.loads(resp['Body'].read().decode('utf-8')) if resp else {}

        printer_config = pricing_data.get('printer_configuration', {}) or {}
        printer_counts = pricing_data.get('printer_counts', {}) or {}

        selected_name = _select_printer_by_lowest_count(printer_config, printer_counts, service_type)
        if not selected_name:
            return ''

        # find field to bump
        st = (service_type or '').lower().strip()
        reverse_fields = {
            'digital_print': ['digital_printer_1', 'digital_printer_2', 'a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
            'project_binding': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
            'gloss_printing': ['gloss_printer_1', 'gloss_printer_2', 'a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
            'jumbo_printing': ['jumbo_printer_1', 'jumbo_printer_2', 'a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
            'regular_print': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
            'regular print': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
            'passport_print': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3'],
            'photo_print': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3'],
            'passport_photo': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3']
        }
        for field in reverse_fields.get(st, ['a4_printer_1', 'a4_printer_2', 'a4_printer_3']):
            if printer_config.get(field) == selected_name:
                count_key = f"{field}_count"
                try:
                    cur = int(pricing_data.setdefault('printer_counts', {}).get(count_key, 0))
                except Exception:
                    cur = 0
                pricing_data['printer_counts'][count_key] = cur + 1
                s3.put_object(Bucket=settings.R2_BUCKET,
                              Key=key,
                              Body=json.dumps(pricing_data, indent=4),
                              ContentType='application/json')
                break
        return selected_name
    except Exception as e:
        print(f"⚠️ assign_printer_and_increment_count error: {str(e)}")
        return ''
def store_vendor_pending_jobs_snapshot(request=None):
    """
    Store vendor pending jobs snapshot for the logged-in vendor only in D1 database.
    This function should be called every 2 minutes via a scheduled task.
    If request is provided, only processes the logged-in vendor from session.
    If request is None (scheduled task), skips execution (no session available).
    """
    try:
        # If called from scheduled task (no request), skip - we need session to identify vendor
        if request is None:
            print(f"⚠️ Scheduled snapshot skipped - requires active vendor session")
            return False
        
        # Get logged-in vendor from session
        vendor_email = request.session.get('vendor_email', '').strip().lower()
        vendor_id = request.session.get('vendor_id', '').strip()
        
        if not vendor_email and not vendor_id:
            print(f"⚠️ No vendor session found, skipping snapshot")
            return False
        
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            print(f"⚠️ Worker API not configured, skipping vendor pending jobs snapshot")
            return False
        
        # Process only the logged-in vendor
        vendors = [{'vendor_id': vendor_id, 'email': vendor_email}]
        print(f"📋 Processing snapshot for logged-in vendor: {vendor_email or vendor_id}")
        
        snapshot_timestamp = datetime.datetime.now().isoformat()
        success_count = 0
        error_count = 0
        
        # Process the logged-in vendor only
        for vendor in vendors:
            try:
                vendor_id = vendor.get('vendor_id', '').strip()
                vendor_email = vendor.get('email', '').strip().lower()
                
                if not vendor_id and not vendor_email:
                    continue
                
                # Get pending jobs for this vendor
                pending_jobs = get_vendor_jobs_from_d1(vendor_id=vendor_id, vendor_email=vendor_email, job_status='NO') or []
                
                # Store snapshot via worker API
                worker_endpoint = build_worker_endpoint('/store-vendor-pending-jobs-snapshot')
                if not worker_endpoint:
                    print(f"⚠️ Could not build worker endpoint")
                    continue
                
                payload = {
                    'vendor_id': vendor_id or '',
                    'vendor_email': vendor_email or '',
                    'pending_jobs': pending_jobs,
                    'snapshot_timestamp': snapshot_timestamp
                }
                
                snapshot_resp = requests.post(
                    worker_endpoint,
                    json=payload,
                    headers={
                        'Content-Type': 'application/json',
                        'x-api-key': api_key
                    },
                    timeout=10
                )
                
                if snapshot_resp.status_code == 200:
                    success_count += 1
                    print(f"✅ Stored {len(pending_jobs)} pending jobs snapshot for vendor {vendor_id or vendor_email}")
                else:
                    error_count += 1
                    error_text = snapshot_resp.text[:200] if snapshot_resp.text else 'Unknown error'
                    print(f"⚠️ Failed to store snapshot for vendor {vendor_id or vendor_email}: {snapshot_resp.status_code} - {error_text}")
                    
            except Exception as vendor_error:
                error_count += 1
                print(f"❌ Error processing vendor {vendor.get('vendor_id', 'unknown')}: {str(vendor_error)}")
                continue
        
        print(f"✅ Vendor pending jobs snapshot completed: {success_count} successful, {error_count} errors")
        return True
        
    except Exception as e:
        print(f"❌ Error storing vendor pending jobs snapshot: {str(e)}")
        return False


@csrf_exempt
def store_vendor_pending_jobs_snapshot_endpoint(request):
    """
    Endpoint to trigger vendor pending jobs snapshot storage for logged-in vendor only.
    This can be called when vendor accesses their dashboard.
    """
    if request.method == 'POST':
        try:
            result = store_vendor_pending_jobs_snapshot(request)
            if result:
                return JsonResponse({'success': True, 'message': 'Vendor pending jobs snapshot stored successfully'}, status=200)
            else:
                return JsonResponse({'success': False, 'error': 'Failed to store vendor pending jobs snapshot'}, status=500)
        except Exception as e:
            print(f"❌ Error in store_vendor_pending_jobs_snapshot_endpoint: {str(e)}")
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    else:
        return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)


def get_vendor_jobs_from_d1(vendor_id=None, vendor_email=None, job_status='NO'):
    """
    Fetch vendor print jobs from the Vendor_print_jobs D1 table via Worker API.
    Supports filtering by vendor_id, vendor_email and job completion status.
    Only uses R2 to generate presigned URLs for the actual files.
    """
    try:
        vendor_id = str(vendor_id).strip() if vendor_id else ''
        vendor_email = (vendor_email or '').strip().lower()
        status_filter = (job_status or '').strip().upper()

        if not vendor_id and not vendor_email:
            print("⚠️ Cannot fetch vendor jobs without vendor_id or vendor_email")
            return []

        if not status_filter:
            status_filter = 'NO'

        worker_payload = {
            'job_completed': status_filter,
        }
        if vendor_id:
            worker_payload['vendor_id'] = vendor_id
        if vendor_email:
            worker_payload['vendor_email'] = vendor_email

        endpoint, resp = post_to_worker('/get-vendor-print-jobs', worker_payload)
        if resp.status_code != 200:
            print(f"⚠️ Worker API request failed ({resp.status_code}) via {endpoint}: {resp.text[:300]}")
            return []

        data = resp.json()
        if not data.get('success'):
            print(f"⚠️ Worker API returned error: {data.get('error', 'Unknown error')}")
            return []

        jobs = data.get('data') or []
        print(f"🔍 get_vendor_jobs_from_d1 - Worker returned {len(jobs)} jobs for vendor_id={vendor_id}, vendor_email={vendor_email}, job_status={status_filter}")
        if jobs:
            print(f"📋 Sample job service_types: {[j.get('service_type', 'N/A') for j in jobs[:5]]}")
        if not jobs:
            return []

        # Process jobs to generate presigned URLs from R2 paths stored in database
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        filtered_jobs = []
        for job in jobs:
            storage_folder = job.get('storage_folder') or 'vendor_print_jobs'
            job['storage_folder'] = storage_folder
            job.setdefault('vendor_id', vendor_id)
            job.setdefault('vendor', job.get('vendor_id', vendor_id))
            job.setdefault('user', job.get('user_email', ''))

            # Get R2 path from database (r2_path field) - ONLY use R2 for file URL
            r2_path = _normalize_r2_key(job.get('r2_path'))
            if not r2_path and job.get('filename'):
                r2_path = f"{storage_folder.rstrip('/')}/{job.get('vendor_id', vendor_id)}/{job.get('filename')}"
                job['r2_path'] = r2_path
            elif r2_path:
                job['r2_path'] = r2_path

            if r2_path:
                try:
                    # Generate presigned URL from R2 path - ONLY for file rendering
                    url = s3.generate_presigned_url(
                        ClientMethod='get_object',
                        Params={
                            'Bucket': settings.R2_BUCKET,
                            'Key': r2_path
                        },
                        ExpiresIn=3600
                    )
                    job['url'] = url
                    job['preview_url'] = url
                    job['download_url'] = url
                except Exception as e:
                    print(f"⚠️ Error generating presigned URL for {r2_path}: {e}")
                    job['url'] = ''
                    job['preview_url'] = ''
                    job['download_url'] = ''
            else:
                job['url'] = job['preview_url'] = job['download_url'] = ''
            
            job.setdefault('preview_url', job.get('url', ''))
            job.setdefault('download_url', job.get('url', ''))
            
            # Map ALL metadata from D1 database (not from R2)
            # All fields come from D1, only file URL comes from R2
            # This ensures we're using D1 as the single source of truth for all metadata
            job_completed_source = job.get('job_completed')
            if not job_completed_source:
                job_completed_source = job.get('job_completed_status')
            if job_completed_source is None or str(job_completed_source).strip() == '':
                job_completed_source = 'NO'
            job_completed_value = str(job_completed_source).strip().upper()
            job['job_completed'] = job_completed_value
            
            # All metadata fields explicitly from D1 database
            job['service_type'] = job.get('service_type', '') or ''
            job['pages'] = str(job.get('page_count', job.get('pages', '0')) or '0')
            job['rendered_status'] = job.get('rendered_status', 'NO') or 'NO'
            job['vendor_status'] = job.get('vendor_status', 'not sended') or 'not sended'
            job['token'] = job.get('token', '') or ''
            job['job_id'] = job.get('job_id', '') or ''
            copies_value = job.get('copies', job.get('num_copies'))
            job['copies'] = str(copies_value) if copies_value not in [None, ''] else '1'
            job['color'] = job.get('color', '') or ''
            job['orientation'] = job.get('orientation', '') or ''
            job['pageSize'] = job.get('pageSize', '') or ''
            job['timestamp'] = job.get('timestamp', '') or ''
            job['uploaded_at'] = job.get('timestamp', '') or ''  # Alias for compatibility
            job['completion_time'] = job.get('completion_time', '') or ''
            job['total_price'] = job.get('total_price', 0) or 0
            job['final_amount'] = job.get('final_amount', 0) or 0
            job['user_email'] = job.get('user_email', '') or ''
            job['vendor_email'] = job.get('vendor_email', '') or ''
            job['status'] = job.get('status', 'pending') or 'pending'
            job['feedback'] = job.get('feedback', '') or ''
            job['quality'] = job.get('quality', '') or ''
            job['thickness'] = job.get('thickness', '') or ''
            job['service_name'] = job.get('service_name', '') or ''
            job['filename'] = job.get('filename', '') or ''  # Ensure filename is from D1

            pricing_details_raw = job.get('pricing_details')
            parsed_pricing = None
            if isinstance(pricing_details_raw, str) and pricing_details_raw.strip():
                try:
                    parsed_pricing = json.loads(pricing_details_raw)
                except json.JSONDecodeError:
                    parsed_pricing = None
            elif isinstance(pricing_details_raw, dict):
                parsed_pricing = pricing_details_raw
            if parsed_pricing is not None:
                job['pricing_details'] = parsed_pricing

            # Create metadata structure from D1 fields ONLY (not from R2)
            # All metadata comes from D1 database, R2 is ONLY used for file URL generation
            job['metadata'] = {
                'status': job.get('status', 'pending') or 'pending',
                'job_completed': job_completed_value,
                'copies': job.get('copies', '1') or '1',
                'color': job.get('color', '') or '',
                'orientation': job.get('orientation', '') or '',
                'page_size': job.get('pageSize', '') or '',
                'pages': str(job.get('page_count', job.get('pages', '0')) or '0'),
                'timestamp': job.get('timestamp', '') or '',
                'vendor': job.get('vendor', vendor_id) or vendor_id,
                'user': job.get('user_email', '') or '',
                'service_type': job.get('service_type', '') or '',
                'job_id': job.get('job_id', '') or '',
                'token': job.get('token', '') or '',
                'vendor_id': job.get('vendor_id', vendor_id) or vendor_id,
                'rendered_status': job.get('rendered_status', 'NO') or 'NO',
                'vendor_status': job.get('vendor_status', 'not sended') or 'not sended',
                'total_price': job.get('total_price', 0) or 0,
                'final_amount': job.get('final_amount', 0) or 0,
                'feedback': job.get('feedback', '') or '',
                'quality': job.get('quality', '') or '',
                'thickness': job.get('thickness', '') or '',
                'service_name': job.get('service_name', '') or '',
                'completion_time': job.get('completion_time', '') or '',
                'platform_profit': job.get('platform_profit', 0) or 0,
                'pageRange': job.get('pageRange', '') or '',
                'bwPageRangeValue': job.get('bwPageRangeValue', '') or '',
                'colorPageRangeValue': job.get('colorPageRangeValue', '') or '',
            }
            
            # Add pricing details from D1 if available
            job['metadata']['pricing_details'] = parsed_pricing
            
            # Enforce local status filtering as a safety net
            # Only filter if status is explicitly different (not just missing/empty)
            normalized_status_raw = job.get('job_completed') or job.get('job_completed_status') or 'NO'
            normalized_status = str(normalized_status_raw).strip().upper()
            
            # If status is empty or None, default to 'NO' and include it if we're looking for 'NO'
            if not normalized_status_raw or normalized_status_raw == '':
                normalized_status = 'NO'
            
            # Include job if status matches, or if status is empty/None and we're looking for 'NO'
            if normalized_status == status_filter or (not normalized_status_raw and status_filter == 'NO'):
                filtered_jobs.append(job)
            else:
                print(f"⚠️ Job filtered out: {job.get('filename', 'unknown')} (token: {job.get('token', 'N/A')}) - normalized_status={normalized_status}, expected={status_filter}, service_type={job.get('service_type', 'N/A')}")

        print(f"✅ get_vendor_jobs_from_d1 - Returning {len(filtered_jobs)} filtered jobs")
        return filtered_jobs
            
    except Exception as e:
        print(f"⚠️ Error fetching vendor jobs from D1: {str(e)}")
        return []


def get_print_requests(request):
    try:
        # Get vendor details from session to filter jobs
        vendor_email = request.session.get('vendor_email')
        vendor_id = request.session.get('vendor_id')  # Get vendor_id directly from session
        vendor_status = request.session.get('vendor_status', 'pending').strip().lower()
        
        print(f"🔍 get_print_requests - Session data - Email: {vendor_email}, Vendor ID: {vendor_id}, Status: {vendor_status}")
        
        if not vendor_id and vendor_email:
            # Try to get vendor_id from vendor details
            vendor_details = get_vendor_details_by_email(vendor_email)
            if vendor_details:
                if vendor_details.get('data_source') == 'd1' and vendor_details.get('vendor_id'):
                    vendor_id = vendor_details.get('vendor_id')
                    request.session['vendor_id'] = vendor_id
                    print(f"🔍 Retrieved vendor ID from D1 details: {vendor_id}")
                else:
                    print("⚠️ Vendor ID unavailable from D1; keeping session vendor_id empty")
        
        if not vendor_id:
            print("❌ No vendor ID found in session - returning empty job list")
            return JsonResponse({"print_requests": []}, status=200)
        
        # Check vendor status - only return jobs if status is 'verified'
        if vendor_status != 'verified':
            print(f"⚠️ Vendor status is '{vendor_status}', not verified. Returning empty job list.")
            return JsonResponse({"print_requests": []}, status=200)
        
        # Fetch pending jobs from D1 database using vendor email fallback
        pending_files = get_vendor_jobs_from_d1(vendor_id=vendor_id, vendor_email=vendor_email, job_status='NO') or []
        
        # Also fetch completed jobs to check if any were completed today
        completed_files = get_vendor_jobs_from_d1(vendor_id=vendor_id, vendor_email=vendor_email, job_status='YES') or []
        
        # Filter completed jobs: only include if completed today
        today = datetime.datetime.now().date()
        today_completed_files = []
        for job in completed_files:
            completion_time = job.get('completion_time') or job.get('timestamp')
            if completion_time:
                try:
                    # Try parsing as ISO string
                    if isinstance(completion_time, str):
                        if 'T' in completion_time or ' ' in completion_time:
                            completion_date = datetime.datetime.fromisoformat(completion_time.replace('Z', '+00:00')).date()
                        else:
                            # Try parsing as timestamp
                            completion_date = datetime.datetime.fromtimestamp(float(completion_time)).date()
                    else:
                        completion_date = datetime.datetime.fromtimestamp(float(completion_time)).date()
                    
                    if completion_date == today:
                        today_completed_files.append(job)
                except Exception as e:
                    print(f"⚠️ Error parsing completion_time for job {job.get('filename', 'unknown')}: {e}")
        
        # Combine pending and today's completed jobs
        files = pending_files + today_completed_files
        
        if not files:
            print(f"ℹ️ No jobs returned from D1 for vendor {vendor_id}")
        
        # Ensure rendered_status present by default
        for job in files:
            if 'rendered_status' not in job:
                job['rendered_status'] = job.get('metadata', {}).get('rendered_status', 'NO')
        
        # Set default printer assignment (no longer using printer selection)
        for job in files:
            job['assigned_printer'] = 'Not Assigned'
            job['printer_name'] = 'Not Assigned'
        
        # Define service types for categorization (same as vendordashboard)
        manual_services = [
            'digital_print', 'project_binding', 'gloss_printing', 'jumbo_printing', 'golden_embossing'
        ]
        regular_services = [
            'regular_print', 'passport_print', 'photo_print', 'regular print', 'passport_photo'
        ]
        
        # Categorize jobs the same way as vendordashboard
        categorized_jobs = {
            'manual': [],
            'requests': [],
            'completed': []
        }
        
        for job in files:
            # Get job status - prioritize job_completed over job_completed_status
            job_completed_source = job.get('job_completed')
            if not job_completed_source:
                job_completed_source = job.get('job_completed_status')
            if job_completed_source is None or str(job_completed_source).strip() == '':
                job_completed_source = 'NO'
            job_completed = str(job_completed_source).strip().upper()
            
            # Normalize job_completed in the job object to ensure frontend receives consistent format
            job['job_completed'] = job_completed
            
            # Ensure service_type is preserved in job object for frontend categorization
            # Frontend needs the original service_type to categorize into passport, golden, etc. sections
            original_service_type = job.get('service_type', '').strip()
            service_type = (original_service_type or '').strip().lower()
            job['service_type'] = original_service_type or service_type  # Preserve original for frontend
            vendor_status = job.get('vendor_status', 'not sended').lower()
            
            # Skip cancelled jobs - they should not be displayed
            if job_completed == 'CANCELLED':
                continue
            
            if job_completed == 'YES':
                # Completed jobs - regardless of service type
                categorized_jobs['completed'].append(job)
            elif job_completed == 'NO':
                # Pending jobs - categorize by service type
                if service_type in manual_services:
                    # Manual print jobs (digital_print, project_binding, etc.)
                    categorized_jobs['manual'].append(job)
                elif service_type in regular_services or service_type == '' or service_type == 'regular print':
                    # Regular print jobs - show both pending and accepted
                    categorized_jobs['requests'].append(job)
                else:
                    # Default to print requests for unknown service types
                    categorized_jobs['requests'].append(job)
        
        print(f"📊 AJAX Job categorization - Manual: {len(categorized_jobs['manual'])}, Requests: {len(categorized_jobs['requests'])}, Completed: {len(categorized_jobs['completed'])}")
        
        # Return pending jobs + today's completed jobs in print_requests
        pending_jobs = categorized_jobs['manual'] + categorized_jobs['requests']
        # Include today's completed jobs in print_requests
        all_print_requests = pending_jobs + categorized_jobs['completed']
        
        return JsonResponse({
            "print_requests": all_print_requests,  # Pending jobs + today's completed jobs
            "categorized": categorized_jobs,
            "counts": {
                "manual": len(categorized_jobs['manual']),
                "requests": len(categorized_jobs['requests']),
                "completed": len(categorized_jobs['completed']),
                "total": len(files)
            }
        }, status=200)
    except Exception as e:
        print(f"Error in get_print_requests: {str(e)}")
        return JsonResponse({"error": str(e), "print_requests": []}, status=500)

# ─────────────────────────────────────────────────────────────
# AUTO PRINT ENDPOINT FOR WEBSOCKET INTEGRATION
# ─────────────────────────────────────────────────────────────


@csrf_exempt
def auto_print_documents(request):
    """
    Get pending print jobs and send them to connected vendor clients via WebSocket
    (Now only for admin or dashboard, not for vendor client polling)
    """
    if request.method == 'POST':
        try:
            # Get all files with job_completed = 'NO' (from vendor folders only)
            pending_jobs = get_pending_print_jobs()

            if not pending_jobs:
                return JsonResponse({
                    'success': True, 
                    'message': 'No pending print jobs found',
                    'jobs_sent': 0
                })

            print(f"🖨️  Auto-print triggered: Found {len(pending_jobs)} pending jobs")
            for job in pending_jobs:
                print(f"   - {job['filename']} (status: {job['metadata']['status']}, completed: {job['metadata']['job_completed']})")

            return JsonResponse({
                'success': True,
                'message': f'Found {len(pending_jobs)} pending print jobs ready for processing',
                'jobs_sent': len(pending_jobs),
                'jobs': pending_jobs
            })

        except Exception as e:
            print(f"Error in auto_print_documents: {str(e)}")
            return JsonResponse({'success': False, 'error': str(e)}, status=500)

    return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)

from django.views.decorators.csrf import csrf_exempt
from django.shortcuts import render, redirect
from django.http import JsonResponse
from django.contrib import messages
from django.urls import reverse

def vendor_documents(request):
    return render(request, 'vendor_documents.html', {
        'vendor_email': request.session.get('vendor_email', ''),
    })

def _sanitize_email_for_r2(email):
    import re as _re
    if not email:
        return ''
    return _re.sub(r'[^a-zA-Z0-9_]', '', email.lower().replace('@', '_at_').replace('.', '_dot_'))

@csrf_exempt  # Form includes CSRF; keep exempt to tolerate AJAX/form variations
def upload_vendor_documents(request):
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)

    # Inline helpers and imports to avoid global import collisions
    import boto3, json, datetime, uuid, re
    from django.conf import settings as _settings

    def _is_allowed_file(file_obj):
        try:
            if not file_obj:
                return False, 'Missing file'
            allowed_mime = {'application/pdf', 'image/jpeg', 'image/png'}
            name_ok = any(file_obj.name.lower().endswith(ext) for ext in ['.pdf', '.jpg', '.jpeg', '.png'])
            type_ok = (getattr(file_obj, 'content_type', '') in allowed_mime) or name_ok
            if not type_ok:
                return False, 'Invalid file type. Only PDF, JPG, PNG allowed.'
            if getattr(file_obj, 'size', 0) > 5 * 1024 * 1024:
                return False, 'File too large. Maximum 5 MB allowed.'
            return True, None
        except Exception:
            return False, 'File validation error'

    def _generate_unique_filename(base_label, original_name):
        safe_label = re.sub(r'[^a-zA-Z0-9_]+', '_', (base_label or 'document').lower()).strip('_')
        ext = ''
        if '.' in original_name:
            ext = '.' + original_name.split('.')[-1].lower()
        timestamp = datetime.datetime.now().strftime('%Y-%m-%d_%H%M%S')
        short_id = uuid.uuid4().hex[:8]
        return f"{safe_label}_{timestamp}_{short_id}{ext}"

    try:
        vendor_email = (request.POST.get('email') or request.session.get('vendor_email') or '').strip()
        if not vendor_email:
            err = 'Email is required'
            messages.error(request, err)
            try:
                return redirect('vendor_documents')
            except Exception:
                return JsonResponse({'success': False, 'error': err}, status=400)

        # Validate required bank fields
        beneficiary_name = (request.POST.get('beneficiaryName') or '').strip()
        account_number = (request.POST.get('accountNumber') or '').strip()
        ifsc_code = (request.POST.get('ifscCode') or '').strip()
        if not (beneficiary_name and account_number and ifsc_code):
            err = 'Beneficiary name, account number and IFSC code are required'
            messages.error(request, err)
            try:
                return redirect('vendor_documents')
            except Exception:
                return JsonResponse({'success': False, 'error': err}, status=400)

        # Files (AADHAR, PAN, Cheque, GST optional)
        candidate_files = [
            ('aadharFile', 'aadhaar'),
            ('qrFile', 'qr_code'),
            ('panFile', 'pan_card'),
            ('chequeFile', 'cheque'),
            ('gstFile', 'gst_certificate'),
        ]

        files_to_upload = []
        uploaded_files = []
        for field_name, label in candidate_files:
            f = request.FILES.get(field_name)
            if not f:
                continue
            ok, reason = _is_allowed_file(f)
            if not ok:
                err = f"{label.replace('_', ' ').title()}: {reason}"
                messages.error(request, err)
                try:
                    return redirect('vendor_documents')
                except Exception:
                    return JsonResponse({'success': False, 'error': err}, status=400)
            unique_name = _generate_unique_filename(label, f.name)
            files_to_upload.append((unique_name, f))
            uploaded_files.append(unique_name)

        if not uploaded_files:
            err = 'Please upload at least one document (Aadhaar/QR/PAN/Cheque/GST).'
            messages.error(request, err)
            try:
                return redirect('vendor_documents')
            except Exception:
                return JsonResponse({'success': False, 'error': err}, status=400)

        sanitized = _sanitize_email_for_r2(vendor_email)
        base_prefix = f"vendor_register_details/{sanitized}/vendor_documents/"

        s3 = boto3.client(
            's3',
            aws_access_key_id=_settings.R2_ACCESS_KEY,
            aws_secret_access_key=_settings.R2_SECRET_KEY,
            endpoint_url=_settings.R2_ENDPOINT.rstrip('/') if getattr(_settings, 'R2_ENDPOINT', None) else None,
            region_name='auto'
        )

        # Upload files
        for unique_name, f in files_to_upload:
            key = base_prefix + unique_name
            s3.put_object(
                Bucket=_settings.R2_BUCKET,
                Key=key,
                Body=f.read(),
                ContentType=getattr(f, 'content_type', 'application/octet-stream') or 'application/octet-stream',
                Metadata={
                    'uploaded_at': datetime.datetime.now().isoformat(),
                    'original_filename': f.name,
                }
            )

        # Build and store registration_details.json (append/update minimal fields)
        submission_time = datetime.datetime.now().isoformat()
        metadata = {
            'submission_time': submission_time,
            'vendor_email': vendor_email,
            'vendor_name': (request.POST.get('vendor_name') or ''),
            'phone_number': (request.POST.get('phone_number') or ''),
            'bank_details': {
                'beneficiary_name': beneficiary_name,
                'beneficiary_bank': (request.POST.get('beneficiaryBank') or ''),
                'account_number': account_number,
                'ifsc_code': ifsc_code,
                'branch_code': (request.POST.get('branchCode') or ''),
                'bank_address': (request.POST.get('bankAddress') or ''),
                'city': (request.POST.get('city') or ''),
                'postal_code': (request.POST.get('postalCode') or ''),
                'country': (request.POST.get('country') or ''),
            },
            'uploaded_files': uploaded_files,
        }

        metadata_key = base_prefix + 'registration_details.json'
        s3.put_object(
            Bucket=_settings.R2_BUCKET,
            Key=metadata_key,
            Body=json.dumps(metadata, ensure_ascii=False, indent=2),
            ContentType='application/json'
        )

        try:
            url = reverse('vendor_documents') + '?success=1'
        except Exception:
            url = '/vendor_documents?success=1'
        return redirect(url)

    except Exception as e:
        messages.error(request, 'Failed to upload documents. Please try again.')
        try:
            return redirect('vendor_documents')
        except Exception:
            return JsonResponse({'success': False, 'error': 'Upload failed'}, status=500)


@csrf_exempt
def get_vendor_print_jobs(request):
    """
    Fetch print jobs for a specific vendor from Vendor_print_jobs table (D1) only.
    Uses Worker API for job list; R2 only for generating download URLs.
    """
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_id = data.get('vendor_id')
            if not vendor_id:
                return JsonResponse({'success': False, 'error': 'Missing vendor_id'})

            vendor_id = str(vendor_id).strip()

            # Treat successful polling as connection activity to prevent UI flapping
            try:
                if get_vendor_auth_status(vendor_id):
                    update_vendor_connection_status(vendor_id, 'connected')
            except Exception:
                pass

            try:
                vendor_email = get_vendor_email_by_vendor_id(vendor_id)
            except Exception:
                vendor_email = None
            printer_config = {}
            if vendor_email:
                printer_config = get_vendor_printer_configuration(vendor_email)

            # Fetch jobs from Vendor_print_jobs table (D1) only via Worker API
            d1_jobs = get_vendor_jobs_from_d1(
                vendor_id=vendor_id,
                vendor_email=vendor_email,
                job_status='NO'
            )

            # Filter to accepted jobs only (same behavior as previous R2-based filter)
            allowed_services = {
                'regular print', 'regular_print',
                'passport photo', 'passport_photo', 'passport print', 'passport_print',
                'photo print', 'photo_print'
            }
            jobs = []
            for job in d1_jobs:
                service_type = (job.get('service_type') or '').strip().lower()
                vendor_status = (job.get('vendor_status') or 'not sended').strip().lower()
                if service_type not in allowed_services or vendor_status != 'accepted':
                    continue
                pages_value = job.get('metadata', {}).get('pages', job.get('pages', '0'))
                service_value = job.get('service_type', 'regular print')
                assigned = _assign_printer_alternating(vendor_id, printer_config, service_value, pages_value)
                meta = job.get('metadata', {})
                meta['vendor_status'] = 'sended'
                meta['assigned_printer'] = assigned
                job_info = {
                    'filename': job.get('filename', ''),
                    'download_url': job.get('download_url', ''),
                    'r2_path': job.get('r2_path', ''),
                    'metadata': meta,
                    'assigned_printer': assigned
                }
                jobs.append(job_info)

            print(f"📋 get_vendor_print_jobs (D1 only): {len(jobs)} jobs for vendor {vendor_id}")
            return JsonResponse({'success': True, 'jobs': jobs})

        except Exception as e:
            print(f"❌ Error fetching vendor jobs: {str(e)}")
            traceback.print_exc()
            return JsonResponse({'success': False, 'error': str(e)})
    return JsonResponse({'success': False, 'error': 'Invalid method'})


def get_vendor_specific_print_jobs(vendor_id):
    """Get pending print jobs from vendor-specific folder in R2 storage"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        pending_jobs = []
        vendor_folder_path = f'vendor_print_jobs/{vendor_id}'
        manual_vendor_folder_path = f'vendor_manual_print_jobs/{vendor_id}'

        # Check both vendor print jobs and vendor manual print jobs folders for documents
        try:
            vendor_objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=vendor_folder_path)
            manual_objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=manual_vendor_folder_path)

            # Combine both object lists
            all_objects = []
            if vendor_objects.get("Contents"):
                all_objects.extend(vendor_objects.get("Contents", []))
            if manual_objects.get("Contents"):
                all_objects.extend(manual_objects.get("Contents", []))

            for obj in all_objects:
                key = obj["Key"]
                filename = key.split("/")[-1]

                # Skip folder itself and non-document files
                if not filename or filename.endswith('.json'):
                    continue

                try:
                    # Get object metadata
                    head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                    metadata = head_response.get('Metadata', {})

                    # Check if job is pending (job_completed = 'NO')
                    job_completed = metadata.get('job_completed', 'NO').upper()
                    status = metadata.get('status', 'pending').lower()

                    if job_completed == 'NO' or status == 'pending':
                        # Generate download URL
                        download_url = s3.generate_presigned_url(
                            ClientMethod='get_object',
                            Params={
                                'Bucket': settings.R2_BUCKET,
                                'Key': key
                            },
                            ExpiresIn=3600
                        )

                        # Build job info
                        job_info = {
                            'filename': filename,
                            'download_url': download_url,
                            'r2_path': key,
                            'user_email': metadata.get('user', ''),
                            'metadata': {
                                'status': 'no',  # Set to 'no' for pending jobs
                                'job_completed': job_completed,
                                'copies': metadata.get('copies', '1'),
                                'color': metadata.get('color', 'Black and White'),
                                'orientation': metadata.get('orientation', 'portrait'),
                                'page_size': metadata.get('pagesize', 'A4'),
                                'pages': metadata.get('pages', '1'),
                                'timestamp': metadata.get('timestamp', obj["LastModified"].isoformat()),
                                'vendor': metadata.get('vendor', vendor_id),
                                'user': metadata.get('user', 'Unknown'),
                                'service_type': metadata.get('service_type', ''),
                                'job_id': metadata.get('job_id', ''),
                                'token': metadata.get('token', ''),
                                'feedback': metadata.get('feedback', ''),
                                'quality': metadata.get('quality', ''),
                                'thickness': metadata.get('thickness', ''),
                                'service_name': metadata.get('service_name', '')
                            }
                        }

                        pending_jobs.append(job_info)
                        print(f"✅ Found pending job for vendor {vendor_id}: {filename}")
                        
                        # Update print_round from 1 to 2 when sending to vendor client
                        current_print_round = metadata.get('print_round', '0')
                        if current_print_round == '1':
                            try:
                                # Update metadata to set print_round to 2
                                updated_metadata = metadata.copy()
                                updated_metadata['print_round'] = '2'
                                
                                # Copy object with updated metadata
                                copy_source = {'Bucket': settings.R2_BUCKET, 'Key': key}
                                s3.copy_object(
                                    CopySource=copy_source,
                                    Bucket=settings.R2_BUCKET,
                                    Key=key,
                                    Metadata=updated_metadata,
                                    MetadataDirective='REPLACE'
                                )
                                print(f"🔄 Updated print_round from 1 to 2 for {filename}")
                            except Exception as update_error:
                                print(f"⚠️ Failed to update print_round for {filename}: {str(update_error)}")

                except Exception as e:
                    print(f"Error processing vendor file {key}: {str(e)}")
                    continue

        except Exception as e:
            print(f"Error accessing vendor folder {vendor_folder_path}: {str(e)}")

        print(f"📋 Total pending jobs found for vendor {vendor_id}: {len(pending_jobs)}")
        return pending_jobs

    except Exception as e:
        print(f"Error getting vendor-specific jobs: {e}")
        return []
def update_job_status_comprehensive(filename, job_completed_status, vendor_id, completion_time):
    """
    Update job status in all relevant folders: user folder, vendor_print_jobs, vendor_manual_print_jobs
    Returns (success, user_email)
    """
    s3 = boto3.client('s3',
                      aws_access_key_id=settings.R2_ACCESS_KEY,
                      aws_secret_access_key=settings.R2_SECRET_KEY,
                      endpoint_url=settings.R2_ENDPOINT,
                      region_name='auto')
    
    try:
        user_email = None
        updated_folders = []
        
        # Define all possible folder paths to check
        folder_paths = [
            f'users/',
            f'vendor_print_jobs/{vendor_id}/',
            f'vendor_manual_print_jobs/{vendor_id}/'
        ]
        
        # Find the file in all folders and update metadata
        for folder_path in folder_paths:
            try:
                # List objects in the folder
                response = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=folder_path)
                
                if 'Contents' in response:
                    for obj in response['Contents']:
                        if obj['Key'].endswith(filename):
                            # Get current metadata
                            head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=obj['Key'])
                            current_metadata = head_response.get('Metadata', {})
                            
                            # Extract user email from user folder
                            if folder_path.startswith('users/') and not user_email:
                                user_email = obj['Key'].split('/')[1]  # Extract email from path
                            
                            # Update metadata
                            updated_metadata = current_metadata.copy()
                            updated_metadata['job_completed'] = job_completed_status
                            if completion_time:
                                updated_metadata['completion_time'] = str(completion_time)
                            
                            # Ensure pricing_details contains platform_profit for user modals
                            try:
                                pricing_details_raw = current_metadata.get('pricing_details')
                                if pricing_details_raw:
                                    # Parse pricing_details from metadata (stringified JSON)
                                    pricing_details_obj = json.loads(pricing_details_raw) if isinstance(pricing_details_raw, str) else pricing_details_raw
                                    if isinstance(pricing_details_obj, dict):
                                        breakdown_obj = pricing_details_obj.get('pricing_breakdown') or {}
                                        # Try to compute platform profit if missing
                                        if 'platform_profit' not in pricing_details_obj:
                                            # Prefer explicit vendor_price from breakdown; fallback to total/total_price
                                            vendor_price_num = None
                                            try:
                                                vendor_price_num = float(breakdown_obj.get('vendor_price', breakdown_obj.get('total_price', pricing_details_obj.get('vendor_price', pricing_details_obj.get('total_price', 0)))))
                                            except Exception:
                                                vendor_price_num = 0.0
                                            # Total charged to user: prefer pricing_details.total_price/total, fallback to metadata total_price
                                            total_user_amount = None
                                            try:
                                                total_user_amount = float(pricing_details_obj.get('total_price', pricing_details_obj.get('total', current_metadata.get('total_price', 0))))
                                            except Exception:
                                                total_user_amount = 0.0
                                            platform_profit_val = max(0.0, (total_user_amount or 0.0) - (vendor_price_num or 0.0))
                                            pricing_details_obj['platform_profit'] = round(platform_profit_val, 2)
                                        # Re-store back into metadata as string
                                        updated_metadata['pricing_details'] = json.dumps(pricing_details_obj)
                            except Exception as _e:
                                # Non-fatal: skip profit enrichment if parsing fails
                                pass
                            
                            # Copy object with updated metadata
                            copy_source = {
                                'Bucket': settings.R2_BUCKET,
                                'Key': obj['Key']
                            }
                            
                            s3.copy_object(
                                CopySource=copy_source,
                                Bucket=settings.R2_BUCKET,
                                Key=obj['Key'],
                                Metadata=updated_metadata,
                                MetadataDirective='REPLACE'
                            )
                            
                            updated_folders.append(obj['Key'])
                            print(f"✅ Updated metadata in {obj['Key']}")
                            
            except Exception as e:
                print(f"⚠️ Error checking folder {folder_path}: {str(e)}")
                continue
        
        if updated_folders:
            print(f"✅ Successfully updated {len(updated_folders)} file(s) for {filename}")
            return True, user_email
        else:
            print(f"❌ No files found to update for {filename}")
            return False, None
            
    except Exception as e:
        print(f"❌ Error in comprehensive job status update: {str(e)}")
        return False, None




@csrf_exempt
def update_job_status(request):
    """
    Update job completion status when vendor client completes printing
    """
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            filename = data.get('filename')
            status = data.get('status', 'completed')
            vendor_id = data.get('vendor_id')
            completion_time = data.get('completion_time')

            if not filename:
                return JsonResponse({'success': False, 'error': 'Filename required'})

            # If vendor_id not provided in request, try to get it from session
            if not vendor_id:
                vendor_email = request.session.get('vendor_email')
                if vendor_email:
                    # Get vendor_id from vendor details
                    vendor_details = get_vendor_details_by_email(vendor_email)
                    if vendor_details:
                        if vendor_details.get('data_source') == 'd1' and vendor_details.get('vendor_id'):
                            vendor_id = vendor_details.get('vendor_id')
                        else:
                            # Legacy fallback for vendors that still exist only in R2
                            try:
                                s3 = boto3.client('s3',
                                                aws_access_key_id=settings.R2_ACCESS_KEY,
                                                aws_secret_access_key=settings.R2_SECRET_KEY,
                                                endpoint_url=settings.R2_ENDPOINT,
                                                region_name='auto')

                                reg_key = f'vendor_register_details/{sanitize_email(vendor_email)}/registration_details.json'
                                response = s3.get_object(Bucket=settings.R2_BUCKET, Key=reg_key)
                                vendor_data = json.loads(response['Body'].read().decode('utf-8'))
                                vendor_id = vendor_data.get('vendor_id', 'vendor1')
                                print(f"⚠️ Using legacy R2 vendor_id {vendor_id} for session vendor {vendor_email}")
                            except:
                                vendor_id = 'vendor1'
                else:
                    vendor_id = 'vendor1'

            status_lower = str(status).lower()

            # When a job is accepted or retried from the dashboard, copy it into
            # vendor_print_jobs/<vendor_id>/ so the vendor client (which polls this
            # folder) can immediately see and process it.
            if status_lower in ['accepted', 'retry'] and vendor_id:
                try:
                    s3 = boto3.client('s3',
                                      aws_access_key_id=settings.R2_ACCESS_KEY,
                                      aws_secret_access_key=settings.R2_SECRET_KEY,
                                      endpoint_url=settings.R2_ENDPOINT,
                                      region_name='auto')

                    # Search common locations for the file by filename
                    search_prefixes = ['users/', f'vendor_manual_print_jobs/{vendor_id}/', f'vendor_print_jobs/{vendor_id}/']
                    source_key = None
                    src_metadata = {}
                    for pref in search_prefixes:
                        resp = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=pref)
                        for obj in resp.get('Contents', []) if resp.get('Contents') else []:
                            if obj['Key'].endswith(filename):
                                source_key = obj['Key']
                                try:
                                    head = s3.head_object(Bucket=settings.R2_BUCKET, Key=source_key)
                                    src_metadata = head.get('Metadata', {})
                                except Exception:
                                    src_metadata = {}
                                break
                        if source_key:
                            break

                    if source_key is None:
                        return JsonResponse({'success': False, 'error': 'Source file not found'}, status=404)

                    # Destination in vendor_print_jobs
                    dest_key = f'vendor_print_jobs/{vendor_id}/{filename}'
                    updated_metadata = src_metadata.copy()
                    updated_metadata['status'] = 'accepted'
                    updated_metadata['job_completed'] = 'NO'

                    s3.copy_object(
                        CopySource={'Bucket': settings.R2_BUCKET, 'Key': source_key},
                        Bucket=settings.R2_BUCKET,
                        Key=dest_key,
                        Metadata=updated_metadata,
                        MetadataDirective='REPLACE'
                    )

                    return JsonResponse({'success': True, 'message': 'Job accepted and queued for vendor client'})
                except Exception as e:
                    return JsonResponse({'success': False, 'error': f'Accept failed: {str(e)}'}, status=500)

            # Convert status to job_completed format for completion updates
            job_completed_status = 'YES' if status_lower in ['completed', 'yes'] else 'NO'

            # Enhanced: Update metadata in all relevant folders and send notification
            success, user_email = update_job_status_comprehensive(filename, job_completed_status, vendor_id, completion_time)

            if success:
                print(f"✅ Job status updated by vendor {vendor_id}: {filename} -> {job_completed_status}")
                
                # Send notification to user if job is completed
                if job_completed_status == 'YES' and user_email:
                    # Get the token associated with this job
                    job_token = get_token_from_file_metadata(filename, vendor_id)
                    if job_token:
                        print(f"🔍 Found token {job_token} for job {filename}")
                    else:
                        print(f"⚠️ No token found for job {filename}")
                    
                    send_job_completion_notification(user_email, filename, vendor_id, 'completed', datetime.datetime.now().isoformat(), job_token)
                
                return JsonResponse({
                    'success': True,
                    'message': f'Job status updated for {filename}',
                    'status': job_completed_status,
                    'user_notified': user_email is not None
                })
            else:
                return JsonResponse({'success': False, 'error': 'Failed to update job status'})

        except json.JSONDecodeError:
            return JsonResponse({'success': False, 'error': 'Invalid JSON'}, status=400)
        except Exception as e:
            print(f"Error updating job status: {str(e)}")
            return JsonResponse({'success': False, 'error': str(e)}, status=500)

    return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)


def get_pending_print_jobs():
    """Get pending print jobs from R2 storage with enhanced validation (for admin or dashboard only)"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        pending_jobs = []

        # Check both vendor print jobs and vendor manual print jobs folders for documents
        try:
            vendor_objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_print_jobs/')
            manual_objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_manual_print_jobs/')

            # Combine both object lists
            all_objects = []
            if vendor_objects.get("Contents"):
                all_objects.extend(vendor_objects.get("Contents", []))
            if manual_objects.get("Contents"):
                all_objects.extend(manual_objects.get("Contents", []))

            for obj in all_objects:
                key = obj["Key"]
                filename = key.split("/")[-1]

                # Skip folder itself and metadata files
                if not filename or filename.lower().endswith('.json'):
                    continue

                # Process files that are in either vendor_print_jobs or vendor_manual_print_jobs folders
                path_parts = key.split('/')
                # Expected structure: vendor_print_jobs/{vendor_id}/{filename} or vendor_manual_print_jobs/{vendor_id}/{filename}
                if len(path_parts) >= 3 and (path_parts[0] == 'vendor_print_jobs' or path_parts[0] == 'vendor_manual_print_jobs'):
                    try:
                        # Get object metadata
                        head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                        metadata = head_response.get('Metadata', {})

                        # Check if job is pending (job_completed = 'NO')
                        job_completed = metadata.get('job_completed', 'NO').upper()
                        status = metadata.get('status', 'pending').lower()
                        service_type = metadata.get('service_type', '').strip().lower()

                        if (job_completed == 'NO' or status == 'pending'):
                            # Generate actual presigned URL for downloading
                            actual_download_url = s3.generate_presigned_url(
                                ClientMethod='get_object',
                                Params={
                                    'Bucket': settings.R2_BUCKET,
                                    'Key': key
                                },
                                ExpiresIn=3600
                            )

                            # Extract vendor info from path
                            vendor_id = path_parts[1] if len(path_parts) > 1 else 'vendor1'

                            # Build job info with proper R2 structure
                            job_info = {
                                'filename': filename,
                                'download_url': actual_download_url,  # Use actual presigned URL for download
                                'r2_path': key,  # Use actual key path
                                'user_email': metadata.get('user', ''),
                                'metadata': {
                                    'status': 'no',  # Set to 'no' for pending jobs
                                    'job_completed': job_completed,
                                    'copies': metadata.get('copies', '1'),
                                    'color': metadata.get('color', 'Black and White'),
                                    'orientation': metadata.get('orientation', 'portrait'),
                                    'page_size': metadata.get('pagesize', 'A4'),
                                    'pages': metadata.get('pages', '1'),
                                    'timestamp': metadata.get('timestamp', obj["LastModified"].isoformat()),
                                    'vendor': metadata.get('vendor', vendor_id),
                                    'user': metadata.get('user', 'Unknown'),
                                    'service_type': metadata.get('service_type', ''),
                                    'job_id': metadata.get('job_id', ''),
                                    'token': metadata.get('token', ''),
                                    'vendor_id': vendor_id,
                                    'feedback': metadata.get('feedback', ''),
                                    'quality': metadata.get('quality', ''),
                                    'thickness': metadata.get('thickness', ''),
                                    'service_name': metadata.get('service_name', '')
                                }
                            }

                            pending_jobs.append(job_info)
                            print(f"✅ Found pending print job for vendor {vendor_id}: {filename} (status: {status}, completed: {job_completed})")

                    except Exception as e:
                        print(f"Error processing vendor file {key}: {str(e)}")
                        continue

        except Exception as e:
            print(f"Error accessing vendor bucket: {str(e)}")

        print(f"📋 Total pending jobs found: {len(pending_jobs)}")
        return pending_jobs

    except Exception as e:
        print(f"Error getting pending jobs: {e}")
        return []

def update_job_status_in_r2(filename, status, vendor_id, user_email, r2_folder_structure):
    """Update job status in R2 storage with enhanced folder structure validation"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        job_completed_status = 'YES' if status.upper() == 'YES' else 'NO'
        updated_files = []
        # Only update vendor-specific folders (no testshop)
        if vendor_id and filename:
            vendor_key = f'vendor_register_details/{vendor_id}/firozshop/{filename}'
            try:
                head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=vendor_key)
                current_metadata = head_response.get('Metadata', {})
                current_metadata['job_completed'] = job_completed_status
                current_metadata['completion_time'] = datetime.datetime.now().isoformat()
                current_metadata['completed_by_vendor'] = vendor_id
                if job_completed_status == 'YES':
                    current_metadata['status'] = 'completed'
                copy_source = {'Bucket': settings.R2_BUCKET, 'Key': vendor_key}
                s3.copy_object(
                    CopySource=copy_source,
                    Bucket=settings.R2_BUCKET,
                    Key=vendor_key,
                    Metadata=current_metadata,
                    MetadataDirective='REPLACE'
                )
                updated_files.append(vendor_key)
                print(f"✅ Updated vendor job status: {vendor_key} -> {job_completed_status}")
            except Exception as e:
                print(f"⚠️  Vendor file {vendor_key} not found or error updating: {str(e)}")
        print(f"📋 Updated {len(updated_files)} file(s) in R2 storage")
        return len(updated_files) > 0
    except Exception as e:
        print(f"❌ Error updating R2 job status: {e}")
        traceback.print_exc()
        return False

def track_job_failure(filename, vendor_id, error_message, user_email):
    """Track job failures with enhanced logging"""
    try:
        # Log failure details
        print(f"Job failure tracked: {filename} by {vendor_id} - {error_message}")

        # Add your failure tracking logic here

        return True

    except Exception as e:
        print(f"Error tracking job failure: {e}")
        traceback.print_exc()
        return False

def update_vendor_status(vendor_id, status, details):
    """Update vendor status with enhanced tracking"""
    try:
        # Update vendor status
        print(f"Vendor status updated: {vendor_id} -> {status}")

        # Add your vendor status update logic here

        return True

    except Exception as e:
        print(f"Error updating vendor status: {e}")
        traceback.print_exc()
        return False

def update_printer_status(vendor_id, printer_stats):
    """Update printer status with enhanced tracking"""
    try:
        # Update printer status
        print(f"Printer status updated for vendor {vendor_id}: {printer_stats}")

        # Add your printer status update logic here

        return True

    except Exception as e:
        print(f"Error updating printer status: {e}")
        traceback.print_exc()
        return False


def update_file_job_status(filename, status='YES', vendor_id=None, completion_time=None):
    """
    Update the job_completed metadata for a specific file in both vendor and user folders
    """
    s3 = boto3.client('s3',
                      aws_access_key_id=settings.R2_ACCESS_KEY,
                      aws_secret_access_key=settings.R2_SECRET_KEY,
                      endpoint_url=settings.R2_ENDPOINT,
                      region_name='auto')

    updated_files = []

    try:
        # Search for the file in vendor folders and user folders
        prefixes_to_search = [
            'vendor_print_jobs/',
            'vendor_manual_print_jobs/',
            'users/'
        ]

        for prefix in prefixes_to_search:
            try:
                # List objects with the prefix
                response = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=prefix)

                for obj in response.get('Contents', []):
                    key = obj['Key']
                    # Check if this is the file we're looking for
                    if key.endswith(filename):
                        try:
                            # Get current object metadata
                            head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                            current_metadata = head_response.get('Metadata', {})

                            # Update job_completed status
                            current_metadata['job_completed'] = status.upper()
                            current_metadata['completion_time'] = datetime.datetime.now().isoformat()

                            # Add vendor information if provided
                            if vendor_id:
                                current_metadata['completed_by_vendor'] = vendor_id

                            # Use provided completion time if available
                            if completion_time:
                                try:
                                    # Convert timestamp to ISO format
                                    completion_dt = datetime.datetime.fromtimestamp(float(completion_time))
                                    current_metadata['completion_time'] = completion_dt.isoformat()
                                except (ValueError, TypeError):
                                    pass  # Use default timestamp if conversion fails

                            # Update status for better tracking
                            if status.upper() == 'YES':
                                current_metadata['status'] = 'completed'

                                # Create notification for job completion
                                user_email = current_metadata.get('user', '')
                                token = current_metadata.get('token', '')
                                service_type = current_metadata.get('service_type', '')

                                if user_email and token:
                                    # Get vendor name
                                    vendor_name = 'PrintMax Vendor'
                                    if vendor_id:
                                        try:
                                            # Try to get vendor details by vendor_id first
                                            vendor_details = get_vendor_details_by_email(vendor_id)
                                            if vendor_details:
                                                vendor_name = vendor_details.get('vendor_name', 'PrintMax Vendor')
                                            else:
                                                # If vendor_id is not an email, try to find vendor by vendor_id
                                                s3 = boto3.client('s3',
                                                                  aws_access_key_id=settings.R2_ACCESS_KEY,
                                                                  aws_secret_access_key=settings.R2_SECRET_KEY,
                                                                  endpoint_url=settings.R2_ENDPOINT,
                                                                  region_name='auto')

                                                # Search for vendor registration with this vendor_id
                                                objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_register_details/')
                                                for obj in objects.get("Contents", []):
                                                    if obj["Key"].endswith('/registration_details.json'):
                                                        try:
                                                            response = s3.get_object(Bucket=settings.R2_BUCKET, Key=obj["Key"])
                                                            vendor_data = json.loads(response['Body'].read().decode('utf-8'))
                                                            if vendor_data.get('vendor_id') == vendor_id:
                                                                vendor_name = vendor_data.get('vendor_name', 'PrintMax Vendor')
                                                                break
                                                        except:
                                                            continue
                                        except:
                                            pass

                                    # Create notification
                                    create_job_completion_notification(
                                        user_email=user_email,
                                        filename=filename,
                                        token=token,
                                        vendor_name=vendor_name,
                                        service_type=service_type,
                                        completion_time=current_metadata.get('completion_time', datetime.datetime.now().isoformat())
                                    )
                            else:
                                current_metadata['status'] = current_metadata.get('status', 'pending')

                            # Copy object with updated metadata
                            copy_source = {'Bucket': settings.R2_BUCKET, 'Key': key}

                            s3.copy_object(
                                CopySource=copy_source,
                                Bucket=settings.R2_BUCKET,
                                Key=key,
                                Metadata=current_metadata,
                                MetadataDirective='REPLACE'
                            )

                            updated_files.append(key)
                            print(f"✅ Updated job status for {key}: {status}")

                        except Exception as e:
                            print(f"❌ Error updating file {key}: {str(e)}")
                            traceback.print_exc()
                            continue

            except Exception as e:
                print(f"❌ Error searching in {prefix}: {str(e)}")
                traceback.print_exc()
                continue

        if updated_files:
            print(f"📋 Successfully updated {len(updated_files)} file(s) in R2 storage")
            return True
        else:
            print(f"⚠️ No files found with filename: {filename}")
            return False

    except Exception as e:
        print(f"❌ Error updating job status for {filename}: {str(e)}")
        traceback.print_exc()
        return False
# ─────────────────────────────────────────────────────────────
# FILE UPLOAD TO CLOUDFLARE R2
# ─────────────────────────────────────────────────────────────
@csrf_exempt  # Use proper CSRF protection in production!
def upload_to_r2(request):
    """
    Upload files to R2 storage with service-based folder routing:
    
    Manual Job Services (stored in vendor_manual_print_jobs/):
    - digital_print: Digital Document Printing
    - gloss_printing: Gloss Print  
    - jumbo_printing: Jumbo Paper Printing
    - golden_embossing: Golden Embossing
    - project_binding: Project Binding
    
    Other Services (stored in vendor_print_jobs/):
    - photo_print: Photo Print
    - passport_photo: Passport Photo
    - Any other services
    """
    if request.method == 'POST':
        try:
            files_uploaded = 0
            file_count = int(request.POST.get('file_count', 0))
            selected_vendor = request.POST.get('selected_vendor', 'firozshop')
            vendor_id = request.POST.get('vendor_id') or get_vendor_id_by_shop_folder(selected_vendor)
            db_storage_failed = False  # Initialize flag for database storage status

            # Initialize S3 client
            s3 = boto3.client('s3',
                              aws_access_key_id=settings.R2_ACCESS_KEY,
                              aws_secret_access_key=settings.R2_SECRET_KEY,
                              endpoint_url=settings.R2_ENDPOINT,
                              region_name='auto')

            # Get user details for folder creation (do not allow anonymous folders)
            user_email = (request.POST.get('user_email') or '').strip()
            if not user_email and request.user.is_authenticated:
                user_email = (request.user.email or '').strip()
            if not user_email:
                return JsonResponse({
                    'success': False,
                    'error': 'User email is required to submit a print job'
                }, status=400)
            user_id = str(request.user.id) if request.user.is_authenticated else ''

            # Get vendor_email from form data (preferred) or fallback to shop folder lookup
            vendor_email = (request.POST.get('vendor_email') or '').strip()
            if not vendor_email and selected_vendor:
                try:
                    # Get vendor email from shop folder as fallback
                    vendor_email = get_vendor_email_by_shop_folder(selected_vendor)
                    print(f"✅ Got vendor_email from shop folder: {vendor_email}")
                except Exception as e:
                    print(f"⚠️ Could not get vendor email for {selected_vendor}: {str(e)}")
            elif vendor_email:
                print(f"✅ Got vendor_email from form data: {vendor_email}")

            # Initialize token variable - will be assigned once per request (not per file)
            token = None

            # Process each file with its corresponding settings
            for i in range(file_count):
                file_key = f'file_{i}'
                settings_key = f'settings_{i}'

                if file_key in request.FILES and settings_key in request.POST:
                    # Get the file
                    file = request.FILES[file_key]
                    file_content = file.read()

                    # Get and parse the settings JSON
                    settings_json = request.POST.get(settings_key)
                    print_settings = json.loads(settings_json)
                    
                    # Debug: Log the print settings to see what's being sent
                    print(f"📋 Print settings for {file.name}:")
                    print(f"   Service type: {print_settings.get('service_type')}")
                    print(f"   Pricing details: {print_settings.get('pricing_details')}")
                    print(f"   All settings keys: {list(print_settings.keys())}")

                    # Get vendor email from settings_json if not already set (fallback)
                    if not vendor_email:
                        vendor_email = print_settings.get('vendor_email', '').strip()
                        if vendor_email:
                            print(f"✅ Got vendor_email from settings_json: {vendor_email}")
                    
                    # Final fallback: Get vendor email from shop folder if still not set
                    if not vendor_email and selected_vendor:
                        try:
                            vendor_email = get_vendor_email_by_shop_folder(selected_vendor)
                            print(f"✅ Got vendor_email from shop folder (final fallback): {vendor_email}")
                        except Exception as e:
                            print(f"⚠️ Could not get vendor email for {selected_vendor}: {str(e)}")
                    
                    # Assign token from vendor pool if vendor email is available
                    # IMPORTANT: Only assign token once per request (not per file) to prevent duplicate token assignments
                    if not token:  # Only assign if token hasn't been assigned yet in this request
                        if vendor_email:
                            token = assign_token_from_vendor_pool(vendor_email)
                            if token is None:
                                # Fallback to sequential token if vendor pool is empty
                                token = get_next_sequential_token()
                                print(f"⚠️ Vendor token pool empty, using fallback token: {token}")
                            else:
                                print(f"✅ Assigned token {token} from vendor pool for {vendor_email}")
                        else:
                            # Fallback to sequential token if no vendor email
                            token = get_next_sequential_token()
                            print(f"⚠️ No vendor email available, using fallback token: {token}")
                    else:
                        # Reuse the same token for all files in this request
                        print(f"✅ Reusing token {token} for file {file.name} (same request)")

                    # Generate a unique job_id for this file (use original_filename + timestamp for idempotency)
                    job_id = print_settings.get('job_id')
                    if not job_id:
                        job_id = str(uuid.uuid4())
                        print_settings['job_id'] = job_id

                    # Determine content type
                    content_type = file.content_type or 'application/octet-stream'

                    # Get file extension for better content type detection
                    file_extension = file.name.split('.')[-1].lower() if '.' in file.name else ''

                    # Get points information from form data
                    points_applied = request.POST.get('points_applied', 'false').lower() == 'true'
                    points_used = request.POST.get('points_used', '0')
                    final_amount = request.POST.get('final_amount', '0')
                    
                    # Build file metadata
                    file_metadata = {
                        'copies': str(print_settings.get("copies", "1")),
                        'color': print_settings.get("color", "Black and White"),
                        'color_mode': str(print_settings.get("color", "Black and White")),
                        'print_type': str(print_settings.get("print_type", "single_side")),
                        'orientation': print_settings.get("orientation", "portrait"),
                        'pageRange': str(print_settings.get("pageRange", "")),
                        'specificPages': str(print_settings.get("specificPages", "")),
                        'pageSize': str(print_settings.get("pageSize", "A4")),
                        'paper_type': str(print_settings.get("paper_type", print_settings.get("pageSize", "A4"))),
                        'layout_type': str(print_settings.get("layout", "single")),
                        'spiralBinding': str(print_settings.get("spiralBinding", "No")),
                        'lamination': str(print_settings.get("lamination", "No")),
                        'timestamp': get_ist_timestamp(),
                        'status': 'pending',
                        'job_completed': 'NO',
                        'vendor_status': 'not sended',
                        'trash': 'NO',
                        'user': user_email,
                        'user_id': user_id,
                        'vendor': vendor_id,
                        'vendor_id': vendor_id,  # Explicitly include vendor_id for User_print_jobs table
                        'vendor_email': vendor_email or '',  # Explicitly include vendor_email
                        'shop_id': str(vendor_id or ''),
                        'job_id': job_id,  # Explicitly include job_id
                        'service_type': print_settings.get('service_type', 'regular print'),  # Explicitly include service_type
                        'token': token,  # Explicitly include token
                        'feedback': print_settings.get('feedback', ''),
                        'quality': print_settings.get('quality', ''),
                        'thickness': print_settings.get('thickness', ''),
                        'service_name': print_settings.get('service_name', ''),
                        'points_applied': str(points_applied),
                        'points_used': str(points_used),
                        'final_amount': str(final_amount),
                        'user_email': user_email,
                        'page_count': str(print_settings.get("page_count", print_settings.get("pages", ""))),
                        'pages': str(print_settings.get("page_count", print_settings.get("pages", "")))
                    }
                    
                    # Add pricing details to metadata if available
                    pricing_details = print_settings.get('pricing_details')
                    if pricing_details:
                        # Create a compact summary for metadata display
                        total_price = pricing_details.get('total_price', 0)
                        breakdown = pricing_details.get('pricing_breakdown', {})
                        
                        # Handle different breakdown formats
                        price_per_page = 0
                        page_count = 0
                        num_copies = 0
                        pricing_key = ''
                        
                        # Check if breakdown is a dictionary (gloss print format)
                        if isinstance(breakdown, dict):
                            price_per_page = breakdown.get('price_per_page', 0)
                            page_count = breakdown.get('page_count', 0)
                            num_copies = breakdown.get('num_copies', 0)
                            pricing_key = breakdown.get('pricing_key_used', '')
                        # Check if breakdown is a list (golden embossing format)
                        elif isinstance(breakdown, list):
                            # Extract information from list format
                            for item in breakdown:
                                if isinstance(item, dict):
                                    label = item.get('label', '')
                                    value = item.get('value', '₹0')
                                    # Extract numeric value from ₹ format
                                    if value.startswith('₹'):
                                        value = value[1:]
                                    try:
                                        numeric_value = int(value)
                                        if 'page' in label.lower():
                                            page_count = numeric_value
                                        elif 'copy' in label.lower():
                                            num_copies = numeric_value
                                        elif 'per page' in label.lower():
                                            price_per_page = numeric_value
                                    except (ValueError, TypeError):
                                        pass
                            
                            # Check if structured data is available for golden embossing
                            if pricing_details.get('structured_data'):
                                structured = pricing_details.get('structured_data', {})
                                price_per_page = structured.get('price_per_page', price_per_page)
                                page_count = structured.get('page_count', page_count)
                                num_copies = structured.get('num_copies', num_copies)
                                pricing_key = f"golden_emboss_{structured.get('paper_type', 'unknown')}"
                        
                        # Store pricing details as compact JSON (ASCII only)
                        compact_pricing = {
                            "total": total_price,
                            "per_page": price_per_page,
                            "pages": page_count,
                            "copies": num_copies,
                            "key": pricing_key,
                            "quality": breakdown.get('quality_upgrade', 0) if isinstance(breakdown, dict) else 0
                        }
                        
                        # Include platform_profit if available in original pricing_details
                        if 'platform_profit' in pricing_details:
                            compact_pricing['platform_profit'] = pricing_details['platform_profit']
                        file_metadata['pricing_details'] = json.dumps(compact_pricing, separators=(',', ':'))
                        
                        # Also store individual fields for better visibility
                        file_metadata['total_price'] = str(total_price)
                        file_metadata['price_per_page'] = str(price_per_page)
                        file_metadata['page_count'] = str(page_count)
                        file_metadata['num_copies'] = str(num_copies)
                        file_metadata['pricing_key'] = pricing_key
                        
                        # Validate that all metadata values are ASCII-compatible
                        for key, value in file_metadata.items():
                            try:
                                value.encode('ascii')
                            except UnicodeEncodeError:
                                print(f"⚠️ Warning: Non-ASCII character found in metadata key '{key}': {value}")
                                # Replace non-ASCII characters with ASCII equivalents
                                file_metadata[key] = value.encode('ascii', errors='replace').decode('ascii')
                        
                        # Final validation - ensure all values are ASCII
                        for key, value in file_metadata.items():
                            if not isinstance(value, str):
                                file_metadata[key] = str(value)
                            # Ensure ASCII compatibility
                            file_metadata[key] = value.encode('ascii', errors='replace').decode('ascii')
                        
                        print(f"💰 Pricing details added to metadata: Rs{total_price}")
                    else:
                        print("⚠️ No pricing details found in print settings")

                    # Check if this is a photo print service
                    service_type = print_settings.get('service_type', '')
                    service_type_lc = service_type.lower() if service_type else ''
                    
                    # Store every user dashboard job under vendor_print_jobs to keep metadata consistent with D1
                    storage_folder = 'vendor_print_jobs'
                    vendor_file_key = f'{storage_folder}/{vendor_id}/{file.name}'
                    user_file_key = f'users/{user_email}/{file.name}'
                    file_metadata['storage_folder'] = storage_folder
                    print(f"📁 Storing {service_type or 'regular print'} job in {storage_folder} folder")

                    # Ensure a default render flag so dashboard can avoid double-rendering
                    if 'rendered_status' not in file_metadata:
                        file_metadata['rendered_status'] = 'NO'

                    # ALL userdashboard modals should NOT store in database here - they will be stored AFTER successful payment
                    # in verify_razorpay_payment (same pattern as jumbo_printing)
                    # This ensures data is only stored after payment is verified, preventing orphaned records
                    document_print_services = ['regular_print', 'regular print', 'document_print']
                    passport_photo_services = ['passport_photo', 'passport_print', 'photo_print']
                    digital_services = ['digital_print']
                    golden_services = ['golden_embossing', 'golden_emboss']
                    gloss_services = ['gloss_printing', 'gloss_print']
                    jumbo_services = ['jumbo_printing', 'jumbo_print']
                    
                    # All userdashboard modal services - skip database storage before payment
                    all_payment_required_services = (document_print_services + passport_photo_services + 
                                                      digital_services + golden_services + gloss_services + jumbo_services)
                    is_payment_required_service = service_type_lc in [s.lower() for s in all_payment_required_services]

                    if service_type in ['photo_print', 'passport_photo']:
                        # If the uploaded file is a PDF, just upload it directly (from jsPDF frontend)
                        if file.name.lower().endswith('.pdf') or file.content_type == 'application/pdf':
                            # Use the same keys as above
                            # Store only the binary file in R2; keep all metadata in database tables
                            s3.put_object(
                                Bucket=settings.R2_BUCKET,
                                Key=vendor_file_key,
                                Body=file_content,
                                ContentType='application/pdf'
                            )
                            s3.put_object(
                                Bucket=settings.R2_BUCKET,
                                Key=user_file_key,
                                Body=file_content,
                                ContentType='application/pdf'
                            )
                            print(f"✅ PDF uploaded directly: {file.name}")
                        else:
                            # Handle photo print processing (backend layout generation)
                            print(f"📸 Processing {service_type} service...")
                            # Collect all image files for this job
                            image_files_data = []
                            for j in range(file_count):
                                if f'file_{j}' in request.FILES:
                                    temp_file = request.FILES[f'file_{j}']
                                    if temp_file.content_type and temp_file.content_type.startswith('image/'):
                                        image_files_data.append(temp_file.read())
                            if not image_files_data:
                                image_files_data = [file_content]  # Use current file if no other images
                            # Create layout configuration
                            layout_config = {
                                'photo_count': int(print_settings.get('photo_count', 1)),
                                'layout': print_settings.get('layout', '1x1'),
                                'image_mode': print_settings.get('image_mode', 'same'),
                                'color': print_settings.get('color', 'Color'),
                                'paper_size': print_settings.get('paper_size', 'A4')
                            }
                            # Create photo layout PDF
                            if service_type == 'passport_photo':
                                # Get country and package from settings
                                country = print_settings.get('country', 'India')
                                total_prints = int(print_settings.get("copies", 8))
                                
                                print(f"🔍 Processing passport photo: Country={country}, Prints={total_prints}")

                                # Country-specific passport photo processing
                                pdf_data = create_passport_photo_layout(file_content, total_prints, country)
                            else:
                                # New photo print layout
                                pdf_data = create_photo_print_layout(image_files_data, layout_config)
                            if pdf_data:
                                # Update file metadata for photo service
                                if service_type == 'passport_photo':
                                    file_metadata.update({
                                        'service_type': service_type,
                                        'photo_count': str(total_prints),
                                        'country': country,
                                        'layout_created': 'YES',
                                        'original_filename': file.name,
                                        'paper_size': 'A4',
                                        'photo_dimensions': get_passport_photo_dimensions(country)
                                    })
                                else:
                                    file_metadata.update({
                                        'service_type': service_type,
                                        'photo_count': str(layout_config['photo_count']),
                                        'layout': layout_config['layout'],
                                        'image_mode': layout_config['image_mode'],
                                        'layout_created': 'YES',
                                        'original_filename': file.name,
                                        'paper_size': layout_config['paper_size']
                                    })
                                # Store only the PDF; persist metadata solely in the database
                                s3.put_object(
                                    Bucket=settings.R2_BUCKET,
                                    Key=vendor_file_key,
                                    Body=pdf_data,
                                    ContentType='application/pdf'
                                )
                                s3.put_object(
                                    Bucket=settings.R2_BUCKET,
                                    Key=user_file_key,
                                    Body=pdf_data,
                                    ContentType='application/pdf'
                                )
                                print(f"✅ Photo layout saved as PDF: {file.name}")
                                
                            else:
                                print(f"❌ Failed to create {service_type} layout")
                                return JsonResponse({'success': False, 'error': f'Failed to create {service_type} layout'}, status=500)
                    else:
                        # Regular file upload for non-passport services (no metadata stored in R2)
                        s3.put_object(
                            Bucket=settings.R2_BUCKET,
                            Key=vendor_file_key,
                            Body=file_content,
                            ContentType=content_type
                        )
                        s3.put_object(
                            Bucket=settings.R2_BUCKET,
                            Key=user_file_key,
                            Body=file_content,
                            ContentType=content_type
                        )

                    # Store print job in D1 database (R2 storage is already done above)
                    # SKIP database storage for ALL userdashboard modal services - they will be stored AFTER successful payment
                    # in verify_razorpay_payment (same pattern as jumbo_printing)
                    # This ensures data is only stored after payment verification, preventing orphaned records
                    if not is_payment_required_service:
                        try:
                            store_vendor_print_job_in_db(
                                vendor_id=vendor_id,
                                vendor_email=vendor_email,
                                user_email=user_email,
                                filename=file.name,
                                storage_folder=file_metadata.get('storage_folder', 'vendor_print_jobs'),
                                r2_path=vendor_file_key,
                                metadata=file_metadata,
                                pricing_details=pricing_details,
                                user_id=user_id,
                                shop_id=vendor_id
                            )
                        except Exception as db_err:
                            print(f"⚠️ Error storing print job in database: {db_err}")
                            # Don't fail the upload if database storage fails

                        try:
                            user_metadata = dict(file_metadata)
                            user_metadata['storage_folder'] = 'users'
                            store_user_print_job_in_db(
                                vendor_id=vendor_id,
                                vendor_email=vendor_email,
                                user_email=user_email,
                                filename=file.name,
                                storage_folder='users',
                                r2_path=user_file_key,
                                metadata=user_metadata,
                                pricing_details=pricing_details,
                                user_id=user_id,
                                shop_id=vendor_id
                            )
                        except Exception as user_db_err:
                            print(f"⚠️ Error storing user print job in database: {user_db_err}")
                    else:
                        print(f"⏭️ Skipping database storage for {service_type} - will be stored after successful payment (same as jumbo_printing)")
                        print(f"   ✅ All userdashboard modals now follow the same pattern: store only after payment verification")

                    files_uploaded += 1

            if files_uploaded > 0:
                # If database storage failed but files uploaded to R2, inform the user explicitly
                if 'db_storage_failed' in locals() and db_storage_failed:
                    return JsonResponse({
                        'success': False,
                        'error': 'Your file was uploaded, but we could not store the print job in our database. Please try again or contact support.'
                    }, status=500)
                return JsonResponse({
                    'success': True,
                    'message': f'{files_uploaded} file(s) uploaded successfully'
                })
            else:
                return JsonResponse({'success': False, 'error': 'No files uploaded'}, status=400)

        except json.JSONDecodeError as e:
            return JsonResponse({'success': False, 'error': f'Invalid JSON in settings: {str(e)}'}, status=400)
        except Exception as e:
            print(f"Upload error: {str(e)}")
            traceback.print_exc()
            return JsonResponse({'success': False, 'error': str(e)}, status=500)

    return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)


# ─────────────────────────────────────────────────────────────
# LIST OBJECTS IN CLOUDFLARE R2
# ─────────────────────────────────────────────────────────────


def list_r2_files():
    s3 = boto3.client('s3',
                      aws_access_key_id=settings.R2_ACCESS_KEY,
                      aws_secret_access_key=settings.R2_SECRET_KEY,
                      endpoint_url=settings.R2_ENDPOINT,
                      region_name='auto')
    try:
        file_data = []
        vendor_objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_print_jobs/')
        manual_objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_manual_print_jobs/')
        all_objects = []
        if vendor_objects.get("Contents"):
            all_objects.extend(vendor_objects.get("Contents", []))
        if manual_objects.get("Contents"):
            all_objects.extend(manual_objects.get("Contents", []))
        for obj in all_objects:
            key = obj["Key"]
            filename = key.split("/")[-1]
            if filename.lower().endswith('.json') or not filename:
                continue
            path_parts = key.split('/')
            if len(path_parts) >= 3 and (path_parts[0] == 'vendor_print_jobs' or path_parts[0] == 'vendor_manual_print_jobs'):
                try:
                    head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                    metadata = head_response.get('Metadata', {})
                    job_completed = metadata.get('job_completed', 'NO').upper()
                    if job_completed not in ['NO', 'YES']:
                        continue
                    url = s3.generate_presigned_url(
                        ClientMethod='get_object',
                        Params={
                            'Bucket': settings.R2_BUCKET,
                            'Key': key
                        },
                        ExpiresIn=3600
                    )
                    download_url = s3.generate_presigned_url(
                        ClientMethod='get_object',
                        Params={
                            'Bucket': settings.R2_BUCKET,
                            'Key': key,
                            'ResponseContentDisposition': f'inline; filename="{filename}"'
                        },
                        ExpiresIn=3600
                    )
                    file_extension = filename.split('.')[-1].lower() if '.' in filename else ''
                    file_type = get_file_type(file_extension)
                    pages = metadata.get('pages', estimate_pages_from_size(obj.get('Size', 0), file_extension))
                    vendor_id = path_parts[1] if len(path_parts) > 1 else 'vendor1'
                    file_info = {
                        "filename": filename,
                        "job_id": metadata.get('job_id', ''),
                        "preview_url": url,
                        "download_url": download_url,
                        "file_type": file_type,
                        "file_extension": file_extension,
                        "size": format_file_size(obj.get('Size', 0)),
                        "user": metadata.get('user', 'Auto User'),
                        "pages": pages,
                        "status": metadata.get('status', 'pending').title(),
                        "uploaded_at": obj["LastModified"].strftime("%Y-%m-%d %H:%M"),
                        "priority": metadata.get('priority', 'Medium'),
                        "copies": metadata.get('copies', '1'),
                        "color": metadata.get('color', 'Black and White'),
                        "orientation": metadata.get('orientation', 'portrait'),
                        "pageRange": metadata.get('pagerange', 'all'),
                        "specificPages": metadata.get('specificpages', ''),
                        "pageSize": metadata.get('pagesize', 'A4'),
                        "spiralBinding": metadata.get('spiralbinding', 'No'),
                        "lamination": metadata.get('lamination', 'No'),
                        "job_completed": metadata.get('job_completed', 'NO'),
                        "trash": metadata.get('trash', 'NO'),
                        "timestamp": metadata.get('timestamp', obj["LastModified"].isoformat()),
                        "service_type": metadata.get('service_type', ''),
                        "service_name": metadata.get('service_name', ''),
                        "token": metadata.get('token', ''),
                        "vendor_id": vendor_id,
                        "feedback": metadata.get('feedback', ''),
                        "quality": metadata.get('quality', ''),
                        "thickness": metadata.get('thickness', '')
                    }
                    file_info["print_options"] = f"{file_info['copies']} copies, {file_info['color']}, {file_info['orientation']}"
                    file_data.append(file_info)
                except Exception as e:
                    print(f"Error processing vendor file {key}: {str(e)}")
                    continue
        return file_data
    except Exception as e:
        print(f"Error listing R2 files: {str(e)}")
        return []

def get_file_type(extension):
    """Get file type based on extension"""
    file_types = {
        'pdf': 'PDF Document',
        'doc': 'Word Document',
        'docx': 'Word Document',
        'txt': 'Text Document',
        'ppt': 'PowerPoint Presentation',
        'pptx': 'PowerPoint Presentation',
        'xls': 'Excel Spreadsheet',
        'xlsx': 'Excel Spreadsheet',
        'jpg': 'JPEG Image',
        'jpeg': 'JPEG Image',
        'png': 'PNG Image',
        'gif': 'GIF Image',
        'bmp': 'BMP Image',
        'tiff': 'TIFF Image',
        'svg': 'SVG Image'
    }
    return file_types.get(extension, 'Document')

def estimate_pages_from_size(file_size, file_extension):
    """Estimate number of pages based on file size and type with improved accuracy"""
    # Convert bytes to KB
    size_kb = file_size / 1024

    # Different estimation for different file types
    if file_extension.lower() in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'svg']:
        return 1  # Images are typically 1 page
    elif file_extension.lower() == 'pdf':
        # PDFs: More accurate estimation based on typical PDF compression
        if size_kb < 50:
            return 1
        elif size_kb < 200:
            return max(1, round(size_kb / 50))  # Small PDFs have less compression
        elif size_kb < 1000:
            return max(1, round(size_kb / 80))  # Medium PDFs
        else:
            return max(1, round(size_kb / 120))  # Large PDFs have better compression
    elif file_extension.lower() in ['doc', 'docx']:
        # Word docs: More accurate estimation
        if size_kb < 100:
            return max(1, round(size_kb / 30))
        else:
            return max(1, round(size_kb / 60))
    elif file_extension.lower() in ['ppt', 'pptx']:
        # PowerPoint: More conservative estimation
        return max(1, round(size_kb / 150))
    elif file_extension.lower() in ['xls', 'xlsx']:
        # Excel: Better estimation based on typical spreadsheet size
        return max(1, round(size_kb / 60))
    elif file_extension.lower() == 'txt':
        # Text files: Very accurate estimation
        return max(1, round(size_kb / 3))  # Assuming ~3KB per page of text
    else:
        # Other files: conservative estimate
        return max(1, round(size_kb / 50))

def format_file_size(size_bytes):
    """Format file size in human readable format"""
    if size_bytes == 0:
        return "0 B"
    size_names = ["B", "KB", "MB", "GB"]
    i = 0
    while size_bytes >= 1024 and i < len(size_names) - 1:
        size_bytes /= 1024.0
        i += 1
    return f"{size_bytes:.1f} {size_names[i]}"


def create_photo_print_layout(input_images_data, layout_config):
    """
    Create a photo print layout with dynamic grid arrangements on A4 page as PDF.
    Args:
        input_images_data (list): List of input image data (bytes)
        layout_config (dict): Layout configuration with grid info and settings
    Returns:
        bytes: PDF data of the layout, None if failed
    """
    try:
        total_prints = layout_config.get('photo_count', 1)
        layout_type = str(layout_config.get('layout', '1'))
        image_mode = layout_config.get('image_mode', 'same')

        print(f"📸 Creating photo print layout for {total_prints} photos in {layout_type} arrangement...")

        # A4 dimensions at 300 DPI for high quality printing
        A4_WIDTH = 2480   # 210mm at 300 DPI
        A4_HEIGHT = 3508  # 297mm at 300 DPI
        MARGIN = 118      # 10mm margins
        SPACING = 59      # 5mm spacing between photos

        # Determine grid layout based on photo count
        if total_prints == 1:
            cols, rows = 1, 1
        elif total_prints == 2:
            cols, rows = 1, 2
        elif total_prints == 4:
            cols, rows = 2, 2
        elif total_prints == 6:
            cols, rows = 2, 3
        elif total_prints == 9:
            cols, rows = 3, 3
        else:
            cols, rows = 1, 1

        actual_photos = min(total_prints, cols * rows)

        # Calculate photo dimensions based on grid layout
        available_width = A4_WIDTH - (2 * MARGIN) - ((cols - 1) * SPACING)
        available_height = A4_HEIGHT - (2 * MARGIN) - ((rows - 1) * SPACING)
        photo_width = available_width // cols
        photo_height = available_height // rows

        # Load and process images
        print("📂 Loading and processing images...")
        processed_images = []

        for i, image_data in enumerate(input_images_data):
            if i >= actual_photos:
                break

            original_image = Image.open(io.BytesIO(image_data))
            if original_image.mode != 'RGB':
                original_image = original_image.convert('RGB')

            # Resize image to fit photo dimensions while maintaining aspect ratio
            original_width, original_height = original_image.size
            scale_width = photo_width / original_width
            scale_height = photo_height / original_height
            scale = min(scale_width, scale_height)

            new_width = int(original_width * scale)
            new_height = int(original_height * scale)
            resized_image = original_image.resize((new_width, new_height), Image.Resampling.LANCZOS)

            # Create photo with white background and centered image
            photo = Image.new('RGB', (photo_width, photo_height), 'white')
            x_offset = (photo_width - new_width) // 2
            y_offset = (photo_height - new_height) // 2
            photo.paste(resized_image, (x_offset, y_offset))
            processed_images.append(photo)

        # If same image mode and we have only one image, replicate it
        if image_mode == 'same' and len(processed_images) == 1:
            single_image = processed_images[0]
            processed_images = [single_image.copy() for _ in range(actual_photos)]

        # Ensure we have enough images for the layout
        while len(processed_images) < actual_photos:
            if len(processed_images) > 0:
                processed_images.append(processed_images[0].copy())
            else:
                # Create a placeholder image if no images available
                placeholder = Image.new('RGB', (photo_width, photo_height), 'lightgray')
                processed_images.append(placeholder)

        # Create the A4 layout with white background
        print(f"📄 Creating A4 layout with {actual_photos} photos in {cols}x{rows} grid...")
        layout = Image.new('RGB', (A4_WIDTH, A4_HEIGHT), 'white')

        # Calculate positioning to center the grid on the page
        total_width = cols * photo_width + (cols - 1) * SPACING
        total_height = rows * photo_height + (rows - 1) * SPACING
        start_x = max(MARGIN, (A4_WIDTH - total_width) // 2)
        start_y = max(MARGIN, (A4_HEIGHT - total_height) // 2)

        # Place photos on the layout
        photo_index = 0
        for row in range(rows):
            for col in range(cols):
                if photo_index >= actual_photos or photo_index >= len(processed_images):
                    break
                x = start_x + col * (photo_width + SPACING)
                y = start_y + row * (photo_height + SPACING)
                layout.paste(processed_images[photo_index], (x, y))
                photo_index += 1

        # Add subtle corner marks for cutting guidance
        draw = ImageDraw.Draw(layout)
        mark_length = 20  # Corner mark length
        mark_color = 'lightgray'  # Light gray for subtle guides
        mark_width = 1  # Thin lines

        for row in range(rows):
            for col in range(cols):
                if (row * cols + col) >= actual_photos:
                    break
                x = start_x + col * (photo_width + SPACING)
                y = start_y + row * (photo_height + SPACING)

                # Corner marks for cutting guidance
                offset = 4  # Distance from photo edge

                # Top-left corner
                draw.line([(x-offset, y-offset), (x-offset+mark_length, y-offset)], fill=mark_color, width=mark_width)
                draw.line([(x-offset, y-offset), (x-offset, y-offset+mark_length)], fill=mark_color, width=mark_width)

                # Top-right corner
                draw.line([(x+photo_width+offset-mark_length, y-offset), (x+photo_width+offset, y-offset)], fill=mark_color, width=mark_width)
                draw.line([(x+photo_width+offset, y-offset), (x+photo_width+offset, y-offset+mark_length)], fill=mark_color, width=mark_width)

                # Bottom-left corner
                draw.line([(x-offset, y+photo_height+offset-mark_length), (x-offset, y+photo_height+offset)], fill=mark_color, width=mark_width)
                draw.line([(x-offset, y+photo_height+offset), (x-offset+mark_length, y+photo_height+offset)], fill=mark_color, width=mark_width)

                # Bottom-right corner
                draw.line([(x+photo_width+offset, y+photo_height+offset-mark_length), (x+photo_width+offset, y+photo_height+offset)], fill=mark_color, width=mark_width)
                draw.line([(x+photo_width+offset-mark_length, y+photo_height+offset), (x+photo_width+offset, y+photo_height+offset)], fill=mark_color, width=mark_width)

        # Convert to high-quality PDF with proper settings
        print("💾 Converting layout to PDF...")
        pdf_buffer = io.BytesIO()

        # Save as PDF with high quality settings for printing
        layout.save(
            pdf_buffer, 
            'PDF', 
            quality=95,  # High quality
            resolution=300.0,  # 300 DPI for print quality
            optimize=False  # Don't optimize to maintain quality
        )

        pdf_data = pdf_buffer.getvalue()
        pdf_buffer.close()

        print(f"✅ Photo print layout created successfully!")
        print(f"   📄 {actual_photos} photos arranged on A4 page")
        print(f"   📐 Grid layout: {cols}x{rows}")
        print(f"   🖼️ Image mode: {image_mode}")
        print(f"   🎨 High quality 300 DPI PDF ready for printing")

        return pdf_data

    except Exception as e:
        print(f"❌ Error creating photo print layout: {e}")
        traceback.print_exc()
        return None
def create_passport_photo_layout(input_image_data, total_prints=8, country='India'):
    """
    Creates passport photo layout based on selected country and specified print count
    """
    try:
        # Country-specific configurations with exact dimensions
        country_config = {
            'India': {'size': (35, 45), 'unit': 'mm', 'prints': [8, 16, 30]},
            'United Arab Emirates (UAE)': {'size': (51, 51), 'unit': 'mm', 'prints': [8]},
            'Saudi Arabia': {'size': (51, 51), 'unit': 'mm', 'prints': [8]},
            'United States': {'size': (51, 51), 'unit': 'mm', 'prints': [8]},
            'Singapore': {'size': (35, 45), 'unit': 'mm', 'prints': [8]},
            'Thailand': {'size': (35, 45), 'unit': 'mm', 'prints': [8]},
            'United Kingdom': {'size': (35, 45), 'unit': 'mm', 'prints': [8]},
            'Qatar': {'size': (51, 51), 'unit': 'mm', 'prints': [8]},
            'Kuwait': {'size': (51, 51), 'unit': 'mm', 'prints': [8]},
            'Canada': {'size': (50, 70), 'unit': 'mm', 'prints': [8]},
            'Australia': {'size': (35, 45), 'unit': 'mm', 'prints': [8]},
            'Maldives': {'size': (51, 51), 'unit': 'mm', 'prints': [8]},
            'Nepal': {'size': (35, 45), 'unit': 'mm', 'prints': [8]},
            'Sri Lanka': {'size': (35, 45), 'unit': 'mm', 'prints': [8]},
            'Malaysia': {'size': (35, 45), 'unit': 'mm', 'prints': [8]},
            'Indonesia': {'size': (51, 51), 'unit': 'mm', 'prints': [8]},
            'Switzerland': {'size': (35, 45), 'unit': 'mm', 'prints': [8]},
            'Bhutan': {'size': (35, 45), 'unit': 'mm', 'prints': [8]},
            'Mauritius': {'size': (51, 51), 'unit': 'mm', 'prints': [8]},
            'France': {'size': (35, 45), 'unit': 'mm', 'prints': [8]},
            'Germany': {'size': (35, 45), 'unit': 'mm', 'prints': [8]},
        }

        selected_config = country_config.get(country, country_config['India'])
        photo_width_mm, photo_height_mm = selected_config['size']

        # Validate print count based on country
        if country == 'India':
            if total_prints not in selected_config['prints']:
                total_prints = 8
        else:
            total_prints = 8

        print(f"📸 Creating passport photo layout for {country}")
        print(f"   📏 Photo size: {photo_width_mm}x{photo_height_mm}mm")
        print(f"   📊 Total prints: {total_prints}")

        # Convert mm to pixels at 300 DPI (11.811 pixels per mm)
        DPI_CONVERSION = 11.811
        photo_width_px = int(photo_width_mm * DPI_CONVERSION)
        photo_height_px = int(photo_height_mm * DPI_CONVERSION)

        # A4 dimensions at 300 DPI
        A4_WIDTH = 2480   # 210mm at 300 DPI
        A4_HEIGHT = 3508  # 297mm at 300 DPI
        MARGIN = 118      # 10mm margins
        SPACING = 59      # 5mm spacing between photos

        # Calculate optimal grid layout for 8 photos (always 2x4 for consistent layout)
        cols, rows = 2, 4

        print(f"📄 Creating {cols}x{rows} grid layout")

        # Load and process the input image
        original_image = Image.open(io.BytesIO(input_image_data))
        if original_image.mode != 'RGB':
            original_image = original_image.convert('RGB')

        # Resize image to passport photo dimensions while maintaining aspect ratio
        original_width, original_height = original_image.size
        scale_width = photo_width_px / original_width
        scale_height = photo_height_px / original_height
        scale = min(scale_width, scale_height)

        new_width = int(original_width * scale)
        new_height = int(original_height * scale)
        resized_image = original_image.resize((new_width, new_height), Image.Resampling.LANCZOS)

        # Create passport photo with white background and centered image
        passport_photo = Image.new('RGB', (photo_width_px, photo_height_px), 'white')
        x_offset = (photo_width_px - new_width) // 2
        y_offset = (photo_height_px - new_height) // 2
        passport_photo.paste(resized_image, (x_offset, y_offset))

        # Create the A4 layout with white background
        layout = Image.new('RGB', (A4_WIDTH, A4_HEIGHT), 'white')

        # Calculate positioning to center the grid on the page
        total_width = cols * photo_width_px + (cols - 1) * SPACING
        total_height = rows * photo_height_px + (rows - 1) * SPACING
        start_x = max(MARGIN, (A4_WIDTH - total_width) // 2)
        start_y = max(MARGIN, (A4_HEIGHT - total_height) // 2)

        # Place photos on the layout
        photo_count = 0
        for row in range(rows):
            for col in range(cols):
                if photo_count >= total_prints:
                    break
                x = start_x + col * (photo_width_px + SPACING)
                y = start_y + row * (photo_height_px + SPACING)
                layout.paste(passport_photo, (x, y))
                photo_count += 1

        # Add cutting guides (corner marks)
        draw = ImageDraw.Draw(layout)
        mark_length = 20
        mark_color = 'lightgray'
        mark_width = 1

        for row in range(rows):
            for col in range(cols):
                if (row * cols + col) >= total_prints:
                    break
                x = start_x + col * (photo_width_px + SPACING)
                y = start_y + row * (photo_height_px + SPACING)

                # Corner marks for cutting guidance
                offset = 4

                # Top-left corner
                draw.line([(x-offset, y-offset), (x-offset+mark_length, y-offset)], fill=mark_color, width=mark_width)
                draw.line([(x-offset, y-offset), (x-offset, y-offset+mark_length)], fill=mark_color, width=mark_width)

                # Top-right corner
                draw.line([(x+photo_width_px+offset-mark_length, y-offset), (x+photo_width_px+offset, y-offset)], fill=mark_color, width=mark_width)
                draw.line([(x+photo_width_px+offset, y-offset), (x+photo_width_px+offset, y-offset+mark_length)], fill=mark_color, width=mark_width)

                # Bottom-left corner
                draw.line([(x-offset, y+photo_height_px+offset-mark_length), (x-offset, y+photo_height_px+offset)], fill=mark_color, width=mark_width)
                draw.line([(x-offset, y+photo_height_px+offset), (x-offset+mark_length, y+photo_height_px+offset)], fill=mark_color, width=mark_width)

                # Bottom-right corner
                draw.line([(x+photo_width_px+offset, y+photo_height_px+offset-mark_length), (x+photo_width_px+offset, y+photo_height_px+offset)], fill=mark_color, width=mark_width)
                draw.line([(x+photo_width_px+offset-mark_length, y+photo_height_px+offset), (x+photo_width_px+offset, y+photo_height_px+offset)], fill=mark_color, width=mark_width)

        # Convert to PDF
        pdf_buffer = io.BytesIO()
        layout.save(
            pdf_buffer, 
            'PDF', 
            quality=95,
            resolution=300.0,
            optimize=False
        )

        pdf_data = pdf_buffer.getvalue()
        pdf_buffer.close()

        print(f"✅ Passport photo layout created successfully!")
        print(f"   📄 {total_prints} photos of {photo_width_mm}x{photo_height_mm}mm each")
        print(f"   🏁 Country: {country}")
        print(f"   🎨 High quality 300 DPI PDF ready for printing")

        return pdf_data

    except Exception as e:
        print(f"❌ Error creating passport photo layout: {e}")
        traceback.print_exc()
        return None


# ─────────────────────────────────────────────────────────────
# HANDLE 'PROCEED TO PRINT' – FILE + SETTINGS
# ─────────────────────────────────────────────────────────────


@csrf_exempt
def process_print_request(request):
    if request.method == 'POST':
        try:
            file_count = int(request.POST.get('file_count', 0))
            files_processed = 0

            # Process each file with its corresponding settings
            for i in range(file_count):
                file_key = f'file_{i}'
                settings_key = f'settings_{i}'

                if file_key in request.FILES and settings_key in request.POST:
                    # Get the file
                    file = request.FILES[file_key]
                    file_content = file.read()

                    # Get and parse the settings JSON
                    settings_json = request.POST.get(settings_key)
                    print_settings = json.loads(settings_json)

                    # Use settings from the parsed JSON for metadata
                    file_name = file.name

                    # Initialize S3 client
                    s3 = boto3.client('s3',
                                      aws_access_key_id=settings.R2_ACCESS_KEY,
                                      aws_secret_access_key=settings.R2_SECRET_KEY,
                                      endpoint_url=settings.R2_ENDPOINT,
                                      region_name='auto')

                    # Upload the original file with metadata
                    s3.put_object(Bucket=settings.R2_BUCKET,
                                  Key=file_name,
                                  Body=file_content,
                                  ContentType=file.content_type,
                                  Metadata={
                                      'copies': str(print_settings.get("copies", "1")),
                                      'color': print_settings.get("color", "Black and White"),
                                      'orientation': print_settings.get("orientation", "portrait"),
                                      'pageRange': str(print_settings.get("pageRange", "")),
                                      'specificPages': str(print_settings.get("specificPages", "")),
                                      'pageSize': str(print_settings.get("pageSize", "A4")),
                                      'spiralBinding': str(print_settings.get("spiralBinding", "No")),
                                      'lamination': str(print_settings.get("lamination", "No")),
                                      'timestamp': get_ist_timestamp(),
                                      'status': 'pending',
                                      'job_completed': 'NO',
                                      'vendor_status': 'not sended',
                                      'trash': 'NO'
                                  })

                    files_processed += 1

            return JsonResponse({'success': True})
        except Exception as e:
            traceback.print_exc()
            return JsonResponse({'success': False, 'error': str(e)})

    return JsonResponse({'success': False, 'error': 'Invalid request method'})


# ─────────────────────────────────────────────────────────────
# CONVERT WORD / PPT TO PDF (for document, gloss, jumbo, golden emboss, digital)
# Try Docling first (modern, no LibreOffice); fall back to LibreOffice
# ─────────────────────────────────────────────────────────────

ALLOWED_CONVERT_EXTENSIONS = ('.doc', '.docx', '.ppt', '.pptx')
MAX_CONVERT_SIZE = 50 * 1024 * 1024  # 50 MB

PY_CONVERT_FORMATS = ('.docx', '.pptx')


def _html_to_pdf(html_content):
    """Convert HTML string to PDF bytes. Returns bytes or None."""
    try:
        from xhtml2pdf import pisa
        import io
        html_full = f'''<!DOCTYPE html><html><head><meta charset="utf-8">
        <style>body{{font-family:Segoe UI,Arial,sans-serif;padding:24px;line-height:1.5;color:#333;}}
        table{{border-collapse:collapse;margin:1em 0;width:100%;}}
        td,th{{border:1px solid #ddd;padding:8px;}}
        pre{{background:#f5f5f5;padding:10px;overflow-x:auto;}}
        h1,h2,h3{{margin-top:1em;}} img{{max-width:100%;}}</style></head>
        <body>{html_content}</body></html>'''
        pdf_buffer = io.BytesIO()
        pisa_status = pisa.CreatePDF(html_full, dest=pdf_buffer, encoding='utf-8')
        return pdf_buffer.getvalue() if not pisa_status.err else None
    except Exception:
        return None


def _convert_docx_with_mammoth(input_path):
    """Convert DOCX to PDF using mammoth (fast, pure Python). Returns PDF bytes or None."""
    try:
        import mammoth
        with open(input_path, 'rb') as f:
            result = mammoth.convert_to_html(f)
        html = result.value
        if not html or not html.strip():
            return None
        return _html_to_pdf(html)
    except Exception:
        return None


def _convert_pptx_with_pptx(input_path):
    """Convert PPTX to PDF using python-pptx (fast, pure Python). Returns PDF bytes or None."""
    try:
        from pptx import Presentation
        import html
        prs = Presentation(input_path)
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            parts.append(f'<div class="slide"><h3>Slide {slide_num}</h3>')
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        escaped = html.escape(text).replace('\n', '<br>')
                        parts.append(f'<p>{escaped}</p>')
            parts.append('</div>')
        html_content = ''.join(parts) if parts else '<p>No content</p>'
        return _html_to_pdf(html_content)
    except Exception:
        return None


def _convert_with_python(ext, input_path):
    """Try Python-based conversion. Returns PDF bytes or None."""
    if ext == '.docx':
        return _convert_docx_with_mammoth(input_path)
    if ext == '.pptx':
        return _convert_pptx_with_pptx(input_path)
    return None


def _find_libreoffice():
    """Return path to LibreOffice executable or None."""
    names = ['soffice', 'soffice.exe', 'libreoffice', 'libreoffice.exe']
    for name in names:
        path = shutil.which(name)
        if path:
            return path
    # Common Windows paths
    if os.name == 'nt':
        for base in [os.environ.get('ProgramFiles', 'C:\\Program Files'),
                     os.environ.get('ProgramFiles(x86)', 'C:\\Program Files (x86)')]:
            for sub in ['LibreOffice\\program', 'OpenOffice 4\\program']:
                exe = os.path.join(base, sub, 'soffice.exe')
                if os.path.isfile(exe):
                    return exe
    return None


@csrf_exempt
@require_POST
def convert_to_pdf(request):
    """Accept a single doc/docx/ppt/pptx file; return PDF bytes. Used by document, gloss, jumbo, golden emboss, digital modals."""
    if 'file' not in request.FILES:
        return JsonResponse({'success': False, 'error': 'No file provided'}, status=400)
    file = request.FILES['file']
    name = (file.name or '').lower()
    ext = os.path.splitext(name)[1]
    if ext not in ALLOWED_CONVERT_EXTENSIONS:
        return JsonResponse({'success': False, 'error': 'Invalid file type. Allowed: .doc, .docx, .ppt, .pptx'}, status=400)
    if file.size > MAX_CONVERT_SIZE:
        return JsonResponse({'success': False, 'error': 'File too large (max 50 MB)'}, status=400)

    tmpdir = None
    try:
        tmpdir = tempfile.mkdtemp()
        input_path = os.path.join(tmpdir, file.name.replace('/', '_').replace('\\', '_'))
        with open(input_path, 'wb') as f:
            for chunk in file.chunks():
                f.write(chunk)

        base = os.path.splitext(os.path.basename(input_path))[0]
        pdf_data = None

        # Try Python converters first (mammoth for DOCX, python-pptx for PPTX - fast, no external deps)
        if ext in PY_CONVERT_FORMATS:
            pdf_data = _convert_with_python(ext, input_path)

        # Fall back to LibreOffice if Docling failed or format is .doc/.ppt
        if not pdf_data:
            soffice = _find_libreoffice()
            if not soffice:
                return JsonResponse({
                    'success': False,
                    'error': 'Server conversion not available. Please upload a PDF or install LibreOffice.'
                }, status=501)
            outdir = os.path.join(tmpdir, 'out')
            os.makedirs(outdir, exist_ok=True)
            cmd = [
                soffice,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', outdir,
                input_path
            ]
            proc = subprocess.run(
                cmd,
                capture_output=True,
                timeout=120,
                cwd=tmpdir
            )
            if proc.returncode != 0:
                return JsonResponse({
                    'success': False,
                    'error': 'Conversion failed. Please try uploading a PDF instead.'
                }, status=500)
            pdf_path = os.path.join(outdir, base + '.pdf')
            if not os.path.isfile(pdf_path):
                return JsonResponse({'success': False, 'error': 'Conversion produced no PDF'}, status=500)
            with open(pdf_path, 'rb') as f:
                pdf_data = f.read()

        from django.http import HttpResponse
        response = HttpResponse(pdf_data, content_type='application/pdf')
        response['Content-Disposition'] = 'inline; filename="%s.pdf"' % (base.replace('"', '_'),)
        return response
    except subprocess.TimeoutExpired:
        return JsonResponse({'success': False, 'error': 'Conversion timed out'}, status=504)
    except Exception as e:
        traceback.print_exc()
        return JsonResponse({'success': False, 'error': str(e)}, status=500)
    finally:
        if tmpdir and os.path.isdir(tmpdir):
            try:
                shutil.rmtree(tmpdir, ignore_errors=True)
            except Exception:
                pass


from django.shortcuts import render
from django.conf import settings
from django.views.decorators.csrf import csrf_exempt
import hmac
import hashlib
from django.urls import reverse
from django.http import JsonResponse, HttpResponse, HttpResponseRedirect

# Google Drive imports (server-side OAuth and Drive access)
try:
    from google.oauth2.credentials import Credentials
    from googleapiclient.discovery import build
    from google_auth_oauthlib.flow import Flow
    from googleapiclient.http import MediaIoBaseDownload
    import io
except Exception:
    Credentials = None
    build = None
    Flow = None
    MediaIoBaseDownload = None
    io = None

# ─────────────────────────────────────────────────────────────
# Google Drive OAuth + File Proxy (for mobile Drive import)
# ─────────────────────────────────────────────────────────────

GOOGLE_DRIVE_SCOPES = ['https://www.googleapis.com/auth/drive.readonly']

def _get_google_oauth_client_config():
    client_id = getattr(settings, 'GOOGLE_OAUTH_CLIENT_ID', None) or getattr(settings, 'GOOGLE_CLIENT_ID', None)
    client_secret = getattr(settings, 'GOOGLE_OAUTH_CLIENT_SECRET', None)
    if not client_id or not client_secret:
        return None
    return {
        'web': {
            'client_id': client_id,
            'project_id': 'smartprint',
            'auth_uri': 'https://accounts.google.com/o/oauth2/auth',
            'token_uri': 'https://oauth2.googleapis.com/token',
            'auth_provider_x509_cert_url': 'https://www.googleapis.com/oauth2/v1/certs',
            'client_secret': client_secret,
            'redirect_uris': [],
            'javascript_origins': []
        }
    }

def drive_oauth_start(request):
    if Flow is None:
        return JsonResponse({'error': 'Google API libraries not installed'}, status=500)

    client_config = _get_google_oauth_client_config()
    if not client_config:
        return JsonResponse({'error': 'Google OAuth not configured'}, status=500)

    redirect_uri = request.build_absolute_uri(reverse('drive_oauth_callback'))
    client_config['web']['redirect_uris'] = [redirect_uri]
    client_config['web']['javascript_origins'] = [request.build_absolute_uri('/')[:-1]]

    flow = Flow.from_client_config(client_config, scopes=GOOGLE_DRIVE_SCOPES, redirect_uri=redirect_uri)
    authorization_url, state = flow.authorization_url(access_type='offline', include_granted_scopes='true', prompt='consent')
    request.session['google_drive_oauth_state'] = state
    return HttpResponseRedirect(authorization_url)

def drive_oauth_callback(request):
    if Flow is None:
        return JsonResponse({'error': 'Google API libraries not installed'}, status=500)

    client_config = _get_google_oauth_client_config()
    if not client_config:
        return JsonResponse({'error': 'Google OAuth not configured'}, status=500)

    redirect_uri = request.build_absolute_uri(reverse('drive_oauth_callback'))
    state = request.session.get('google_drive_oauth_state')
    flow = Flow.from_client_config(client_config, scopes=GOOGLE_DRIVE_SCOPES, state=state, redirect_uri=redirect_uri)
    try:
        flow.fetch_token(authorization_response=request.build_absolute_uri())
        creds = flow.credentials
        request.session['google_credentials'] = {
            'token': creds.token,
            'refresh_token': getattr(creds, 'refresh_token', None),
            'token_uri': creds.token_uri,
            'client_id': creds.client_id,
            'client_secret': creds.client_secret,
            'scopes': creds.scopes,
        }
        # Redirect back to user dashboard (print modal will retry list)
        return HttpResponseRedirect(reverse('userdashboard'))
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)

def drive_list_files(request):
    if not request.user.is_authenticated:
        return JsonResponse({'error': 'Not authenticated'}, status=401)
    creds_data = request.session.get('google_credentials')
    if not creds_data or Credentials is None or build is None:
        return JsonResponse({'error': 'Not authenticated with Google'}, status=401)
    try:
        creds = Credentials.from_authorized_user_info(creds_data)
        service = build('drive', 'v3', credentials=creds)
        # Limit to common document types
        q = "mimeType!='application/vnd.google-apps.folder' and trashed=false"
        results = service.files().list(
            q=q,
            orderBy='modifiedTime desc',
            pageSize=20,
            fields="files(id, name, mimeType, size)"
        ).execute()
        items = results.get('files', [])
        return JsonResponse({'files': items})
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)

def drive_download_file(request):
    if not request.user.is_authenticated:
        return JsonResponse({'error': 'Not authenticated'}, status=401)
    file_id = request.GET.get('file_id')
    if not file_id:
        return JsonResponse({'error': 'file_id required'}, status=400)
    creds_data = request.session.get('google_credentials')
    if not creds_data or Credentials is None or build is None or MediaIoBaseDownload is None:
        return JsonResponse({'error': 'Not authenticated with Google'}, status=401)
    try:
        creds = Credentials.from_authorized_user_info(creds_data)
        service = build('drive', 'v3', credentials=creds)
        request_media = service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request_media)
        done = False
        while not done:
            status, done = downloader.next_chunk()
        fh.seek(0)
        # Try to get file metadata (name, mimeType)
        meta = service.files().get(fileId=file_id, fields='name, mimeType').execute()
        filename = meta.get('name', 'downloaded_file')
        content_type = meta.get('mimeType', 'application/octet-stream')
        resp = HttpResponse(fh.read(), content_type=content_type)
        resp['Content-Disposition'] = f'attachment; filename="{filename}"'
        return resp
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
def drive_fetch_file(request):
    """
    Mobile-friendly endpoint: accepts a Google Drive fileId (from Picker)
    and downloads the file server-side using stored OAuth credentials.
    The file is stored temporarily in R2 under temp_drive_uploads/<session>/<uuid>_<filename>.
    Returns a JSON payload with a temporary key and basic metadata for later finalize.
    """
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Method not allowed'}, status=405)

    try:
        body = json.loads(request.body or '{}')
        file_id = body.get('file_id')
        if not file_id:
            return JsonResponse({'success': False, 'error': 'file_id is required'}, status=400)

        creds_data = request.session.get('google_credentials')
        if not creds_data or Credentials is None or build is None or MediaIoBaseDownload is None:
            return JsonResponse({'success': False, 'error': 'Not authenticated with Google'}, status=401)

        # Ensure session key exists for namespacing temp uploads
        if not request.session.session_key:
            request.session.save()
        session_key = request.session.session_key

        creds = Credentials.from_authorized_user_info(creds_data)
        service = build('drive', 'v3', credentials=creds)

        # Fetch metadata first
        meta = service.files().get(fileId=file_id, fields='name, mimeType, size').execute()
        filename = meta.get('name', 'document')
        content_type = meta.get('mimeType', 'application/octet-stream')

        # Stream download
        request_media = service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request_media)
        done = False
        while not done:
            _, done = downloader.next_chunk()
        fh.seek(0)

        # Store to R2 temporary location
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        unique_id = uuid.uuid4().hex
        safe_filename = re.sub(r'[^A-Za-z0-9._-]', '_', filename)
        temp_key = f"temp_drive_uploads/{session_key}/{unique_id}_{safe_filename}"

        s3.put_object(
            Bucket=settings.R2_BUCKET,
            Key=temp_key,
            Body=fh.getvalue(),
            ContentType=content_type,
            Metadata={
                'source': 'google_drive_picker',
                'original_filename': safe_filename,
            }
        )

        return JsonResponse({
            'success': True,
            'temp_key': temp_key,
            'filename': filename,
            'mimeType': content_type,
            'size': int(meta.get('size') or 0)
        })
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)


@csrf_exempt
def finalize_drive_upload(request):
    """
    After successful payment, move the temporary object into a permanent location.
    Expects JSON: { temp_key, target_path (optional), filename (optional) }
    Returns the final key.
    """
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Method not allowed'}, status=405)

    try:
        body = json.loads(request.body or '{}')
        temp_key = body.get('temp_key')
        filename = body.get('filename')
        if not temp_key:
            return JsonResponse({'success': False, 'error': 'temp_key is required'}, status=400)

        # Determine final storage path. Prefer user context if available.
        user_email = request.user.email if getattr(request, 'user', None) and request.user.is_authenticated else None
        vendor_id = request.session.get('vendor_id')

        safe_name = re.sub(r'[^A-Za-z0-9._-]', '_', (filename or temp_key.rsplit('/', 1)[-1]))
        if user_email:
            final_key = f"user_drive_uploads/{sanitize_email(user_email)}/{safe_name}"
        elif vendor_id:
            final_key = f"vendor_print_jobs/{vendor_id}/{safe_name}"
        else:
            # Fallback into a generic bucket path
            final_key = f"user_drive_uploads/anonymous/{safe_name}"

        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        # Copy then delete temp
        copy_source = {'Bucket': settings.R2_BUCKET, 'Key': temp_key}
        s3.copy_object(CopySource=copy_source, Bucket=settings.R2_BUCKET, Key=final_key, MetadataDirective='COPY')
        s3.delete_object(Bucket=settings.R2_BUCKET, Key=temp_key)

        return JsonResponse({'success': True, 'final_key': final_key})
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)

def sign_in(request):
    client_id = settings.GOOGLE_CLIENT_ID
    print(f"🔍 Debug: Google Client ID loaded: {client_id[:20] if client_id else 'None'}...")
    return render(request, 'login.html', {'client_id': client_id})


# ─────────────────────────────────────────────────────────────
# Razorpay: Create Order
# ─────────────────────────────────────────────────────────────
def _ensure_pkg_resources_available_for_razorpay():
    """
    Razorpay Python SDK (razorpay==1.4.2) hard-depends on `pkg_resources`
    (from setuptools). Some environments (minimal installs / certain packagers)
    omit it and crash on import.

    This function preserves existing behavior when `pkg_resources` exists,
    and provides a tiny compatibility shim when it doesn't, using
    `importlib.metadata` for the version lookup Razorpay uses.
    """
    try:
        import pkg_resources  # noqa: F401
        return
    except ModuleNotFoundError:
        pass

    import sys
    import types
    from importlib import metadata as importlib_metadata

    class DistributionNotFound(Exception):
        """Compat: mirrors pkg_resources.DistributionNotFound"""

    def require(dist_name: str):
        try:
            version = importlib_metadata.version(dist_name)
        except importlib_metadata.PackageNotFoundError as e:
            raise DistributionNotFound(str(e))
        # Razorpay accesses: pkg_resources.require("razorpay")[0].version
        return [types.SimpleNamespace(version=version)]

    shim = types.ModuleType("pkg_resources")
    shim.DistributionNotFound = DistributionNotFound
    shim.require = require
    sys.modules["pkg_resources"] = shim

@csrf_exempt
def create_razorpay_order(request):
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Invalid method'}, status=405)

    try:
        _ensure_pkg_resources_available_for_razorpay()
        # Lazy import so that deploys/admin startup don't fail
        # if Razorpay or its transitive dependencies are missing.
        import razorpay

        body = json.loads(request.body.decode('utf-8')) if request.body else {}
        amount_paise = int(body.get('amount_paise'))  # amount in paise
        receipt = body.get('receipt', f"rcpt_{int(time.time())}")

        if not settings.RAZORPAY_KEY_ID or not settings.RAZORPAY_KEY_SECRET:
            return JsonResponse({'success': False, 'error': 'Razorpay keys not configured'}, status=500)

        client = razorpay.Client(auth=(settings.RAZORPAY_KEY_ID, settings.RAZORPAY_KEY_SECRET))
        order = client.order.create({
            'amount': amount_paise,
            'currency': 'INR',
            'receipt': receipt,
            'payment_capture': 1
        })

        return JsonResponse({'success': True, 'order': order, 'key_id': settings.RAZORPAY_KEY_ID})
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)


# ─────────────────────────────────────────────────────────────
# Razorpay: Verify Payment and Persist Job
# ─────────────────────────────────────────────────────────────
@csrf_exempt
def verify_razorpay_payment(request):
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Invalid method'}, status=405)

    try:
        # Expect multipart: payment details + files + settings
        payment_id = request.POST.get('razorpay_payment_id')
        order_id = request.POST.get('razorpay_order_id')
        signature = request.POST.get('razorpay_signature')
        
        # Normal payment verification flow - all services require payment
        if not (payment_id and order_id and signature):
            return JsonResponse({'success': False, 'error': 'Missing payment details'}, status=400)

        generated_signature = hmac.new(
            settings.RAZORPAY_KEY_SECRET.encode('utf-8'),
            f"{order_id}|{payment_id}".encode('utf-8'),
            hashlib.sha256
        ).hexdigest()

        if generated_signature != signature:
            return JsonResponse({'success': False, 'error': 'Signature verification failed'}, status=400)

        # Signature valid → proceed to store files (reuse existing logic from process_print_request)
        file_count = int(request.POST.get('file_count', 0))
        
        # Validate file_count - must be greater than 0
        if file_count <= 0:
            return JsonResponse({
                'success': False, 
                'error': 'No files provided. Please upload at least one file.'
            }, status=400)
        
        # Debug: Log file count and available files
        print(f"📦 Processing {file_count} file(s)")
        print(f"📋 Available files in request.FILES: {list(request.FILES.keys())}")
        print(f"📋 Available POST keys: {list(request.POST.keys())[:10]}...")  # First 10 keys
        
        files_processed = 0
        files_failed = 0
        token_value = ''
        failed_files = []
        total_payment_amount = 0
        compensation_points_awarded = 0
        points_refunded = False  # CRITICAL: Single authoritative guard flag to prevent double refunding
        response_data = {
            'success': True,
            'files_processed': 0,
            'files_failed': 0,
            'failed_files': [],
            'token': '',
            'printer_name': '',
            'points_allotted': 0  # Only set to non-zero when upload fails (see below)
        }

        def safe_float(value, default=0.0):
            try:
                return float(value)
            except (TypeError, ValueError):
                return default

        user_email = (request.POST.get('user_email') or '').strip()
        if not user_email and request.user.is_authenticated:
            user_email = (request.user.email or '').strip()
        if not user_email:
            return JsonResponse({
                'success': False,
                'error': 'User email is required to finalize the print job'
            }, status=400)

        # Get vendor_email from form data (preferred) or fallback to shop folder lookup
        selected_vendor = request.POST.get('selected_vendor') or ''
        vendor_email = (request.POST.get('vendor_email') or '').strip()
        if not vendor_email and selected_vendor:
            try:
                vendor_email = get_vendor_email_by_shop_folder(selected_vendor)
                print(f"✅ Got vendor_email from shop folder: {vendor_email}")
            except Exception as e:
                print(f"⚠️ Could not get vendor email for {selected_vendor}: {str(e)}")
        elif vendor_email:
            print(f"✅ Got vendor_email from form data: {vendor_email}")

        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        for i in range(file_count):
            file_key = f'file_{i}'
            settings_key = f'settings_{i}'
            
            print(f"🔍 Processing file {i}: file_key={file_key}, settings_key={settings_key}")
            print(f"   Available FILES keys: {list(request.FILES.keys())}")
            print(f"   Available POST keys (first 20): {list(request.POST.keys())[:20]}")
            
            # Validate that both file and settings exist
            if file_key not in request.FILES:
                print(f"❌ Missing file at index {i} (file_key: {file_key})")
                print(f"   Available FILES: {list(request.FILES.keys())}")
                files_failed += 1
                failed_files.append({
                    'filename': f'file_{i}',
                    'error': f'File {i} not found in request.FILES. Available keys: {list(request.FILES.keys())}'
                })
                continue
                
            if settings_key not in request.POST:
                print(f"❌ Missing settings at index {i} (settings_key: {settings_key})")
                files_failed += 1
                failed_files.append({
                    'filename': request.FILES[file_key].name if file_key in request.FILES else f'file_{i}',
                    'error': f'Settings {i} not found in request'
                })
                continue
            
            try:
                fobj = request.FILES[file_key]
                file_content = fobj.read()
                settings_json = request.POST.get(settings_key)
                
                if not settings_json:
                    raise ValueError(f"Settings JSON is empty for file {i}")
                    
                print_settings = json.loads(settings_json)
            except json.JSONDecodeError as e:
                print(f"❌ Invalid JSON in settings for file {i}: {str(e)}")
                files_failed += 1
                failed_files.append({
                    'filename': fobj.name if 'fobj' in locals() else f'file_{i}',
                    'error': f'Invalid settings JSON: {str(e)}'
                })
                continue
            except Exception as e:
                print(f"❌ Error processing file {i}: {str(e)}")
                files_failed += 1
                failed_files.append({
                    'filename': fobj.name if 'fobj' in locals() else f'file_{i}',
                    'error': str(e)
                })
                continue

            # Resolve user and vendor context (ensure vendor_id is present for Worker validation)
            vendor_id = request.POST.get('vendor_id') or get_vendor_id_by_shop_folder(selected_vendor)

            # Get vendor email from settings_json if not already set (fallback)
            if not vendor_email:
                vendor_email = print_settings.get('vendor_email', '').strip()
                if vendor_email:
                    print(f"✅ Got vendor_email from settings_json: {vendor_email}")
            
            # Final fallback: Get vendor email from shop folder if still not set
            if not vendor_email and selected_vendor:
                try:
                    vendor_email = get_vendor_email_by_shop_folder(selected_vendor)
                    print(f"✅ Got vendor_email from shop folder (final fallback): {vendor_email}")
                except Exception as e:
                    print(f"⚠️ Could not get vendor email for {selected_vendor}: {str(e)}")

            # If vendor_id is still missing but vendor_email is available, derive vendor_id for Worker storage
            if (not vendor_id or vendor_id.strip() == '') and vendor_email:
                try:
                    vendor_id_lookup = get_vendor_id_by_vendor_email(vendor_email)
                    if vendor_id_lookup:
                        vendor_id = vendor_id_lookup
                        print(f"✅ Resolved vendor_id {vendor_id} from vendor_email {vendor_email}")
                except Exception as e:
                    print(f"⚠️ Could not resolve vendor_id from vendor_email {vendor_email}: {str(e)}")

            # Store every paid job inside vendor_print_jobs
            # Normalize service_type to ensure consistent handling (same logic as jumbo_printing)
            service_type_raw = (print_settings.get('service_type') or '').strip()
            service_type_lc = service_type_raw.lower()
            
            # Check service_name to distinguish between photoprint and passport photo
            service_name = (print_settings.get('service_name') or '').strip().lower()
            
            # Normalize service_type: convert 'regular print' to 'regular_print' for consistency with jumbo_printing
            # This ensures regular print is handled EXACTLY like jumbo_printing
            # CRITICAL: Document print modal must use 'regular_print' to match jumbo_printing storage pattern
            if service_type_lc in ['regular print', 'regular_print', 'document_print', 'document print']:
                service_type = 'regular_print'  # Normalize to match jumbo_printing pattern
                print(f"📝 Normalized service_type '{service_type_raw}' -> 'regular_print' (same as jumbo_printing)")
            elif service_type_lc in ['jumbo_printing', 'jumbo_print']:
                service_type = 'jumbo_printing'  # Keep jumbo_printing as-is
                print(f"📝 Service type '{service_type_raw}' -> 'jumbo_printing'")
            elif service_type_lc in ['photo_print'] and service_name == 'photprint':
                # CRITICAL: Keep 'photo_print' as-is when service_name is 'photprint' (photoprint service)
                service_type = 'photo_print'
                print(f"📝 Keeping service_type 'photo_print' for photoprint service (service_name: '{service_name}')")
            elif service_type_lc in ['passport_photo', 'passport_print', 'photo_print', 'passport photo', 'passport print']:
                service_type = 'passport_photo'  # Normalize all passport photo variants to 'passport_photo'
                print(f"📝 Normalized service_type '{service_type_raw}' -> 'passport_photo'")
            elif service_type_lc in ['digital_print']:
                service_type = 'digital_print'
            elif service_type_lc in ['golden_embossing', 'golden_emboss']:
                service_type = 'golden_embossing'
            elif service_type_lc in ['gloss_printing', 'gloss_print']:
                service_type = 'gloss_printing'
            else:
                service_type = service_type_raw or 'regular_print'  # Default to regular_print if not specified
                if not service_type_raw:
                    print(f"⚠️ No service_type specified, defaulting to 'regular_print'")
            
            # Define service types that should be stored in vendor_print_jobs with consistent R2 path pattern
            # Document print model, passport photo model, digital, golden, gloss, jumbo print model
            document_print_services = ['regular_print', 'regular print', 'document_print']
            passport_photo_services = ['passport_photo', 'passport_print', 'photo_print']
            digital_services = ['digital_print']
            golden_services = ['golden_embossing', 'golden_emboss']
            gloss_services = ['gloss_printing', 'gloss_print']
            jumbo_services = ['jumbo_printing', 'jumbo_print']
            
            # All these service types use vendor_print_jobs storage with pattern: {storage_folder}/{vendor_id}/{filename}
            all_special_services = (document_print_services + passport_photo_services + 
                                   digital_services + golden_services + gloss_services + jumbo_services)
            
            # Ensure consistent storage folder and R2 path for all service types (SAME as jumbo_printing)
            storage_folder = 'vendor_print_jobs'
            
            # Always construct R2 path as: {storage_folder}/{vendor_id}/{filename} (SAME as jumbo_printing)
            # This ensures consistent storage in R2 per vendor_id
            vendor_file_key = f'{storage_folder}/{vendor_id}/{fobj.name}'
            user_file_key = f'users/{user_email}/{fobj.name}'
            
            # Log service type for debugging (ensure regular print is logged same as jumbo_printing)
            # Always log for regular_print and jumbo_printing to ensure they're processed identically
            if service_type in ['regular_print', 'jumbo_printing'] or service_type_lc in [s.lower() for s in all_special_services]:
                print(f"📦 Storing {service_type} service (normalized from '{service_type_raw}') with consistent R2 path: {vendor_file_key}")
                print(f"   ✅ Service type normalized: '{service_type_raw}' -> '{service_type}' (same pattern as jumbo_printing)")
                print(f"   🔍 Will store in D1 database after payment verification (same as jumbo_printing)")
                print(f"   📋 File: {fobj.name}, Vendor: {vendor_email}, User: {user_email}")
            else:
                print(f"⚠️ Service type '{service_type}' may not be handled correctly - check storage logic")
                print(f"   📋 File: {fobj.name}, Vendor: {vendor_email}, User: {user_email}")

            # Ensure vendor_id is populated before DB calls (Worker requires it)
            if (not vendor_id or str(vendor_id).strip() == ''):
                if selected_vendor:
                    vendor_id = selected_vendor
                    print(f"✅ Fallback vendor_id from selected_vendor: {vendor_id}")
                elif vendor_email:
                    try:
                        vendor_id_lookup = get_vendor_id_by_vendor_email(vendor_email)
                        if vendor_id_lookup:
                            vendor_id = vendor_id_lookup
                            print(f"✅ Fallback vendor_id from vendor_email lookup: {vendor_id}")
                    except Exception as e:
                        print(f"⚠️ Failed vendor_id lookup from vendor_email {vendor_email}: {e}")

            # Assign token from vendor pool if available
            # CRITICAL: Only assign token once per payment request to prevent duplicate token assignments
            # All files in the same payment should share the same token
            if not token_value:  # Only generate once per request
                try:
                    # Assign token from vendor pool if vendor email is available
                    # The worker API ensures atomic token assignment (checks for 'free' status and updates to 'busy')
                    if vendor_email:
                        assigned_token = assign_token_from_vendor_pool(vendor_email)
                        if assigned_token is None:
                            # Fallback to sequential token if vendor pool is empty
                            token_value = get_next_sequential_token()
                            print(f"⚠️ Vendor token pool empty for {vendor_email}, using fallback token: {token_value}")
                        else:
                            token_value = str(assigned_token)
                            print(f"✅ Assigned token {token_value} from vendor pool for {vendor_email} (service: {service_type})")
                    else:
                        # Fallback to sequential token if no vendor email
                        token_value = get_next_sequential_token()
                        print(f"⚠️ No vendor email available, using fallback token: {token_value}")
                except Exception as e:
                    print(f"❌ Error in token assignment: {str(e)}")
                    token_value = get_next_sequential_token()  # Use proper sequential token instead of random
            else:
                # Reuse the same token for all files in this payment request
                print(f"✅ Reusing token {token_value} for file {fobj.name} (same payment request, service: {service_type})")

            # Generate a unique job_id if not provided
            job_id = print_settings.get('job_id', '').strip()
            if not job_id:
                import uuid
                job_id = str(uuid.uuid4())
                print(f"✅ Generated job_id: {job_id} for file {fobj.name}")
            
            # Build metadata (extend base with payment details)
            # Ensure ALL required fields are included: vendor_id, vendor_email, service_type, token, job_id
            # Handle Mixed color with page ranges
            color_value = print_settings.get('color', 'Black and White')
            page_range_value = str(print_settings.get('pageRange', ''))
            specific_pages_value = str(print_settings.get('specificPages', ''))
            
            # For Mixed color, combine bw and color page ranges into pageRange field
            bw_range_value = str(print_settings.get('bwPageRangeValue', '') or '').strip()
            color_range_value = str(print_settings.get('colorPageRangeValue', '') or '').strip()
            if color_value == 'Mixed':
                bw_range = (print_settings.get('bwPageRange') or 'all').lower()
                color_range = (print_settings.get('colorPageRange') or 'all').lower()
                use_bw_val = bw_range in ('range', 'specific') or bool(bw_range_value)
                use_color_val = color_range in ('range', 'specific') or bool(color_range_value)
                if not use_bw_val and not use_color_val:
                    page_range_value = 'BW: all | Color: all'
                else:
                    bw_str = f"BW: {bw_range_value if use_bw_val else 'all'}"
                    color_str = f"Color: {color_range_value if use_color_val else 'all'}"
                    page_range_value = f"{bw_str} | {color_str}"
            
            # Store "Mix" in DB when user selects Both (Mixed); keep Mixed elsewhere for compatibility
            color_for_db = 'Mix' if color_value == 'Mixed' else color_value
            metadata = {
                'copies': str(print_settings.get('copies', '1')),
                'color': color_for_db,
                'print_type': str(print_settings.get('print_type', 'single_side')),
                'orientation': print_settings.get('orientation', 'portrait'),
                'pageRange': page_range_value,
                'specificPages': specific_pages_value,
                'pageSize': str(print_settings.get('pageSize', 'A4')),
                'spiralBinding': str(print_settings.get('spiralBinding', 'No')),
                'lamination': str(print_settings.get('lamination', 'No')),
                'timestamp': get_ist_timestamp(),
                'status': 'pending',
                'job_completed': 'NO',
                'vendor_status': 'not sended',
                'trash': 'NO',
                'user': user_email,
                'vendor': vendor_id,
                'vendor_id': vendor_id,  # Explicitly include vendor_id
                'vendor_email': vendor_email or '',  # Explicitly include vendor_email
                'job_id': job_id,  # Use generated or provided job_id
                'service_type': service_type,  # Use normalized service_type (same as jumbo_printing)
                'service_name': str(print_settings.get('service_name', '')),
                'token': token_value,  # Explicitly include token
                'printer_name': '',
                'payment_id': payment_id,
                'order_id': order_id
            }
            
            # Always store B&W and Color page range values for DB (Mixed and non-Mixed)
            metadata['bwPageRange'] = str(print_settings.get('bwPageRange', 'all'))
            metadata['bwPageRangeValue'] = bw_range_value
            metadata['colorPageRange'] = str(print_settings.get('colorPageRange', 'all'))
            metadata['colorPageRangeValue'] = color_range_value

            # Ensure default rendered_status
            if 'rendered_status' not in metadata:
                metadata['rendered_status'] = 'NO'

            # Add shop_address and shop_name to metadata (for database storage)
            shop_address = print_settings.get('shop_address', '')
            shop_name = print_settings.get('shop_name', '')
            if (not shop_address or not shop_name) and vendor_email:
                try:
                    vendor_data = get_vendor_coordinates_from_email(vendor_email)
                    if vendor_data:
                        if not shop_address:
                            shop_address = vendor_data.get('shop_address', '')
                        if not shop_name:
                            shop_name = vendor_data.get('vendor_name', vendor_data.get('shop_name', ''))
                except Exception as e:
                    print(f"⚠️ Could not get shop address/name from vendor email: {str(e)}")
            
            metadata['shop_address'] = shop_address
            metadata['shop_name'] = shop_name

            # Resolve vendor email for printer assignment and assign by lowest count
            assigned_printer_name = ''
            try:
                # Use already-extracted vendor_email for printer assignment
                if vendor_email:
                    assigned_printer_name = assign_printer_and_increment_count(vendor_email, service_type)
                    if assigned_printer_name:
                        metadata['printer_name'] = assigned_printer_name
            except Exception as e:
                print(f"⚠️ Printer assignment failed: {str(e)}")

            # Include pricing details compactly if present (reuse logic from upload_to_r2 when possible)
            pricing_details = print_settings.get('pricing_details')
            pricing_details_serialized = None
            if pricing_details:
                if isinstance(pricing_details, dict):
                    try:
                        pricing_details_serialized = json.dumps(pricing_details)
                    except Exception:
                        pricing_details_serialized = None
                try:
                    breakdown = pricing_details.get('pricing_breakdown', {})
                    price_per_page = 0
                    page_count = 0
                    num_copies = 0
                    pricing_key = ''
                    if isinstance(breakdown, dict):
                        price_per_page = breakdown.get('price_per_page', 0)
                        page_count = breakdown.get('page_count', 0)
                        num_copies = breakdown.get('num_copies', 0)
                        pricing_key = breakdown.get('pricing_key_used', '')
                    # Also try to get from print_settings directly if not in breakdown
                    if page_count == 0:
                        page_count = print_settings.get('page_count', print_settings.get('pages', 0))
                    if num_copies == 0:
                        num_copies = print_settings.get('copies', 1)
                    
                    # Store only base_price in pricing_details (not full JSON)
                    base_price = breakdown.get('base_price', 0) if isinstance(breakdown, dict) else 0
                    if base_price == 0:
                        base_price = pricing_details.get('base_price', 0)
                    
                    # Store only base_price value, not full JSON
                    metadata['pricing_details'] = str(base_price) if base_price else None
                    metadata['total_price'] = str(pricing_details.get('total_price', 0))
                    # Explicitly store page_count and num_copies in metadata for database storage
                    metadata['page_count'] = str(page_count)
                    metadata['num_copies'] = str(num_copies)
                    metadata['price_per_page'] = str(price_per_page)
                    if 'platform_profit' in pricing_details:
                        metadata['platform_profit'] = str(pricing_details['platform_profit'])
                except Exception as e:
                    print(f"⚠️ Error processing pricing details: {e}")
                    pass
            if pricing_details_serialized and not metadata.get('pricing_details_raw'):
                metadata['pricing_details_raw'] = pricing_details_serialized
            if pricing_details and not metadata.get('final_amount'):
                final_amount_value = pricing_details.get('total_price') or pricing_details.get('final_amount')
                if final_amount_value is not None:
                    metadata['final_amount'] = str(final_amount_value)
            else:
                # Even without pricing_details, try to extract page_count and num_copies from print_settings
                page_count = print_settings.get('page_count', print_settings.get('pages', 0))
                num_copies = print_settings.get('copies', 1)
                if page_count:
                    metadata['page_count'] = str(page_count)
                if num_copies:
                    metadata['num_copies'] = str(num_copies)

            # Persist points usage/allocation and pricing into metadata for D1 storage
            metadata['points_applied'] = request.POST.get('points_applied', 'false')
            metadata['points_used'] = request.POST.get('points_used', '0')
            if not metadata.get('final_amount'):
                fallback_final_amount = request.POST.get('final_amount') or print_settings.get('final_amount')
                if fallback_final_amount is not None:
                    metadata['final_amount'] = str(fallback_final_amount)
            if pricing_details and not metadata.get('platform_profit'):
                platform_profit_val = pricing_details.get('platform_profit') or pricing_details.get('platform_commission')
                if platform_profit_val is not None:
                    metadata['platform_profit'] = str(platform_profit_val)

            # ATOMIC STORAGE: Store files in R2 and both database tables - all must succeed
            vendor_stored = False
            user_stored = False
            vendor_r2_stored = False
            user_r2_stored = False
            upload_error_msg = None
            db_storage_failed = False
            
            try:
                # Step 1: Store to vendor folder in R2
                # Store only the file bytes in R2 (metadata stays in DB tables)
                s3.put_object(
                    Bucket=settings.R2_BUCKET,
                    Key=vendor_file_key,
                    Body=file_content,
                    ContentType=fobj.content_type
                )
                vendor_r2_stored = True
                print(f"✅ Stored {fobj.name} to vendor R2 folder")

                # Step 2: Store a copy under the user's folder in R2
                s3.put_object(
                    Bucket=settings.R2_BUCKET,
                    Key=user_file_key,
                    Body=file_content,
                    ContentType=fobj.content_type
                )
                user_r2_stored = True
                print(f"✅ Stored {fobj.name} to user R2 folder")
                
                # Step 3: Store in vendor_print_jobs table (CRITICAL - must succeed)
                if not vendor_email and vendor_id:
                    try:
                        vendor_email = get_vendor_email_by_vendor_id(vendor_id)
                        print(f"✅ Got vendor_email from vendor_id: {vendor_email}")
                    except Exception as e:
                        print(f"⚠️ Could not get vendor email from vendor_id {vendor_id}: {str(e)}")
                
                # Prefer full pricing details JSON when present for accurate D1 storage
                pricing_details_for_db = metadata.get('pricing_details_raw') or pricing_details
                
                try:
                    vendor_stored = store_vendor_print_job_in_db(
                        vendor_id=vendor_id,
                        vendor_email=vendor_email,
                        user_email=user_email,
                        filename=fobj.name,
                        storage_folder=storage_folder,
                        r2_path=vendor_file_key,
                        metadata=metadata,
                        pricing_details=pricing_details_for_db,
                        user_id=str(request.user.id) if request.user.is_authenticated else None,
                        shop_id=vendor_id
                    )
                    
                    if not vendor_stored:
                        print(f"❌ Failed to store {fobj.name} in vendor_print_jobs table for service_type: {service_type}")
                        print(f"   🔍 Debug: vendor_id={vendor_id}, vendor_email={vendor_email}, user_email={user_email}")
                        db_storage_failed = True
                        raise Exception(f"Failed to store in vendor_print_jobs table for {service_type}")
                    else:
                        print(f"✅ Successfully stored {fobj.name} in vendor_print_jobs table for service_type: {service_type}")
                        print(f"   🔍 Token: {metadata.get('token')}, Job ID: {metadata.get('job_id')}, R2 Path: {vendor_file_key}")
                except Exception as vendor_db_error:
                    print(f"❌ Database error storing {fobj.name} in vendor_print_jobs: {str(vendor_db_error)}")
                    db_storage_failed = True
                    raise Exception(f"Database error: Failed to store in vendor_print_jobs table - {str(vendor_db_error)}")
                
                # Step 4: Store in user_print_jobs table (CRITICAL - must succeed)
                user_metadata = dict(metadata)
                user_metadata['storage_folder'] = 'users'
                try:
                    user_stored = store_user_print_job_in_db(
                        vendor_id=vendor_id,
                        vendor_email=vendor_email,
                        user_email=user_email,
                        filename=fobj.name,
                        storage_folder='users',
                        r2_path=user_file_key,
                        metadata=user_metadata,
                        pricing_details=pricing_details_for_db,
                        user_id=str(request.user.id) if request.user.is_authenticated else None,
                        shop_id=vendor_id
                    )
                    
                    if not user_stored:
                        print(f"❌ Failed to store {fobj.name} in user_print_jobs table for service_type: {service_type}")
                        print(f"   🔍 Debug: vendor_id={vendor_id}, vendor_email={vendor_email}, user_email={user_email}")
                        db_storage_failed = True
                        raise Exception(f"Failed to store in user_print_jobs table for {service_type}")
                    else:
                        print(f"✅ Successfully stored {fobj.name} in user_print_jobs table for service_type: {service_type}")
                        print(f"   🔍 Token: {user_metadata.get('token')}, Job ID: {user_metadata.get('job_id')}, R2 Path: {user_file_key}")
                except Exception as user_db_error:
                    print(f"❌ Database error storing {fobj.name} in user_print_jobs: {str(user_db_error)}")
                    db_storage_failed = True
                    raise Exception(f"Database error: Failed to store in user_print_jobs table - {str(user_db_error)}")
                
                # All steps succeeded - mark as processed
                files_processed += 1
                print(f"✅ Successfully stored file {fobj.name} (service_type: {service_type}) in R2 and both database tables")
                
            except Exception as upload_error:
                upload_error_msg = str(upload_error)
                print(f"❌ Failed to store file {fobj.name}: {upload_error_msg}")
                
                # ROLLBACK: Delete R2 files if they were stored but database failed
                try:
                    if vendor_r2_stored:
                        s3.delete_object(Bucket=settings.R2_BUCKET, Key=vendor_file_key)
                        print(f"🔄 Rolled back vendor R2 file: {fobj.name}")
                    if user_r2_stored:
                        s3.delete_object(Bucket=settings.R2_BUCKET, Key=user_file_key)
                        print(f"🔄 Rolled back user R2 file: {fobj.name}")
                except Exception as rollback_err:
                    print(f"⚠️ Error during rollback: {rollback_err}")
                
                # Track failed file and calculate refund amount
                files_failed += 1
                failed_files.append({
                    'filename': fobj.name,
                    'error': upload_error_msg,
                    'pricing_details': pricing_details,
                    'service_type': service_type  # Include service_type for debugging
                })
                
                # Calculate payment amount for this failed file
                if pricing_details:
                    file_price = safe_float(pricing_details.get('total_price', 0))
                    total_payment_amount += file_price
                    print(f"💰 Added ₹{file_price} to refund amount for failed {service_type} file: {fobj.name}")
                else:
                    # If no pricing_details, try to get from metadata
                    file_price = metadata.get('total_price') or metadata.get('final_amount')
                    if file_price:
                        file_price_float = safe_float(file_price, 0.0)
                        total_payment_amount += file_price_float
                        print(f"💰 Added ₹{file_price_float} to refund amount for failed {service_type} file: {fobj.name} (from metadata)")
                    else:
                        print(f"⚠️ No pricing found for failed {service_type} file: {fobj.name}")

        # CRITICAL: If any files failed (including database storage failures), calculate total refund
        # This includes files that failed R2 upload OR database storage
        if files_failed > 0:
            print(f"⚠️ {files_failed} file(s) failed to store. Raw refund amount before capping: ₹{total_payment_amount}")
            print(f"⚠️ Failed files: {[f['filename'] for f in failed_files]}")

            # IMPORTANT: Cap the refund/compensation amount to the ACTUAL amount charged after points
            # Frontend always sends 'final_amount' which already includes any point discounts.
            try:
                final_amount_str = request.POST.get('final_amount')
                final_amount_paid = safe_float(final_amount_str, 0.0) if final_amount_str else 0.0
            except Exception:
                final_amount_paid = 0.0

            if final_amount_paid > 0 and total_payment_amount > final_amount_paid:
                print(f"⚖️ Capping total_payment_amount from ₹{total_payment_amount} to actual paid amount ₹{final_amount_paid}")
                total_payment_amount = final_amount_paid

        # Deduct points immediately after successful payment verification
        try:
            points_applied = request.POST.get('points_applied', 'false').lower() == 'true'
            points_used = request.POST.get('points_used', '0')

            # Preserve decimals from the client but never allow negatives or overdrafts
            try:
                points_used_numeric = float(points_used)
            except (TypeError, ValueError):
                points_used_numeric = 0.0

            if points_applied and points_used_numeric > 0 and user_email:
                # Fetch current balance to ensure we never over-deduct
                current_balance = get_total_user_points(user_email)
                safe_balance = max(0.0, float(current_balance or 0))
                # Deduct only what is available and cap at the invoiced amount
                points_to_deduct = min(points_used_numeric, safe_balance)

                if points_to_deduct <= 0:
                    print(f"⚠️ Points deduction skipped for {user_email}: no available balance.")
                else:
                    # Store back to the response for transparency
                    response_data['points_deducted'] = round(points_to_deduct, 1)
                    response_data['points_balance_before'] = round(safe_balance, 1)

                    success = deduct_user_points(
                        user_email,
                        points_to_deduct,
                        f'Points used for payment {payment_id}'
                    )
                    if success:
                        print(f"💰 Deducted {points_to_deduct} points from {user_email} for payment {payment_id}")
                    else:
                        print(f"❌ Failed to deduct {points_to_deduct} points from {user_email}")
        except Exception as e:
            print(f"⚠️ Error deducting points after payment: {str(e)}")

        # NOTE: Points are NOT automatically earned after successful payments
        # Points are only refunded when payment fails (see refund logic below)

        # REFUND LOGIC: If files failed to upload, refund the payment amount AND return points used
        if files_failed > 0 and total_payment_amount > 0:
            try:
                # CRITICAL: Return points that were used if files failed to store
                if points_applied and points_used_numeric > 0 and user_email:
                    try:
                        # Return the points that were deducted
                        points_returned = add_user_points(
                            user_email,
                            points_used_numeric,
                            f'Points returned - {files_failed} file(s) failed to store for payment {payment_id if payment_id else "no_payment"}'
                        )
                        if points_returned:
                            print(f"✅ Returned {points_used_numeric} points to {user_email} due to file storage failure")
                            response_data['points_returned'] = points_used_numeric
                        else:
                            print(f"❌ Failed to return {points_used_numeric} points to {user_email}")
                    except Exception as points_return_error:
                        print(f"⚠️ Error returning points after file storage failure: {str(points_return_error)}")
                
                # Convert amount to paise (Razorpay uses paise)
                refund_amount_paise = int(float(total_payment_amount) * 100)
                
                # Attempt Razorpay refund (only if payment_id exists - skip for document print without payment)
                refund_success = False
                refund_id = None
                if payment_id:  # Only attempt refund if payment was made
                    try:
                        _ensure_pkg_resources_available_for_razorpay()
                        # Lazy import to avoid startup-time failures if Razorpay
                        # isn't available; this code path only runs after payment.
                        import razorpay

                        if settings.RAZORPAY_KEY_ID and settings.RAZORPAY_KEY_SECRET:
                            client = razorpay.Client(auth=(settings.RAZORPAY_KEY_ID, settings.RAZORPAY_KEY_SECRET))
                            refund_data = {
                                'amount': refund_amount_paise,
                                'notes': {
                                    'reason': 'Document upload failed',
                                    'failed_files': ', '.join([f['filename'] for f in failed_files]),
                                    'payment_id': payment_id
                                }
                            }
                            refund_response = client.payment.refund(payment_id, refund_data)
                            if refund_response and refund_response.get('id'):
                                refund_id = refund_response.get('id')
                                refund_success = True
                                print(f"✅ Successfully refunded ₹{total_payment_amount} (Payment ID: {payment_id}, Refund ID: {refund_id})")
                            else:
                                print(f"⚠️ Razorpay refund API returned unexpected response: {refund_response}")
                        else:
                            print(f"⚠️ Razorpay keys not configured, cannot process refund")
                    except Exception as refund_err:
                        print(f"❌ Razorpay refund failed: {str(refund_err)}")
                        # Fallback to points if refund fails
                        refund_success = False
                else:
                    print(f"⚠️ Skipping Razorpay refund - no payment_id (document print without payment)")
                    refund_success = False
                
                # If refund failed, compensate with points (ONLY if not already refunded)
                if not refund_success and not points_refunded:
                    user_email_for_refund = request.user.email if request.user.is_authenticated else 'anonymous'
                    if user_email_for_refund != 'anonymous':
                        # Convert payment amount to points (1 rupee = 1 point) - preserve decimal values
                        points_to_assign = float(total_payment_amount)
                        payment_ref = payment_id if payment_id else 'no_payment'
                        success = add_user_points(user_email_for_refund, points_to_assign, f"Compensation for {files_failed} failed upload(s) - refund failed for payment {payment_ref}")
                        if success:
                            compensation_points_awarded = points_to_assign
                            points_refunded = True  # Mark as refunded to prevent double refunding
                            response_data['compensation_points_awarded'] = points_to_assign  # Add for frontend compatibility
                            response_data['points_compensated'] = points_to_assign
                            print(f"💰 Assigned {points_to_assign} points to {user_email_for_refund} as compensation (refund failed)")
                        else:
                            print(f"❌ Failed to assign points to {user_email_for_refund} for failed uploads")
            except Exception as e:
                print(f"⚠️ Error processing refund/compensation for failed uploads: {str(e)}")

        # Prepare response with refund information
        # Include refund details in response for ALL service types (regular_print, jumbo_printing, etc.)
        refund_message = None
        refund_amount = 0
        if files_failed > 0 and total_payment_amount > 0:
            refund_amount = total_payment_amount
            if compensation_points_awarded > 0:
                refund_message = f"Payment of ₹{refund_amount:.2f} has been compensated with {compensation_points_awarded:.1f} points due to upload failure."
            else:
                refund_message = f"Payment of ₹{refund_amount:.2f} has been refunded due to document upload failure."
        
        response_data.update({
            'files_processed': files_processed,
            'files_failed': files_failed,
            'failed_files': failed_files,
            'upload_failed': files_failed > 0,
            'token': token_value,
            'printer_name': locals().get('metadata', {}).get('printer_name', '') if files_processed else response_data.get('printer_name', ''),
            'refund_message': refund_message,
            'refund_amount': refund_amount
        })

        # Always store points in user_points when uploads fail so they are instantly available
        # This ensures compensation for ALL service types (document print, passport photo, jumbo print, etc.)
        # CRITICAL: Only award points if not already refunded (prevents double refunding)
        if files_failed > 0 and total_payment_amount > 0 and not points_refunded:
            try:
                user_email_for_points = user_email or (request.user.email if request.user.is_authenticated else None)
                if user_email_for_points:
                    points_to_assign = float(total_payment_amount)
                    payment_ref = payment_id if payment_id else 'no_payment'
                    # Get service types of failed files for better logging
                    failed_service_types = [f.get('service_type', 'unknown') for f in failed_files]
                    service_types_str = ', '.join(set(failed_service_types))
                    success = add_user_points(
                        user_email_for_points,
                        points_to_assign,
                        f"Compensation points for {files_failed} failed upload(s) ({service_types_str}) - payment {payment_ref}"
                    )
                    if success:
                        compensation_points_awarded = points_to_assign
                        points_refunded = True  # Mark as refunded to prevent double refunding
                        response_data['points_compensated'] = points_to_assign
                        response_data['compensation_points_awarded'] = points_to_assign  # Add for frontend compatibility
                        print(f"💰 Assigned {points_to_assign} points to {user_email_for_points} for failed uploads (compensation) - Service types: {service_types_str}")
                    else:
                        print(f"❌ Failed to assign compensation points to {user_email_for_points}")
            except Exception as e:
                print(f"⚠️ Error assigning compensation points after failed uploads: {str(e)}")
                import traceback
                traceback.print_exc()
        elif files_failed > 0 and total_payment_amount > 0 and points_refunded:
            # Points already refunded - ensure response includes the compensation info
            print(f"✅ Compensation points already awarded: {compensation_points_awarded} points (prevented double refund)")
            if compensation_points_awarded > 0:
                response_data['compensation_points_awarded'] = compensation_points_awarded
                response_data['points_compensated'] = compensation_points_awarded
        
        # CRITICAL: Set points_allotted only when upload fails (for modal message display)
        # points_allotted should match compensation_points_awarded when files fail
        if files_failed > 0 and compensation_points_awarded > 0:
            response_data['points_allotted'] = compensation_points_awarded
            print(f"💰 Set points_allotted to {compensation_points_awarded} for failed uploads")
        else:
            # Ensure points_allotted is 0 on success
            response_data['points_allotted'] = 0
        
        # Add refund information if files failed
        if files_failed > 0 and total_payment_amount > 0:
            response_data['refund_amount'] = float(f"{total_payment_amount:.2f}")
            response_data['refund_message'] = (
                f"Payment of ₹{total_payment_amount:.2f} has been refunded due to document upload failure. "
                "Please check your payment method for the refund."
            )
        
        # Final validation: Ensure response includes all required fields
        if 'files_processed' not in response_data:
            response_data['files_processed'] = files_processed
        if 'files_failed' not in response_data:
            response_data['files_failed'] = files_failed
        if 'failed_files' not in response_data:
            response_data['failed_files'] = failed_files
        if 'upload_failed' not in response_data:
            response_data['upload_failed'] = files_failed > 0

        print(f"📊 Final response: files_processed={files_processed}, files_failed={files_failed}, token={token_value}")
        return JsonResponse(response_data)
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"❌ Exception in verify_razorpay_payment: {str(e)}")
        print(f"📋 Traceback:\n{error_trace}")
        
        # If payment verification fails but payment was attempted, return points if they were deducted
        try:
            points_applied = request.POST.get('points_applied', 'false').lower() == 'true'
            points_used = int(request.POST.get('points_used', '0'))
            user_email = request.user.email if request.user.is_authenticated else None
            
            # Check if we're on localhost
            is_localhost = request.get_host() in ['127.0.0.1:8000', 'localhost:8000', '127.0.0.1', 'localhost']
            
            # Return points if payment verification failed and points were used
            if points_applied and points_used > 0 and user_email:
                # Check if points were already deducted (they shouldn't be, but check anyway)
                # Actually, points are only deducted after successful verification, so this is a safety check
                # But if verification fails, we should return points if they were somehow deducted
                if is_localhost:
                    # Always return points on localhost if payment fails
                    success = add_user_points(user_email, points_used, f'Points returned - payment verification failed on localhost: {str(e)}')
                    if success:
                        print(f"✅ Returned {points_used} points to {user_email} due to payment verification failure on localhost")
        except Exception as return_error:
            print(f"⚠️ Error returning points after verification failure: {str(return_error)}")
        
        traceback.print_exc()
        return JsonResponse({'success': False, 'error': str(e)}, status=500)
from django.http import JsonResponse
from django.contrib.auth import login
from django.contrib.auth.models import User
import requests
def auth_receiver(request):
    try:
        if request.method == 'POST':
            token = request.POST.get('credential')
            
            if not token:
                return JsonResponse({'status': 'error', 'message': 'No credential provided'}, status=400)
            
            # ULTRA-FAST: Local JWT decode only (no network calls)
            try:
                # Only use local JWT decode for maximum speed
                payload = jwt.decode(token, options={"verify_signature": False})
                data = payload
                print(f"⚡ INSTANT: Local token decode for {payload.get('email', 'unknown')}")
            except Exception as local_error:
                print(f"❌ Local decode failed: {str(local_error)}")
                return JsonResponse({'status': 'error', 'message': 'Invalid token format'}, status=400)
            
            if 'sub' in data:  # 'sub' is the unique Google user ID
                email = data.get('email')
                if not email:
                    return JsonResponse({'status': 'error', 'message': 'Email not found in token'}, status=400)
                
                google_user_id = data['sub']
                
                # ✅ STRICT: Verify email is verified by Google
                email_verified = data.get('email_verified', False)
                if not email_verified:
                    print(f"❌ Email {email} not verified by Google")
                    return JsonResponse({'status': 'error', 'message': 'Email not verified by Google'}, status=403)

                # ✅ STRICT: Store signup details in D1 FIRST (synchronously) before allowing login
                # This ensures user exists in D1 before they can access dashboard
                try:
                    signup_payload = {
                        'email': email,
                        'google_user_id': google_user_id,
                        'name': data.get('name', ''),
                        'given_name': data.get('given_name', ''),
                        'family_name': data.get('family_name', ''),
                        'picture': data.get('picture', ''),
                        'email_verified': bool(email_verified),
                        'signup_timestamp': data.get('signup_timestamp') or timezone.now().isoformat(),
                        'last_login': timezone.now().isoformat(),
                        'is_active': True
                    }
                    
                    # Store signup to D1 synchronously to ensure it exists
                    endpoint, resp = post_to_worker('/add-user-signup', signup_payload)
                    if resp.status_code != 200:
                        print(f"⚠️ Failed to store signup in D1 ({resp.status_code}): {resp.text[:300]}")
                        # Still allow login but log warning
                    else:
                        response_json = resp.json()
                        if response_json.get('success'):
                            print(f"✅ User signup stored in D1 for {email}")
                        else:
                            print(f"⚠️ D1 signup response indicates failure: {response_json}")
                except Exception as d1_error:
                    print(f"❌ Error storing signup to D1: {str(d1_error)}")
                    # Don't block login if D1 storage fails, but log it

                # ✅ Find or create user with enhanced details (FAST - database only)
                try:
                    user, created = User.objects.get_or_create(
                        username=email,
                        defaults={
                            'email': email,
                            'first_name': data.get('given_name', ''),
                            'last_name': data.get('family_name', ''),
                        }
                    )
                    
                    # ✅ Update user details if they exist but are incomplete
                    if not created:
                        if not user.first_name and data.get('given_name'):
                            user.first_name = data.get('given_name')
                        if not user.last_name and data.get('family_name'):
                            user.last_name = data.get('family_name')
                        user.save()
                    
                    # ✅ Set up persistent session
                    # Explicitly specify backend since multiple authentication backends are configured.
                    # Google users should always use Django's default ModelBackend.
                    login(request, user, backend='django.contrib.auth.backends.ModelBackend')
                    
                    # ✅ Store additional user info in session for quick access
                    request.session['user_email'] = email
                    request.session['user_name'] = data.get('name', '')
                    request.session['user_picture'] = data.get('picture', '')
                    request.session['google_user_id'] = google_user_id
                    request.session['auth_method'] = 'google'  # Mark authentication method
                    
                    print(f"✅ User {email} logged in successfully with persistent session")
                    print(f"🔍 Session after login: {request.session.keys()}")
                    print(f"🔍 User authenticated after login: {request.user.is_authenticated}")
                    print(f"🔍 Google User ID: {google_user_id}")
                    
                    return JsonResponse({'status': 'success', 'email': email, 'redirect': '/userdashboard/'})
                except Exception as db_error:
                    print(f"❌ Database error during user creation/login: {str(db_error)}")
                    traceback.print_exc()
                    return JsonResponse({'status': 'error', 'message': f'Database error: {str(db_error)}'}, status=500)
            return JsonResponse({'status': 'error', 'message': 'Invalid token: missing sub field'}, status=400)
        return JsonResponse({'status': 'error', 'message': 'Invalid request method'}, status=405)
    except Exception as e:
        print(f"❌ Unexpected error in auth_receiver: {str(e)}")
        traceback.print_exc()
        return JsonResponse({'status': 'error', 'message': f'Internal server error: {str(e)}'}, status=500)


def get_passport_photo_dimensions(country):
    """Get passport photo dimensions for a specific country"""
    country_config = {
        'India': '35x45mm',
        'United Arab Emirates (UAE)': '51x51mm',
        'Saudi Arabia': '51x51mm',
        'United States': '51x51mm',
        'Singapore': '35x45mm',
        'Thailand': '35x45mm',
        'United Kingdom': '35x45mm',
        'Qatar': '51x51mm',
        'Kuwait': '51x51mm',
        'Canada': '50x70mm',
        'Australia': '35x45mm',
        'Maldives': '51x51mm',
        'Nepal': '35x45mm',
        'Sri Lanka': '35x45mm',
        'Malaysia': '35x45mm',
        'Indonesia': '51x51mm',
        'Switzerland': '35x45mm',
        'Bhutan': '35x45mm',
        'Mauritius': '51x51mm',
        'France': '35x45mm',
        'Germany': '35x45mm',
    }
    return country_config.get(country, '35x45mm')


def photoprint(request):
    """
    Render the photo print page
    """
    return render(request, 'photoprint.html')


LOCATION_SESSION_TTL_MINUTES = 15
LOCATION_SESSION_MAX_AGE_MINUTES = 60


def _cleanup_stale_location_sessions():
    """
    Remove stale vendor location sessions older than an hour to keep table lean.
    """
    expiry_threshold = timezone.now() - datetime.timedelta(minutes=LOCATION_SESSION_MAX_AGE_MINUTES)
    VendorLocationSession.objects.filter(created_at__lt=expiry_threshold).delete()


def _mark_session_expired_if_needed(session):
    """
    Mark pending sessions as expired when they cross the TTL.
    """
    if not session:
        return None

    if session.status == VendorLocationSession.STATUS_COMPLETED:
        return session

    age = timezone.now() - session.created_at
    if age > datetime.timedelta(minutes=LOCATION_SESSION_TTL_MINUTES):
        if session.status != VendorLocationSession.STATUS_EXPIRED:
            session.status = VendorLocationSession.STATUS_EXPIRED
            session.save(update_fields=['status', 'updated_at'])
    return session


def _parse_geocode_components(components):
    """
    Extract city, state, locality, and pincode from Google Geocoder components.
    """
    result = {
        'city': '',
        'state': '',
        'locality': '',
        'pincode': '',
    }

    for comp in components or []:
        types = comp.get('types', [])
        long_name = comp.get('long_name', '')
        if 'postal_code' in types and not result['pincode']:
            result['pincode'] = long_name
        if 'administrative_area_level_1' in types and not result['state']:
            result['state'] = long_name
        if 'locality' in types:
            result['city'] = long_name or result['city']
        if 'administrative_area_level_2' in types and not result['city']:
            result['city'] = long_name
        if ('sublocality' in types or 'sublocality_level_1' in types) and not result['locality']:
            result['locality'] = long_name

    return result


def _reverse_geocode_coordinates(latitude, longitude):
    """
    Use Google Maps Geocoding API to resolve a full address from coordinates.
    """
    api_key = getattr(settings, 'GOOGLE_MAPS_API', '') or getattr(settings, 'GOOGLE_DEVELOPER_KEY', '')
    if not api_key:
        raise RuntimeError('Google Maps API key is not configured')

    endpoint = 'https://maps.googleapis.com/maps/api/geocode/json'
    params = {
        'latlng': f'{latitude},{longitude}',
        'key': api_key,
    }

    response = requests.get(endpoint, params=params, timeout=10)
    response.raise_for_status()
    data = response.json()

    if data.get('status') != 'OK' or not data.get('results'):
        raise RuntimeError(f"Google Geocoding failed: {data.get('status')}")

    primary = data['results'][0]
    components = _parse_geocode_components(primary.get('address_components', []))
    components['full_address'] = primary.get('formatted_address', '')
    return components


def _to_decimal_or_none(value):
    try:
        return Decimal(str(value))
    except (InvalidOperation, TypeError, ValueError):
        return None


# ─────────────────────────────────────────────────────────────
# VENDOR REGISTRATION AND PRICING VIEWS
# ─────────────────────────────────────────────────────────────

def vendor_register(request):
    """
    Render the vendor registration page
    """
    return render(request, 'vendor_register.html')


@require_GET
def vendor_location_mobile(request):
    """
    Lightweight mobile-only view where vendors approve precise GPS coordinates.
    """
    session_id = request.GET.get('session_id', '').strip()
    return render(request, 'vendor_location_mobile.html', {
        'session_id': session_id,
    })


@require_GET
def user_location_mobile(request):
    """
    Mobile-first view for end users to share precise GPS coordinates via phone.
    """
    session_id = request.GET.get('session_id', '').strip()
    return render(request, 'user_location_mobile.html', {
        'session_id': session_id,
    })


def _build_location_session_response(request, mobile_view_name):
    """
    Shared helper to create QR-based location sessions for vendors or users.
    """
    try:
        _cleanup_stale_location_sessions()
        session = VendorLocationSession.objects.create(session_id=uuid.uuid4().hex)
        mobile_url = request.build_absolute_uri(
            f"{reverse(mobile_view_name)}?session_id={session.session_id}"
        )
        return JsonResponse({
            'success': True,
            'session_id': session.session_id,
            'mobile_url': mobile_url,
            'expires_in': LOCATION_SESSION_TTL_MINUTES * 60,
        })
    except Exception as exc:
        print(f"❌ Failed to create {mobile_view_name} session: {exc}")
        return JsonResponse({
            'success': False,
            'message': 'Could not start location verification. Please try again.',
        }, status=500)


@require_POST
def create_vendor_location_session(request):
    """
    Generate a unique session ID + mobile deep link used for vendor QR-based location capture.
    """
    return _build_location_session_response(request, 'vendor_location_mobile')


@require_POST
def create_user_location_session(request):
    """
    Generate a unique session ID + mobile deep link used for user QR-based location capture.
    """
    return _build_location_session_response(request, 'user_location_mobile')


@csrf_exempt  # Mobile deep link already carries random session_id token
@require_POST
def vendor_set_location(request):
    """
    Called from the vendor mobile page after GPS permission is granted.
    Saves coordinates + reverse geocoded address against the session.
    """
    return _handle_location_submission(request, context_label='vendor')


@csrf_exempt  # Mobile deep link already carries random session_id token
@require_POST
def user_set_location(request):
    """
    Called from the user mobile page after GPS permission is granted.
    Saves coordinates + reverse geocoded address against the session.
    """
    return _handle_location_submission(
        request,
        context_label='user',
        include_location_payload=True,
    )


def _handle_location_submission(request, *, context_label='vendor', include_location_payload=False):
    """
    Shared handler for vendor/user mobile GPS submission.
    """
    try:
        payload = json.loads(request.body or '{}')
    except json.JSONDecodeError:
        return JsonResponse({'success': False, 'message': 'Invalid JSON'}, status=400)

    session_id = str(payload.get('session_id', '')).strip()
    latitude = payload.get('latitude')
    longitude = payload.get('longitude')

    if not session_id or latitude is None or longitude is None:
        return JsonResponse({'success': False, 'message': 'Missing session or coordinates'}, status=400)

    session = VendorLocationSession.objects.filter(session_id=session_id).first()
    if not session:
        return JsonResponse({'success': False, 'message': 'Session not found'}, status=404)

    _mark_session_expired_if_needed(session)
    if session.status == VendorLocationSession.STATUS_EXPIRED:
        return JsonResponse({'success': False, 'message': 'Session expired'}, status=410)
    if session.status == VendorLocationSession.STATUS_COMPLETED:
        response_payload = {'success': True, 'message': 'Location already captured'}
        if include_location_payload:
            response_payload['location'] = {
                'latitude': float(session.latitude) if session.latitude is not None else None,
                'longitude': float(session.longitude) if session.longitude is not None else None,
                'full_address': session.full_address,
                'locality': session.locality,
                'city': session.city,
                'state': session.state,
                'pincode': session.pincode,
                'captured_at': session.updated_at.isoformat() if session.updated_at else '',
            }
        return JsonResponse(response_payload)

    try:
        lat_float = float(latitude)
        lng_float = float(longitude)
    except (TypeError, ValueError):
        return JsonResponse({'success': False, 'message': 'Invalid coordinates'}, status=400)

    try:
        geo_details = _reverse_geocode_coordinates(lat_float, lng_float)
    except Exception as exc:
        print(f"❌ Reverse geocoding failed ({context_label}): {exc}")
        return JsonResponse({
            'success': False,
            'message': 'Unable to verify address from GPS. Please retry.',
        }, status=502)

    session.latitude = _to_decimal_or_none(lat_float)
    session.longitude = _to_decimal_or_none(lng_float)
    session.full_address = geo_details.get('full_address', '')
    session.locality = geo_details.get('locality', '')
    session.city = geo_details.get('city', '')
    session.state = geo_details.get('state', '')
    session.pincode = geo_details.get('pincode', '')
    session.status = VendorLocationSession.STATUS_COMPLETED
    session.save(update_fields=[
        'latitude', 'longitude', 'full_address', 'locality',
        'city', 'state', 'pincode', 'status', 'updated_at'
    ])

    response_payload = {'success': True, 'message': 'Location captured successfully'}
    if include_location_payload:
        response_payload['location'] = {
            'latitude': lat_float,
            'longitude': lng_float,
            'full_address': session.full_address,
            'locality': session.locality,
            'city': session.city,
            'state': session.state,
            'pincode': session.pincode,
            'captured_at': session.updated_at.isoformat() if session.updated_at else '',
        }
    return JsonResponse(response_payload)


@require_GET
def vendor_location_status(request):
    """
    Laptop polls this endpoint every 2 seconds to check if the phone reported coordinates.
    """
    session_id = request.GET.get('session_id', '').strip()
    if not session_id:
        return JsonResponse({'success': False, 'status': 'missing_session'}, status=400)

    try:
        session = VendorLocationSession.objects.get(session_id=session_id)
    except VendorLocationSession.DoesNotExist:
        return JsonResponse({'success': False, 'status': 'not_found'}, status=404)

    session = _mark_session_expired_if_needed(session)

    if session.status == VendorLocationSession.STATUS_PENDING:
        return JsonResponse({'success': False, 'status': 'pending'})

    if session.status == VendorLocationSession.STATUS_EXPIRED:
        return JsonResponse({'success': False, 'status': 'expired'}, status=410)

    location_payload = {
        'latitude': float(session.latitude) if session.latitude is not None else None,
        'longitude': float(session.longitude) if session.longitude is not None else None,
        'full_address': session.full_address,
        'locality': session.locality,
        'city': session.city,
        'state': session.state,
        'pincode': session.pincode,
        'captured_at': session.updated_at.isoformat() if session.updated_at else '',
    }

    return JsonResponse({'success': True, 'location': location_payload})


@require_GET
def user_location_status(request):
    """
    User dashboard polling endpoint (delegates to vendor handler).
    """
    return vendor_location_status(request)


def vendor_documents(request):
    """
    Render the vendor documents upload page
    """
    return render(request, 'vendor_documents.html', {
        'vendor_email': request.session.get('vendor_email', ''),
    })


def _is_allowed_file(file_obj):
    """Validate file type and size (<= 5 MB)."""
    try:
        if not file_obj:
            return False, 'Missing file'
        allowed_mime = {'application/pdf', 'image/jpeg', 'image/png'}
        name_ok = any(file_obj.name.lower().endswith(ext) for ext in ['.pdf', '.jpg', '.jpeg', '.png'])
        type_ok = (getattr(file_obj, 'content_type', '') in allowed_mime) or name_ok
        if not type_ok:
            return False, 'Invalid file type. Only PDF, JPG, PNG allowed.'
        # 5 MB limit
        if getattr(file_obj, 'size', 0) > 5 * 1024 * 1024:
            return False, 'File too large. Maximum 5 MB allowed.'
        return True, None
    except Exception:
        return False, 'File validation error'


def _parse_page_range(range_str):
    """Parse page range string (e.g., '1-5, 8, 10-12') into list of page numbers"""
    pages = set()
    parts = range_str.split(',')
    for part in parts:
        trimmed = part.strip()
        if '-' in trimmed:
            start_end = trimmed.split('-')
            if len(start_end) == 2:
                try:
                    start = int(start_end[0].strip())
                    end = int(start_end[1].strip())
                    pages.update(range(start, end + 1))
                except (ValueError, TypeError):
                    pass
        else:
            try:
                page_num = int(trimmed)
                pages.add(page_num)
            except (ValueError, TypeError):
                pass
    return sorted(list(pages))

def _generate_unique_filename(base_label, original_name):
    """Generate a unique filename using label + timestamp + short uuid + original extension."""
    import re as _re
    import uuid as _uuid
    import datetime as _dt
    safe_label = _re.sub(r'[^a-zA-Z0-9_]+', '_', (base_label or 'document').lower()).strip('_')
    ext = ''
    if '.' in (original_name or ''):
        ext = '.' + original_name.split('.')[-1].lower()
    timestamp = _dt.datetime.now().strftime('%Y-%m-%d_%H%M%S')
    short_id = _uuid.uuid4().hex[:8]
    return f"{safe_label}_{timestamp}_{short_id}{ext}"


def _collect_vendor_form_data(request):
    """Collect vendor details from POST and session with sensible fallbacks."""
    data = {}
    # Accept email from POST, session, or querystring (?email=) to be resilient
    email_qs = request.GET.get('email') or ''
    data['vendor_email'] = (request.POST.get('email') or request.session.get('vendor_email') or email_qs).strip()
    data['vendor_name'] = (request.POST.get('vendor_name') or request.POST.get('beneficiaryName') or '').strip()
    data['phone_number'] = (request.POST.get('phone_number') or '').strip()
    data['beneficiary_name'] = (request.POST.get('beneficiaryName') or '').strip()
    data['beneficiary_bank'] = (request.POST.get('beneficiaryBank') or '').strip()
    data['account_number'] = (request.POST.get('accountNumber') or '').strip()
    data['ifsc_code'] = (request.POST.get('ifscCode') or '').strip()
    data['branch_code'] = (request.POST.get('branchCode') or '').strip()
    data['bank_address'] = (request.POST.get('bankAddress') or '').strip()
    data['city'] = (request.POST.get('city') or '').strip()
    data['postal_code'] = (request.POST.get('postalCode') or '').strip()
    data['country'] = (request.POST.get('country') or '').strip()
    return data


def _wants_json_response(request):
    requested_with = request.headers.get('x-requested-with', '')
    accepts = request.headers.get('accept', '')
    return 'xmlhttprequest' in requested_with.lower() or 'application/json' in accepts.lower()


def _build_s3_client():
    return boto3.client(
        's3',
        aws_access_key_id=settings.R2_ACCESS_KEY,
        aws_secret_access_key=settings.R2_SECRET_KEY,
        endpoint_url=settings.R2_ENDPOINT.rstrip('/') if getattr(settings, 'R2_ENDPOINT', None) else None,
        region_name='auto'
    )


@csrf_exempt  # Template includes CSRF; keep exempt in case of AJAX from external client
def upload_vendor_documents(request):
    """
    Handle vendor documents form submission: validate, upload files to Cloudflare R2, and write metadata JSON.

    R2 key structure (updated):
    vendor_register_details/<sanitized_email>/vendor_document/<files>
    vendor_register_details/<sanitized_email>/vendor_document/document.json
    """
    from django.urls import reverse as _reverse
    import json as _json
    import datetime as _dt
    import os as _os
    import traceback as _tb

    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)

    try:
        form_data = _collect_vendor_form_data(request)

        # Basic validations
        if not form_data['vendor_email']:
            err = 'Email is required'
            if _wants_json_response(request):
                return JsonResponse({'success': False, 'error': err}, status=400)
            messages.error(request, err)
            return redirect('vendor_documents')

        if not form_data['beneficiary_name'] or not form_data['account_number'] or not form_data['ifsc_code']:
            err = 'Beneficiary name, account number and IFSC code are required'
            if _wants_json_response(request):
                return JsonResponse({'success': False, 'error': err}, status=400)
            messages.error(request, err)
            return redirect('vendor_documents')

        # Files (AADHAR, PAN, Cheque, GST optional)
        candidate_files = [
            ('aadharFile', 'aadhaar'),
            ('qrFile', 'qr_code'),
            ('panFile', 'pan_card'),
            ('chequeFile', 'cheque'),
            ('gstFile', 'gst_certificate'),
        ]

        uploaded_files = []
        files_to_upload = []
        for field_name, label in candidate_files:
            f = request.FILES.get(field_name)
            if not f:
                continue
            ok, reason = _is_allowed_file(f)
            if not ok:
                err = f"{label.replace('_', ' ').title()}: {reason}"
                if _wants_json_response(request):
                    return JsonResponse({'success': False, 'error': err}, status=400)
                messages.error(request, err)
                return redirect('vendor_documents')
            unique_name = _generate_unique_filename(label, f.name)
            files_to_upload.append((unique_name, f))
            uploaded_files.append(unique_name)

        if not uploaded_files:
            err = 'Please upload at least one document (Aadhaar/QR/PAN/Cheque/GST).'
            if _wants_json_response(request):
                return JsonResponse({'success': False, 'error': err}, status=400)
            messages.error(request, err)
            return redirect('vendor_documents')

        # Prepare paths and client
        sanitized = sanitize_email(form_data['vendor_email'])
        # Store under a dedicated vendor_document folder
        base_prefix = f"vendor_register_details/{sanitized}/vendor_document/"

        s3 = _build_s3_client()

        # Upload files
        for unique_name, f in files_to_upload:
            key = base_prefix + unique_name
            s3.put_object(
                Bucket=settings.R2_BUCKET,
                Key=key,
                Body=f.read(),
                ContentType=getattr(f, 'content_type', 'application/octet-stream'),
                Metadata={
                    'uploaded_at': _dt.datetime.now().isoformat(),
                    'original_filename': f.name,
                }
            )

        # Build metadata JSON
        import uuid as _uuid
        vendor_id = str(_uuid.uuid4())
        submission_time = _dt.datetime.now().isoformat()
        metadata = {
            'vendor_id': vendor_id,
            'submission_time': submission_time,
            'vendor_email': form_data['vendor_email'],
            'vendor_name': form_data.get('vendor_name', ''),
            'phone_number': form_data.get('phone_number', ''),
            'bank_details': {
                'beneficiary_name': form_data['beneficiary_name'],
                'beneficiary_bank': form_data['beneficiary_bank'],
                'account_number': form_data['account_number'],
                'ifsc_code': form_data['ifsc_code'],
                'branch_code': form_data['branch_code'],
                'bank_address': form_data['bank_address'],
                'city': form_data['city'],
                'postal_code': form_data['postal_code'],
                'country': form_data['country'],
            },
            'uploaded_files': uploaded_files,
        }

        metadata_key = base_prefix + 'document.json'
        s3.put_object(
            Bucket=settings.R2_BUCKET,
            Key=metadata_key,
            Body=_json.dumps(metadata, ensure_ascii=False, indent=2),
            ContentType='application/json'
        )

        if _wants_json_response(request):
            return JsonResponse({'success': True, 'uploaded_files': uploaded_files, 'metadata_key': metadata_key})

        messages.success(request, 'Documents uploaded successfully!')
        try:
            url = _reverse('vendor_documents') + '?success=1'
        except Exception:
            url = '/vendor_documents?success=1'
        return redirect(url)

    except Exception as e:
        print(f"Error in upload_vendor_documents: {str(e)}")
        if _wants_json_response(request):
            return JsonResponse({'success': False, 'error': 'Upload failed'}, status=500)
        messages.error(request, 'Failed to upload documents. Please try again.')
        try:
            url = _reverse('vendor_documents')
        except Exception:
            url = '/vendor_documents'
        return redirect(url)
@csrf_exempt
def vendor_pricing(request):
    """
    Render the pricing form on GET, handle pricing submission on POST.
    """
    if request.method == 'GET':
        load_pricing = request.GET.get('load_pricing')

        # Prefer session-scoped vendor identity; fall back to URL param only if necessary
        vendor_email = (request.session.get('vendor_email') or '').strip()
        url_vendor_email = (request.GET.get('vendorEmail') or request.GET.get('email') or '').strip()
        if url_vendor_email:
            vendor_email = url_vendor_email

        if not vendor_email:
            if load_pricing:
                return JsonResponse({
                    'success': False,
                    'message': 'Vendor email missing in session'
                }, status=400)
            messages.error(request, 'Please complete registration or login again to continue pricing setup.')
            return redirect('vendor_register')

        # Fetch vendor profile with D1-first strategy and R2 fallback
        vendor_details = None
        try:
            vendor_details = get_vendor_details_by_email(vendor_email)
            if not vendor_details:
                # Newly created vendors might still be propagating; retry once
                time.sleep(1)
                vendor_details = get_vendor_details_by_email(vendor_email)
        except Exception as e:
            print(f"❌ Failed to fetch vendor details for pricing: {e}")

        # Load existing pricing data when requested
        pricing_data = None
        if vendor_email:
            try:
                s3 = boto3.client(
                    's3',
                    aws_access_key_id=settings.R2_ACCESS_KEY,
                    aws_secret_access_key=settings.R2_SECRET_KEY,
                    endpoint_url=settings.R2_ENDPOINT,
                    region_name='auto'
                )
            except Exception as s3_error:
                s3 = None
                print(f"⚠️ Unable to initialise R2 client for pricing fetch: {s3_error}")
        else:
            s3 = None

        if load_pricing:
            if s3:
                try:
                    pricing_key = f'vendor_register_details/{sanitize_email(vendor_email)}/pricing.json'
                    response = s3.get_object(Bucket=settings.R2_BUCKET, Key=pricing_key)
                    pricing_data = json.loads(response['Body'].read().decode('utf-8'))
                except Exception as e:
                    print(f"ℹ️ No existing pricing data found for {vendor_email}: {e}")
                    pricing_data = None
            if pricing_data:
                return JsonResponse({
                    'success': True,
                    'pricing_data': pricing_data.get('pricing_data', {}),
                    'categorized_pricing': pricing_data.get('categorized_pricing', {}),
                    'services_summary': pricing_data.get('services_summary', {})
                })
            return JsonResponse({
                'success': False,
                'message': 'No pricing data found'
            })

        context = {
            'vendor_email': vendor_email,
            'vendor': vendor_details or {},
            'vendor_details': vendor_details or {},
            'pricing_data': pricing_data
        }
        return render(request, 'vendor_pricing.html', context)
    elif request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_email = data.get('vendor_email') or data.get('email') or data.get('vendor_id')
            pricing_entries = data.get('pricing_entries', [])

            if not vendor_email:
                return JsonResponse({'success': False, 'message': 'Vendor email required'})

            # Initialize empty structures (R2 storage removed - all data now in D1)
            existing_pricing_data = {}
            existing_categorized_pricing = {
                'digital_print': {},
                'jumbo_print': {},
                'gloss_print': {},
                'photo_print': {},
                'golden_embossing': {},
                'passport_photo': {},
                'a4_print': {},
                'lamination': {},
                'binding': {}
            }
            existing_data = {}

            # Categorize pricing data by service type - merge with existing data
            pricing_entries = data.get('pricing_entries', [])
            from_dashboard = data.get('from_dashboard', False)
            categorized_pricing = existing_categorized_pricing.copy()  # Start with existing data
            
            # Update existing pricing_data with new values
            updated_pricing_data = existing_pricing_data.copy()
            
            # Organize pricing entries by service category - only update non-empty values
            for entry in pricing_entries:
                service_type = entry.get('service_type', '')
                price = entry.get('price', 0)
                
                # Convert price to float if it's a string
                try:
                    price = float(price) if price else 0
                except (ValueError, TypeError):
                    price = 0
                
                # Only update if price is provided and greater than 0
                if price and price > 0:
                    # Update the main pricing_data structure
                    updated_pricing_data[service_type] = price
                    
                    # Categorize based on service type prefix
                    if service_type.startswith('digital_print'):
                        categorized_pricing['digital_print'][service_type] = price
                    elif service_type.startswith('jumbo_print'):
                        categorized_pricing['jumbo_print'][service_type] = price
                    elif service_type.startswith('gloss_print'):
                        categorized_pricing['gloss_print'][service_type] = price
                    elif service_type.startswith('photo_print'):
                        categorized_pricing['photo_print'][service_type] = price
                    elif service_type.startswith('regular_print'):
                        categorized_pricing['a4_print'][service_type] = price
                    elif service_type.startswith('golden_emboss'):
                        categorized_pricing['golden_embossing'][service_type] = price
                    elif service_type.startswith('passport_print'):
                        categorized_pricing['passport_photo'][service_type] = price
                    elif service_type.startswith('lamination'):
                        categorized_pricing['lamination'][service_type] = price
                    elif service_type.startswith('tape_binding') or service_type.startswith('spiral_binding'):
                        categorized_pricing['binding'][service_type] = price
                    else:
                        # Fallback to general category
                        if 'general' not in categorized_pricing:
                            categorized_pricing['general'] = {}
                        categorized_pricing['general'][service_type] = price
            
            # Calculate services summary based on all pricing data (existing + new)
            all_pricing_entries = list(updated_pricing_data.values())
            # Convert all prices to float for comparison
            def safe_float_compare(price):
                try:
                    return float(price) > 0 if price else False
                except (ValueError, TypeError):
                    return False
            
            total_services = len([price for price in all_pricing_entries if safe_float_compare(price)])
            available_services = len([price for price in all_pricing_entries if safe_float_compare(price)])
            not_available_services = len([price for price in all_pricing_entries if not safe_float_compare(price)])
            
            # Save to D1 database via Worker API (REQUIRED - no R2 storage)
            api_url = getattr(settings, 'WORKER_API_URL', '')
            api_key = getattr(settings, 'WORKER_API_KEY', '')
            
            if api_url and api_key:
                # Construct the Worker API endpoint for vendor pricing
                if '/add-contact' in api_url:
                    worker_endpoint = api_url.replace('/add-contact', '/add-vendor-pricing')
                elif '/add-vendor-register' in api_url:
                    worker_endpoint = api_url.replace('/add-vendor-register', '/add-vendor-pricing')
                else:
                    worker_endpoint = api_url.rstrip('/') + '/add-vendor-pricing'
                
                # Helper function to safely convert price to float
                def safe_float_price(price):
                    try:
                        return float(price) if price else 0.0
                    except (ValueError, TypeError):
                        return 0.0
                
                # Prepare pricing data for D1 database - ONLY include fields that are actually provided
                # This prevents overwriting existing values with zeros for fields not in the form
                worker_payload = {
                    'vendor_email': vendor_email,
                    'last_updated': datetime.datetime.now().isoformat(),
                    'is_active': 'yes'
                }
                
                # Only add pricing fields that are actually in updated_pricing_data (i.e., were provided in the form)
                # Map form field names to database column names
                field_mapping = {
                    # Digital Print
                    'digital_print_a4_color': 'digital_print_a4_color',
                    'digital_print_a3_color': 'digital_print_a3_color',
                    'digital_print_12x18_color': 'digital_print_12x18_color',
                    'digital_print_a2_color': 'digital_print_a2_color',
                    'digital_print_a1_color': 'digital_print_a1_color',
                    'digital_print_a0_color': 'digital_print_a0_color',
                    # Regular Print
                    'regular_print_a4_bw': 'regular_print_a4_bw',
                    'regular_print_a4_color': 'regular_print_a4_color',
                    # Document Print (Letter)
                    'doc_letter_bw': 'doc_letter_bw',
                    'doc_letter_color': 'doc_letter_color',
                    # Photo Print
                    'photo_print_a4_bw': 'photo_print_a4_bw',
                    'photo_print_a4_color': 'photo_print_a4_color',
                    # Gloss Print
                    'gloss_print_a4_color': 'gloss_print_a4_color',
                    'gloss_print_a3_color': 'gloss_print_a3_color',
                    'gloss_print_a2_color': 'gloss_print_a2_color',
                    'gloss_print_a1_color': 'gloss_print_a1_color',
                    'gloss_print_a0_color': 'gloss_print_a0_color',
                    # Jumbo Print
                    'jumbo_print_a3_bw': 'jumbo_print_a3_bw',
                    'jumbo_print_a3_color': 'jumbo_print_a3_color',
                    'jumbo_print_a2_bw': 'jumbo_print_a2_bw',
                    'jumbo_print_a2_color': 'jumbo_print_a2_color',
                    'jumbo_print_a1_bw': 'jumbo_print_a1_bw',
                    'jumbo_print_a1_color': 'jumbo_print_a1_color',
                    'jumbo_print_a0_bw': 'jumbo_print_a0_bw',
                    'jumbo_print_a0_color': 'jumbo_print_a0_color',
                    # Passport Photo
                    'passport_print_8': 'passport_print_8',
                    'passport_print_16': 'passport_print_16',
                    'passport_print_30': 'passport_print_30',
                    # Golden Embossing
                    'golden_emboss_cover': 'golden_emboss_cover',
                    'golden_emboss_a4_color': 'golden_emboss_a4_color',
                    'golden_emboss_bond_color': 'golden_emboss_bond_color',
                    # Lamination
                    'lamination_a4_standard': 'lamination_a4_standard',
                    'lamination_a4_glossy': 'lamination_a4_glossy',
                    'lamination_a3_standard': 'lamination_a3_standard',
                    'lamination_a3_glossy': 'lamination_a3_glossy',
                    'lamination_a2_standard': 'lamination_a2_standard',
                    'lamination_a2_glossy': 'lamination_a2_glossy',
                    'lamination_a1_standard': 'lamination_a1_standard',
                    'lamination_a1_glossy': 'lamination_a1_glossy',
                    'lamination_a0_standard': 'lamination_a0_standard',
                    'lamination_a0_glossy': 'lamination_a0_glossy',
                    # Binding
                    'tape_binding_a4_100': 'tape_binding_a4_100',
                    'tape_binding_a4_200': 'tape_binding_a4_200',
                    'tape_binding_a3_100': 'tape_binding_a3_100',
                    'tape_binding_a3_200': 'tape_binding_a3_200',
                    'spiral_binding_a4_100': 'spiral_binding_a4_100',
                    'spiral_binding_a4_200': 'spiral_binding_a4_200',
                    'spiral_binding_a3_100': 'spiral_binding_a3_100',
                    'spiral_binding_a3_200': 'spiral_binding_a3_200',
                }
                
                # Only include fields that are actually in updated_pricing_data
                for form_field, db_field in field_mapping.items():
                    if form_field in updated_pricing_data:
                        worker_payload[db_field] = safe_float_price(updated_pricing_data[form_field])
                
                print(f"💾 Saving vendor pricing to D1 database via Worker API...")
                print(f"🔗 Worker endpoint: {worker_endpoint}")
                
                try:
                    import requests
                    resp = requests.post(
                        worker_endpoint,
                        json=worker_payload,
                        headers={
                            'Content-Type': 'application/json',
                            'x-api-key': api_key
                        },
                        timeout=15
                    )
                    
                    print(f"📡 Worker API Response Status: {resp.status_code}")
                    
                    if resp.status_code == 200:
                        try:
                            response_data = resp.json()
                            if response_data.get('success'):
                                print(f"✅ Vendor pricing saved to D1 database successfully")
                            else:
                                error_msg = response_data.get('error', 'Unknown error from Worker API')
                                print(f"❌ Worker API returned error: {error_msg}")
                        except json.JSONDecodeError:
                            print(f"❌ Invalid JSON response from Worker API: {resp.text[:200]}")
                    else:
                        error_text = resp.text[:500] if resp.text else 'No error message'
                        print(f"⚠️ D1 database save failed with status {resp.status_code}: {error_text}")
                except requests.exceptions.RequestException as e:
                    print(f"⚠️ Failed to connect to Worker API: {str(e)}")
                except Exception as e:
                    print(f"⚠️ Unexpected error saving to D1 database: {str(e)}")
            else:
                print(f"❌ Worker API not configured - cannot save pricing data")
                return JsonResponse({
                    'success': False,
                    'error': 'Server configuration error: Worker API not configured. Please contact support.'
                }, status=500)
            
            if not from_dashboard:
                # Create 300 tokens and save to D1 database via Worker API
                print(f"🎫 Creating 300 tokens for vendor {vendor_email}...")
                
                # Prepare token data for D1 database
                tokens_payload = {
                    'vendor_email': vendor_email,
                    'tokens': []
                }
                
                # Generate 300 tokens (1-300) all set to "free"
                for i in range(1, 301):
                    tokens_payload['tokens'].append({
                        'token_number': i,
                        'status': 'free'
                    })
                
                # Construct the Worker API endpoint for tokens
                if '/add-contact' in api_url:
                    token_endpoint = api_url.replace('/add-contact', '/add-vendor-tokens')
                elif '/add-vendor-register' in api_url:
                    token_endpoint = api_url.replace('/add-vendor-register', '/add-vendor-tokens')
                elif '/add-vendor-pricing' in api_url:
                    token_endpoint = api_url.replace('/add-vendor-pricing', '/add-vendor-tokens')
                else:
                    token_endpoint = api_url.rstrip('/') + '/add-vendor-tokens'
                
                print(f"🔗 Token endpoint: {token_endpoint}")
                
                try:
                    import requests
                    resp = requests.post(
                        token_endpoint,
                        json=tokens_payload,
                        headers={
                            'Content-Type': 'application/json',
                            'x-api-key': api_key
                        },
                        timeout=15
                    )
                    
                    print(f"📡 Token API Response Status: {resp.status_code}")
                    
                    if resp.status_code == 200:
                        try:
                            response_data = resp.json()
                            if response_data.get('success'):
                                print(f"✅ Successfully created 300 tokens for vendor {vendor_email} in D1 database")
                            else:
                                error_msg = response_data.get('error', 'Unknown error from Worker API')
                                print(f"❌ Token API returned error: {error_msg}")
                        except json.JSONDecodeError:
                            print(f"❌ Invalid JSON response from Token API: {resp.text[:200]}")
                    else:
                        error_text = resp.text[:500] if resp.text else 'No error message'
                        print(f"⚠️ Token save failed with status {resp.status_code}: {error_text}")
                except requests.exceptions.RequestException as e:
                    print(f"⚠️ Failed to connect to Token API: {str(e)}")
                except Exception as e:
                    print(f"⚠️ Unexpected error saving tokens: {str(e)}")
            else:
                print("ℹ️ Dashboard update detected - skipping token regeneration to speed up response.")

            # Calculate detailed statistics for success message
            # Convert all prices to float for comparison
            def safe_float_compare(price):
                try:
                    return float(price) > 0 if price else False
                except (ValueError, TypeError):
                    return False
            
            total_services = len([price for price in updated_pricing_data.values() if safe_float_compare(price)])
            available_services = len([price for price in updated_pricing_data.values() if safe_float_compare(price)])
            not_available_services = len([price for price in updated_pricing_data.values() if not safe_float_compare(price)])
            
            # Create detailed success message based on update type
            if from_dashboard:
                success_message = f"✅ Pricing Updated Successfully!\n\n"
                success_message += f"📊 Services Updated: {len(pricing_entries)}\n\n"
                success_message += f"Your pricing changes are now live!"
            else:
                success_message = f"✅ Pricing Configuration Updated Successfully!\n\n"
                success_message += f"📊 Services Configured: {total_services}\n"
                success_message += f"✅ Available Services: {available_services}\n"
                success_message += f"❌ Not Available: {not_available_services}\n"
                success_message += f"🎫 Token File Created: 300 tokens generated and set to 'free'\n\n"
                success_message += f"Your pricing is now live and customers can place orders!"
            
            return JsonResponse({
                'success': True,
                'message': 'Pricing configuration saved successfully',
                'detailed_message': success_message,
                'total_services': total_services,
                'available_services_count': available_services,
                'not_available_services_count': not_available_services,
                'categorized_pricing': categorized_pricing,
                'updated_at': datetime.datetime.now().isoformat()
            })

        except Exception as e:
            print(f"❌ Error saving pricing data: {str(e)}")
            return JsonResponse({
                'success': False,
                'message': f'Error saving pricing: {str(e)}'
            })
    else:
        return JsonResponse({
            'success': False,
            'message': 'Invalid request method'
        })


def vendor_info(request, vendor_id):
    """
    Get vendor information by vendor ID
    """
    try:
        if vendor_id:
            try:
                endpoint, resp = post_to_worker('/get-vendor-by-id', {'vendor_id': vendor_id})
                if resp.status_code == 200:
                    payload = resp.json()
                    if payload.get('success'):
                        vendor = payload.get('vendor') or {}
                        return JsonResponse({
                            'success': True,
                            'vendor': {
                                'vendor_id': vendor.get('vendor_id', vendor_id),
                                'vendor_name': vendor.get('vendor_name', ''),
                                'email': vendor.get('email') or vendor.get('vendor_email', ''),
                                'phone_number': vendor.get('phone_number', '')
                            }
                        })
                elif resp.status_code != 404:
                    print(f"⚠️ Worker vendor lookup failed ({resp.status_code}): {resp.text[:300]}")
            except Exception as worker_error:
                print(f"⚠️ Worker vendor lookup error: {worker_error}")

        return JsonResponse({
            'success': False,
            'message': 'Vendor not found in D1 database'
        }, status=404)

    except Exception as e:
        print(f"❌ Error fetching vendor info: {str(e)}")
        return JsonResponse({
            'success': False,
            'message': f'Error fetching vendor info: {str(e)}'
        })
# Add vendor login endpoint
@csrf_exempt
def vendor_login(request):
    """
    Handle vendor login by email using D1 database (vendor_register_details table)
    Uses the same authentication system as R2 storage (pbkdf2_sha256 password hashing)
    """
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            email = data.get('email')  # frontend now sends email as 'email'
            password = data.get('password')

            if not email or not password:
                return JsonResponse({
                    'success': False,
                    'message': 'Email and password are required'
                })

            # Get vendor from D1 database via Worker API
            api_url = getattr(settings, 'WORKER_API_URL', '')
            api_key = getattr(settings, 'WORKER_API_KEY', '')
            
            if not api_url or not api_key:
                print(f"❌ Worker API not configured")
                return JsonResponse({
                    'success': False,
                    'message': 'Server configuration error: Worker API not configured. Please contact support.'
                }, status=500)

            # Construct the Worker API endpoint for getting vendor
            if '/add-contact' in api_url:
                worker_endpoint = api_url.replace('/add-contact', '/get-vendor-by-email')
            elif '/add-vendor-register' in api_url:
                worker_endpoint = api_url.replace('/add-vendor-register', '/get-vendor-by-email')
            else:
                worker_endpoint = api_url.rstrip('/') + '/get-vendor-by-email'
            
            try:
                resp = requests.post(
                    worker_endpoint,
                    json={'email': email},
                    headers={
                        'Content-Type': 'application/json',
                        'x-api-key': api_key
                    },
                    timeout=10
                )
                
                if resp.status_code == 404:
                    return JsonResponse({
                        'success': False,
                        'message': 'Vendor not found with this email address'
                    })

                if resp.status_code != 200:
                    error_text = resp.text[:500] if resp.text else 'Unknown error'
                    print(f"❌ Worker API error: {error_text}")
                    return JsonResponse({
                        'success': False,
                        'message': 'Error finding vendor account'
                    })
                
                response_data = resp.json()
                if not response_data.get('success'):
                    return JsonResponse({
                        'success': False,
                        'message': response_data.get('error', 'Vendor not found')
                    })
                
                vendor_data = response_data.get('vendor')
                if not vendor_data:
                    return JsonResponse({
                        'success': False,
                        'message': 'Vendor not found with this email address'
                    })
                
                # Get password_hash from database
                password_hash = vendor_data.get('password_hash')
                if not password_hash:
                    print(f"❌ No password_hash found for vendor: {email}")
                    return JsonResponse({
                        'success': False,
                        'message': 'Invalid vendor account configuration'
                    })
                
                # Check password using Django's check_password (supports pbkdf2_sha256 format)
                if check_password(password, password_hash):
                    vendor_id = str(vendor_data.get('vendor_id') or '').strip()
                    vendor_name = vendor_data.get('vendor_name', '')
                    vendor_email_db = (vendor_data.get('email') or vendor_data.get('vendor_email') or email or '').strip()
                    vendor_status = (vendor_data.get('status') or 'pending').strip().lower()
                    
                    # Persist canonical vendor identifiers in the session
                    request.session['vendor_email'] = vendor_email_db
                    request.session['vendor_name'] = vendor_name
                    request.session['vendor_status'] = vendor_status
                    if vendor_id:
                        request.session['vendor_id'] = vendor_id
                    else:
                        request.session.pop('vendor_id', None)
                    
                    # Ensure Vendor_service_availability row exists on login without resetting saved values
                    _sync_vendor_service_on_login(vendor_email_db, vendor_id)
                    
                    # Keep service reminder disabled so existing values are not interrupted
                    request.session['service_update_needed'] = False

                    print(f"✅ Vendor login successful: {email} (ID: {vendor_id}, Status: {vendor_status})")
                    
                    return JsonResponse({
                        'success': True,
                        'message': 'Login successful',
                        'vendor': {
                            'vendor_id': vendor_id,
                            'vendor_name': vendor_name,
                            'email': vendor_email_db,
                            'status': vendor_status
                        }
                    })
                else:
                    return JsonResponse({
                        'success': False,
                        'message': 'Invalid password'
                    })

            except requests.exceptions.RequestException as e:
                print(f"❌ Failed to connect to Worker API: {str(e)}")
                return JsonResponse({
                    'success': False,
                    'message': f'Failed to connect to database service: {str(e)}. Please try again later.'
                }, status=500)
            except Exception as e:
                print(f"❌ Error during vendor login: {str(e)}")
                import traceback
                print(f"❌ Traceback: {traceback.format_exc()}")
                return JsonResponse({
                    'success': False,
                    'message': f'Login error: {str(e)}'
                })

        except Exception as e:
            print(f"❌ Error during vendor login: {str(e)}")
            import traceback
            print(f"❌ Traceback: {traceback.format_exc()}")
            return JsonResponse({
                'success': False,
                'message': f'Login error: {str(e)}'
            })

    return JsonResponse({
        'success': False,
        'message': 'Invalid request method'
    })


# Add vendor registration endpoint
@csrf_exempt
def vendor_register_api(request):
    """
    Handle vendor registration API
    """
    if request.method == 'POST':
        try:
            print(f"🔍 Starting vendor registration process...")
            
            # Parse JSON data with better error handling
            try:
                if not request.body:
                    print(f"❌ Empty request body")
                    return JsonResponse({
                        'success': False,
                        'message': 'Empty request body. Please check your request.'
                    }, status=400)
                
                data = json.loads(request.body)
                print(f"📝 Received registration data: {list(data.keys())}")
            except json.JSONDecodeError as e:
                print(f"❌ JSON decode error: {str(e)}")
                print(f"❌ Request body: {request.body}")
                return JsonResponse({
                    'success': False,
                    'message': 'Invalid JSON data. Please check your request.'
                }, status=400)
            except Exception as e:
                print(f"❌ Error parsing request: {str(e)}")
                return JsonResponse({
                    'success': False,
                    'message': 'Error processing request. Please try again.'
                }, status=400)
            
            # Extract form data
            email = data.get('email', '').strip()
            password = data.get('password', '')
            vendor_name = data.get('vendor_name', '').strip()
            phone_number = data.get('phone_number', '').strip()
            state = data.get('state', '').strip()
            city = data.get('city', '').strip()
            locality = data.get('locality', '').strip()
            shop_address = data.get('shop_address', '').strip()
            full_address = data.get('full_address', '').strip()
            if full_address:
                shop_address = full_address
            pincode = data.get('pincode', '').strip()
            # Support both old (latitude/longitude) and new (vendor_lat/vendor_lng) field names
            latitude = data.get('vendor_lat') or data.get('latitude', '0')
            longitude = data.get('vendor_lng') or data.get('longitude', '0')

            print(f"📧 Processing registration for: {email}")
            print(f"📍 Location data: lat={latitude}, lng={longitude}")

            # Validate location coordinates
            if not latitude or not longitude or latitude == '0' or longitude == '0':
                print(f"❌ Location fetch failed: lat={latitude}, lng={longitude}")
                return JsonResponse({
                    'success': False,
                    'message': 'Location could not be determined. Please try again.'
                }, status=400)
            
            # Validate latitude and longitude are valid numbers
            try:
                lat_float = float(latitude)
                lng_float = float(longitude)
            except (ValueError, TypeError):
                print(f"❌ Invalid coordinate format: lat={latitude}, lng={longitude}")
                return JsonResponse({
                    'success': False,
                    'message': 'Invalid location coordinates. Please refresh the page and try again.'
                }, status=400)
            
            # Validate coordinate ranges
            if not (-90 <= lat_float <= 90) or not (-180 <= lng_float <= 180):
                print(f"❌ Invalid coordinate range: lat={lat_float}, lng={lng_float}")
                return JsonResponse({
                    'success': False,
                    'message': 'Invalid location coordinates. Please refresh the page and try again.'
                }, status=400)

            # Validate required fields
            required_fields = {
                'email': email,
                'password': password,
                'vendor_name': vendor_name,
                'phone_number': phone_number,
                'state': state,
                'city': city,
                'locality': locality,
                'shop_address': shop_address,
                'full_address': shop_address,
                'pincode': pincode
            }
            
            missing_fields = [field for field, value in required_fields.items() if not value]
            if missing_fields:
                print(f"❌ Missing required fields: {missing_fields}")
                return JsonResponse({
                    'success': False,
                    'message': f'Missing required fields: {", ".join(missing_fields)}'
                }, status=400)

            # Validate email format
            email_regex = r'^[^\s@]+@[^\s@]+\.[^\s@]+$'
            if not re.match(email_regex, email):
                print(f"❌ Invalid email format: {email}")
                return JsonResponse({
                    'success': False,
                    'message': 'Please enter a valid email address'
                }, status=400)

            # Validate password strength
            if len(password) < 8:
                print(f"❌ Password too short: {len(password)} characters")
                return JsonResponse({
                    'success': False,
                    'message': 'Password must be at least 8 characters long'
                }, status=400)

            if not re.search(r'[a-zA-Z]', password) or not re.search(r'\d', password):
                print(f"❌ Password doesn't meet complexity requirements")
                return JsonResponse({
                    'success': False,
                    'message': 'Password must contain at least one letter and one number'
                }, status=400)

            # Validate phone number (10 digits)
            phone_clean = re.sub(r'\D', '', phone_number)
            if len(phone_clean) != 10:
                print(f"❌ Invalid phone number: {phone_number} (cleaned: {phone_clean})")
                return JsonResponse({
                    'success': False,
                    'message': 'Please enter a valid 10-digit phone number'
                }, status=400)

            # Validate PIN code (6 digits)
            if not re.match(r'^\d{6}$', pincode):
                print(f"❌ Invalid PIN code: {pincode}")
                return JsonResponse({
                    'success': False,
                    'message': 'Please enter a valid 6-digit PIN code'
                }, status=400)

            # Generate unique 10-digit vendor ID and token
            vendor_id = str(random.randint(1000000000, 9999999999))
            vendor_token = str(random.randint(1000000000, 9999999999))
            print(f"🆔 Generated vendor ID: {vendor_id}")

            # Hash password
            password_hash = make_password(password)

            # Save to D1 database via Worker API (REQUIRED - no fallback to R2)
            api_url = getattr(settings, 'WORKER_API_URL', '')
            api_key = getattr(settings, 'WORKER_API_KEY', '')
            
            if not api_url or not api_key:
                print(f"❌ Worker API not configured")
                return JsonResponse({
                    'success': False,
                    'message': 'Server configuration error: Worker API not configured. Please contact support.'
                }, status=500)

            # Construct the Worker API endpoint for vendor registration
            # Handle both cases: URL with /add-contact or base URL
            if '/add-contact' in api_url:
                worker_endpoint = api_url.replace('/add-contact', '/add-vendor-register')
            else:
                # If it's a base URL, append the endpoint
                worker_endpoint = api_url.rstrip('/') + '/add-vendor-register'
            
            worker_payload = {
                'email': email,
                'password_hash': password_hash,
                'vendor_name': vendor_name,
                'phone_number': phone_number,
                'state': state,
                'city': city,
                'locality': locality,
                'shop_address': shop_address,
                'full_address': shop_address,
                'pincode': pincode,
                # Store coordinates as proper numeric values
                'latitude': lat_float,
                'longitude': lng_float,
                # Also expose vendor_lat/vendor_lng for downstream consumers
                'vendor_lat': lat_float,
                'vendor_lng': lng_float,
                'vendor_id': vendor_id,
                'vendor_token': vendor_token,
                'status': 'pending'
            }
            
            print(f"💾 Saving vendor registration to D1 database via Worker API...")
            print(f"🔗 Worker endpoint: {worker_endpoint}")
            
            try:
                resp = requests.post(
                    worker_endpoint,
                    json=worker_payload,
                    headers={
                        'Content-Type': 'application/json',
                        'x-api-key': api_key
                    },
                    timeout=10
                )
                
                print(f"📡 Worker API Response Status: {resp.status_code}")
                
                if resp.status_code == 200:
                    try:
                        response_data = resp.json()
                        if response_data.get('success'):
                            print(f"✅ Vendor registration saved to D1 database successfully")
                        else:
                            error_msg = response_data.get('error', 'Unknown error from Worker API')
                            print(f"❌ Worker API returned error: {error_msg}")
                            return JsonResponse({
                                'success': False,
                                'message': f'Failed to save registration: {error_msg}'
                            }, status=500)
                    except json.JSONDecodeError:
                        print(f"❌ Invalid JSON response from Worker API: {resp.text[:200]}")
                        return JsonResponse({
                            'success': False,
                            'message': 'Invalid response from database service. Please try again.'
                        }, status=500)
                else:
                    error_text = resp.text[:500] if resp.text else 'No error message'
                    print(f"❌ D1 database save failed with status {resp.status_code}: {error_text}")
                    
                    # Provide helpful error message
                    if resp.status_code == 404:
                        error_msg = 'Database endpoint not found. Please ensure the Worker API is deployed with the /add-vendor-register endpoint.'
                    elif resp.status_code == 401:
                        error_msg = 'Database authentication failed. Please check API key configuration.'
                    else:
                        error_msg = f'Database error (status {resp.status_code}): {error_text}'
                    
                    return JsonResponse({
                        'success': False,
                        'message': error_msg
                    }, status=500)
                    
            except requests.exceptions.RequestException as e:
                print(f"❌ Failed to connect to Worker API: {str(e)}")
                return JsonResponse({
                    'success': False,
                    'message': f'Failed to connect to database service: {str(e)}. Please try again later.'
                }, status=500)
            except Exception as e:
                print(f"❌ Unexpected error saving to D1 database: {str(e)}")
                import traceback
                print(f"❌ Traceback: {traceback.format_exc()}")
                return JsonResponse({
                    'success': False,
                    'message': f'Unexpected error occurred while saving registration. Please try again.'
                }, status=500)

            # Registration saved to D1 database successfully above
            # Now proceed with non-critical operations (email, etc.)
            
            email_sent = False
            try:  # Welcome email (fire-and-forget to avoid blocking request)
                print(f"📧 Sending welcome email to: {email}")
                threading.Thread(target=send_welcome_email, args=(email, vendor_name, password, vendor_id), daemon=True).start()
                email_sent = True
            except Exception as e:
                print(f"⚠️ Welcome email failed to dispatch: {e}")

            # Optional: send a one-time SMTP test email if configured
            try:
                test_to = getattr(settings, 'EMAIL_TEST_TO', None)
                if test_to:
                    from django.core.mail import send_mail as _send_mail
                    threading.Thread(
                        target=_send_mail,
                        args=(
                            'PrintMax SMTP Test',
                            'This is a test email sent after vendor registration to verify SMTP.',
                            settings.DEFAULT_FROM_EMAIL,
                            [test_to],
                        ),
                        kwargs={
                            'fail_silently': True
                        },
                        daemon=True
                    ).start()
                    print(f"📨 SMTP test email queued to {test_to}")
            except Exception as e:
                print(f"⚠️ Could not queue SMTP test email: {e}")

            shop_folder_name = sanitize_shop_name(vendor_name)

            print(f"🎉 Registration completed successfully for {email}")

            # Persist newly registered vendor in session so downstream steps don't need URL params
            request.session['vendor_email'] = email
            request.session['vendor_name'] = vendor_name
            request.session['vendor_id'] = vendor_id
            request.session['vendor_token'] = vendor_token
            request.session.modified = True

            return JsonResponse({
                'success': True,
                'message': 'Registration successful' + (" (Welcome email sent)" if email_sent else " (Welcome email will be sent later)"),
                'vendor_email': email,
                'vendor_id': vendor_id,
                'vendor_token': vendor_token,
                'shop_folder': shop_folder_name
            })

        except Exception as e:
            print(f"❌ Unexpected error during vendor registration: {str(e)}")
            print(f"❌ Error type: {type(e).__name__}")
            print(f"❌ Error details: {str(e)}")
            import traceback
            print(f"❌ Traceback: {traceback.format_exc()}")
            
            # Check if this is a critical error that prevents registration
            error_message = str(e).lower()
            if any(keyword in error_message for keyword in ['json', 'decode', 'parse', 'invalid', 'malformed']):
                return JsonResponse({
                    'success': False,
                    'message': 'Invalid request data. Please check your information and try again.'
                }, status=400)
            else:
                return JsonResponse({
                    'success': False,
                    'message': f'Registration failed due to a server error: {str(e)}'
                }, status=500)

    return JsonResponse({
        'success': False,
        'message': 'Invalid request method'
    })

@csrf_exempt
def vendor_authenticate(request):
    """
    Authenticate vendor using vendor_id and vendor_token (hashed in shop_info.json)
    """
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_email = data.get('vendor_email')
            vendor_id = data.get('vendor_id')
            vendor_token = data.get('vendor_token')
            shop_name = data.get('shop_name')

            if not all([vendor_email, vendor_id, vendor_token, shop_name]):
                return JsonResponse({'success': False, 'error': 'Missing credentials'}, status=400)

            s3 = boto3.client('s3',
                aws_access_key_id=settings.R2_ACCESS_KEY,
                aws_secret_access_key=settings.R2_SECRET_KEY,
                endpoint_url=settings.R2_ENDPOINT,
                region_name='auto'
            )
            shop_folder = sanitize_shop_name(shop_name)
            shop_info_key = f'vendor_register_details/{sanitize_email(vendor_email)}/{shop_folder}/shop_info.json'
            try:
                response = s3.get_object(Bucket=settings.R2_BUCKET, Key=shop_info_key)
                shop_info = json.loads(response['Body'].read().decode('utf-8'))
                vendor_id_hash = shop_info.get('vendor_id_hash')
                vendor_token_hash = shop_info.get('vendor_token_hash')
                if check_password(vendor_id, vendor_id_hash) and check_password(vendor_token, vendor_token_hash):
                    # Update authentication timestamp for 8-hour tracking
                    update_vendor_auth_timestamp(vendor_id)
                    return JsonResponse({'success': True, 'message': 'Authenticated'})
                else:
                    return JsonResponse({'success': False, 'error': 'Invalid credentials'}, status=401)
            except Exception as e:
                return JsonResponse({'success': False, 'error': f'Shop info not found: {str(e)}'}, status=404)
        except Exception as e:
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)

def sanitize_email(email):
    # Lowercase, replace @ with _at_, . with _dot_, and remove other special chars
    return re.sub(r'[^a-zA-Z0-9_]', '', email.lower().replace('@', '_at_').replace('.', '_dot_'))

def sanitize_shop_name(shop_name):
    # Convert to lowercase, replace spaces with underscores, remove special chars except underscores
    sanitized = re.sub(r'[^a-zA-Z0-9_\s]', '', shop_name.lower())
    sanitized = re.sub(r'\s+', '_', sanitized.strip())
    return sanitized


def _categorize_vendor_pricing_row(row_dict):
    """
    Build a categorized_pricing structure similar to the one returned by the Worker API
    using a row from the Vendor_pricing table.
    """
    if not row_dict:
        return {}

    prefix_rules = {
        'digital_print': ('digital_print_',),
        'regular_print': ('regular_print_',),
        'photo_print': ('photo_print_',),
        'gloss_print': ('gloss_print_',),
        'jumbo_print': ('jumbo_print_',),
        'passport_photo': ('passport_print_',),
        'golden_embossing': ('golden_emboss_',),
        'lamination': ('lamination_',),
        'binding': ('tape_binding_', 'spiral_binding_'),
    }

    categorized = {}
    for category, prefixes in prefix_rules.items():
        category_data = {}
        for key, value in row_dict.items():
            if value in (None, ''):
                continue
            if any(key.startswith(prefix) for prefix in prefixes):
                category_data[key] = value
        if category_data:
            categorized[category] = category_data

    # Create an a4_print alias so existing UI code gets the expected keys
    regular = categorized.get('regular_print', {})
    a4_alias = {}
    bw_price = regular.get('regular_print_a4_bw')
    color_price = regular.get('regular_print_a4_color')
    if bw_price is not None:
        a4_alias['a4_print_single_bw'] = bw_price
    if color_price is not None:
        a4_alias['a4_print_single_color'] = color_price
    if a4_alias:
        categorized['a4_print'] = a4_alias

    return categorized


def normalize_color_key(value, default='color'):
    """
    Normalize color strings coming from various parts of the product UI.
    """
    normalized = (value or '').strip().lower()
    color_aliases = {
        'single_color': 'color',
        'colour': 'color',
        'color': 'color',
        'colour_print': 'color',
        'color_print': 'color',
        'single_colour': 'color',
        'single_bw': 'bw',
        'single_black_white': 'bw',
        'bw': 'bw',
        'black_white': 'bw',
        'black & white': 'bw',
        'black and white': 'bw',
        'mono': 'bw',
    }
    return color_aliases.get(normalized, normalized or default)


def safe_price(value, default=0.0):
    """
    Convert pricing values from DB/API to float safely.
    """
    try:
        if value in (None, ''):
            return default
        return float(value)
    except (TypeError, ValueError):
        return default


def get_vendor_pricing_from_local_db(vendor_email):
    """
    Attempt to read vendor pricing from the local Vendor_pricing table when the Worker API
    is unreachable. This keeps pricing-dependent flows functional for on-prem deployments.
    """
    if not vendor_email:
        return None

    try:
        with connection.cursor() as cursor:
            cursor.execute(
                """
                SELECT *
                FROM Vendor_pricing
                WHERE LOWER(vendor_email) = LOWER(%s)
                ORDER BY COALESCE(last_updated, '') DESC, id DESC
                LIMIT 1
                """,
                [vendor_email]
            )
            row = cursor.fetchone()
            if not row:
                return None

            columns = [col[0] for col in cursor.description]
            row_dict = {columns[idx]: value for idx, value in enumerate(row)}
            categorized = _categorize_vendor_pricing_row(row_dict)
            return {
                'pricing_data': row_dict,
                'categorized_pricing': categorized,
                'services_summary': {}
            }
    except Exception as exc:
        print(f"⚠️ Unable to fetch vendor pricing from local DB for {vendor_email}: {exc}")

    return None

def get_vendor_pricing_from_d1(vendor_email):
    """
    Helper function to get vendor pricing from d1 database via Worker API
    Returns pricing data dict with categorized_pricing, or None if failed
    """
    api_url = getattr(settings, 'WORKER_API_URL', '')
    api_key = getattr(settings, 'WORKER_API_KEY', '')
    
    if not api_url or not api_key:
        print(f"⚠️ Worker API not configured - cannot fetch pricing from d1 database")
        return None
    
    # Construct the Worker API endpoint for getting vendor pricing
    if '/add-contact' in api_url:
        worker_endpoint = api_url.replace('/add-contact', '/get-vendor-pricing')
    elif '/add-vendor-register' in api_url:
        worker_endpoint = api_url.replace('/add-vendor-register', '/get-vendor-pricing')
    elif '/add-vendor-pricing' in api_url:
        worker_endpoint = api_url.replace('/add-vendor-pricing', '/get-vendor-pricing')
    else:
        worker_endpoint = api_url.rstrip('/') + '/get-vendor-pricing'
    
    # Add vendor_email as query parameter (URL encode it)
    from urllib.parse import quote
    worker_endpoint = f"{worker_endpoint}?vendor_email={quote(vendor_email)}"
    
    try:
        resp = requests.get(
            worker_endpoint,
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=15
        )
        
        if resp.status_code == 200:
            response_data = resp.json()
            if response_data.get('success'):
                print(f"✅ Successfully fetched vendor pricing from d1 database for {vendor_email}")
                # Return pricing data in the format expected by pricing calculation functions
                pricing_data = {
                    'pricing_data': response_data.get('pricing', {}),
                    'categorized_pricing': response_data.get('categorized_pricing', {}),
                    'services_summary': response_data.get('services_summary', {})
                }
                return pricing_data
            else:
                error_msg = response_data.get('error', 'Unknown error from Worker API')
                print(f"⚠️ Worker API returned error: {error_msg}")
        else:
            error_text = resp.text[:500] if resp.text else 'No error message'
            print(f"⚠️ Failed to fetch pricing from d1 database with status {resp.status_code}: {error_text}")
    except requests.exceptions.RequestException as e:
        print(f"⚠️ Failed to connect to Worker API: {str(e)}")
    except Exception as e:
        print(f"⚠️ Unexpected error fetching pricing from d1 database: {str(e)}")
    
    # Fallback to local Vendor_pricing table so pricing-dependent flows still work
    local_pricing = get_vendor_pricing_from_local_db(vendor_email)
    if local_pricing:
        print(f"ℹ️ Falling back to local Vendor_pricing table for {vendor_email}")
        return local_pricing
    
    return None

@csrf_exempt
def get_vendor_pricing(request):
    """
    Get vendor pricing data from d1 database via Worker API
    """
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_email = data.get('vendor_email')
            
            if not vendor_email:
                return JsonResponse({'success': False, 'error': 'Vendor email required'})

            # Try to fetch from d1 database via Worker API
            api_url = getattr(settings, 'WORKER_API_URL', '')
            api_key = getattr(settings, 'WORKER_API_KEY', '')
            
            if api_url and api_key:
                # Construct the Worker API endpoint for getting vendor pricing
                if '/add-contact' in api_url:
                    worker_endpoint = api_url.replace('/add-contact', '/get-vendor-pricing')
                elif '/add-vendor-register' in api_url:
                    worker_endpoint = api_url.replace('/add-vendor-register', '/get-vendor-pricing')
                elif '/add-vendor-pricing' in api_url:
                    worker_endpoint = api_url.replace('/add-vendor-pricing', '/get-vendor-pricing')
                else:
                    worker_endpoint = api_url.rstrip('/') + '/get-vendor-pricing'
                
                # Add vendor_email as query parameter (URL encode it)
                from urllib.parse import quote
                worker_endpoint = f"{worker_endpoint}?vendor_email={quote(vendor_email)}"
                
                try:
                    resp = requests.get(
                        worker_endpoint,
                        headers={
                            'Content-Type': 'application/json',
                            'x-api-key': api_key
                        },
                        timeout=15
                    )
                    
                    if resp.status_code == 200:
                        response_data = resp.json()
                        if response_data.get('success'):
                            print(f"✅ Successfully fetched vendor pricing from d1 database for {vendor_email}")
                            # Combine pricing and categorized_pricing into a single object for frontend compatibility
                            pricing_obj = response_data.get('pricing', {})
                            pricing_obj['categorized_pricing'] = response_data.get('categorized_pricing', {})
                            pricing_obj['services_summary'] = response_data.get('services_summary', {})
                            return JsonResponse({
                                'success': True,
                                'pricing': pricing_obj,
                                'categorized_pricing': response_data.get('categorized_pricing', {}),
                                'services_summary': response_data.get('services_summary', {})
                            })
                        else:
                            error_msg = response_data.get('error', 'Unknown error from Worker API')
                            print(f"⚠️ Worker API returned error: {error_msg}")
                    else:
                        error_text = resp.text[:500] if resp.text else 'No error message'
                        print(f"⚠️ Failed to fetch pricing from d1 database with status {resp.status_code}: {error_text}")
                except requests.exceptions.RequestException as e:
                    print(f"⚠️ Failed to connect to Worker API: {str(e)}")
                except Exception as e:
                    print(f"⚠️ Unexpected error fetching from d1 database: {str(e)}")
            else:
                print(f"⚠️ Worker API not configured - cannot fetch pricing from d1 database")
            
            # Fallback: Return default pricing if d1 fetch fails
            print(f"⚠️ Falling back to default pricing for {vendor_email}")
            # Return default pricing if vendor pricing not found
            default_pricing = {
                'digital_print_a4_color': 5,
                'digital_print_a3_color': 8,
                'digital_print_12x18_color': 12,
                'digital_print_a2_color': 15,
                'digital_print_a1_color': 22,
                'digital_print_a0_color': 30,
                'regular_print_a4_bw': 2,
                'regular_print_a4_color': 5,
                'photo_print_a4_bw': 6,
                'photo_print_a4_color': 12,
                'gloss_print_a4_color': 8,
                'gloss_print_a3_color': 12,
                'gloss_print_a2_color': 18,
                'gloss_print_a1_color': 24,
                'gloss_print_a0_color': 32,
                'jumbo_print_a3_bw': 10,
                'jumbo_print_a3_color': 16,
                'jumbo_print_a2_bw': 18,
                'jumbo_print_a2_color': 24,
                'jumbo_print_a1_bw': 26,
                'jumbo_print_a1_color': 32,
                'jumbo_print_a0_bw': 34,
                'jumbo_print_a0_color': 40,
                'passport_print_8': 40,
                'passport_print_16': 70,
                'passport_print_30': 120,
                'golden_emboss_cover': 50,
                'golden_emboss_a4_color': 10,
                'golden_emboss_bond_color': 10,
                'lamination_a4_standard': 30,
                'lamination_a4_glossy': 35,
                'lamination_a3_standard': 45,
                'lamination_a3_glossy': 55,
                'lamination_a2_standard': 65,
                'lamination_a2_glossy': 75,
                'lamination_a1_standard': 85,
                'lamination_a1_glossy': 95,
                'lamination_a0_standard': 105,
                'lamination_a0_glossy': 115,
                'tape_binding_a4_100': 40,
                'tape_binding_a4_200': 60,
                'tape_binding_a3_100': 70,
                'tape_binding_a3_200': 90,
                'spiral_binding_a4_100': 45,
                'spiral_binding_a4_200': 65,
                'spiral_binding_a3_100': 75,
                'spiral_binding_a3_200': 95
            }
            
            # Create categorized default pricing
            default_categorized = {
                'digital_print': {
                    'digital_print_a4_color': 5,
                    'digital_print_a3_color': 8,
                    'digital_print_12x18_color': 12,
                    'digital_print_a2_color': 15,
                    'digital_print_a1_color': 22,
                    'digital_print_a0_color': 30
                },
                'a4_print': {
                    'regular_print_a4_bw': 2,
                    'regular_print_a4_color': 5
                },
                'photo_print': {
                    'photo_print_a4_bw': 6,
                    'photo_print_a4_color': 12
                },
                'gloss_print': {
                    'gloss_print_a4_color': 8,
                    'gloss_print_a3_color': 12,
                    'gloss_print_a2_color': 18,
                    'gloss_print_a1_color': 24,
                    'gloss_print_a0_color': 32
                },
                'jumbo_print': {
                    'jumbo_print_a3_bw': 10,
                    'jumbo_print_a3_color': 16,
                    'jumbo_print_a2_bw': 18,
                    'jumbo_print_a2_color': 24,
                    'jumbo_print_a1_bw': 26,
                    'jumbo_print_a1_color': 32,
                    'jumbo_print_a0_bw': 34,
                    'jumbo_print_a0_color': 40
                },
                'passport_photo': {
                    'passport_print_8': 40,
                    'passport_print_16': 70,
                    'passport_print_30': 120
                },
                'golden_embossing': {
                    'golden_emboss_cover': 50,
                    'golden_emboss_a4_color': 10,
                    'golden_emboss_bond_color': 10
                },
                'lamination': {
                    'lamination_a4_standard': 30,
                    'lamination_a4_glossy': 35,
                    'lamination_a3_standard': 45,
                    'lamination_a3_glossy': 55,
                    'lamination_a2_standard': 65,
                    'lamination_a2_glossy': 75,
                    'lamination_a1_standard': 85,
                    'lamination_a1_glossy': 95,
                    'lamination_a0_standard': 105,
                    'lamination_a0_glossy': 115
                },
                'binding': {
                    'tape_binding_a4_100': 40,
                    'tape_binding_a4_200': 60,
                    'tape_binding_a3_100': 70,
                    'tape_binding_a3_200': 90,
                    'spiral_binding_a4_100': 45,
                    'spiral_binding_a4_200': 65,
                    'spiral_binding_a3_100': 75,
                    'spiral_binding_a3_200': 95
                }
            }
            
            # Combine pricing and categorized_pricing into a single object for frontend compatibility
            default_pricing['categorized_pricing'] = default_categorized
            default_pricing['services_summary'] = {
                'total_services': len(default_pricing),
                'available_services_count': len(default_pricing),
                'not_available_services_count': 0
            }
            
            return JsonResponse({
                'success': True,
                'pricing': default_pricing,
                'categorized_pricing': default_categorized,
                'services_summary': {
                    'total_services': len(default_pricing),
                    'available_services_count': len(default_pricing),
                    'not_available_services_count': 0
                }
            })
                
        except Exception as e:
            print(f"Error in get_vendor_pricing: {str(e)}")
            return JsonResponse({'success': False, 'error': str(e)})
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'})
@csrf_exempt
def calculate_gloss_print_pricing(request):
    """Calculate pricing for gloss print service based on vendor pricing.json"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_email = data.get('vendor_email')
            print_size = data.get('print_size', 'A4')
            gloss_type = data.get('gloss_type', 'standard_gloss')
            print_color = data.get('print_color', 'Color')
            page_count = data.get('page_count', 1)
            num_copies = data.get('num_copies', 1)
            
            print(f"Received gloss print data: vendor_email={vendor_email}, print_size={print_size}, gloss_type={gloss_type}, print_color={print_color}, page_count={page_count}, num_copies={num_copies}")
            print(f"Page count type: {type(page_count)}, value: {page_count}")
            print(f"Num copies type: {type(num_copies)}, value: {num_copies}")
            
            # Convert page_count to int
            try:
                page_count = int(page_count)
                print(f"Page count after conversion: {page_count}")
            except (ValueError, TypeError):
                page_count = 1
                print("Invalid page count, using default: 1")
            
            # Ensure page_count is at least 1
            if page_count < 1:
                page_count = 1
            
            # Convert num_copies to int
            try:
                num_copies = int(num_copies)
                print(f"Num copies after conversion: {num_copies}")
            except (ValueError, TypeError):
                num_copies = 1
                print("Invalid num copies, using default: 1")
            
            # Ensure num_copies is at least 1
            if num_copies < 1:
                num_copies = 1
            
            # Get vendor pricing from d1 database
            print(f"🔍 Gloss print - Vendor email: {vendor_email}")
            pricing_data = get_vendor_pricing_from_d1(vendor_email)
            
            if pricing_data is None:
                print(f"❌ Could not load vendor pricing from d1 database for vendor: {vendor_email}")
                return JsonResponse({
                    'success': False,
                    'error': f'Unable to load pricing data for vendor. Please contact the vendor.'
                })
            
            # Construct pricing key based on user selections
            size_key = (print_size or 'A4').strip().lower()
            size_candidates = [size_key]
            if size_key == 'a1':
                size_candidates.append('al')
            
            gloss_type_key = (gloss_type or 'standard_gloss').strip().lower()
            color_key = normalize_color_key(print_color, default='color')
            
            # For gloss printing, the pricing is nested under "categorized_pricing.gloss_print" object
            # Access the nested gloss_print pricing data
            categorized_pricing = pricing_data.get('categorized_pricing', {})
            print(f"📊 Categorized pricing keys: {list(categorized_pricing.keys())}")
            
            gloss_pricing = categorized_pricing.get('gloss_print', {})
            print(f"📊 Gloss pricing keys: {list(gloss_pricing.keys())}")
            
            # Construct candidate pricing keys to support legacy schemas
            lookup_candidates = []
            for size_variant in size_candidates:
                lookup_candidates.extend([
                    f'gloss_print_{size_variant}_{gloss_type_key}_{color_key}',
                    f'gloss_print_{size_variant}_{gloss_type_key}',
                    f'gloss_print_{size_variant}_{color_key}',
                    f'gloss_print_{size_variant}'
                ])
            base_price, resolved_key = resolve_pricing_key(gloss_pricing, *lookup_candidates)
            
            # Check if pricing is available
            if base_price in (None, ''):
                available_keys = list(gloss_pricing.keys())
                print(f"❌ Pricing not found for keys: {lookup_candidates}")
                print(f"Available gloss pricing keys: {available_keys}")
                return JsonResponse({
                    'success': False,
                    'error': f'Pricing not available for {print_size} {gloss_type}. Please ask the vendor to configure gloss pricing in Dashboard.'
                })
            
            print(f"✅ Using gloss pricing key: {resolved_key}")
            base_price = safe_price(base_price, default=0.0)
            
            # Calculate price per page
            price_per_page = base_price
            
            # Calculate total price (price per page * number of pages * number of copies)
            total_price = price_per_page * page_count * num_copies
            
            print(f"Calculation: base_price={base_price}, price_per_page={price_per_page}, page_count={page_count}, num_copies={num_copies}, total_price={total_price}")
            
            # Prepare pricing breakdown
            pricing_breakdown = {
                'base_price': base_price,
                'price_per_page': price_per_page,
                'page_count': page_count,
                'num_copies': num_copies,
                'total_price': total_price,
                'pricing_key_used': resolved_key or lookup_candidates[0]
            }
            
            return JsonResponse({
                'success': True,
                'pricing_breakdown': pricing_breakdown,
                'total_price': total_price
            })
            
        except Exception as e:
            print(f"Error calculating gloss print pricing: {e}")
            traceback.print_exc()
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    
    return JsonResponse({
        'success': False,
        'error': 'Invalid request method'
    })
@csrf_exempt
def calculate_golden_emboss_pricing(request):
    """Calculate pricing for golden embossing service based on vendor pricing.json"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_email = data.get('vendor_email')
            print_color = data.get('print_color', 'Color')
            paper_type = data.get('paper_type', 'A4')
            num_copies = data.get('num_copies', 1)
            page_count = data.get('page_count', 1)
            
            print(f"Received golden emboss data: vendor_email={vendor_email}, print_color={print_color}, paper_type={paper_type}, num_copies={num_copies}, page_count={page_count}")
            
            # Convert page_count to int
            try:
                page_count = int(page_count)
                print(f"Page count after conversion: {page_count}")
            except (ValueError, TypeError):
                page_count = 1
                print("Invalid page count, using default: 1")
            
            # Ensure page_count is at least 1
            if page_count < 1:
                page_count = 1
            
            # Convert num_copies to int
            try:
                num_copies = int(num_copies)
                print(f"Num copies after conversion: {num_copies}")
            except (ValueError, TypeError):
                num_copies = 1
                print("Invalid num copies, using default: 1")
            
            # Ensure num_copies is at least 1
            if num_copies < 1:
                num_copies = 1
            
            # Get vendor pricing from d1 database
            print(f"🔍 Golden emboss - Vendor email: {vendor_email}")
            pricing_data = get_vendor_pricing_from_d1(vendor_email)
            
            if pricing_data is None:
                print(f"❌ Could not load vendor pricing from d1 database for vendor: {vendor_email}")
                return JsonResponse({
                    'success': False,
                    'error': f'Unable to load pricing data for vendor. Please contact the vendor.'
                })
            
            # Get the categorized pricing for golden embossing
            categorized_pricing = pricing_data.get('categorized_pricing', {})
            print(f"📊 Categorized pricing keys: {list(categorized_pricing.keys())}")
            
            golden_emboss_pricing = categorized_pricing.get('golden_embossing', {})
            print(f"📊 Golden emboss pricing keys: {list(golden_emboss_pricing.keys())}")
            
            cover_price = safe_price(golden_emboss_pricing.get('golden_emboss_cover'))
            a4_color_price = safe_price(golden_emboss_pricing.get('golden_emboss_a4_color'))
            bond_color_price = safe_price(golden_emboss_pricing.get('golden_emboss_bond_color'))

            if cover_price <= 0 and a4_color_price <= 0 and bond_color_price <= 0:
                print("❌ Golden emboss pricing table missing cover/A4/bond entries")
                return JsonResponse({
                    'success': False,
                    'error': 'Golden emboss pricing not configured for this vendor. Please ask the vendor to set cover and paper rates.'
                })

            normalized_paper_type = str(paper_type or 'A4').strip().lower()
            if normalized_paper_type in ('bond', 'bond paper', 'bond_paper'):
                selected_paper_type = 'Bond'
                selected_pricing_key = 'golden_emboss_bond_color'
                selected_per_page_price = bond_color_price
                selected_label = 'Bond Color'
            else:
                selected_paper_type = 'A4'
                selected_pricing_key = 'golden_emboss_a4_color'
                selected_per_page_price = a4_color_price
                selected_label = 'A4 Color'

            # Fallback safely to the other paper key when the selected one is not configured.
            if selected_per_page_price <= 0:
                fallback_price = bond_color_price if selected_paper_type == 'A4' else a4_color_price
                if fallback_price > 0:
                    print(
                        f"⚠️ {selected_pricing_key} not configured, falling back to "
                        f"{'golden_emboss_bond_color' if selected_paper_type == 'A4' else 'golden_emboss_a4_color'}"
                    )
                    selected_per_page_price = fallback_price
                else:
                    return JsonResponse({
                        'success': False,
                        'error': f'{selected_paper_type} pricing is not configured for this vendor.'
                    })

            cover_cost_per_book = max(cover_price, 0.0)
            color_cost_per_book = max(selected_per_page_price, 0.0) * page_count
            print(
                f"Golden emboss cover cost: {cover_cost_per_book}, "
                f"{selected_label} charge per book: {color_cost_per_book}"
            )
            
            total_price_per_book = cover_cost_per_book + color_cost_per_book
            if total_price_per_book <= 0:
                return JsonResponse({
                    'success': False,
                    'error': 'Golden emboss pricing resulted in zero total. Please ask vendor to configure valid rates.'
                })
            
            # Calculate total price for all copies
            total_price = total_price_per_book * num_copies
            print(f"Total for {num_copies} copies: {total_price_per_book} * {num_copies} = {total_price}")
            
            # Prepare pricing breakdown
            pricing_breakdown = [
                {
                    'label': 'Golden Emboss Cover (per book)',
                    'value': f'₹{cover_cost_per_book:.2f}'
                },
                {
                    'label': f'{selected_label} ({page_count} pages)',
                    'value': f'₹{color_cost_per_book:.2f}'
                }
            ]
            
            pricing_breakdown.append({
                'label': 'Total per book',
                'value': f'₹{total_price_per_book:.2f}'
            })
            pricing_breakdown.append({
                'label': 'Number of copies',
                'value': f'{num_copies}'
            })
            
            # Also provide structured data for metadata storage
            structured_breakdown = {
                'pricing_breakdown': pricing_breakdown,
                'total_price': total_price,
                'price_per_page': selected_per_page_price,
                'page_count': page_count,
                'num_copies': num_copies,
                'paper_type': selected_paper_type,
                'emboss_cost': cover_cost_per_book,
                'paper_cost': color_cost_per_book,
                'total_per_book': total_price_per_book,
                'total_pages': page_count * num_copies,  # Total pages across all copies
                'pricing_key_used': selected_pricing_key
            }
            
            return JsonResponse({
                'success': True,
                'pricing_breakdown': pricing_breakdown,
                'total_price': total_price,
                'structured_data': structured_breakdown
            })
            
        except Exception as e:
            print(f"Error calculating golden emboss pricing: {e}")
            traceback.print_exc()
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    
    return JsonResponse({
        'success': False,
        'error': 'Invalid request method'
    })

def resolve_pricing_key(pricing_dict, *candidates):
    """
    Helper to find a pricing value using multiple key variations (case-insensitive).
    Returns tuple of (value, matched_key) or (None, None).
    """
    if not pricing_dict:
        return None, None

    for candidate in candidates:
        if not candidate:
            continue
        value = pricing_dict.get(candidate)
        if value in (None, ''):
            lower_candidate = candidate.lower()
            for existing_key, existing_value in pricing_dict.items():
                if existing_key.lower() == lower_candidate and existing_value not in (None, ''):
                    return existing_value, existing_key
        else:
            return value, candidate
    return None, None

@csrf_exempt
def calculate_photo_print_pricing(request):
    """Calculate pricing for photo print service based on vendor pricing.json"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_email = data.get('vendor_email')
            photo_count = data.get('photo_count', 1)
            layout_slots = data.get('layout_slots', 1)
            orientation = data.get('orientation', 'portrait')
            color = data.get('color', 'Color')
            copies = data.get('copies', 1)
            
            print(f"Received photo print data: vendor_email={vendor_email}, photo_count={photo_count}, layout_slots={layout_slots}, orientation={orientation}, color={color}, copies={copies}")
            
            # Convert values to int
            try:
                photo_count = int(photo_count)
                layout_slots = int(layout_slots)
                copies = int(copies)
            except (ValueError, TypeError):
                photo_count = 1
                layout_slots = 1
                copies = 1
            
            # Ensure values are at least 1
            if photo_count < 1: photo_count = 1
            if layout_slots < 1: layout_slots = 1
            if copies < 1: copies = 1
            
            # Get vendor pricing from d1 database
            print(f"🔍 Photo print - Vendor email: {vendor_email}")
            pricing_data = get_vendor_pricing_from_d1(vendor_email)
            
            if pricing_data is None:
                print(f"❌ Could not load vendor pricing from d1 database for vendor: {vendor_email}")
                return JsonResponse({
                    'success': False,
                    'error': f'Unable to load pricing data for vendor. Please contact the vendor.'
                })
            
            # For photo printing, the pricing is nested under "categorized_pricing.photo_print" object
            # Access the nested photo_print pricing data
            categorized_pricing = pricing_data.get('categorized_pricing', {})
            print(f"📊 Categorized pricing keys: {list(categorized_pricing.keys())}")
            
            photo_pricing = categorized_pricing.get('photo_print', {})
            print(f"📊 Photo pricing keys: {list(photo_pricing.keys())}")
            
            # Map color to the correct pricing key format
            # The pricing.json has: photo_print_a4_bw, photo_print_a4_color
            color = data.get('color', 'Color')
            if color == 'Color' or color == 'color':
                pricing_key_name = 'photo_print_a4_color'
            else:
                pricing_key_name = 'photo_print_a4_bw'
            
            print(f"🔍 Looking for pricing key: {pricing_key_name}")
            
            # Get base price from the nested photo_print object
            base_price = photo_pricing.get(pricing_key_name)
            
            # Check if pricing is available
            if base_price is None:
                print(f"❌ Pricing not found for key: {pricing_key_name}")
                print(f"Available photo pricing keys: {list(photo_pricing.keys())}")
                return JsonResponse({
                    'success': False,
                    'error': f'Pricing not available for photo print. Please contact the vendor.'
                })
            
            # Calculate price per layout (1 layout = 1 A4 page)
            price_per_layout = base_price
            
            # Calculate total price (price per layout * number of copies)
            total_price = price_per_layout * copies
            
            # Calculate platform profit (20% commission for photo print, same as A4 print)
            platform_commission_percent = 20  # 20% commission for photo print
            platform_profit = (total_price * platform_commission_percent) / 100
            
            # Calculate final amount (total_price + platform_profit)
            final_amount = total_price + platform_profit
            
            print(f"Calculation: base_price={base_price}, price_per_layout={price_per_layout}, copies={copies}, total_price={total_price}, platform_profit={platform_profit}, final_amount={final_amount}")
            
            # Prepare pricing breakdown
            pricing_breakdown = {
                'base_price': base_price,
                'price_per_layout': price_per_layout,
                'photo_count': photo_count,
                'layout_slots': layout_slots,
                'copies': copies,
                'total_price': total_price,
                'total_bill_amount': total_price,  # Same as total_price before commission
                'platform_profit': round(platform_profit, 2),
                'platform_commission': platform_commission_percent,
                'final_amount': round(final_amount, 2),
                'pricing_key_used': pricing_key_name,
                'total_pages': copies  # For photo print, total pages = number of copies (layouts)
            }
            
            return JsonResponse({
                'success': True,
                'pricing_breakdown': pricing_breakdown,
                'total_price': total_price
            })
            
        except Exception as e:
            print(f"Error calculating photo print pricing: {e}")
            traceback.print_exc()
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    
    return JsonResponse({
        'success': False,
        'error': 'Invalid request method'
    })

@csrf_exempt
def calculate_digital_print_pricing(request):
    """
    Calculate digital print pricing based on user selections and vendor pricing
    """
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_email = data.get('vendor_email')
            paper_size = data.get('paper_size')
            print_color = data.get('print_color')
            num_copies = data.get('num_copies', 1)
            print_quality = data.get('print_quality')
            page_count = data.get('page_count', 1)
            
            print(f"Received data: vendor_email={vendor_email}, paper_size={paper_size}, print_color={print_color}, num_copies={num_copies}, print_quality={print_quality}, page_count={page_count}")
            print(f"Page count type: {type(page_count)}, value: {page_count}")
            
            # Ensure page_count is an integer
            try:
                page_count = int(page_count)
                print(f"Page count after conversion: {page_count}")
            except (ValueError, TypeError):
                print(f"Error converting page_count to int, using default: {page_count}")
                page_count = 1
            
            # Ensure page_count is at least 1
            if page_count < 1:
                print(f"Page count was {page_count}, setting to 1")
                page_count = 1
            
            if not all([vendor_email, paper_size, print_color, print_quality]):
                return JsonResponse({'success': False, 'error': 'Missing required parameters'})

            # Get vendor pricing from d1 database
            pricing_data = get_vendor_pricing_from_d1(vendor_email)
            
            if pricing_data is None:
                print(f"❌ Could not load vendor pricing from d1 database for vendor: {vendor_email}")
                return JsonResponse({
                    'success': False,
                    'error': f'Unable to load pricing data for vendor. Please contact the vendor.'
                })
            
            # Get digital print pricing from categorized pricing
            digital_pricing = pricing_data.get('categorized_pricing', {}).get('digital_print', {})
            
            # Construct the pricing key based on user selections
            paper_size_key = (paper_size or 'A4').strip().lower()
            size_candidates = [paper_size_key]
            if paper_size_key == 'a1':
                size_candidates.append('al')
            
            color_key_raw = (print_color or '').strip().lower()
            color_key_normalized = normalize_color_key(print_color, default='color')
            color_candidates = []
            for candidate in [color_key_raw, color_key_normalized, 'color']:
                if candidate:
                    color_candidates.append(candidate)
            if color_key_normalized == 'color':
                color_candidates.append('single_color')
            elif color_key_normalized == 'bw':
                color_candidates.append('single_bw')
            # Preserve order but remove duplicates
            seen = set()
            color_candidates = [c for c in color_candidates if not (c in seen or seen.add(c))]
            
            lookup_candidates = []
            for size_variant in size_candidates:
                for color_variant in color_candidates:
                    lookup_candidates.append(f"digital_print_{size_variant}_{color_variant}")
                lookup_candidates.append(f"digital_print_{size_variant}")
            base_price, resolved_key = resolve_pricing_key(digital_pricing, *lookup_candidates)
            
            if base_price in (None, ''):
                print(f"❌ Digital pricing not found for keys: {lookup_candidates}")
                print(f"Available digital pricing keys: {list(digital_pricing.keys())}")
                return JsonResponse({
                    'success': False,
                    'error': f'Pricing not available for {paper_size} {print_color}. Please contact the vendor.'
                })
            
            base_price = safe_price(base_price, default=0.0)
            
            # Apply quality upgrade if high quality is selected
            quality_upgrade = 0
            if print_quality == 'high_quality':
                quality_upgrade = safe_price(digital_pricing.get('digital_print_high_quality', 0), default=0.0)
            
            # Calculate total price (price per page * number of pages * number of copies)
            price_per_page = base_price + quality_upgrade
            total_price = price_per_page * page_count * num_copies
            
            print(f"Calculation: base_price={base_price}, quality_upgrade={quality_upgrade}, price_per_page={price_per_page}, page_count={page_count}, num_copies={num_copies}, total_price={total_price}")
            
            # Prepare pricing breakdown
            pricing_breakdown = {
                'base_price': base_price,
                'quality_upgrade': quality_upgrade,
                'price_per_page': price_per_page,
                'page_count': page_count,
                'num_copies': num_copies,
                'total_price': total_price,
                'pricing_key_used': resolved_key or lookup_candidates[0]
            }
            
            return JsonResponse({
                'success': True,
                'pricing_breakdown': pricing_breakdown,
                'total_price': total_price
            })
            
        except Exception as e:
                print(f"Error calculating pricing for {vendor_email}: {str(e)}")
                # Return default pricing calculation based on actual pricing.json structure
                default_pricing = {
                    'digital_print_a4_single_color': 5,
                    'digital_print_a3_single_color': 10,
                    'digital_print_a2_single_color': 15,
                    'digital_print_al_single_color': 20,  # Note: A1 is stored as 'al' in pricing.json
                    'digital_print_12x18_single_color': 26,
                    'digital_print_standard_quality': 22,
                    'digital_print_high_quality': 19
                }
                
                # Construct the pricing key
                # Handle special case for A1 (it's stored as 'al' in the pricing.json)
                paper_size_key = paper_size.lower()
                if paper_size_key == 'a1':
                    paper_size_key = 'al'
                
                pricing_key_name = f"digital_print_{paper_size_key}_{print_color}"
                base_price = default_pricing.get(pricing_key_name, 5)
                
                quality_upgrade = 0
                if print_quality == 'high_quality':
                    quality_upgrade = default_pricing.get('digital_print_high_quality', 19)
                
                price_per_page = base_price + quality_upgrade
                total_price = price_per_page * page_count * num_copies
                
                print(f"Default calculation: base_price={base_price}, quality_upgrade={quality_upgrade}, price_per_page={price_per_page}, page_count={page_count}, num_copies={num_copies}, total_price={total_price}")
                
                pricing_breakdown = {
                    'base_price': base_price,
                    'quality_upgrade': quality_upgrade,
                    'price_per_page': price_per_page,
                    'page_count': page_count,
                    'num_copies': num_copies,
                    'total_price': total_price,
                    'pricing_key_used': pricing_key_name
                }
                
                return JsonResponse({
                    'success': True,
                    'pricing_breakdown': pricing_breakdown,
                    'total_price': total_price,
                    'using_default_pricing': True
                })
                
        except Exception as e:
            print(f"Error in calculate_digital_print_pricing: {str(e)}")
            return JsonResponse({'success': False, 'error': str(e)})
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'})
@csrf_exempt
def calculate_jumbo_print_pricing(request):
    """Calculate pricing for jumbo print service based on vendor pricing.json"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_email = data.get('vendor_email')
            print_size = data.get('print_size', 'A3')
            print_color = data.get('print_color', 'Color')
            print_quality = data.get('print_quality', 'standard')
            page_count = data.get('page_count', 1)
            num_copies = data.get('num_copies', 1)
            
            print(f"Received jumbo print data: vendor_email={vendor_email}, print_size={print_size}, print_color={print_color}, print_quality={print_quality}, page_count={page_count}, num_copies={num_copies}")
            print(f"Page count type: {type(page_count)}, value: {page_count}")
            print(f"Num copies type: {type(num_copies)}, value: {num_copies}")
            
            # Convert page_count to int
            try:
                page_count = int(page_count)
                print(f"Page count after conversion: {page_count}")
            except (ValueError, TypeError):
                page_count = 1
                print("Invalid page count, using default: 1")
            
            # Ensure page_count is at least 1
            if page_count < 1:
                page_count = 1
            
            # Convert num_copies to int
            try:
                num_copies = int(num_copies)
                print(f"Num copies after conversion: {num_copies}")
            except (ValueError, TypeError):
                num_copies = 1
                print("Invalid num copies, using default: 1")
            
            # Ensure num_copies is at least 1
            if num_copies < 1:
                num_copies = 1
            
            # Get vendor pricing from d1 database
            print(f"🔍 Jumbo print - Vendor email: {vendor_email}")
            pricing_data = get_vendor_pricing_from_d1(vendor_email)
            
            if pricing_data is None:
                print(f"❌ Could not load vendor pricing from d1 database for vendor: {vendor_email}")
                return JsonResponse({
                    'success': False,
                    'error': f'Unable to load pricing data for vendor. Please contact the vendor.'
                })
            
            # Construct pricing key based on user selections
            size_key = (print_size or 'A3').strip().lower()
            size_key = size_key.replace(' ', '')
            size_candidates = [size_key]
            if size_key == 'a1':
                size_candidates.append('al')
            
            # For jumbo printing, the pricing is nested under "categorized_pricing.jumbo_print" object
            # Access the nested jumbo_print pricing data
            categorized_pricing = pricing_data.get('categorized_pricing', {})
            print(f"📊 Categorized pricing keys: {list(categorized_pricing.keys())}")
            
            jumbo_pricing = categorized_pricing.get('jumbo_print', {})
            print(f"📊 Jumbo pricing keys: {list(jumbo_pricing.keys())}")
            
            # Map print color to pricing key format
            color_key = normalize_color_key(print_color, default='color')
            color_candidates = []
            for candidate in [color_key, 'single_color' if color_key == 'color' else None, 'single_bw' if color_key == 'bw' else None]:
                if candidate:
                    color_candidates.append(candidate)
            if color_key == 'bw':
                color_candidates.append('bw')
            # Deduplicate while preserving order
            seen_colors = set()
            color_candidates = [c for c in color_candidates if not (c in seen_colors or seen_colors.add(c))]
            
            lookup_candidates = []
            for size_variant in size_candidates:
                for color_variant in color_candidates:
                    lookup_candidates.append(f'jumbo_print_{size_variant}_{color_variant}')
                lookup_candidates.append(f'jumbo_print_{size_variant}')
            
            # Get base price from the nested jumbo_print object
            base_price, resolved_key = resolve_pricing_key(jumbo_pricing, *lookup_candidates)
            
            # Check if pricing is available
            if base_price in (None, ''):
                print(f"❌ Pricing not found for keys: {lookup_candidates}")
                print(f"Available jumbo pricing keys: {list(jumbo_pricing.keys())}")
                return JsonResponse({
                    'success': False,
                    'error': f'Pricing not available for {print_size} {print_color}. Please contact the vendor.'
                })
            
            base_price = safe_price(base_price, default=0.0)
            
            # Calculate price per page
            price_per_page = base_price
            
            # Calculate total price (price per page * number of pages * number of copies)
            total_price = price_per_page * page_count * num_copies
            
            print(f"Calculation: base_price={base_price}, price_per_page={price_per_page}, page_count={page_count}, num_copies={num_copies}, total_price={total_price}")
            
            # Prepare pricing breakdown
            pricing_breakdown = {
                'base_price': base_price,
                'price_per_page': price_per_page,
                'page_count': page_count,
                'num_copies': num_copies,
                'total_price': total_price,
                'pricing_key_used': resolved_key or lookup_candidates[0]
            }
            
            return JsonResponse({
                'success': True,
                'pricing_breakdown': pricing_breakdown,
                'total_price': total_price
            })
            
        except Exception as e:
            print(f"Error calculating jumbo print pricing: {e}")
            traceback.print_exc()
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    
    return JsonResponse({
        'success': False,
        'error': 'Invalid request method'
    })
@csrf_exempt
def calculate_passport_photo_pricing(request):
    """Calculate pricing for passport photo service based on vendor pricing.json"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_email = data.get('vendor_email')
            country = data.get('country', 'india')
            photo_package = data.get('photo_package', '8')
            color_photo = data.get('color_photo', 'color_photo')
            paper_size = data.get('paper_size', 'A4')
            
            print(f"Received passport photo data: vendor_email={vendor_email}, country={country}, photo_package={photo_package}, color_photo={color_photo}, paper_size={paper_size}")
            
            # Convert photo_package to int
            try:
                photo_package = int(photo_package)
                print(f"Photo package after conversion: {photo_package}")
            except (ValueError, TypeError):
                photo_package = 8
                print("Invalid photo package, using default: 8")
            
            # Ensure photo_package is valid
            if photo_package not in [8, 16, 30]:
                photo_package = 8
                print(f"Invalid photo package {photo_package}, using default: 8")
            
            # Get vendor pricing from d1 database
            print(f"🔍 Passport photo - Vendor email: {vendor_email}")
            pricing_data = get_vendor_pricing_from_d1(vendor_email)
            
            if pricing_data is None:
                print(f"❌ Could not load vendor pricing from d1 database for vendor: {vendor_email}")
                return JsonResponse({
                    'success': False,
                    'error': f'Unable to load pricing data for vendor. Please contact the vendor.'
                })
            
            # For passport photo, the pricing is nested under "categorized_pricing.passport_photo" object
            # Access the nested passport_photo pricing data
            categorized_pricing = pricing_data.get('categorized_pricing', {})
            print(f"📊 Categorized pricing keys: {list(categorized_pricing.keys())}")
            
            passport_pricing = categorized_pricing.get('passport_photo', {})
            print(f"📊 Passport photo pricing keys: {list(passport_pricing.keys())}")
            
            # Construct the pricing key format: passport_print_{photo_package}
            # The pricing.json has: passport_print_8, passport_print_16, passport_print_30
            pricing_key_name = f'passport_print_{photo_package}'
            print(f"🔍 Looking for pricing key: {pricing_key_name}")
            
            # Get base price from the nested passport_photo object
            base_price = passport_pricing.get(pricing_key_name)
            
            # Check if pricing is available
            if base_price is None:
                print(f"❌ Pricing not found for key: {pricing_key_name}")
                print(f"Available passport photo pricing keys: {list(passport_pricing.keys())}")
                return JsonResponse({
                    'success': False,
                    'error': f'Pricing not available for {photo_package} photos. Please contact the vendor.'
                })
            
            # Calculate total price (base price for the package)
            total_price = base_price
            
            print(f"Calculation: base_price={base_price}, total_price={total_price}")
            
            # Prepare pricing breakdown
            pricing_breakdown = {
                'base_price': base_price,
                'total_price': total_price,
                'pricing_key_used': pricing_key_name,
                'country': country,
                'photo_package': photo_package,
                'color_photo': color_photo,
                'paper_size': paper_size,
                'total_pages': 1  # For passport photo, it's always 1 page (A4 sheet)
            }
            
            return JsonResponse({
                'success': True,
                'pricing_breakdown': pricing_breakdown,
                'total_price': total_price
            })
            
        except Exception as e:
            print(f"Error calculating passport photo pricing: {e}")
            traceback.print_exc()
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    
    return JsonResponse({
        'success': False,
        'error': 'Invalid request method'
    })

@csrf_exempt
def calculate_a4_print_pricing(request):
    """Calculate pricing for document print service (A4 / Letter / A3) based on vendor pricing.json
    Formula: Copies × [(B&W Pages × B&W Rate) + (Color Pages × Color Rate) + 
             (IF Lamination Selected → Lamination Rate) + 
             (IF Spiral Binding Selected → Spiral Binding Rate for page based of document ELSE → 0) + 
             (IF Tape Binding Selected → Tape Binding Rate as per page of document ELSE → 0)]
    Final Amount Customer Pays = Total Bill Amount × 1.20
    """
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_email = data.get('vendor_email')
            print_type = data.get('print_type', 'single_bw')
            total_pages = data.get('total_pages', 1)
            total_copies = data.get('total_copies', 1)
            bw_pages = data.get('bw_pages', 0)  # Black and white pages
            color_pages = data.get('color_pages', 0)  # Color pages
            lamination = data.get('lamination', False)  # Lamination option
            spiral_binding = data.get('spiral_binding', False)  # Spiral binding option
            tape_binding = data.get('tape_binding', False)  # Tape binding option
            page_size_raw = data.get('page_size', 'A4')
            page_size = str(page_size_raw).strip().upper() if page_size_raw else 'A4'
            
            print(f"Received A4 print data: vendor_email={vendor_email}, print_type={print_type}, total_pages={total_pages}, total_copies={total_copies}, bw_pages={bw_pages}, color_pages={color_pages}, lamination={lamination}, spiral_binding={spiral_binding}, tape_binding={tape_binding}, page_size={page_size}")
            
            # Convert values to int/float
            try:
                total_pages = int(total_pages)
                total_copies = int(total_copies)
                bw_pages = int(bw_pages) if bw_pages else 0
                color_pages = int(color_pages) if color_pages else 0
                lamination = bool(lamination) if lamination is not None else False
                spiral_binding = bool(spiral_binding) if spiral_binding is not None else False
                tape_binding = bool(tape_binding) if tape_binding is not None else False
            except (ValueError, TypeError):
                total_pages = 1
                total_copies = 1
                bw_pages = 0
                color_pages = 0
                lamination = False
                spiral_binding = False
                tape_binding = False
            
            # Ensure values are at least 0/1
            if total_pages < 1: total_pages = 1
            if total_copies < 1: total_copies = 1
            if bw_pages < 0: bw_pages = 0
            if color_pages < 0: color_pages = 0
            
            # If bw_pages and color_pages are not provided, calculate from print_type
            if bw_pages == 0 and color_pages == 0:
                if print_type == 'single_color':
                    color_pages = total_pages
                else:
                    bw_pages = total_pages
            
            # Get vendor pricing from d1 database - STRICTLY use vendor pricing only
            print(f"🔍 A4 print - Vendor email: {vendor_email}")
            pricing_data = get_vendor_pricing_from_d1(vendor_email)
            
            if pricing_data is None:
                print(f"❌ Could not load vendor pricing from d1 database for vendor: {vendor_email}")
                return JsonResponse({
                    'success': False,
                    'error': f'Unable to load pricing data for vendor. Please contact the vendor.'
                })
            
            # Access categorized pricing and flat pricing
            categorized_pricing = pricing_data.get('categorized_pricing', {})
            flat_pricing = pricing_data.get('pricing_data') or pricing_data.get('pricing') or {}
            print(f"📊 Categorized pricing keys: {list(categorized_pricing.keys())}")
            
            a4_pricing = categorized_pricing.get('a4_print', {})
            lamination_pricing = categorized_pricing.get('lamination', {})
            binding_pricing = categorized_pricing.get('binding', {})
            
            print(f"📊 A4 print pricing keys: {list(a4_pricing.keys())}")
            print(f"📊 Lamination pricing keys: {list(lamination_pricing.keys())}")
            print(f"📊 Binding pricing keys: {list(binding_pricing.keys())}")
            
            # Get B&W and Color rates - STRICTLY from vendor pricing, no defaults.
            if page_size == 'LETTER':
                # Use dedicated Letter document print pricing from Vendor_pricing
                bw_rate = flat_pricing.get('doc_letter_bw')
                color_rate = flat_pricing.get('doc_letter_color')
            elif page_size == 'A3':
                # For A3 document printing, reuse Jumbo A3 rates
                jumbo_pricing = categorized_pricing.get('jumbo_print', {})
                bw_rate = jumbo_pricing.get('jumbo_print_a3_bw')
                color_rate = jumbo_pricing.get('jumbo_print_a3_color')
            else:
                # Default and backward‑compatible behaviour: A4 pricing
                bw_rate = a4_pricing.get('regular_print_a4_bw')
                color_rate = a4_pricing.get('regular_print_a4_color')
            
            # Check if pricing is available - STRICTLY require vendor pricing
            if bw_rate is None and bw_pages > 0:
                print(f"❌ B&W pricing not found for page_size={page_size}")
                print(f"Available A4 print pricing keys: {list(a4_pricing.keys())}")
                return JsonResponse({
                    'success': False,
                    'error': f'Pricing not available for Black & White printing. Please contact the vendor.'
                })
            
            if color_rate is None and color_pages > 0:
                print(f"❌ Color pricing not found for page_size={page_size}")
                print(f"Available A4 print pricing keys: {list(a4_pricing.keys())}")
                return JsonResponse({
                    'success': False,
                    'error': f'Pricing not available for Color printing. Please contact the vendor.'
                })
            
            # Convert rates to float, default to 0 if None
            bw_rate = safe_price(bw_rate, default=0.0)
            color_rate = safe_price(color_rate, default=0.0)
            
            # Calculate base print cost: (B&W Pages × B&W Rate) + (Color Pages × Color Rate)
            base_print_cost = (bw_pages * bw_rate) + (color_pages * color_rate)
            
            # Calculate lamination cost (per page)
            lamination_cost = 0.0
            if lamination:
                # Try to get A4 lamination rate (standard or glossy)
                lamination_key = 'lamination_a4_standard'  # Default to standard
                lamination_rate = lamination_pricing.get(lamination_key)
                if lamination_rate is None:
                    # Try glossy
                    lamination_key = 'lamination_a4_glossy'
                    lamination_rate = lamination_pricing.get(lamination_key)
                
                if lamination_rate is None:
                    print(f"⚠️ Lamination pricing not found. Available keys: {list(lamination_pricing.keys())}")
                    return JsonResponse({
                        'success': False,
                        'error': f'Lamination pricing not available. Please contact the vendor.'
                    })
                
                lamination_rate = safe_price(lamination_rate, default=0.0)
                # Lamination is per page
                lamination_cost = lamination_rate * total_pages
            
            # Calculate spiral binding cost (per document, based on page count)
            spiral_binding_cost = 0.0
            if spiral_binding:
                # Determine binding key based on page count
                # Spiral binding: a4_100 (up to 100 pages) or a4_200 (101-200 pages, or >200 pages)
                if total_pages <= 100:
                    spiral_key = 'spiral_binding_a4_100'
                else:
                    spiral_key = 'spiral_binding_a4_200'
                
                print(f"🔍 Spiral binding: total_pages={total_pages}, selected_key={spiral_key}")
                spiral_rate = binding_pricing.get(spiral_key)
                print(f"🔍 Spiral binding rate retrieved: {spiral_rate} (type: {type(spiral_rate)})")
                
                if spiral_rate is None or spiral_rate == '':
                    print(f"⚠️ Spiral binding pricing not found for key: {spiral_key}")
                    print(f"Available binding keys: {list(binding_pricing.keys())}")
                    return JsonResponse({
                        'success': False,
                        'error': f'Spiral binding pricing not available. Please contact the vendor.'
                    })
                
                spiral_rate = safe_price(spiral_rate, default=0.0)
                if spiral_rate <= 0:
                    print(f"⚠️ Spiral binding rate is invalid (0 or negative): {spiral_rate} for key: {spiral_key}")
                    print(f"Available binding keys with values: {[(k, v) for k, v in binding_pricing.items()]}")
                    return JsonResponse({
                        'success': False,
                        'error': f'Spiral binding pricing is not configured properly. Please contact the vendor.'
                    })
                
                print(f"✅ Spiral binding rate: {spiral_rate} for {total_pages} pages (key: {spiral_key})")
                # Spiral binding is per copy (document)
                spiral_binding_cost = spiral_rate * total_copies
                print(f"💰 Spiral binding cost: {spiral_rate} × {total_copies} = {spiral_binding_cost}")
            
            # Calculate tape binding cost (per document, based on page count)
            tape_binding_cost = 0.0
            if tape_binding:
                # Determine binding key based on page count
                # Tape binding: a4_100 (up to 100 pages) or a4_200 (101-200 pages, or >200 pages)
                if total_pages <= 100:
                    tape_key = 'tape_binding_a4_100'
                else:
                    tape_key = 'tape_binding_a4_200'
                
                print(f"🔍 Tape binding: total_pages={total_pages}, selected_key={tape_key}")
                tape_rate = binding_pricing.get(tape_key)
                print(f"🔍 Tape binding rate retrieved: {tape_rate} (type: {type(tape_rate)})")
                
                if tape_rate is None or tape_rate == '':
                    print(f"⚠️ Tape binding pricing not found for key: {tape_key}")
                    print(f"Available binding keys: {list(binding_pricing.keys())}")
                    return JsonResponse({
                        'success': False,
                        'error': f'Tape binding pricing not available. Please contact the vendor.'
                    })
                
                tape_rate = safe_price(tape_rate, default=0.0)
                if tape_rate <= 0:
                    print(f"⚠️ Tape binding rate is invalid (0 or negative): {tape_rate} for key: {tape_key}")
                    print(f"Available binding keys with values: {[(k, v) for k, v in binding_pricing.items()]}")
                    return JsonResponse({
                        'success': False,
                        'error': f'Tape binding pricing is not configured properly. Please contact the vendor.'
                    })
                
                print(f"✅ Tape binding rate: {tape_rate} for {total_pages} pages (key: {tape_key})")
                # Tape binding is per copy (document)
                tape_binding_cost = tape_rate * total_copies
                print(f"💰 Tape binding cost: {tape_rate} × {total_copies} = {tape_binding_cost}")
            
            # Calculate total bill amount (before commission) per copy
            # Formula: (B&W Pages × B&W Rate) + (Color Pages × Color Rate) + 
            #          (IF Lamination Selected → Lamination Rate) + 
            #          (IF Spiral Binding Selected → Spiral Binding Rate) + 
            #          (IF Tape Binding Selected → Tape Binding Rate)
            cost_per_copy = base_print_cost + lamination_cost + spiral_binding_cost + tape_binding_cost
            
            # Total bill amount (before commission) = cost_per_copy × copies
            total_bill_amount = cost_per_copy * total_copies
            
            # Final Amount Customer Pays = Total Bill Amount × 1.20 (20% commission)
            final_amount = total_bill_amount * 1.20
            
            print(f"💰 Calculation: bw_pages={bw_pages}, color_pages={color_pages}, bw_rate={bw_rate}, color_rate={color_rate}")
            print(f"💰 Base print cost={base_print_cost}, lamination_cost={lamination_cost}, spiral_binding_cost={spiral_binding_cost}, tape_binding_cost={tape_binding_cost}")
            print(f"💰 Cost per copy={cost_per_copy}, total_copies={total_copies}, total_bill_amount={total_bill_amount}, final_amount={final_amount}")
            
            # Prepare pricing breakdown
            pricing_breakdown = {
                'vendor_email': vendor_email,  # Include vendor email for shop lookup
                'bw_pages': bw_pages,
                'color_pages': color_pages,
                'bw_rate': bw_rate,
                'color_rate': color_rate,
                'base_print_cost': base_print_cost,
                'lamination_cost': lamination_cost,
                'spiral_binding_cost': spiral_binding_cost,
                'tape_binding_cost': tape_binding_cost,
                'cost_per_copy': cost_per_copy,
                'total_pages': total_pages,
                'total_copies': total_copies,
                'total_bill_amount': total_bill_amount,  # Before commission
                'final_amount': final_amount,  # After 20% commission
                'commission_percentage': 20,
                'commission_amount': final_amount - total_bill_amount,
                'lamination': lamination,
                'spiral_binding': spiral_binding,
                'tape_binding': tape_binding
            }
            
            return JsonResponse({
                'success': True,
                'pricing_breakdown': pricing_breakdown,
                'total_price': final_amount  # Return final amount with commission
            })
            
        except Exception as e:
            print(f"Error calculating A4 print pricing: {e}")
            traceback.print_exc()
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    
    return JsonResponse({
        'success': False,
        'error': 'Invalid request method'
    })

@csrf_exempt
def get_available_shops(request):
    """
    Get all available shops from d1 database Vendor_register table with caching
    """
    import time
    
    # Check cache first
    cache_key = "available_shops_cache"
    current_time = time.time()
    
    if (hasattr(get_available_shops, '_cache') and 
        cache_key in get_available_shops._cache and 
        cache_key in get_available_shops._cache_timestamps and 
        current_time - get_available_shops._cache_timestamps[cache_key] < 300):  # 5 minutes
        
        print(f"⚡ CACHE HIT: Using cached available shops")
        return JsonResponse({
            'success': True,
            'shops': get_available_shops._cache[cache_key]
        })
    
    try:
        # Try to fetch from d1 database via Worker API
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if api_url and api_key:
            # Construct the Worker API endpoint for getting all vendors
            # Remove any existing endpoint paths and add /get-all-vendors
            base_url = api_url.rstrip('/')
            # Remove common endpoint paths if present
            for endpoint in ['/add-contact', '/add-vendor-register', '/add-vendor-pricing', '/get-all-vendors']:
                if base_url.endswith(endpoint):
                    base_url = base_url[:-len(endpoint)]
            worker_endpoint = base_url.rstrip('/') + '/get-all-vendors'
            
            try:
                resp = requests.get(
                    worker_endpoint,
                    headers={
                        'Content-Type': 'application/json',
                        'x-api-key': api_key
                    },
                    timeout=15
                )
                
                if resp.status_code == 200:
                    response_data = resp.json()
                    if response_data.get('success'):
                        vendors = response_data.get('vendors', [])
                        shops = []
                        
                        for vendor in vendors:
                            vendor_name = vendor.get('vendor_name', '')
                            vendor_email = vendor.get('email', '')
                            
                            if vendor_name and vendor_email:
                                shop_folder = sanitize_shop_name(vendor_name)
                                service_availability = vendor.get('service_availability') or {}
                                shop_info = {
                                    'shop_name': vendor_name,
                                    'shop_folder': shop_folder,
                                    'vendor_email': vendor_email,
                                    'phone_number': vendor.get('phone_number', ''),
                                    'shop_address': vendor.get('shop_address', ''),
                                    'city': vendor.get('city', ''),
                                    'pincode': vendor.get('pincode', ''),
                                    'latitude': vendor.get('latitude', ''),
                                    'longitude': vendor.get('longitude', ''),
                                    'status': vendor.get('status', 'Available'),
                                    'vendor_id': vendor.get('vendor_id', ''),
                                    'vendor_token': vendor.get('vendor_token', ''),
                                    'pending_jobs_count': vendor.get('pending_jobs_count', 0),
                                    'service_availability': service_availability,
                                    # Flatten service flags - only 1/true = available, else NOT available
                                    'digital_print': service_availability.get('service_data', {}).get('digital_print', False),
                                    'project_binding': service_availability.get('service_data', {}).get('project_binding', False),
                                    'gloss_printing': service_availability.get('service_data', {}).get('gloss_printing', False),
                                    'jumbo_printing': service_availability.get('service_data', {}).get('jumbo_printing', False),
                                    'regular_print': service_availability.get('service_data', {}).get('regular_print', False),
                                    'passport_print': service_availability.get('service_data', {}).get('passport_print', False),
                                    'photo_print': service_availability.get('service_data', {}).get('photo_print', False),
                                }
                                if not any(s['shop_folder'] == shop_folder for s in shops):
                                    shops.append(shop_info)
                        
                        if len(shops) > 0:
                            print(f"✅ Successfully fetched {len(shops)} shops from d1 database")
                            
                            # Only cache successful results with shops
                            if not hasattr(get_available_shops, '_cache'):
                                get_available_shops._cache = {}
                                get_available_shops._cache_timestamps = {}
                            get_available_shops._cache[cache_key] = shops
                            get_available_shops._cache_timestamps[cache_key] = current_time
                            
                            return JsonResponse({
                                'success': True,
                                'shops': shops,
                                'total_shops': len(shops)
                            })
                        else:
                            error_msg = "No shops found in database"
                            print(f"⚠️ {error_msg}")
                            # Clear cache if we get empty results
                            if hasattr(get_available_shops, '_cache') and cache_key in get_available_shops._cache:
                                del get_available_shops._cache[cache_key]
                                if cache_key in get_available_shops._cache_timestamps:
                                    del get_available_shops._cache_timestamps[cache_key]
                    else:
                        error_msg = response_data.get('error', 'Unknown error from Worker API')
                        print(f"⚠️ Worker API returned error: {error_msg}")
                        # Clear cache on error
                        if hasattr(get_available_shops, '_cache') and cache_key in get_available_shops._cache:
                            del get_available_shops._cache[cache_key]
                            if cache_key in get_available_shops._cache_timestamps:
                                del get_available_shops._cache_timestamps[cache_key]
                else:
                    error_text = resp.text[:500] if resp.text else 'No error message'
                    print(f"⚠️ Failed to fetch shops from d1 database with status {resp.status_code}: {error_text}")
                    # Clear cache on error
                    if hasattr(get_available_shops, '_cache') and cache_key in get_available_shops._cache:
                        del get_available_shops._cache[cache_key]
                        if cache_key in get_available_shops._cache_timestamps:
                            del get_available_shops._cache_timestamps[cache_key]
            except requests.exceptions.RequestException as e:
                print(f"⚠️ Failed to connect to Worker API: {str(e)}")
                # Clear cache on connection error
                if hasattr(get_available_shops, '_cache') and cache_key in get_available_shops._cache:
                    del get_available_shops._cache[cache_key]
                    if cache_key in get_available_shops._cache_timestamps:
                        del get_available_shops._cache_timestamps[cache_key]
            except Exception as e:
                print(f"⚠️ Unexpected error fetching shops from d1 database: {str(e)}")
                # Clear cache on unexpected error
                if hasattr(get_available_shops, '_cache') and cache_key in get_available_shops._cache:
                    del get_available_shops._cache[cache_key]
                    if cache_key in get_available_shops._cache_timestamps:
                        del get_available_shops._cache_timestamps[cache_key]
        else:
            print(f"⚠️ Worker API not configured - cannot fetch shops from d1 database")
        
        # Fallback: Return empty shops list if d1 fetch fails (but don't cache it)
        print(f"⚠️ Falling back to empty shops list")
        shops = []
        
        return JsonResponse({
            'success': True,
            'shops': shops,
            'total_shops': len(shops),
            'message': 'No shops available. Please check database connection.'
        })
    except Exception as e:
        print(f"Error getting available shops: {str(e)}")
        return JsonResponse({
            'success': False,
            'error': str(e),
            'shops': []
        })

def vendor_email_folder(email):
    return f'vendor_register_details/{sanitize_email(email)}'


@csrf_exempt
def get_available_printers(request):
    """
    Get available printers for a specific service type from vendor pricing data
    """
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Method not allowed'}, status=405)
    
    try:
        data = json.loads(request.body)
        service_type = (data.get('service_type') or '').strip().lower()
        
        if not service_type:
            return JsonResponse({'success': False, 'error': 'Service type is required'}, status=400)
        
        # Get vendor email from session
        vendor_email = request.session.get('vendor_email')
        if not vendor_email:
            return JsonResponse({'success': False, 'error': 'Vendor not authenticated'}, status=401)
        
        # Get vendor pricing data
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        pricing_key = f"vendor_register_details/{sanitize_email(vendor_email)}/pricing.json"
        
        try:
            response = s3.get_object(Bucket=settings.R2_BUCKET, Key=pricing_key)
            pricing_data = json.loads(response['Body'].read().decode('utf-8'))
        except Exception as e:
            return JsonResponse({'success': False, 'error': f'Failed to load pricing data: {str(e)}'}, status=404)
        
        printer_config = pricing_data.get('printer_configuration') or pricing_data.get('printers') or {}
        printers = []
        
        # Normalize service type variants used across UI/backend
        normalized = service_type.replace(' ', '_')
        if normalized == 'gloss_printing':
            normalized = 'gloss_print'
        if normalized == 'jumbo_printing':
            normalized = 'jumbo_print'
        if normalized in ('regular_print', 'regular'):
            normalized = 'a4_print'
        if not normalized:
            normalized = 'a4_print'

        # Map service types to printer configuration keys
        service_to_printer_mapping = {
            'digital_print': ['digital_printer_1', 'digital_printer_2'],
            'project_binding': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
            'jumbo_print': ['jumbo_printer_1', 'jumbo_printer_2'],
            'gloss_print': ['gloss_printer_1', 'gloss_printer_2'],
            'photo_print': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
            'passport_print': ['passport_printer_1', 'passport_printer_2', 'passport_printer_3'],
            'a4_print': ['a4_printer_1', 'a4_printer_2', 'a4_printer_3'],
        }
        
        printer_keys = service_to_printer_mapping.get(normalized, [])
        
        for key in printer_keys:
            printer_name = (printer_config.get(key) if isinstance(printer_config, dict) else None) or 'NA'
            if printer_name and printer_name != 'NA':
                printers.append({'name': printer_name, 'type': key, 'key': key})

        # De-duplicate by name while preserving order
        seen = set()
        unique_printers = []
        for p in printers:
            if p['name'] not in seen:
                seen.add(p['name'])
                unique_printers.append(p)
        
        return JsonResponse({
            'success': True,
            'printers': unique_printers,
            'service_type': normalized
        })
        
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)


@csrf_exempt
def assign_printer_to_job(request):
    """
    Assign a printer to a specific job and update the job's metadata in R2
    """
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Method not allowed'}, status=405)
    
    try:
        data = json.loads(request.body)
        filename = data.get('filename')
        service_type = data.get('service_type')
        printer_name = data.get('printer_name')
        
        if not all([filename, service_type, printer_name]):
            return JsonResponse({'success': False, 'error': 'filename, service_type, and printer_name are required'}, status=400)
        
        # Get vendor email from session
        vendor_email = request.session.get('vendor_email')
        if not vendor_email:
            return JsonResponse({'success': False, 'error': 'Vendor not authenticated'}, status=401)
        
        # Initialize S3 client
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        # Also get vendor_id from session, as jobs are rendered under vendor_id-based folders
        vendor_id = request.session.get('vendor_id') or request.session.get('vendorId')

        # Define possible paths where the print job might be stored
        possible_paths = [
            f'vendor_print_jobs/{vendor_email}/{filename}',
            f'vendor_manual_print_jobs/{vendor_email}/{filename}',
            f'vendor_register_details/{vendor_email.replace("@", "_at_").replace(".", "_dot_")}/firozshop/{filename}',
            f'vendor_print_jobs/{vendor_email.replace("@", "_at_").replace(".", "_dot_")}/{filename}',
            f'vendor_manual_print_jobs/{vendor_email.replace("@", "_at_").replace(".", "_dot_")}/{filename}',
        ]

        # Add vendor_id based paths (these are used by get_print_requests)
        if vendor_id:
            possible_paths = [
                f'vendor_print_jobs/{vendor_id}/{filename}',
                f'vendor_manual_print_jobs/{vendor_id}/{filename}',
            ] + possible_paths
        
        updated = False
        updated_path = None
        
        # Search for the file in all possible paths
        for path in possible_paths:
            try:
                print(f"🔍 Searching for {filename} at {path}")
                head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=path)
                current_metadata = head_response.get('Metadata', {})
                
                # Do not store metadata on R2 for vendor/user print jobs; rely on DB tables instead
                if path.startswith('vendor_print_jobs/') or path.startswith('users/'):
                    print(f"ℹ️ Skipping R2 metadata write for {path}; DB holds metadata.")
                    updated = True
                    updated_path = path
                else:
                    # Update metadata for non-print-job assets
                    current_metadata['assigned_printer'] = printer_name
                    # Avoid R2 metadata writes for print jobs; rely on DB tables instead
                    if path.startswith('vendor_print_jobs/') or path.startswith('users/'):
                        print(f"ℹ️ Skipping printer metadata write for {path}; DB holds metadata.")
                        updated = True
                        updated_path = path
                    else:
                        current_metadata['printer_name'] = printer_name
                        current_metadata['printer_assigned_at'] = datetime.datetime.now().isoformat()
                        current_metadata['service_type'] = service_type
                        
                        copy_source = {'Bucket': settings.R2_BUCKET, 'Key': path}
                        s3.copy_object(
                            CopySource=copy_source,
                            Bucket=settings.R2_BUCKET,
                            Key=path,
                            Metadata=current_metadata,
                            MetadataDirective='REPLACE'
                        )
                        updated = True
                        updated_path = path
                print(f"✅ Updated printer assignment for {filename} at {path}")
                break
                
            except Exception as e:
                print(f"   ⚠️ Not found at {path}: {str(e)}")
                continue
        
        if not updated:
            return JsonResponse({
                'success': False, 
                'error': f'Print job {filename} not found in any vendor folder'
            }, status=404)
        
        return JsonResponse({
            'success': True,
            'message': f'Printer {printer_name} assigned to job {filename}',
            'filename': filename,
            'service_type': service_type,
            'printer_name': printer_name,
            'updated_path': updated_path
        })
        
    except Exception as e:
        print(f"❌ Error assigning printer to job: {str(e)}")
        return JsonResponse({'success': False, 'error': str(e)}, status=500)

def get_vendor_email_by_shop_folder(shop_folder):
    """Get vendor email by shop folder name from R2 storage"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        # Search through vendor registration details to find matching shop folder
        objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_register_details/')
        for obj in objects.get("Contents", []):
            if obj["Key"].endswith('/registration_details.json'):
                try:
                    response = s3.get_object(Bucket=settings.R2_BUCKET, Key=obj["Key"])
                    vendor_data = json.loads(response['Body'].read().decode('utf-8'))
                    vendor_name = vendor_data.get('vendor_name', '')
                    vendor_email = vendor_data.get('vendor_email', '')

                    # Check if this vendor's sanitized shop name matches
                    if sanitize_shop_name(vendor_name) == shop_folder:
                        return vendor_email
                except Exception as e:
                    print(f"Error reading vendor data from {obj['Key']}: {str(e)}")
                    continue

        # Fallback for firozshop or unknown shops
        return 'firozshop@example.com'

    except Exception as e:
        print(f"Error finding vendor email for shop {shop_folder}: {str(e)}")
        return 'firozshop@example.com'

def get_vendor_id_by_shop_folder(shop_folder):
    """Get vendor_id by shop folder name from R2 storage"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        # Search through vendor registration details to find matching shop folder
        objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_register_details/')
        for obj in objects.get("Contents", []):
            if obj["Key"].endswith('/registration_details.json'):
                try:
                    response = s3.get_object(Bucket=settings.R2_BUCKET, Key=obj["Key"])
                    vendor_data = json.loads(response['Body'].read().decode('utf-8'))
                    vendor_name = vendor_data.get('vendor_name', '')
                    vendor_id = vendor_data.get('vendor_id', '')

                    # Check if this vendor's sanitized shop name matches
                    if sanitize_shop_name(vendor_name) == shop_folder:
                        return vendor_id
                except Exception as e:
                    print(f"Error reading vendor data from {obj['Key']}: {str(e)}")
                    continue

        # Fallback for firozshop or unknown shops
        return 'vendor1'

    except Exception as e:
        print(f"Error finding vendor_id for shop {shop_folder}: {str(e)}")
        return 'vendor1'
def get_vendor_email_by_vendor_id(vendor_id):
    """Get vendor email by vendor_id from R2 storage"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        # Search through vendor registration details to find matching vendor_id
        objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_register_details/')
        for obj in objects.get("Contents", []):
            if obj["Key"].endswith('/registration_details.json'):
                try:
                    response = s3.get_object(Bucket=settings.R2_BUCKET, Key=obj["Key"])
                    vendor_data = json.loads(response['Body'].read().decode('utf-8'))
                    stored_vendor_id = vendor_data.get('vendor_id', '')
                    vendor_email = vendor_data.get('vendor_email', '')

                    # Check if this vendor's ID matches
                    if stored_vendor_id == vendor_id:
                        return vendor_email
                except Exception as e:
                    print(f"Error reading vendor data from {obj['Key']}: {str(e)}")
                    continue

        # Fallback for firozshop or unknown vendors
        return 'firozshop@example.com'

    except Exception as e:
        print(f"Error finding vendor email for vendor_id {vendor_id}: {str(e)}")
        return 'firozshop@example.com'

def get_vendor_id_by_vendor_email(vendor_email):
    """Get vendor_id by vendor_email from R2 storage"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        normalized_email = (vendor_email or '').strip().lower()
        objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_register_details/')
        for obj in objects.get("Contents", []):
            if obj["Key"].endswith('/registration_details.json'):
                try:
                    response = s3.get_object(Bucket=settings.R2_BUCKET, Key=obj["Key"])
                    vendor_data = json.loads(response['Body'].read().decode('utf-8'))
                    stored_email = (vendor_data.get('vendor_email', '') or '').strip().lower()
                    if stored_email and stored_email == normalized_email:
                        return vendor_data.get('vendor_id', '') or 'vendor1'
                except Exception as e:
                    print(f"Error reading vendor data from {obj['Key']}: {str(e)}")
                    continue

        return 'vendor1'

    except Exception as e:
        print(f"Error finding vendor_id for vendor_email {vendor_email}: {str(e)}")
        return 'vendor1'

# This code incorporates address fields into the vendor registration API and updates the pricing structure to handle comprehensive xerox shop pricing.

# --- PicWish Passport Photo Enhancement API ---
@csrf_exempt
def debug_vendor_registrations(request):
    """
    Debug endpoint to list all vendor registrations
    """
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        vendors = []
        try:
            objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_register_details/')
            for obj in objects.get("Contents", []):
                key = obj["Key"]
                if key.endswith('/registration_details.json'):
                    try:
                        response = s3.get_object(Bucket=settings.R2_BUCKET, Key=key)
                        vendor_data = json.loads(response['Body'].read().decode('utf-8'))
                        vendors.append({
                            'vendor_id': vendor_data.get('vendor_id', ''),
                            'vendor_email': vendor_data.get('vendor_email', ''),
                            'vendor_name': vendor_data.get('vendor_name', ''),
                            'latitude': vendor_data.get('latitude', ''),
                            'longitude': vendor_data.get('longitude', ''),
                            'file_path': key
                        })
                    except Exception as e:
                        print(f"Error reading vendor data from {key}: {str(e)}")
                        continue
        except Exception as e:
            print(f"Error listing vendor folders: {str(e)}")
        
        return JsonResponse({
            'success': True,
            'vendors': vendors,
            'total_vendors': len(vendors)
        })
    except Exception as e:
        print(f"Error getting vendor registrations: {str(e)}")
        return JsonResponse({
            'success': False,
            'error': str(e),
            'vendors': []
        })

@csrf_exempt
def enhance_passport_photo(request):
    """
    Accepts an image file, sends it to PicWish API, downloads the enhanced image, and returns it as base64.
    """
    if request.method != 'POST' or 'file' not in request.FILES:
        return JsonResponse({'success': False, 'error': 'No image uploaded.'}, status=400)

    image_file = request.FILES['file']
    api_key = getattr(settings, 'PICWISH_API', None)
    if not api_key:
        return JsonResponse({'success': False, 'error': 'API key not configured'}, status=500)
    
    url = 'https://techhk.aoscdn.com/api/tasks/visual/scale'
    headers = {'X-API-KEY': api_key}
    files = {'image_file': (image_file.name, image_file.read(), image_file.content_type)}
    data = {
        'sync': '1',  # Synchronous
        'type': 'face',  # For passport/portrait
        'scale_factor': '2',  # 2x enhancement
        'return_type': '1'  # Return image URL
    }
    try:
        response = requests.post(url, headers=headers, files=files, data=data, timeout=60)
        result = response.json()
        if result.get('status') == 200 and 'image' in result.get('data', {}):
            enhanced_url = result['data']['image']
            # Download the enhanced image immediately
            img_resp = requests.get(enhanced_url, timeout=60)
            if img_resp.status_code == 200:
                img_b64 = base64.b64encode(img_resp.content).decode('utf-8')
                return JsonResponse({'success': True, 'enhanced_image_b64': img_b64})
            else:
                return JsonResponse({'success': False, 'error': 'Failed to download enhanced image.'}, status=500)
        else:
            return JsonResponse({'success': False, 'error': result.get('error', 'Enhancement failed.')}, status=500)
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)


@csrf_exempt
def enhance_photo(request):
    """
    PicWish Photo Enhancement API endpoint.
    Accepts an image file, sends it to PicWish API, and returns the enhanced image URL.
    Used by the Enhance section in passport-photo-service.
    """
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)
    
    try:
        # Get uploaded image file
        image_file = request.FILES.get('image')
        if not image_file:
            return JsonResponse({'success': False, 'error': 'No image uploaded'}, status=400)
        
        # PicWish API configuration - Load from settings
        api_key = getattr(settings, 'PICWISH_API', None)
        if not api_key:
            print("DEBUG: PICWISH_API not found in settings")
            return JsonResponse({'success': False, 'error': 'API key not configured'}, status=500)
        
        print("DEBUG: Using PicWish API key:", api_key[:10] + "..." if len(api_key) > 10 else api_key)
        
        api_url = 'https://techhk.aoscdn.com/api/tasks/visual/scale'
        
        # Prepare request headers - Use Authorization header format as per PicWish API docs
        headers = {
            "Authorization": f"API-Key {api_key}",
            "Content-Type": "application/json"
        }
        
        # Prepare form data
        data = {
            'sync': 1,  # Synchronous request (1 = sync, 0 = async)
            'type': request.POST.get('type', 'face'),  # Enhancement type: 'face' for passport photos
            'return_type': 1  # Return type: 1 = image URL, 2 = base64
        }
        
        print("DEBUG: Headers:", headers)
        print("DEBUG: Payload:", data)
        
        # Prepare file upload (read file content)
        files = {'image_file': (image_file.name, image_file.read(), image_file.content_type)}
        
        # Make API request to PicWish
        # Note: When using files, requests will set Content-Type to multipart/form-data automatically
        # So we remove Content-Type from headers for file uploads
        headers_for_upload = {
            "Authorization": f"API-Key {api_key}"
        }
        response = requests.post(api_url, headers=headers_for_upload, data=data, files=files, timeout=60)
        
        # If 401 with Authorization header, try X-API-KEY as fallback
        if response.status_code == 401:
            print("DEBUG: 401 with Authorization header, trying X-API-KEY fallback")
            headers_fallback = {'X-API-KEY': api_key}
            # Reset file pointer
            image_file.seek(0)
            files = {'image_file': (image_file.name, image_file.read(), image_file.content_type)}
            response = requests.post(api_url, headers=headers_fallback, data=data, files=files, timeout=60)
        
        print("DEBUG: PicWish response status:", response.status_code)
        print("DEBUG: PicWish response body:", response.text)
        
        # Handle different HTTP status codes
        if response.status_code == 200:
            # Parse JSON response
            result = response.json()
            
            # Check if API returned success status
            if result.get('status') == 200:
                enhanced_url = result.get('data', {}).get('image')
                if enhanced_url:
                    return JsonResponse({
                        'success': True, 
                        'enhanced_image_url': enhanced_url
                    })
                else:
                    return JsonResponse({
                        'success': False, 
                        'error': 'No image URL in API response'
                    })
            else:
                # API returned error in response body
                error_msg = result.get('message') or result.get('msg') or 'Enhancement failed'
                return JsonResponse({
                    'success': False, 
                    'error': error_msg
                })
        
        elif response.status_code == 401:
            # Unauthorized - Invalid or expired API key
            try:
                error_detail = response.json()
                error_msg = error_detail.get('message') or error_detail.get('msg') or 'Unauthorized API Key'
            except:
                error_msg = 'Unauthorized API Key. Please check your key.'
            return JsonResponse({
                'success': False, 
                'error': 'Unauthorized API Key. Please check your key.'
            }, status=401)
        
        elif response.status_code == 429:
            # Rate limit exceeded
            return JsonResponse({
                'success': False, 
                'error': 'Rate limit exceeded. Please try again later.'
            }, status=429)
        
        elif response.status_code == 400:
            # Bad request - invalid parameters or file format
            try:
                error_detail = response.json()
                error_msg = error_detail.get('message') or error_detail.get('msg') or 'Invalid request'
            except:
                error_msg = 'Invalid request. Please check your image file.'
            return JsonResponse({
                'success': False, 
                'error': error_msg
            }, status=400)
        
        else:
            # Other HTTP errors
            try:
                error_detail = response.json()
                error_msg = error_detail.get('message') or error_detail.get('msg') or f'API request failed with status {response.status_code}'
            except:
                error_msg = f'API request failed with status {response.status_code}'
            return JsonResponse({
                'success': False, 
                'error': error_msg
            }, status=response.status_code)
            
    except requests.exceptions.Timeout:
        return JsonResponse({
            'success': False, 
            'error': 'Request timeout. The API took too long to respond.'
        }, status=504)
    
    except requests.exceptions.RequestException as e:
        return JsonResponse({
            'success': False, 
            'error': f'Network error: {str(e)}'
        }, status=500)
    
    except Exception as e:
        # Log the exception for debugging
        import traceback
        print(f"Enhance photo error: {str(e)}")
        print(traceback.format_exc())
        return JsonResponse({
            'success': False, 
            'error': f'Server error: {str(e)}'
        }, status=500)


@csrf_exempt
def enhance_passport_photo_with_picwish(request):
    """
    PicWish Passport Photo Enhancement API endpoint (Async).
    Supports both task creation and task polling.
    - POST with file: Creates async enhancement task, returns task_id
    - GET with task_id: Polls task status, returns enhanced image URL when ready
    """
    # Get API key from settings
    api_key = getattr(settings, 'PICWISH_API', None)
    if not api_key:
        print("DEBUG: PICWISH_API not found in settings")
        return JsonResponse({'success': False, 'error': 'API key not configured'}, status=500)
    
    print("DEBUG: Using PicWish API key:", api_key[:10] + "..." if len(api_key) > 10 else api_key)
    
    api_url = 'https://techhk.aoscdn.com/api/tasks/visual/scale'
    
    if request.method == 'POST':
        # Create enhancement task
        try:
            # Get uploaded image file
            image_file = request.FILES.get('image_file')
            if not image_file:
                return JsonResponse({'success': False, 'error': 'No image uploaded'}, status=400)
            
            # Prepare request headers - Try Authorization format first, fallback to X-API-KEY if needed
            # Note: PicWish API may accept either format, but Authorization: API-Key is the standard
            headers = {
                "Authorization": f"API-Key {api_key}"
            }
            
            # Prepare form data for async task creation
            data = {
                'sync': 0,  # Async request
                'type': request.POST.get('type', 'face'),  # Enhancement type: 'face' for passport photos
                'return_type': 1  # Return type: 1 = image URL, 2 = base64
            }
            
            print("DEBUG: Headers:", headers)
            print("DEBUG: Payload:", data)
            
            # Prepare file upload
            files = {'image_file': (image_file.name, image_file.read(), image_file.content_type)}
            
            # Make API request to PicWish
            response = requests.post(api_url, headers=headers, data=data, files=files, timeout=60)
            
            # If 401 with Authorization header, try X-API-KEY as fallback
            if response.status_code == 401:
                print("DEBUG: 401 with Authorization header, trying X-API-KEY fallback")
                headers_fallback = {'X-API-KEY': api_key}
                # Reset file pointer
                image_file.seek(0)
                files = {'image_file': (image_file.name, image_file.read(), image_file.content_type)}
                response = requests.post(api_url, headers=headers_fallback, data=data, files=files, timeout=60)
            
            print("DEBUG: PicWish response status:", response.status_code)
            print("DEBUG: PicWish response body:", response.text)
            
            if response.status_code == 200:
                result = response.json()
                if result.get('status') == 200 and result.get('data', {}).get('task_id'):
                    task_id = result['data']['task_id']
                    return JsonResponse({
                        'success': True,
                        'data': {
                            'task_id': task_id
                        }
                    })
                else:
                    error_msg = result.get('msg') or result.get('message') or 'Failed to create task'
                    return JsonResponse({'success': False, 'error': error_msg}, status=500)
            elif response.status_code == 401:
                error_detail = response.json() if response.text else {}
                error_msg = error_detail.get('message') or error_detail.get('msg') or 'Unauthorized API Key'
                print(f"DEBUG: 401 Error - {error_msg}")
                return JsonResponse({'success': False, 'error': error_msg}, status=401)
            else:
                try:
                    error_detail = response.json()
                    error_msg = error_detail.get('message') or error_detail.get('msg') or f'API request failed with status {response.status_code}'
                except:
                    error_msg = f'API request failed with status {response.status_code}'
                return JsonResponse({'success': False, 'error': error_msg}, status=response.status_code)
                
        except Exception as e:
            import traceback
            print(f"DEBUG: Exception in enhance_passport_photo_with_picwish (POST): {str(e)}")
            print(traceback.format_exc())
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    
    elif request.method == 'GET':
        # Poll task status
        try:
            task_id = request.GET.get('task_id')
            if not task_id:
                return JsonResponse({'success': False, 'error': 'task_id required'}, status=400)
            
            # Prepare request headers - Try Authorization format first, fallback to X-API-KEY if needed
            headers = {
                "Authorization": f"API-Key {api_key}"
            }
            
            # Poll task status
            poll_url = f"{api_url}/{task_id}"
            response = requests.get(poll_url, headers=headers, timeout=60)
            
            # If 401 with Authorization header, try X-API-KEY as fallback
            if response.status_code == 401:
                print("DEBUG: 401 with Authorization header, trying X-API-KEY fallback for polling")
                headers_fallback = {'X-API-KEY': api_key}
                response = requests.get(poll_url, headers=headers_fallback, timeout=60)
            
            print(f"DEBUG: Polling task {task_id} - status: {response.status_code}")
            print(f"DEBUG: Poll response body: {response.text}")
            
            if response.status_code == 200:
                result = response.json()
                if result.get('status') == 200 and result.get('data'):
                    data = result['data']
                    if data.get('progress') == 100 and data.get('state') == 1 and data.get('image'):
                        # Task completed successfully
                        return JsonResponse({
                            'success': True,
                            'data': {
                                'progress': 100,
                                'state': 1,
                                'image': data['image']
                            }
                        })
                    elif data.get('state') == 2:
                        # Task failed
                        return JsonResponse({
                            'success': False,
                            'error': 'Enhancement failed',
                            'data': data
                        })
                    else:
                        # Task still in progress
                        return JsonResponse({
                            'success': True,
                            'data': data
                        })
                else:
                    error_msg = result.get('msg') or result.get('message') or 'Failed to get task status'
                    return JsonResponse({'success': False, 'error': error_msg}, status=500)
            elif response.status_code == 401:
                error_detail = response.json() if response.text else {}
                error_msg = error_detail.get('message') or error_detail.get('msg') or 'Unauthorized API Key'
                return JsonResponse({'success': False, 'error': error_msg}, status=401)
            else:
                try:
                    error_detail = response.json()
                    error_msg = error_detail.get('message') or error_detail.get('msg') or f'API request failed with status {response.status_code}'
                except:
                    error_msg = f'API request failed with status {response.status_code}'
                return JsonResponse({'success': False, 'error': error_msg}, status=response.status_code)
                
        except Exception as e:
            import traceback
            print(f"DEBUG: Exception in enhance_passport_photo_with_picwish (GET): {str(e)}")
            print(traceback.format_exc())
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    
    else:
        return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)


# ─────────────────────────────────────────────────────────────
# FILE UPLOAD TO CLOUDFLARE R2
# ─────────────────────────────────────────────────────────────

# ─────────────────────────────────────────────────────────────
# FORGOT PASSWORD SERVICE
# ─────────────────────────────────────────────────────────────

def forgot_password_page(request):
    """Render the forgot password page"""
    return render(request, 'forgot_password.html')
@csrf_exempt
def forgot_password(request):
    """
    Legacy endpoint kept for compatibility.
    Now only validates that the vendor exists and returns success;
    password reset itself is handled directly via reset_password without OTP.
    """
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            email = data.get('email')
            
            if not email:
                return JsonResponse({'success': False, 'error': 'Email is required'})
            
            # Use D1 vendor_register_details as the source of truth
            api_url = getattr(settings, 'WORKER_API_URL', '')
            api_key = getattr(settings, 'WORKER_API_KEY', '')

            if not api_url or not api_key:
                print("❌ Worker API not configured for forgot_password")
                return JsonResponse({'success': False, 'error': 'Server configuration error. Please contact support.'}, status=500)

            if '/add-contact' in api_url:
                worker_endpoint = api_url.replace('/add-contact', '/get-vendor-by-email')
            elif '/add-vendor-register' in api_url:
                worker_endpoint = api_url.replace('/add-vendor-register', '/get-vendor-by-email')
            else:
                worker_endpoint = api_url.rstrip('/') + '/get-vendor-by-email'

            resp = requests.post(
                worker_endpoint,
                json={'email': email},
                headers={
                    'Content-Type': 'application/json',
                    'x-api-key': api_key
                },
                timeout=10
            )

            if resp.status_code == 404:
                return JsonResponse({'success': False, 'error': 'Vendor not found with this email'})
            if resp.status_code != 200:
                print(f"❌ Worker error in forgot_password: {resp.status_code} {resp.text[:300]}")
                return JsonResponse({'success': False, 'error': 'Error verifying vendor account'}, status=500)

            payload = resp.json()
            if not payload.get('success'):
                return JsonResponse({'success': False, 'error': payload.get('error', 'Vendor not found')}, status=404)

            return JsonResponse({
                'success': True,
                'message': 'Vendor verified. You can now reset your password directly.'
            })
            
        except Exception as e:
            print(f"Error in forgot_password: {str(e)}")
            return JsonResponse({'success': False, 'error': 'Internal server error'})
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'})

@csrf_exempt
def verify_reset_code(request):
    """
    Legacy endpoint kept for compatibility with older clients.
    Always returns success once the vendor email is valid.
    """
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            email = data.get('email')
            
            if not email:
                return JsonResponse({'success': False, 'error': 'Email is required'})

            # Simply confirm vendor exists in D1
            api_url = getattr(settings, 'WORKER_API_URL', '')
            api_key = getattr(settings, 'WORKER_API_KEY', '')

            if not api_url or not api_key:
                print("❌ Worker API not configured for verify_reset_code")
                return JsonResponse({'success': False, 'error': 'Server configuration error. Please contact support.'}, status=500)

            if '/add-contact' in api_url:
                worker_endpoint = api_url.replace('/add-contact', '/get-vendor-by-email')
            elif '/add-vendor-register' in api_url:
                worker_endpoint = api_url.replace('/add-vendor-register', '/get-vendor-by-email')
            else:
                worker_endpoint = api_url.rstrip('/') + '/get-vendor-by-email'

            resp = requests.post(
                worker_endpoint,
                json={'email': email},
                headers={
                    'Content-Type': 'application/json',
                    'x-api-key': api_key
                },
                timeout=10
            )

            if resp.status_code == 404:
                return JsonResponse({'success': False, 'error': 'Vendor not found with this email'})
            if resp.status_code != 200:
                print(f"❌ Worker error in verify_reset_code: {resp.status_code} {resp.text[:300]}")
                return JsonResponse({'success': False, 'error': 'Error verifying vendor account'}, status=500)

            payload = resp.json()
            if not payload.get('success'):
                return JsonResponse({'success': False, 'error': payload.get('error', 'Vendor not found')}, status=404)

            return JsonResponse({'success': True, 'message': 'Vendor verified successfully'})

        except Exception as e:
            print(f"Error in verify_reset_code: {str(e)}")
            return JsonResponse({'success': False, 'error': 'Internal server error'})
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'})

@csrf_exempt
def reset_password(request):
    """
    Reset vendor password without OTP, updating both D1 vendor_register_details
    and legacy R2 JSON files for backward compatibility.
    """
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            email = data.get('email')
            new_password = data.get('new_password')
            
            if not all([email, new_password]):
                return JsonResponse({'success': False, 'error': 'Email and new password are required'})
            
            if len(new_password) < 8:
                return JsonResponse({'success': False, 'error': 'Password must be at least 8 characters long'})

            # Hash the new password using the same scheme as vendor_login
            hashed_password = make_password(new_password)

            # First, update D1 vendor_register_details via Worker API
            api_url = getattr(settings, 'WORKER_API_URL', '')
            api_key = getattr(settings, 'WORKER_API_KEY', '')

            if not api_url or not api_key:
                print("❌ Worker API not configured for reset_password")
                return JsonResponse({'success': False, 'error': 'Server configuration error. Please contact support.'}, status=500)

            if '/add-contact' in api_url:
                worker_endpoint = api_url.replace('/add-contact', '/update-vendor-password')
            elif '/add-vendor-register' in api_url:
                worker_endpoint = api_url.replace('/add-vendor-register', '/update-vendor-password')
            else:
                worker_endpoint = api_url.rstrip('/') + '/update-vendor-password'

            try:
                resp = requests.post(
                    worker_endpoint,
                    json={
                        'email': email,
                        'new_password_hash': hashed_password,
                        'source': 'django_reset_password'
                    },
                    headers={
                        'Content-Type': 'application/json',
                        'x-api-key': api_key
                    },
                    timeout=10
                )

                if resp.status_code == 404:
                    return JsonResponse({'success': False, 'error': 'Vendor not found with this email'}, status=404)
                if resp.status_code != 200:
                    print(f"❌ Worker error in update-vendor-password: {resp.status_code} {resp.text[:300]}")
                    return JsonResponse({'success': False, 'error': 'Failed to update password in database'}, status=500)

                payload = resp.json()
                if not payload.get('success'):
                    return JsonResponse({'success': False, 'error': payload.get('error', 'Failed to update password')}, status=500)
            except Exception as worker_exc:
                print(f"❌ Exception calling update-vendor-password: {worker_exc}")
                return JsonResponse({'success': False, 'error': 'Failed to update password in database'}, status=500)

            # Best-effort: update legacy R2 JSON so older flows (if any) remain consistent
            try:
                s3 = boto3.client('s3',
                    aws_access_key_id=settings.R2_ACCESS_KEY,
                    aws_secret_access_key=settings.R2_SECRET_KEY,
                    endpoint_url=settings.R2_ENDPOINT,
                    region_name='auto'
                )

                reg_key = f'vendor_register_details/{sanitize_email(email)}/registration_details.json'
                try:
                    response = s3.get_object(Bucket=settings.R2_BUCKET, Key=reg_key)
                    vendor_data = json.loads(response['Body'].read().decode('utf-8'))
                except Exception:
                    vendor_data = {'email': email}

                vendor_data['hashed_password'] = hashed_password
                vendor_data['password_updated_at'] = datetime.datetime.now().isoformat()

                s3.put_object(
                    Bucket=settings.R2_BUCKET,
                    Key=reg_key,
                    Body=json.dumps(vendor_data),
                    ContentType='application/json'
                )

                login_key = f'vendor_register_details/{sanitize_email(email)}/login_details.json'
                try:
                    login_response = s3.get_object(Bucket=settings.R2_BUCKET, Key=login_key)
                    login_data = json.loads(login_response['Body'].read().decode('utf-8'))
                except Exception:
                    login_data = {'email': email}

                login_data['hashed_password'] = hashed_password
                login_data['last_password_reset'] = datetime.datetime.now().isoformat()

                s3.put_object(
                    Bucket=settings.R2_BUCKET,
                    Key=login_key,
                    Body=json.dumps(login_data),
                    ContentType='application/json'
                )
            except Exception as legacy_exc:
                # Log but don't fail the reset if D1 update already succeeded
                print(f"⚠️ Legacy R2 password sync failed for {email}: {legacy_exc}")

            return JsonResponse({
                'success': True,
                'message': 'Password reset successfully'
            })
            
        except Exception as e:
            print(f"Error in reset_password: {str(e)}")
            return JsonResponse({'success': False, 'error': 'Internal server error'})
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'})


# ─────────────────────────────────────────────────────────────
# ENHANCED PHOTO SERVICE
# ─────────────────────────────────────────────────────────────

# ─────────────────────────────────────────────────────────────
# EMAIL UTILITIES
# ─────────────────────────────────────────────────────────────

def send_password_reset_email(email, verification_code, vendor_name):
    """
    Send password reset verification email to vendor
    """
    try:
        subject = 'Password Reset Verification - PrintMax'
        
        # Render HTML email template
        html_message = render_to_string('emails/password_reset.html', {
            'vendor_name': vendor_name,
            'verification_code': verification_code,
            'email': email
        })
        
        # Create plain text version
        plain_message = strip_tags(html_message)
        
        # Send email
        send_mail(
            subject=subject,
            message=plain_message,
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],
            html_message=html_message,
            fail_silently=False,
        )
        
        print(f"✅ Password reset email sent successfully to {email}")
        return True
        
    except Exception as e:
        print(f"❌ Error sending password reset email to {email}: {str(e)}")
        return False

def send_welcome_email(email, vendor_name, password, vendor_id):
    """
    Send welcome email to new vendors with login credentials
    Uses direct SMTP connection with SSL (port 465) and TLS (port 587) only
    """
    subject = 'Welcome to PrintMax - Your Vendor Account is Ready!'
    
    # HTML email template (EXACT CONTENT PRESERVED)
    html_message = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Welcome to PrintMax</title>
            <meta http-equiv="Content-Type" content="text/html; charset=UTF-8">
            <meta name="format-detection" content="telephone=no">
            <meta name="format-detection" content="date=no">
            <meta name="format-detection" content="address=no">
            <meta name="format-detection" content="email=no">
            <!--[if mso]>
            <noscript>
                <xml>
                    <o:OfficeDocumentSettings>
                        <o:PixelsPerInch>96</o:PixelsPerInch>
                    </o:OfficeDocumentSettings>
                </xml>
            </noscript>
            <![endif]-->
            <style>
                body {{
                    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
                    line-height: 1.6;
                    color: #333;
                    max-width: 600px;
                    margin: 0 auto;
                    padding: 20px;
                    background-color: #f8f9fa;
                }}
                /* Prevent email clients from treating content as quoted */
                .email-content {{
                    display: block !important;
                    visibility: visible !important;
                    opacity: 1 !important;
                }}
                /* Override any potential quoted text hiding */
                [data-ogsc] {{
                    display: block !important;
                }}
                /* Ensure all content is visible */
                * {{
                    max-height: none !important;
                    overflow: visible !important;
                }}
                .container {{
                    background: white;
                    border-radius: 10px;
                    padding: 30px;
                    box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
                }}
                .header {{
                    text-align: center;
                    margin-bottom: 30px;
                    padding-bottom: 20px;
                    border-bottom: 3px solid #1976d2;
                }}
                .logo {{
                    font-size: 28px;
                    font-weight: bold;
                    color: #1976d2;
                    margin-bottom: 10px;
                }}
                .credentials-box {{
                    background: linear-gradient(135deg, #e3f2fd 0%, #f3e5f5 100%);
                    border: 2px solid #1976d2;
                    border-radius: 8px;
                    padding: 20px;
                    margin: 20px 0;
                    text-align: center;
                }}
                .credentials-title {{
                    font-size: 18px;
                    font-weight: bold;
                    color: #1976d2;
                    margin-bottom: 15px;
                }}
                .credential-item {{
                    background: white;
                    border: 1px solid #ddd;
                    border-radius: 5px;
                    padding: 12px;
                    margin: 10px 0;
                    font-family: 'Courier New', monospace;
                    font-size: 16px;
                    font-weight: bold;
                    color: #2e7d32;
                }}
                .credential-label {{
                    font-size: 14px;
                    color: #666;
                    margin-bottom: 5px;
                    font-weight: normal;
                }}
                .steps {{
                    background: #f8f9fa;
                    border-radius: 8px;
                    padding: 20px;
                    margin: 20px 0;
                }}
                .step {{
                    margin: 15px 0;
                    padding: 10px;
                    background: white;
                    border-left: 4px solid #4caf50;
                    border-radius: 0 5px 5px 0;
                }}
                .step-number {{
                    font-weight: bold;
                    color: #4caf50;
                    margin-right: 10px;
                }}
                .footer {{
                    text-align: center;
                    margin-top: 30px;
                    padding-top: 20px;
                    border-top: 1px solid #eee;
                    color: #666;
                    font-size: 14px;
                }}
                .button {{
                    display: inline-block;
                    background: linear-gradient(135deg, #1976d2, #1565c0);
                    color: white;
                    padding: 12px 25px;
                    text-decoration: none;
                    border-radius: 5px;
                    font-weight: bold;
                    margin: 10px 5px;
                }}
                .warning {{
                    background: #fff3cd;
                    border: 1px solid #ffeaa7;
                    border-radius: 5px;
                    padding: 15px;
                    margin: 20px 0;
                    color: #856404;
                }}
            </style>
        </head>
        <body>
            <div class="email-content">
            <table width="100%" cellpadding="0" cellspacing="0" border="0" style="background-color: #f8f9fa;">
                <tr>
                    <td align="center" style="padding: 20px;">
                        <table class="container" cellpadding="0" cellspacing="0" border="0" style="background: white; border-radius: 10px; padding: 30px; box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1); max-width: 600px; width: 100%;">
                            <tr>
                                <td>
                                    <div class="header" style="text-align: center; margin-bottom: 30px; padding-bottom: 20px; border-bottom: 3px solid #1976d2;">
                                        <div class="logo" style="font-size: 28px; font-weight: bold; color: #1976d2; margin-bottom: 10px;">🖨️ PrintMax</div>
                                        <h2 style="margin: 0; color: #333;">Welcome to PrintMax, {vendor_name}!</h2>
                                        <p style="margin: 10px 0; color: #666;">Your vendor account has been successfully created and is ready to use.</p>
                                        <div style="background: linear-gradient(135deg, #4caf50 0%, #2e7d32 100%); color: white; padding: 10px 20px; border-radius: 25px; display: inline-block; margin-top: 15px; font-weight: bold;">
                                            ✅ Registration Completed Successfully
                                        </div>
                                        
                                        <!-- Login Credentials directly below success message -->
                                        <div style="background: linear-gradient(135deg, #e3f2fd 0%, #f3e5f5 100%); border: 2px solid #1976d2; border-radius: 8px; padding: 20px; margin: 20px 0; text-align: center;">
                                            <div style="font-size: 18px; font-weight: bold; color: #1976d2; margin-bottom: 15px;">🔐 Your Login Credentials</div>
                                            <p style="margin-bottom: 20px; color: #666;">Please save these credentials securely. You'll need them to access your vendor dashboard.</p>
                                            
                                            <div style="background: white; border: 1px solid #ddd; border-radius: 5px; padding: 12px; margin: 10px 0; font-family: 'Courier New', monospace; font-size: 16px; font-weight: bold; color: #2e7d32;">
                                                <div style="font-size: 14px; color: #666; margin-bottom: 5px; font-weight: normal;">Vendor Email:</div>
                                                {email}
                                            </div>
                                            
                                            <div style="background: white; border: 1px solid #ddd; border-radius: 5px; padding: 12px; margin: 10px 0; font-family: 'Courier New', monospace; font-size: 16px; font-weight: bold; color: #2e7d32;">
                                                <div style="font-size: 14px; color: #666; margin-bottom: 5px; font-weight: normal;">Password:</div>
                                                {password}
                                            </div>
                                        </div>
                                    </div>
                                </td>
                            </tr>

                            <tr>
                                <td>
                                    <div class="warning" style="background: #fff3cd; border: 1px solid #ffeaa7; border-radius: 5px; padding: 15px; margin: 20px 0; color: #856404;">
                                        <strong>⚠️ Important Security Notice:</strong><br>
                                        This is your PrintMax vendor password. Please keep it secure and do not share it with anyone. 
                                        Use these credentials to login to your vendor dashboard.
                                    </div>

                                    <div class="steps" style="background: #f8f9fa; border-radius: 8px; padding: 20px; margin: 20px 0;">
                                        <h3 style="color: #1976d2; margin-top: 0;">Next Steps:</h3>
                                        
                                        <div class="step" style="margin: 15px 0; padding: 10px; background: white; border-left: 4px solid #4caf50; border-radius: 0 5px 5px 0;">
                                            <span class="step-number" style="font-weight: bold; color: #4caf50; margin-right: 10px;">1.</span>
                                            <strong>Access Your Dashboard:</strong> Click the "Vendor" button in the top navigation on the homepage to open your Vendor Dashboard.
                                        </div>
                                        
                                        <div class="step" style="margin: 15px 0; padding: 10px; background: white; border-left: 4px solid #4caf50; border-radius: 0 5px 5px 0;">
                                            <span class="step-number" style="font-weight: bold; color: #4caf50; margin-right: 10px;">2.</span>
                                            <strong>Watch the Demo Video:</strong> Please watch the training video before using the system to print customer documents.
                                        </div>
                                        
                                        <div class="step" style="margin: 15px 0; padding: 10px; background: white; border-left: 4px solid #4caf50; border-radius: 0 5px 5px 0;">
                                            <span class="step-number" style="font-weight: bold; color: #4caf50; margin-right: 10px;">3.</span>
                                            <strong>Start Receiving Orders:</strong> Once your account is activated, you will start receiving customer orders.
                                        </div>
                                    </div>

                                    <div style="text-align: center; margin: 30px 0;">
                                        <a href="https://printmax.onrender.com/" style="display: inline-block; background: linear-gradient(135deg, #1976d2, #1565c0); color: white; padding: 12px 25px; text-decoration: none; border-radius: 5px; font-weight: bold; margin: 10px 5px;">Login to Dashboard</a>
                                    </div>

                                    <div class="footer" style="text-align: center; margin-top: 30px; padding-top: 20px; border-top: 1px solid #eee; color: #666; font-size: 14px;">
                                        <p>Thank you for choosing PrintMax as your printing partner!</p>
                                        <p>If you have any questions, please contact our support team.</p>
                                        <p style="font-size: 12px; color: #999;">
                                            This email was sent to {email}. Please do not reply to this email.
                                        </p>
                                    </div>
                                </td>
                            </tr>
                        </table>
                    </td>
                </tr>
            </table>
            </div>
        </body>
        </html>
        """
    
    # Plain text version (EXACT CONTENT PRESERVED)
    plain_message = f"""
        Welcome to PrintMax, {vendor_name}!

        Your vendor account has been successfully created and is ready to use.

        LOGIN CREDENTIALS:
        ==================
        Vendor Email: {email}
        Password: {password}

        IMPORTANT: This is your PrintMax vendor password. Please keep it secure and do not share it with anyone. 
        Use these credentials to login to your vendor dashboard.

        NEXT STEPS:
        1. Access Your Dashboard: Click the "Vendor" button in the top navigation on the homepage to open your Vendor Dashboard.
        2. Watch the Demo Video: Please watch the training video before using the system to print customer documents.
        3. Start Receiving Orders: Once your account is activated, you will start receiving customer orders.

        Quick Link:
        - Login to Dashboard: https://printmax.onrender.com/

        Thank you for choosing PrintMax as your printing partner!

        If you have any questions, please contact our support team.

        This email was sent to {email}. Please do not reply to this email.
        """
    
    # Validate email configuration first
    if not settings.EMAIL_HOST_USER or not settings.EMAIL_HOST_PASSWORD:
        print(f"⚠️ Email credentials not configured. Skipping email send to {email}")
        print(f"   EMAIL_HOST_USER: {'Set' if settings.EMAIL_HOST_USER else 'Missing'}")
        print(f"   EMAIL_HOST_PASSWORD: {'Set' if settings.EMAIL_HOST_PASSWORD else 'Missing'}")
        return False
    
    # CRITICAL: FROM address MUST match EMAIL_HOST_USER for authentication
    # Extract email address - use authenticated email as FROM
    from_email_addr = settings.EMAIL_HOST_USER
    
    # For display, we can use a formatted version, but SMTP envelope must use authenticated email
    display_from = f'PrintMax <{from_email_addr}>' if from_email_addr else 'PrintMax'
    
    print(f"📧 Email config - Host: {settings.EMAIL_HOST}")
    print(f"📧 Authenticated as: {from_email_addr}")
    print(f"📧 Sending to: {email}")
    
    # Direct SMTP connection using SSL (port 465) and TLS (port 587) only
    try:
        import smtplib
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart
        import ssl
        
        # Create message
        msg = MIMEMultipart('alternative')
        msg['Subject'] = subject
        msg['From'] = display_from  # Display name in header
        msg['To'] = email
        
        # Add headers to prevent email clients from treating content as quoted
        msg['X-Mailer'] = 'PrintMax Vendor System'
        msg['X-Priority'] = '1'
        msg['X-MSMail-Priority'] = 'High'
        msg['Importance'] = 'high'
        msg['X-Original-Sender'] = from_email_addr
        msg['X-Auto-Response-Suppress'] = 'All'
        msg['Precedence'] = 'bulk'
        msg['X-Entity-Ref-ID'] = f'vendor-welcome-{vendor_id}'
        
        # Add text and HTML parts
        part1 = MIMEText(plain_message, 'plain')
        part2 = MIMEText(html_message, 'html')
        msg.attach(part1)
        msg.attach(part2)
        
        # Try SSL first (port 465) - Primary configuration
        try:
            print("🔄 Attempting SMTP connection with SSL (port 465)...")
            context = ssl.create_default_context()
            server = smtplib.SMTP_SSL(
                settings.EMAIL_HOST, 
                465, 
                timeout=getattr(settings, 'EMAIL_TIMEOUT', 30),
                context=context
            )
            server.set_debuglevel(0)
            server.ehlo()
            server.login(settings.EMAIL_HOST_USER, settings.EMAIL_HOST_PASSWORD)
            server.sendmail(from_email_addr, [email], msg.as_string())
            server.quit()
            print(f"✅ Registration email sent successfully via SMTP (SSL port 465) to {email}")
            return True
        except Exception as e_ssl:
            print(f"❌ SSL (465) failed: {str(e_ssl)}")
            
            # Try TLS (port 587) as fallback
            try:
                print("🔄 Attempting SMTP connection with TLS (port 587)...")
                context = ssl.create_default_context()
                server = smtplib.SMTP(settings.EMAIL_HOST, 587, timeout=getattr(settings, 'EMAIL_TIMEOUT', 30))
                server.set_debuglevel(0)
                server.ehlo()
                server.starttls(context=context)
                server.ehlo()
                server.login(settings.EMAIL_HOST_USER, settings.EMAIL_HOST_PASSWORD)
                server.sendmail(from_email_addr, [email], msg.as_string())
                server.quit()
                print(f"✅ Registration email sent successfully via SMTP (TLS port 587) to {email}")
                return True
            except Exception as e_tls:
                print(f"❌ TLS (587) failed: {str(e_tls)}")
                print(f"❌ All email methods failed. Check SMTP credentials and server configuration.")
                return False
                    
    except Exception as e:
        print(f"❌ Email send failed: {str(e)}")
        import traceback
        print(f"❌ Full traceback: {traceback.format_exc()}")
        return False


# ─────────────────────────────────────────────────────────────
# LOGOUT SERVICE
# ─────────────────────────────────────────────────────────────

def logout_view(request):
    """
    Handle logout for both users and vendors
    """
    try:
        # Clear all session data
        request.session.flush()
        
        # Clear any authentication cookies
        response = redirect('home')
        response.delete_cookie('sessionid')
        response.delete_cookie('csrftoken')
        
        print("✅ User logged out successfully")
        return response
        
    except Exception as e:
        print(f"❌ Error during logout: {str(e)}")
        # Even if there's an error, redirect to home
        return redirect('home')


# ─────────────────────────────────────────────────────────────
# FORGOT PASSWORD SERVICE
# ─────────────────────────────────────────────────────────────

def get_vendor_email_by_vendor_id(vendor_id):
    """
    Given a vendor_id, search all vendor registration details in R2 and return the corresponding email.
    """
    import boto3, json
    from django.conf import settings

    s3 = boto3.client('s3',
        aws_access_key_id=settings.R2_ACCESS_KEY,
        aws_secret_access_key=settings.R2_SECRET_KEY,
        endpoint_url=settings.R2_ENDPOINT,
        region_name='auto'
    )

    try:
        # List all registration details
        objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_register_details/')
        for obj in objects.get("Contents", []):
            if obj["Key"].endswith('/registration_details.json'):
                response = s3.get_object(Bucket=settings.R2_BUCKET, Key=obj["Key"])
                vendor_data = json.loads(response['Body'].read().decode('utf-8'))
                if vendor_data.get('vendor_id') == vendor_id:
                    return vendor_data.get('vendor_email')
    except Exception as e:
        print(f"Error finding vendor email for vendor_id {vendor_id}: {str(e)}")
    return None
# ─────────────────────────────────────────────────────────────
# Token management helpers
# ─────────────────────────────────────────────────────────────

def free_token_in_vendor_pool(vendor_email: str, token: str) -> bool:
    """Mark a token as 'free' in Vendor_tokens D1 table based on vendor_email"""
    try:
        if not token:
            return False

        # Free token in D1 Vendor_tokens table (primary storage)
        d1_freed = False
        try:
            api_url = getattr(settings, 'WORKER_API_URL', '')
            api_key = getattr(settings, 'WORKER_API_KEY', '')
            
            if api_url and api_key:
                token_number = int(token) if token.isdigit() else None
                if token_number:
                    worker_endpoint = api_url.rstrip('/') + '/free-vendor-token'
                    worker_payload = {
                        'vendor_email': vendor_email,
                        'token_number': token_number
                    }
                    
                    resp = requests.post(
                        worker_endpoint,
                        json=worker_payload,
                        headers={
                            'x-api-key': api_key,
                            'Content-Type': 'application/json'
                        },
                        timeout=10
                    )
                    
                    if resp.status_code == 200:
                        data = resp.json()
                        if data.get('success'):
                            d1_freed = True
                            print(f"✅ Freed token {token_number} in D1 Vendor_tokens table for vendor {vendor_email}")
                        else:
                            print(f"⚠️ D1 API returned error: {data.get('error', 'Unknown error')}")
                    else:
                        print(f"⚠️ Failed to free token in D1: {resp.status_code} - {resp.text[:200]}")
                else:
                    print(f"⚠️ Token '{token}' is not a valid integer, skipping D1 update")
            else:
                print("⚠️ Worker API not configured - skipping D1 token update")
        except Exception as e:
            print(f"⚠️ Error freeing token in D1 for {vendor_email}: {str(e)}")

        # Return True if D1 update was successful
        return d1_freed
    except Exception as e:
        print(f"❌ Error freeing token for {vendor_email}: {str(e)}")
        return False

def get_token_from_file_metadata(filename: str, vendor_id: str) -> str:
    """Fetch the 'token' value from object metadata for the given filename."""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        possible_paths = [
            f'vendor_print_jobs/{vendor_id}/{filename}',
            f'vendor_manual_print_jobs/{vendor_id}/{filename}',
            f'vendor_register_details/{vendor_id}/firozshop/{filename}',
            f'vendor_register_details/{sanitize_email(vendor_id)}/firozshop/{filename}',
        ]

        for path in possible_paths:
            try:
                head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=path)
                metadata = head_response.get('Metadata', {})
                token = metadata.get('token', '')
                if token:
                    return token
            except Exception:
                continue

        # Search under users as last resort
        users_prefix = 'users/'
        try:
            users_response = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=users_prefix)
            for obj in users_response.get('Contents', []):
                key = obj['Key']
                if key.endswith(f"/{filename}"):
                    head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                    metadata = head_response.get('Metadata', {})
                    token = metadata.get('token', '')
                    if token:
                        return token
        except Exception:
            pass

        return ''
    except Exception as e:
        print(f"❌ Error fetching token for {filename}: {str(e)}")
        return ''

def get_vendor_coordinates_from_email(vendor_email):
    """
    Get vendor coordinates from vendor_register_details/{vendor_email}/registration_details.json
    """
    s3 = boto3.client('s3',
                      aws_access_key_id=settings.R2_ACCESS_KEY,
                      aws_secret_access_key=settings.R2_SECRET_KEY,
                      endpoint_url=settings.R2_ENDPOINT,
                      region_name='auto')
    
    try:
        # Construct the key path
        sanitized_email = sanitize_email(vendor_email)
        key = f'vendor_register_details/{sanitized_email}/registration_details.json'
        
        # Get the vendor registration details
        response = s3.get_object(Bucket=settings.R2_BUCKET, Key=key)
        vendor_data = json.loads(response['Body'].read().decode('utf-8'))
        
        return {
            'latitude': vendor_data.get('latitude'),
            'longitude': vendor_data.get('longitude'),
            'vendor_name': vendor_data.get('vendor_name'),
            'city': vendor_data.get('city'),
            'shop_address': vendor_data.get('shop_address')
        }
    except Exception as e:
        print(f"Error getting vendor coordinates for {vendor_email}: {str(e)}")
        return None
def get_vendor_coordinates(request):
    """
    Return vendor coordinates as JSON for the map, non-blocking for dashboard load with caching.
    """
    import boto3, json
    from django.conf import settings
    import time
    
    # Check cache first
    cache_key = "vendor_coordinates_cache"
    current_time = time.time()
    
    if (hasattr(get_vendor_coordinates, '_cache') and 
        cache_key in get_vendor_coordinates._cache and 
        cache_key in get_vendor_coordinates._cache_timestamps and 
        current_time - get_vendor_coordinates._cache_timestamps[cache_key] < 300):  # 5 minutes
        
        print(f"⚡ CACHE HIT: Using cached vendor coordinates")
        return JsonResponse({'coordinates': get_vendor_coordinates._cache[cache_key]})
    
    coordinates = []
    try:
        s3 = boto3.client('s3',
            aws_access_key_id=settings.R2_ACCESS_KEY,
            aws_secret_access_key=settings.R2_SECRET_KEY,
            endpoint_url=settings.R2_ENDPOINT,
            region_name='auto')
        objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='vendor_register_details/')
        for obj in objects.get("Contents", []):
            if obj["Key"].endswith('/registration_details.json'):
                try:
                    response = s3.get_object(Bucket=settings.R2_BUCKET, Key=obj["Key"])
                    vendor_data = json.loads(response['Body'].read().decode('utf-8'))
                    lat = vendor_data.get('latitude')
                    lng = vendor_data.get('longitude')
                    if lat and lng:
                        coordinates.append({
                            'vendor_id': vendor_data.get('vendor_id'),
                            'vendor_email': vendor_data.get('vendor_email'),
                            'latitude': lat,
                            'longitude': lng,
                            'shop_name': vendor_data.get('vendor_name'),
                            'city': vendor_data.get('city')
                        })
                except Exception as e:
                    continue
        
        # Cache the result
        if not hasattr(get_vendor_coordinates, '_cache'):
            get_vendor_coordinates._cache = {}
            get_vendor_coordinates._cache_timestamps = {}
        get_vendor_coordinates._cache[cache_key] = coordinates
        get_vendor_coordinates._cache_timestamps[cache_key] = current_time
        
        return JsonResponse({'coordinates': coordinates})
    except Exception as e:
        return JsonResponse({'coordinates': [], 'error': str(e)})

@csrf_exempt
def get_distance(request):
    """
    Secure backend proxy for Google Distance Matrix API.
    Uses server-side API key (not referrer-restricted).
    """
    # Validate that server-side API key is configured
    if not settings.GOOGLE_MAPS_API:
        error_msg = "GOOGLE_MAPS_API is not configured. Distance calculation unavailable."
        print(f"⚠️ WARNING: {error_msg}")
        # Return success response but with error status so frontend can handle gracefully
        return JsonResponse({
            'error': error_msg,
            'status': 'ERROR',
            'destination_addresses': [],
            'origin_addresses': [],
            'rows': []
        }, status=200)
    
    user_lat = request.GET.get("user_lat")
    user_lng = request.GET.get("user_lng")
    shop_lat = request.GET.get("shop_lat")
    shop_lng = request.GET.get("shop_lng")
    
    url = "https://maps.googleapis.com/maps/api/distancematrix/json"
    params = {
        "origins": f"{user_lat},{user_lng}",
        "destinations": f"{shop_lat},{shop_lng}",
        "mode": "driving",
        "key": settings.GOOGLE_MAPS_API
    }
    
    try:
        response = requests.get(url, params=params, timeout=10)
        response_data = response.json()
        
        # Temporary logging to verify API response (remove after confirming it works)
        print("GOOGLE RESPONSE:", response_data)
        
        return JsonResponse(response_data)
    except Exception as e:
        print(f"❌ Distance Matrix API Error: {str(e)}")
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
def get_accurate_location(request):
    """
    Backend proxy for Google Geolocation API.
    Called by frontend (userdashboard & vendor_register) to get accurate lat/lng.
    Uses server-side API key (not referrer-restricted).
    """
    if request.method != "POST":
        return JsonResponse({"success": False, "message": "POST required"}, status=405)
    
    # Validate that server-side API key is configured
    if not settings.GOOGLE_MAPS_API:
        error_msg = "GOOGLE_MAPS_API is not configured. Location services will use browser geolocation."
        print(f"⚠️ WARNING: {error_msg}")
        # Return success: false but with 200 status so frontend can handle gracefully
        return JsonResponse({
            "success": False,
            "message": error_msg,
            "fallback_to_browser": True
        }, status=200)
    
    url = f"https://www.googleapis.com/geolocation/v1/geolocate?key={settings.GOOGLE_MAPS_API}"
    
    try:
        google_response = requests.post(url, json={}, timeout=10)
        data = google_response.json()
        
        # 🔥 DEBUG (DO NOT REMOVE) - Print RAW Google API response
        print("📡 RAW GOOGLE RESPONSE:", data)
        
        if "location" not in data:
            print(f"❌ Google Geolocation API failed: {data}")
            return JsonResponse({
                "success": False,
                "message": "Google Geolocation failed",
                "google_error": data
            }, status=400)
        
        lat = data["location"]["lat"]
        lng = data["location"]["lng"]
        accuracy = data.get("accuracy")
        
        print(f"✅ Google Geolocation API success: lat={lat}, lng={lng}, accuracy={accuracy}")
        
        return JsonResponse({
            "success": True,
            "latitude": lat,
            "longitude": lng,
            "accuracy": accuracy,
        })
    except Exception as e:
        print(f"❌ Google Geolocation API Error: {str(e)}")
        return JsonResponse({
            "success": False,
            "message": str(e)
        }, status=500)

def create_job_completion_notification(user_email, filename, token, vendor_name, service_type, completion_time):
    # Stub: implement notification logic if needed
    pass

@csrf_exempt
def mark_notification_read(request):
    # Stub: Mark a notification as read (expand as needed)
    if request.method == 'POST':
        return JsonResponse({'success': True, 'message': 'Notification marked as read'})
    return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)

@csrf_exempt
def get_user_notifications(request):
    # Stub: Return user notifications (expand as needed)
    if request.method == 'POST':
        notifications = [
            {'id': 1, 'message': 'Your print job is ready!', 'read': False}
        ]
        return JsonResponse({'success': True, 'notifications': notifications})
    return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)

def vendor_about(request):
    return render(request, 'vendor_about.html')

def vendor_about(request):
    return render(request, 'vendor_about.html')


def _default_service_availability():
    """
    Returns the default service availability payload used when a vendor
    has not configured their services yet.
    """
    return {
        "digital_print": True,
        "project_binding": True,
        "gloss_printing": True,
        "jumbo_printing": True,
        "regular_print": True,
        "passport_print": True,
        "photo_print": True,
        "vendor_shop_avaliability": "online",
    }


def _reset_service_availability():
    """
    Returns the service availability payload with all services set to False (OFF).
    Used to force vendors to update their services upon login.
    """
    return {
        "digital_print": False,
        "project_binding": False,
        "gloss_printing": False,
        "jumbo_printing": False,
        "regular_print": False,
        "passport_print": False,
        "photo_print": False,
        "vendor_shop_avaliability": "online", # Keep shop online so they can access dashboard
    }

def _sync_vendor_service_on_login(vendor_email, vendor_id):
    """
    Ensure the Vendor_service_availability table has up-to-date rows whenever a vendor logs in.
    Preserves existing settings and only initializes defaults if missing.
    """
    if not vendor_email or not vendor_id:
        return

    try:
        existing_service_data = None
        try:
            endpoint, resp = post_to_worker('/get-vendor-service', {
                'vendor_email': vendor_email,
                'vendor_id': vendor_id
            })
            if resp.status_code == 200:
                payload = resp.json()
                if payload.get('success'):
                    service_payload = payload.get('service') or {}
                    existing_service_data = service_payload.get('service_data') or None
            elif resp.status_code != 404:
                print(f"⚠️ Worker get-vendor-service failed during login sync ({resp.status_code}) via {endpoint}: {resp.text[:200]}")
        except Exception as fetch_err:
            print(f"⚠️ Unable to read vendor service data during login sync: {fetch_err}")

        if existing_service_data is None:
            service_payload = _default_service_availability()
            try:
                endpoint, resp = post_to_worker('/upsert-vendor-service', {
                    'vendor_email': vendor_email,
                    'vendor_id': vendor_id,
                    'service_data': service_payload,
                    'updated_by': vendor_email
                })
                if resp.status_code != 200:
                    print(f"⚠️ Worker upsert-vendor-service failed during login sync ({resp.status_code}) via {endpoint}: {resp.text[:200]}")
            except Exception as upsert_err:
                print(f"⚠️ Unable to upsert vendor service data during login sync: {upsert_err}")
    except Exception as sync_err:
        print(f"⚠️ Vendor service sync on login failed: {sync_err}")


@csrf_exempt
def get_vendor_details(request):
    email = request.GET.get('email')
    if not email:
        return JsonResponse({'error': 'Email required'}, status=400)
    details = get_vendor_details_by_email(email)
    if not details:
        return JsonResponse({'error': 'Vendor not found'}, status=404)
    return JsonResponse(details)

@csrf_exempt
@require_POST
def update_vendor_service_availability(request):
    """
    Persist vendor service availability settings to D1 storage.
    """
    try:
        vendor_email = request.session.get('vendor_email') or request.POST.get('vendor_email') or ''
        vendor_id = request.session.get('vendor_id') or request.POST.get('vendor_id') or ''

        if not vendor_email or not vendor_id:
            return JsonResponse({
                "success": False,
                "error": "vendor_email and vendor_id are required"
            }, status=400)

        try:
            service_data = json.loads(request.body.decode('utf-8')) if request.body else {}
        except json.JSONDecodeError:
            return JsonResponse({"success": False, "error": "Invalid JSON."}, status=400)

        if not isinstance(service_data, dict):
            return JsonResponse({
                "success": False,
                "error": "Service payload must be a JSON object"
            }, status=400)

        updated_by = (
            request.session.get('user_email')
            or (request.user.email if request.user.is_authenticated else None)
            or vendor_email
        )

        endpoint, resp = post_to_worker('/upsert-vendor-service', {
            'vendor_email': vendor_email,
            'vendor_id': vendor_id,
            'service_data': service_data,
            'updated_by': updated_by,
        })

        if resp.status_code != 200:
            print(f"❌ Worker upsert-vendor-service failed ({resp.status_code}) via {endpoint}: {resp.text[:300]}")
            return JsonResponse({
                "success": False,
                "error": "Failed to save service availability"
            }, status=resp.status_code)

        response_data = resp.json()
        if not response_data.get('success'):
            return JsonResponse(response_data, status=400)

        return JsonResponse({
            "success": True,
            "data": response_data.get('service', {}).get('service_data', service_data),
            "source": "D1",
            "message": response_data.get('message', 'Service availability updated')
        })
    except Exception as e:
        return JsonResponse({"success": False, "error": str(e)}, status=500)

@csrf_exempt
def get_vendor_service_availability(request):
    """
    Retrieve vendor service availability from D1 with safe fallback defaults.
    """
    try:
        vendor_email = request.session.get('vendor_email') or request.GET.get('vendor_email') or ''
        vendor_id = request.session.get('vendor_id') or request.GET.get('vendor_id') or ''

        if not vendor_email and not vendor_id:
            return JsonResponse({
                "success": False,
                "error": "vendor_email or vendor_id is required"
            }, status=400)

        payload = {
            'vendor_email': vendor_email,
            'vendor_id': vendor_id
        }
        endpoint, resp = post_to_worker('/get-vendor-service', payload)

        if resp.status_code == 404:
            return JsonResponse({
                "success": True,
                "data": _default_service_availability(),
                "source": "defaults",
                "note": "No service configuration found in D1"
            })

        if resp.status_code != 200:
            print(f"❌ Worker get-vendor-service failed ({resp.status_code}) via {endpoint}: {resp.text[:300]}")
            return JsonResponse({
                "success": False,
                "error": "Unable to load service availability"
            }, status=resp.status_code)

        response_json = resp.json()
        if not response_json.get('success'):
            return JsonResponse(response_json, status=400)

        service_payload = response_json.get('service') or {}
        service_data = service_payload.get('service_data') or _default_service_availability()

        return JsonResponse({
            "success": True,
            "data": service_data,
            "source": "D1",
            "updated_at": service_payload.get('updated_at')
        })
    except Exception as e:
        return JsonResponse({"success": False, "error": str(e)}, status=500)

@csrf_exempt
def test_r2_simple(request):
    """
    Simple R2 connection test
    """
    try:
        endpoint_url = settings.R2_ENDPOINT.rstrip('/') if settings.R2_ENDPOINT else ''
        
        s3 = boto3.client(
            's3',
            aws_access_key_id=settings.R2_ACCESS_KEY,
            aws_secret_access_key=settings.R2_SECRET_KEY,
            endpoint_url=endpoint_url,
            region_name='auto',
        )
        
        # Test bucket access
        s3.head_bucket(Bucket=settings.R2_BUCKET)
        
        # Test simple upload
        test_key = "test_simple.json"
        test_data = {"test": "data", "timestamp": datetime.datetime.now().isoformat()}
        
        s3.put_object(
            Bucket=settings.R2_BUCKET,
            Key=test_key,
            Body=json.dumps(test_data),
            ContentType='application/json'
        )
        
        return JsonResponse({
            "success": True,
            "message": "R2 connection successful",
            "test_file": test_key
        })
        
    except Exception as e:
        return JsonResponse({
            "success": False,
            "error": str(e),
            "endpoint": settings.R2_ENDPOINT,
            "bucket": settings.R2_BUCKET
        })

@csrf_exempt
def list_vendor_folder(request):
    """
    List contents of vendor folder in R2 to debug folder creation
    """
    try:
        # Get vendor email from session or request parameters
        vendor_email = request.session.get('vendor_email')
        if not vendor_email:
            vendor_email = request.GET.get('vendor_email')
        if not vendor_email:
            return JsonResponse({
                "success": False,
                "error": "No vendor email found in session or request parameters"
            })
        
        sanitized_email = sanitize_email(vendor_email)
        folder_prefix = f"vendor_register_details/{sanitized_email}/"
        
        print(f"🔍 DEBUG - Listing folder: {folder_prefix}")
        
        endpoint_url = settings.R2_ENDPOINT.rstrip('/') if settings.R2_ENDPOINT else ''
        
        s3 = boto3.client(
            's3',
            aws_access_key_id=settings.R2_ACCESS_KEY,
            aws_secret_access_key=settings.R2_SECRET_KEY,
            endpoint_url=endpoint_url,
            region_name='auto',
        )
        
        # List objects in the vendor folder
        response = s3.list_objects_v2(
            Bucket=settings.R2_BUCKET,
            Prefix=folder_prefix
        )
        
        files = []
        if 'Contents' in response:
            for obj in response['Contents']:
                files.append({
                    'key': obj['Key'],
                    'size': obj['Size'],
                    'last_modified': obj['LastModified'].isoformat()
                })
        
        return JsonResponse({
            "success": True,
            "vendor_email": vendor_email,
            "sanitized_email": sanitized_email,
            "folder_prefix": folder_prefix,
            "files": files,
            "total_files": len(files)
        })
        
    except Exception as e:
        return JsonResponse({
            "success": False,
            "error": str(e)
        })

@csrf_exempt
def update_vendor_availability(request):
    """
    Update the vendor's availability status (online/offline) in D1 service storage.
    """
    if request.method == 'POST':
        try:
            vendor_email = request.session.get('vendor_email') or request.POST.get('vendor_email') or ''
            vendor_id = request.session.get('vendor_id') or request.POST.get('vendor_id') or ''
            if not vendor_email or not vendor_id:
                return JsonResponse({"success": False, "error": "vendor_email and vendor_id are required"}, status=400)

            try:
                data = json.loads(request.body.decode('utf-8'))
                availability_status = data.get('vendor_shop_avaliability', 'online')
            except Exception:
                return JsonResponse({"success": False, "error": "Invalid JSON."}, status=400)

            service_data = _default_service_availability()

            # Load existing service config from D1 (if present)
            try:
                endpoint, resp = post_to_worker('/get-vendor-service', {
                    'vendor_email': vendor_email,
                    'vendor_id': vendor_id
                })
                if resp.status_code == 200:
                    payload = resp.json()
                    if payload.get('success'):
                        service_payload = payload.get('service') or {}
                        service_data = service_payload.get('service_data') or service_data
            except Exception as fetch_err:
                print(f"⚠️ Unable to read existing vendor availability for {vendor_email}: {fetch_err}")

            service_data['vendor_shop_avaliability'] = availability_status

            endpoint, resp = post_to_worker('/upsert-vendor-service', {
                'vendor_email': vendor_email,
                'vendor_id': vendor_id,
                'service_data': service_data,
                'updated_by': request.session.get('user_email') or vendor_email
            })

            if resp.status_code != 200:
                return JsonResponse({
                    "success": False,
                    "error": "Failed to update availability in D1"
                }, status=resp.status_code)

            response_json = resp.json()
            if not response_json.get('success'):
                return JsonResponse(response_json, status=400)

            return JsonResponse({
                "success": True,
                "message": f"Vendor availability updated to {availability_status}",
                "data": service_data,
                "source": "D1"
            })
        except Exception as e:
            return JsonResponse({"success": False, "error": str(e)}, status=500)
    
    return JsonResponse({"success": False, "error": "Invalid request method"}, status=405)

@csrf_exempt
def test_profile_image_url(request):
    """
    Test if a profile image URL is accessible
    """
    if request.method == 'GET':
        try:
            vendor_email = request.session.get('vendor_email')
            if not vendor_email:
                return JsonResponse({
                    "success": False,
                    "error": "Vendor not authenticated"
                }, status=401)
            
            vendor_details = get_vendor_details_by_email(vendor_email)
            if vendor_details and vendor_details.get('profile_image_url'):
                profile_url = vendor_details['profile_image_url']
                
                # Check if the file actually exists in R2
                try:
                    # Extract key from URL
                    key = None
                    if '?' in profile_url:
                        base_url = profile_url.split('?')[0]
                        if settings.R2_ENDPOINT in base_url:
                            key = base_url.replace(settings.R2_ENDPOINT, '').lstrip('/')
                            if key.startswith(settings.R2_BUCKET + '/'):
                                key = key[len(settings.R2_BUCKET + '/'):]
                    else:
                        if settings.R2_ENDPOINT in profile_url:
                            key = profile_url.replace(settings.R2_ENDPOINT, '').lstrip('/')
                            if key.startswith(settings.R2_BUCKET + '/'):
                                key = key[len(settings.R2_BUCKET + '/'):]
                    
                    file_exists = False
                    if key:
                        try:
                            s3 = boto3.client(
                                's3',
                                aws_access_key_id=settings.R2_ACCESS_KEY,
                                aws_secret_access_key=settings.R2_SECRET_KEY,
                                endpoint_url=settings.R2_ENDPOINT.rstrip('/'),
                                region_name='auto'
                            )
                            s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                            file_exists = True
                            print(f"✅ Profile image file exists in R2: {key}")
                        except Exception as e:
                            print(f"❌ Profile image file does not exist in R2: {key} - {str(e)}")
                            file_exists = False
                except Exception as e:
                    print(f"❌ Error checking file existence: {str(e)}")
                    file_exists = False
                
                # Test if the URL is accessible
                try:
                    response = requests.head(profile_url, timeout=5)
                    is_accessible = response.status_code == 200
                    error_message = None
                except Exception as e:
                    is_accessible = False
                    error_message = str(e)
                
                return JsonResponse({
                    "success": True,
                    "profile_image_url": profile_url,
                    "is_accessible": is_accessible,
                    "status_code": response.status_code if 'response' in locals() else None,
                    "error_message": error_message,
                    "file_exists_in_r2": file_exists if 'file_exists' in locals() else None,
                    "extracted_key": key if 'key' in locals() else None
                })
            else:
                return JsonResponse({
                    "success": True,
                    "profile_image_url": None,
                    "is_accessible": False
                })
                
        except Exception as e:
            print(f"❌ Error testing profile image URL: {str(e)}")
            return JsonResponse({
                "success": False,
                "error": f"Server error: {str(e)}"
            }, status=500)
    
    return JsonResponse({
        "success": False,
        "error": "Invalid request method"
    }, status=405)

@csrf_exempt
def get_vendor_profile_image(request):
    """
    Get vendor profile image URL from Cloudflare R2
    """
    if request.method == 'GET':
        try:
            vendor_email = request.session.get('vendor_email')
            if not vendor_email:
                return JsonResponse({
                    "success": False,
                    "error": "Vendor not authenticated"
                }, status=401)
            
            vendor_details = get_vendor_details_by_email(vendor_email)
            if vendor_details and vendor_details.get('profile_image_url'):
                return JsonResponse({
                    "success": True,
                    "profile_image_url": vendor_details['profile_image_url']
                })
            else:
                return JsonResponse({
                    "success": True,
                    "profile_image_url": None
                })
                
        except Exception as e:
            print(f"❌ Error getting profile image: {str(e)}")
            return JsonResponse({
                "success": False,
                "error": f"Server error: {str(e)}"
            }, status=500)
    
    return JsonResponse({
        "success": False,
        "error": "Invalid request method"
    }, status=405)
def update_vendor_profile(request):
    """
    Update vendor profile details and profile image in Cloudflare R2 under:
    vendor_register_details/{sanitized_email}/profile/
    """
    if request.method == 'POST':
        try:
            vendor_email = request.session.get('vendor_email')
            if not vendor_email:
                return JsonResponse({
                    "success": False,
                    "error": "Vendor not authenticated"
                }, status=401)
            
            sanitized_email = sanitize_email(vendor_email)
            
            # Initialize S3 client
            endpoint_url = settings.R2_ENDPOINT.rstrip('/') if settings.R2_ENDPOINT else ''
            s3 = boto3.client(
                's3',
                aws_access_key_id=settings.R2_ACCESS_KEY,
                aws_secret_access_key=settings.R2_SECRET_KEY,
                endpoint_url=endpoint_url,
                region_name='auto',
                config=boto3.session.Config(
                    signature_version='s3v4',
                    retries={'max_attempts': 3}
                )
            )
            
            # Handle profile image upload
            profile_image_url = None
            if 'profile_image' in request.FILES:
                image_file = request.FILES['profile_image']
                
                # Validate image file
                if not image_file.content_type.startswith('image/'):
                    return JsonResponse({
                        "success": False,
                        "error": "Invalid file type. Please upload an image."
                    }, status=400)
                
                # Generate unique filename
                file_extension = image_file.name.split('.')[-1].lower()
                image_filename = f"profile_image_{int(time.time())}.{file_extension}"
                image_key = f"vendor_register_details/{sanitized_email}/profile/{image_filename}"
                
                # Upload image to R2
                try:
                    s3.put_object(
                        Bucket=settings.R2_BUCKET,
                        Key=image_key,
                        Body=image_file.read(),
                        ContentType=image_file.content_type,
                        Metadata={
                            'uploaded_at': datetime.datetime.now().isoformat(),
                            'original_filename': image_file.name,
                            'file_size': str(image_file.size)
                        }
                    )
                    
                    # Generate public URL for the image (more reliable for R2)
                    try:
                        # Use public URL format with bucket name
                        profile_image_url = f"{settings.R2_ENDPOINT.rstrip('/')}/{settings.R2_BUCKET}/{image_key}"
                        print(f"✅ Generated public URL: {profile_image_url}")
                        
                        # Test if the public URL works
                        try:
                            test_response = requests.head(profile_image_url, timeout=5)
                            if test_response.status_code == 200:
                                print(f"✅ Public URL is accessible")
                            else:
                                print(f"⚠️ Public URL returned status {test_response.status_code}")
                                # Fallback to presigned URL with shorter expiration
                                presigned_url = s3.generate_presigned_url(
                                    'get_object',
                                    Params={'Bucket': settings.R2_BUCKET, 'Key': image_key},
                                    ExpiresIn=3600 * 24 * 7  # 7 days expiration (shorter for testing)
                                )
                                profile_image_url = presigned_url
                                print(f"✅ Fallback to presigned URL: {profile_image_url}")
                        except Exception as test_error:
                            print(f"⚠️ Public URL test failed: {test_error}")
                            # Fallback to presigned URL with shorter expiration
                            presigned_url = s3.generate_presigned_url(
                                'get_object',
                                Params={'Bucket': settings.R2_BUCKET, 'Key': image_key},
                                ExpiresIn=3600 * 24 * 7  # 7 days expiration (shorter for testing)
                            )
                            profile_image_url = presigned_url
                            print(f"✅ Fallback to presigned URL: {profile_image_url}")
                            
                    except Exception as url_error:
                        print(f"❌ Failed to generate URL: {url_error}")
                        # Final fallback
                        profile_image_url = f"{settings.R2_ENDPOINT.rstrip('/')}/{settings.R2_BUCKET}/{image_key}"
                    
                    print(f"✅ Profile image uploaded: {image_key}")
                    print(f"✅ Final profile image URL: {profile_image_url}")
                    print(f"✅ R2_ENDPOINT: {settings.R2_ENDPOINT}")
                    print(f"✅ R2_BUCKET: {settings.R2_BUCKET}")
                    
                except Exception as e:
                    print(f"❌ Error uploading profile image: {str(e)}")
                    return JsonResponse({
                        "success": False,
                        "error": f"Failed to upload profile image: {str(e)}"
                    }, status=500)
            
            # Get form data
            vendor_name = request.POST.get('vendor_name', '').strip()
            remove_profile_image = request.POST.get('remove_profile_image', '').strip() == 'true'
            
            # Validate required fields
            if not vendor_name:
                return JsonResponse({
                    "success": False,
                    "error": "Vendor name is required"
                }, status=400)
            
            # Update registration details with new profile data
            try:
                reg_key = f'vendor_register_details/{sanitized_email}/registration_details.json'
                response = s3.get_object(Bucket=settings.R2_BUCKET, Key=reg_key)
                vendor_data = json.loads(response['Body'].read().decode('utf-8'))
                
                # Get old profile image URL for cleanup
                old_profile_url = vendor_data.get('profile_image_url', '')
                old_image_key = None
                
                # Extract old image key if exists
                if old_profile_url:
                    print(f"🔍 DEBUG - Extracting key from old URL: {old_profile_url}")
                    if '?' in old_profile_url:
                        base_url = old_profile_url.split('?')[0]
                        print(f"🔍 DEBUG - Base URL (after removing query): {base_url}")
                        if settings.R2_ENDPOINT in base_url:
                            old_image_key = base_url.replace(settings.R2_ENDPOINT, '').lstrip('/')
                            print(f"🔍 DEBUG - After removing endpoint: {old_image_key}")
                            if old_image_key.startswith(settings.R2_BUCKET + '/'):
                                old_image_key = old_image_key[len(settings.R2_BUCKET + '/'):]
                                print(f"🔍 DEBUG - After removing bucket: {old_image_key}")
                    else:
                        if settings.R2_ENDPOINT in old_profile_url:
                            old_image_key = old_profile_url.replace(settings.R2_ENDPOINT, '').lstrip('/')
                            print(f"🔍 DEBUG - After removing endpoint (no query): {old_image_key}")
                            if old_image_key.startswith(settings.R2_BUCKET + '/'):
                                old_image_key = old_image_key[len(settings.R2_BUCKET + '/'):]
                                print(f"🔍 DEBUG - After removing bucket (no query): {old_image_key}")
                    
                    print(f"🔍 DEBUG - Final extracted old_image_key: {old_image_key}")
                else:
                    print(f"🔍 DEBUG - No old profile URL found")
                
                # Update vendor data
                vendor_data['vendor_name'] = vendor_name
                
                if remove_profile_image:
                    # Remove profile image
                    vendor_data['profile_image_url'] = ''
                    vendor_data['profile_image_removed_at'] = datetime.datetime.now().isoformat()
                    profile_image_url = None
                    print(f"✅ Profile image removed for vendor {vendor_email}")
                    
                    # Delete the old image from R2 if it exists
                    if old_image_key:
                        try:
                            s3.delete_object(Bucket=settings.R2_BUCKET, Key=old_image_key)
                            print(f"✅ Deleted old profile image from R2: {old_image_key}")
                        except Exception as delete_error:
                            print(f"⚠️ Failed to delete old profile image from R2 {old_image_key}: {delete_error}")
                    else:
                        print(f"⚠️ No old image key found to delete")
                elif profile_image_url:
                    # Update with new profile image
                    vendor_data['profile_image_url'] = profile_image_url
                    vendor_data['profile_image_updated_at'] = datetime.datetime.now().isoformat()
                    
                    # Clean up old image if it exists and is different from new one
                    if old_image_key and old_image_key != image_key:
                        try:
                            s3.delete_object(Bucket=settings.R2_BUCKET, Key=old_image_key)
                            print(f"✅ Deleted old profile image: {old_image_key}")
                        except Exception as delete_error:
                            print(f"⚠️ Failed to delete old profile image {old_image_key}: {delete_error}")
                else:
                    # No new image uploaded, keep existing one if any
                    if old_profile_url:
                        profile_image_url = old_profile_url
                        print(f"✅ Keeping existing profile image: {old_profile_url}")
                
                # Debug: Print the updated vendor data
                print(f"🔍 DEBUG - Updated vendor data:")
                print(f"   - vendor_name: {vendor_data['vendor_name']}")
                print(f"   - profile_image_url: {vendor_data.get('profile_image_url', 'None')}")
                
                # Upload updated registration details
                s3.put_object(
                    Bucket=settings.R2_BUCKET,
                    Key=reg_key,
                    Body=json.dumps(vendor_data, indent=2),
                    ContentType='application/json'
                )
                
                print(f"✅ Profile updated for vendor {vendor_email}")
                
                response_data = {
                    "success": True,
                    "message": "Profile updated successfully",
                    "vendor_name": vendor_name,
                    "updated_at": datetime.datetime.now().isoformat()
                }
                
                if remove_profile_image:
                    response_data["message"] = "Profile image removed successfully"
                    response_data["profile_image_url"] = None
                elif profile_image_url:
                    response_data["profile_image_url"] = profile_image_url
                
                return JsonResponse(response_data)
                
            except Exception as e:
                print(f"❌ Error updating registration details: {str(e)}")
                return JsonResponse({
                    "success": False,
                    "error": f"Failed to update profile details: {str(e)}"
                }, status=500)
                
        except Exception as e:
            print(f"❌ Server error in update_vendor_profile: {str(e)}")
            return JsonResponse({
                "success": False,
                "error": f"Server error: {str(e)}"
            }, status=500)
    
    return JsonResponse({
        "success": False,
        "error": "Invalid request method"
    }, status=405)

@csrf_exempt
def generate_fresh_preview_url(request):
    """
    Generate a fresh preview URL for a file to avoid expiration issues
    """
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            filename = data.get('filename')
            job_data = data.get('job_data', {}) or {}
            
            if not filename:
                return JsonResponse({
                    "success": False,
                    "error": "Filename is required"
                }, status=400)
            
            # Get vendor email from session
            vendor_email = request.session.get('vendor_email')
            if not vendor_email:
                return JsonResponse({
                    "success": False,
                    "error": "Vendor not authenticated"
                }, status=401)
            
            # Initialize S3 client
            s3 = boto3.client('s3',
                            aws_access_key_id=settings.R2_ACCESS_KEY,
                            aws_secret_access_key=settings.R2_SECRET_KEY,
                            endpoint_url=settings.R2_ENDPOINT,
                            region_name='auto')
            
            # Try to find the file in various locations
            possible_keys = []
            r2_path = _normalize_r2_key(job_data.get('r2_path') or job_data.get('metadata', {}).get('r2_path'))
            if r2_path:
                possible_keys.append(r2_path)

            # Try extracting exact object key from already known URLs first
            def _extract_key_from_url(raw_url):
                try:
                    if not raw_url:
                        return None
                    path = _normalize_r2_key(raw_url)
                    return path or None
                except Exception:
                    return None

            for url_field in ('preview_url', 'download_url'):
                extracted = _extract_key_from_url(job_data.get(url_field))
                if extracted:
                    possible_keys.append(extracted)
            
            # Check vendor print jobs (fallback construction)
            vendor_id = (
                job_data.get('vendor_id')
                or job_data.get('vendor')
                or job_data.get('shop_id')
                or job_data.get('metadata', {}).get('vendor_id')
                or job_data.get('metadata', {}).get('vendor')
                or request.session.get('vendor_id')
                or '9080823634'
            )
            storage_folder = (
                job_data.get('storage_folder')
                or job_data.get('metadata', {}).get('storage_folder')
                or 'vendor_print_jobs'
            )
            if vendor_id:
                possible_keys.append(f"{storage_folder.rstrip('/')}/{vendor_id}/{filename}")
            
            # Check user uploads
            user_email = (
                job_data.get('user_email')
                or job_data.get('user')
                or job_data.get('metadata', {}).get('user_email')
                or job_data.get('metadata', {}).get('user')
            )
            if user_email:
                possible_keys.append(f"users/{user_email}/{filename}")
            
            # Check manual jobs
            possible_keys.append(f"manual_print_jobs/{filename}")
            
            # Find the actual file
            found_key = None
            seen_keys = set()
            for key in possible_keys:
                normalized_key = key.lstrip('/')
                if normalized_key in seen_keys:
                    continue
                seen_keys.add(normalized_key)
                try:
                    s3.head_object(Bucket=settings.R2_BUCKET, Key=normalized_key)
                    found_key = normalized_key
                    break
                except:
                    continue

            # Final fallback: scan likely prefixes for exact filename match
            if not found_key:
                candidate_prefixes = []
                if storage_folder:
                    candidate_prefixes.append(f"{storage_folder.rstrip('/')}/")
                if vendor_id:
                    candidate_prefixes.append(f"{storage_folder.rstrip('/')}/{vendor_id}/")
                    candidate_prefixes.append(f"vendor_print_jobs/{vendor_id}/")
                if user_email:
                    candidate_prefixes.append(f"users/{user_email}/")
                candidate_prefixes.extend([
                    "vendor_print_jobs/",
                    "manual_print_jobs/",
                    "users/"
                ])

                prefix_seen = set()
                for prefix in candidate_prefixes:
                    normalized_prefix = (prefix or '').lstrip('/')
                    if not normalized_prefix or normalized_prefix in prefix_seen:
                        continue
                    prefix_seen.add(normalized_prefix)
                    try:
                        continuation_token = None
                        while True:
                            list_args = {
                                'Bucket': settings.R2_BUCKET,
                                'Prefix': normalized_prefix,
                                'MaxKeys': 1000,
                            }
                            if continuation_token:
                                list_args['ContinuationToken'] = continuation_token

                            resp = s3.list_objects_v2(**list_args)
                            for item in resp.get('Contents', []) or []:
                                key = _normalize_r2_key(item.get('Key'))
                                if key and key.split('/')[-1] == filename:
                                    found_key = key
                                    break

                            if found_key:
                                break

                            if not resp.get('IsTruncated'):
                                break

                            continuation_token = resp.get('NextContinuationToken')
                            if not continuation_token:
                                break

                        if found_key:
                            break
                    except Exception:
                        continue

            if not found_key:
                return JsonResponse({
                    "success": False,
                    "error": "File not found"
                }, status=404)
            
            # Generate fresh presigned URL
            fresh_url = s3.generate_presigned_url(
                ClientMethod='get_object',
                Params={
                    'Bucket': settings.R2_BUCKET,
                    'Key': found_key
                },
                ExpiresIn=3600  # 1 hour
            )
            
            return JsonResponse({
                "success": True,
                "preview_url": fresh_url,
                "found_key": found_key
            })
            
        except Exception as e:
            print(f"Error generating fresh preview URL: {str(e)}")
            return JsonResponse({
                "success": False,
                "error": str(e)
            }, status=500)
    
    return JsonResponse({
        "success": False,
        "error": "Invalid request method"
    }, status=405)

@csrf_exempt
def test_r2_url_generation(request):
    """
    Test R2 URL generation for debugging
    """
    if request.method == 'GET':
        try:
            vendor_email = request.session.get('vendor_email')
            if not vendor_email:
                return JsonResponse({
                    "success": False,
                    "error": "Vendor not authenticated"
                }, status=401)
            
            sanitized_email = sanitize_email(vendor_email)
            
            # Initialize S3 client
            endpoint_url = settings.R2_ENDPOINT.rstrip('/') if settings.R2_ENDPOINT else ''
            s3 = boto3.client(
                's3',
                aws_access_key_id=settings.R2_ACCESS_KEY,
                aws_secret_access_key=settings.R2_SECRET_KEY,
                endpoint_url=endpoint_url,
                region_name='auto',
                config=boto3.session.Config(
                    signature_version='s3v4',
                    retries={'max_attempts': 3}
                )
            )
            
            # Test different URL formats
            test_key = f"vendor_register_details/{sanitized_email}/test_image.jpg"
            
            # Format 1: Presigned URL
            try:
                presigned_url = s3.generate_presigned_url(
                    'get_object',
                    Params={'Bucket': settings.R2_BUCKET, 'Key': test_key},
                    ExpiresIn=3600  # 1 hour
                )
                print(f"✅ Generated presigned URL: {presigned_url}")
            except Exception as e:
                presigned_url = f"Error: {str(e)}"
                print(f"❌ Error generating presigned URL: {str(e)}")
            
            # Format 2: Public URL
            public_url = f"{settings.R2_ENDPOINT.rstrip('/')}/{test_key}"
            
            # Format 3: Custom domain URL (if configured)
            custom_url = None
            if hasattr(settings, 'R2_PUBLIC_URL') and settings.R2_PUBLIC_URL:
                custom_url = f"{settings.R2_PUBLIC_URL.rstrip('/')}/{test_key}"
            
            return JsonResponse({
                "success": True,
                "test_key": test_key,
                "presigned_url": presigned_url,
                "public_url": public_url,
                "custom_url": custom_url,
                "r2_endpoint": settings.R2_ENDPOINT,
                "r2_bucket": settings.R2_BUCKET,
                "has_custom_url": hasattr(settings, 'R2_PUBLIC_URL') and settings.R2_PUBLIC_URL
            })
            
        except Exception as e:
            return JsonResponse({
                "success": False,
                "error": str(e)
            }, status=500)
    
    return JsonResponse({
        "success": False,
        "error": "Invalid request method"
    }, status=405)

@csrf_exempt
def accept_job(request):
    """Accept a print job and update vendor status to accepted"""
    try:
        data = json.loads(request.body)
        filename = data.get('filename')
        vendor_id = data.get('vendor_id')
        
        if not filename or not vendor_id:
            return JsonResponse({'error': 'Filename and vendor_id required'}, status=400)
        
        # Update job status to accepted in R2
        success = update_job_vendor_status(filename, vendor_id, 'accepted')
        
        if success:
            # Also update the job in the original location to maintain consistency
            try:
                s3 = boto3.client('s3',
                    aws_access_key_id=settings.R2_ACCESS_KEY,
                    aws_secret_access_key=settings.R2_SECRET_KEY,
                    endpoint_url=settings.R2_ENDPOINT,
                    region_name='auto')
                
                # Update in print_requests folder
                source_key = f'print_requests/{filename}'
                try:
                    response = s3.get_object(Bucket=settings.R2_BUCKET, Key=source_key)
                    src_metadata = json.loads(response['Body'].read().decode('utf-8'))
                    src_metadata['vendor_status'] = 'accepted'
                    src_metadata['vendor_id'] = vendor_id
                    src_metadata['accepted_at'] = datetime.datetime.now().isoformat()
                    
                    s3.put_object(
                        Bucket=settings.R2_BUCKET,
                        Key=source_key,
                        Body=json.dumps(src_metadata, indent=2),
                        ContentType='application/json'
                    )
                except:
                    pass  # If not found in print_requests, that's okay
                
                print(f"✅ Job {filename} accepted by vendor {vendor_id}")
                
            except Exception as e:
                print(f"⚠️ Error updating original job status: {e}")
        
        if success:
            return JsonResponse({
                'success': True,
                'message': f'Job {filename} accepted successfully',
                'status': 'accepted'
            })
        else:
            return JsonResponse({
                'success': False,
                'error': 'Failed to accept job'
            }, status=500)
            
    except Exception as e:
        print(f"Error accepting job: {e}")
        return JsonResponse({'error': 'Internal server error'}, status=500)

def get_2day_period_for_date(date_obj):
    """
    Calculate the 2-day period for a given date.
    Returns (period_start, period_end) as date objects.
    Ensures non-overlapping 2-day buckets: 1-2, 3-4, 5-6, ..., 27-28, 29-30 (not 2-3, 3-4, 4-5).
    """
    # Get the day of month
    day = date_obj.day
    
    # 2-day buckets: (1-2), (3-4), (5-6), ... so one row per two calendar days
    # Odd days (1,3,5,...): start=this day, end=next day. Even days (2,4,6,...): start=prev day, end=this day
    
    if day % 2 == 1:  # Odd day (1, 3, 5, ..., 27, 29, 31)
        period_start = date_obj
        period_end = date_obj + datetime.timedelta(days=1)
    else:  # Even day (2, 4, 6, ..., 28, 30)
        period_start = date_obj - datetime.timedelta(days=1)
        period_end = date_obj
    
    # Handle month boundaries
    if period_end.month != period_start.month:
        # If period_end goes into next month, adjust to last day of current month
        import calendar
        last_day = calendar.monthrange(period_start.year, period_start.month)[1]
        period_end = datetime.date(period_start.year, period_start.month, last_day)
    
    return period_start, period_end

def create_or_update_vendor_transaction(vendor_email, vendor_id, vendor_name, completion_date, filename, total_price, platform_profit):
    """
    Create or update vendor transaction report for the 2-day period containing completion_date.
    This ensures all jobs completed within the same 2-day period are aggregated together.
    """
    try:
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            print("⚠️ Worker API not configured - skipping transaction report creation")
            return False
        
        # Best-effort: resolve a friendly vendor/shop name if we only have a placeholder
        try:
            normalized_name = (vendor_name or '').strip()
            if not normalized_name or normalized_name.lower() in ['printmax vendor', 'unknown vendor']:
                # Try rich vendor details first
                details = get_vendor_details_by_email(vendor_email) or {}
                resolved_name = (
                    details.get('vendor_name')
                    or details.get('shop_name')
                )
                # Fallback to coordinates helper (which also returns vendor_name / shop_name)
                if not resolved_name:
                    coords = get_vendor_coordinates_from_email(vendor_email) or {}
                    resolved_name = (
                        coords.get('vendor_name')
                        or coords.get('shop_name')
                    )
                if resolved_name:
                    vendor_name = resolved_name
        except Exception as name_err:
            # Do not block transaction creation because of name resolution issues
            print(f"⚠️ Unable to resolve vendor_name for transactions: {name_err}")

        # Calculate 2-day period
        period_start, period_end = get_2day_period_for_date(completion_date)
        period_start_str = period_start.strftime('%Y-%m-%d')
        period_end_str = period_end.strftime('%Y-%m-%d')
        
        print(f"📊 Creating/updating transaction report for period: {period_start_str} to {period_end_str}")
        
        # Get job details to calculate totals
        # We'll aggregate all jobs in this period, but for now add this single job
        total_earning = float(total_price) - float(platform_profit)
        
        # Call worker API to create or update transaction
        worker_endpoint = api_url.rstrip('/') + '/aggregate-vendor-transaction'
        
        payload = {
            'vendor_email': vendor_email,
            'vendor_id': vendor_id,
            'vendor_name': vendor_name,
            'current_date': completion_date.strftime('%Y-%m-%d'),
            'period_start': period_start_str,
            'period_end': period_end_str,
            'total_price': float(total_price),
            'platform_profit': float(platform_profit),
            'total_earning': total_earning,
            'total_documents': 1,
            'filename': filename  # Include filename to prevent double counting
        }
        
        response = requests.post(
            worker_endpoint,
            json=payload,
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        if response.status_code == 200:
            print(f"✅ Transaction report created/updated for {vendor_email} (period: {period_start_str} to {period_end_str})")
            return True
        else:
            print(f"⚠️ Failed to create/update transaction report: {response.status_code} - {response.text[:200]}")
            return False
            
    except Exception as e:
        print(f"❌ Error creating/updating vendor transaction: {e}")
        import traceback
        traceback.print_exc()
        return False


def vendor_transactions_history(request):
    """
    Vendor-facing endpoint to fetch transaction reports from the D1 vendor_transaction table.
    Returns 2-day period reports for the authenticated vendor, without requiring admin access.
    """
    try:
        # Prefer vendor email from session to ensure correct vendor is used
        session_vendor_email = (request.session.get('vendor_email') or '').strip().lower()
        requested_vendor_email = (request.GET.get('vendor_email') or '').strip().lower()

        vendor_email = session_vendor_email or requested_vendor_email
        if not vendor_email:
            return JsonResponse({
                'success': False,
                'error': 'Vendor email not found in session'
            }, status=401)

        selected_month = request.GET.get('month')  # format YYYY-MM

        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')

        if not api_url or not api_key:
            return JsonResponse({
                'success': False,
                'error': 'Worker API not configured'
            }, status=500)

        worker_endpoint = api_url.rstrip('/') + '/get-vendor-transactions'

        # Ask the Worker API to filter by this vendor if it supports it
        payload = {
            'vendor_email': vendor_email,
        }
        if selected_month:
            payload['month'] = selected_month

        try:
            resp = requests.post(
                worker_endpoint,
                json=payload,
                headers={
                    'x-api-key': api_key,
                    'Content-Type': 'application/json'
                },
                timeout=20
            )

            if resp.status_code != 200:
                error_text = resp.text[:500] if hasattr(resp, 'text') else 'No error details'
                print(f"⚠️ Failed to get vendor transactions from D1: {resp.status_code} - {error_text}")
                return JsonResponse({
                    'success': False,
                    'error': f'Failed to fetch transactions: {resp.status_code}'
                }, status=500)

            data = resp.json()
            if not data.get('success'):
                error_msg = data.get('error', 'Unknown error')
                print(f"⚠️ Worker API returned error for vendor transactions: {error_msg}")
                return JsonResponse({
                    'success': False,
                    'error': error_msg
                }, status=500)

            transactions = data.get('transactions', []) or []

            print(f"🔍 vendor_transactions_history - raw transactions from worker: {len(transactions)} for vendor_email filter {vendor_email}")

            # Extra safety: only keep rows for this vendor
            vendor_email_lower = vendor_email.lower()
            transactions = [
                t for t in transactions
                if (t.get('vendor_email') or '').strip().lower() == vendor_email_lower
            ]

            reports = []
            available_months = set()

            for trans in transactions:
                period_start = (trans.get('period_start') or '').strip()
                period_end = (trans.get('period_end') or '').strip()

                # Derive display string if not provided
                period_display = trans.get('period_display') or ''
                if not period_display and period_start and period_end:
                    try:
                        period_display = f"{period_start} to {period_end}"
                    except Exception:
                        period_display = ''

                # Extract month for filter dropdown
                if period_start:
                    try:
                        report_month = period_start[:7]  # YYYY-MM
                        available_months.add(report_month)
                    except Exception:
                        pass

                # Normalize numeric fields
                total_documents = (
                    trans.get('total_documents')
                    if trans.get('total_documents') is not None
                    else trans.get('total_docum', 0)
                ) or 0

                try:
                    total_price = float(trans.get('total_price', 0.0) or 0.0)
                except Exception:
                    total_price = 0.0

                try:
                    platform_profit = float(trans.get('platform_profit', 0.0) or 0.0)
                except Exception:
                    platform_profit = 0.0

                # Prefer total_earning if provided, otherwise derive it
                try:
                    total_earning = float(trans.get('total_earning', 0.0) or 0.0)
                except Exception:
                    total_earning = 0.0
                if total_earning == 0.0 and (total_price or platform_profit):
                    total_earning = total_price - platform_profit

                # Choose a reasonable "generated_at" timestamp for sorting & display
                generated_at = (
                    trans.get('updated_at')
                    or trans.get('created_at')
                    or period_end
                    or period_start
                )

                report_data = {
                    'id': trans.get('id'),
                    'transaction_id': trans.get('id'),
                    'vendor_email': vendor_email,
                    'vendor_id': trans.get('vendor_id', ''),
                    'vendor_name': trans.get('vendor_name', 'Unknown Vendor'),
                    'total_documents': total_documents,
                    'total_earning': total_earning,
                    'total_amount': total_earning,  # alias for front-end
                    'total_price': total_price,
                    'platform_profit': platform_profit,
                    'payment_status': trans.get('payment_status', 'not_completed'),
                    'period_start': period_start,
                    'period_end': period_end,
                    'period_display': period_display,
                    'generated_at': generated_at,
                    # Optional extra details if the Worker provides them
                    'service_breakdown': trans.get('service_breakdown'),
                }

                reports.append(report_data)

            print(f"📊 vendor_transactions_history - returning {len(reports)} reports, months={sorted(list(available_months))}")

            # Newest first based on generated_at / period_start
            def _sort_key(r):
                return (r.get('generated_at') or r.get('period_start') or '')

            reports.sort(key=_sort_key, reverse=True)

            return JsonResponse({
                'success': True,
                'reports': reports,
                # Name expected by vendor dashboard JS
                'available_months_list': sorted(list(available_months), reverse=True),
            })

        except requests.exceptions.RequestException as exc:
            print(f"⚠️ Network error fetching vendor transactions from D1: {exc}")
            return JsonResponse({
                'success': False,
                'error': f'Network error: {str(exc)}'
            }, status=500)
        except Exception as exc:
            print(f"⚠️ Unexpected error fetching vendor transactions from D1: {exc}")
            return JsonResponse({
                'success': False,
                'error': f'Unexpected error: {str(exc)}'
            }, status=500)

    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e)
        }, status=500)


@csrf_exempt
def vendor_transaction_details(request):
    """
    Return per-job details for a specific 2-day vendor transaction period.
    Data is sourced from the vendor_notification table via the Worker API.
    """
    if request.method != 'POST':
        return JsonResponse(
            {'success': False, 'error': 'Only POST method is allowed'},
            status=405
        )

    try:
        try:
            payload = json.loads(request.body or '{}')
        except Exception:
            payload = request.POST

        period_start = (payload.get('period_start') or '').strip()
        period_end = (payload.get('period_end') or '').strip()

        if not period_start or not period_end:
            return JsonResponse(
                {'success': False, 'error': 'period_start and period_end are required'},
                status=400
            )

        vendor_email = (request.session.get('vendor_email') or '').strip().lower()
        if not vendor_email:
            return JsonResponse(
                {'success': False, 'error': 'Vendor email not found in session'},
                status=401
            )

        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')

        if not api_url or not api_key:
            return JsonResponse(
                {'success': False, 'error': 'Worker API not configured'},
                status=500
            )

        worker_endpoint = api_url.rstrip('/') + '/get-all-vendor-jobs'
        request_payload = {
            'vendor_email': vendor_email,
            'week_start': period_start,
            'week_end': period_end,
        }

        try:
            resp = requests.post(
                worker_endpoint,
                json=request_payload,
                headers={
                    'x-api-key': api_key,
                    'Content-Type': 'application/json'
                },
                timeout=20
            )

            if resp.status_code != 200:
                error_text = resp.text[:500] if hasattr(resp, 'text') else 'No error details'
                print(f"⚠️ Failed to get vendor jobs for details: {resp.status_code} - {error_text}")
                return JsonResponse(
                    {'success': False, 'error': 'Failed to fetch report details'},
                    status=500
                )

            data = resp.json()
            if not data.get('success'):
                error_msg = data.get('error', 'Unknown error')
                print(f"⚠️ Worker API returned error for vendor jobs details: {error_msg}")
                return JsonResponse(
                    {'success': False, 'error': error_msg},
                    status=500
                )

            jobs_raw = data.get('data', []) or []
            jobs = []

            for job in jobs_raw:
                filename = (job.get('filename') or '').strip() or 'Unknown file'
                token = (job.get('token') or '').strip()

                try:
                    total_price = float(job.get('total_price', 0.0) or 0.0)
                except Exception:
                    total_price = 0.0

                try:
                    platform_profit = float(job.get('platform_profit', 0.0) or 0.0)
                except Exception:
                    platform_profit = 0.0

                payment = total_price - platform_profit

                jobs.append({
                    'filename': filename,
                    'token': token,
                    'payment': round(payment, 2),
                })

            return JsonResponse({'success': True, 'jobs': jobs})

        except requests.exceptions.RequestException as exc:
            print(f"⚠️ Network error fetching vendor jobs details from D1: {exc}")
            return JsonResponse(
                {'success': False, 'error': 'Network error while fetching report details'},
                status=500
            )
        except Exception as exc:
            print(f"⚠️ Unexpected error fetching vendor jobs details from D1: {exc}")
            return JsonResponse(
                {'success': False, 'error': 'Unexpected error while fetching report details'},
                status=500
            )

    except Exception as e:
        return JsonResponse(
            {'success': False, 'error': str(e)},
            status=500
        )

@csrf_exempt
def mark_job_completed(request):
    """Mark a print job as completed by updating job_completed to 'YES' and free the token"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            filename = data.get('filename')
            
            if not filename:
                return JsonResponse({'success': False, 'error': 'Filename required'})
            
            # Get vendor email from session
            vendor_email = request.session.get('vendor_email')
            if not vendor_email:
                return JsonResponse({'success': False, 'error': 'Vendor not authenticated'})
            
            # Get or create a lock for this specific filename to prevent concurrent completions
            with _job_completion_lock:
                if filename not in _job_completion_locks:
                    _job_completion_locks[filename] = threading.Lock()
                file_lock = _job_completion_locks[filename]
            
            # Acquire lock for this specific job to prevent concurrent completions
            if not file_lock.acquire(blocking=False):
                # Another request is already processing this job
                return JsonResponse({
                    'success': False, 
                    'error': 'Job completion already in progress. Please wait...'
                }, status=409)
            
            try:
                # Get vendor details
                vendor_details = get_vendor_details_by_email(vendor_email)
                if not vendor_details:
                    return JsonResponse({'success': False, 'error': 'Vendor details not found'})
                
                vendor_id = vendor_details.get('vendor_id', 'vendor1')
                # Get vendor shop name - try vendor_name first (from database), then shop_name as fallback
                vendor_name = vendor_details.get('vendor_name') or vendor_details.get('shop_name', 'Unknown Vendor')
                
                # Get the token associated with this job before updating status
                job_token = get_token_from_file_metadata(filename, vendor_id)
                if job_token:
                    print(f"🔍 Found token {job_token} for job {filename}")
                else:
                    # Try to get token from Vendor_print_jobs table if not found in metadata
                    job_token = get_token_from_vendor_print_jobs(filename, vendor_id, vendor_email)
                    if job_token:
                        print(f"🔍 Retrieved token {job_token} from Vendor_print_jobs for job {filename}")
                    else:
                        print(f"⚠️ No token found for job {filename}")
                
                # Get user_email from D1 Vendor_print_jobs table FIRST (before updating)
                completion_time = datetime.datetime.now().isoformat()
                user_email = None
                
                # Fetch user_email from D1 Vendor_print_jobs table
                try:
                    api_url_lookup = getattr(settings, 'WORKER_API_URL', '')
                    api_key_lookup = getattr(settings, 'WORKER_API_KEY', '')
                    if api_url_lookup and api_key_lookup:
                        base_url_lookup = api_url_lookup.rstrip('/')
                        if '/add-contact' in base_url_lookup:
                            lookup_endpoint = base_url_lookup.replace('/add-contact', '/get-vendor-print-jobs')
                        elif '/add-vendor-register' in base_url_lookup:
                            lookup_endpoint = base_url_lookup.replace('/add-vendor-register', '/get-vendor-print-jobs')
                        else:
                            lookup_endpoint = base_url_lookup + '/get-vendor-print-jobs'

                        lookup_payload = {
                            'filename': filename,
                            'vendor_id': vendor_id,
                            'vendor_email': vendor_email,
                        }

                        lookup_resp = requests.post(
                            lookup_endpoint,
                            json=lookup_payload,
                            headers={
                                'Content-Type': 'application/json',
                                'x-api-key': api_key_lookup,
                            },
                            timeout=10,
                        )

                        if lookup_resp.status_code == 200:
                            lookup_data = lookup_resp.json()
                            if lookup_data.get('success') and lookup_data.get('data'):
                                jobs = lookup_data.get('data', [])
                                if jobs:
                                    inferred_email = jobs[0].get('user_email') or ''
                                    if inferred_email:
                                        user_email = inferred_email
                                        print(f"📧 Retrieved user email from D1 Vendor_print_jobs for completion: {user_email}")
                                    else:
                                        print(f"⚠️ Job found in Vendor_print_jobs but user_email is empty for {filename}")
                        else:
                            print(f"⚠️ Failed to lookup user_email from Vendor_print_jobs: {lookup_resp.status_code}")
                    else:
                        print("⚠️ Worker API not configured for user_email lookup")
                except Exception as e:
                    print(f"⚠️ Error getting user email from D1 Vendor_print_jobs for completion: {e}")
                
                # Fallback: If user_email is not found in D1, try R2 metadata
                if not user_email:
                    try:
                        print(f"🔍 user_email not found in D1, trying R2 metadata for {filename}")
                        user_email = get_user_email_from_file_metadata(filename, vendor_id)
                        if user_email:
                             print(f"📧 Retrieved user email from R2 metadata: {user_email}")
                    except Exception as e:
                        print(f"⚠️ Error getting user email from R2 metadata: {e}")
                
                # Update D1 database tables synchronously (User_print_jobs and Vendor_print_jobs)
                db_update_success = False
                try:
                    api_url = getattr(settings, 'WORKER_API_URL', '')
                    api_key = getattr(settings, 'WORKER_API_KEY', '')
                    
                    if api_url and api_key:
                        # Use /update-job-completed endpoint which properly UPDATES existing rows
                        base_url = api_url.rstrip('/')
                        if '/add-contact' in base_url:
                            worker_endpoint = base_url.replace('/add-contact', '/update-job-completed')
                        elif '/add-vendor-register' in base_url:
                            worker_endpoint = base_url.replace('/add-vendor-register', '/update-job-completed')
                        else:
                            worker_endpoint = base_url + '/update-job-completed'
                        
                        worker_payload = {
                            'filename': filename,
                            'job_completed': 'YES',
                            'rendered_status': 'YES',
                            'completion_time': completion_time,
                            'vendor_email': vendor_email,
                            'vendor_id': vendor_id,
                            'vendor_name': vendor_name,
                            'user_email': user_email if user_email else ''
                        }
                        print(f"🔄 Updating D1 tables synchronously for {filename} with user_email: {user_email if user_email else 'NOT PROVIDED'}")
                        try:
                            resp = requests.post(
                                worker_endpoint,
                                json=worker_payload,
                                headers={
                                    'x-api-key': api_key,
                                    'Content-Type': 'application/json'
                                },
                                timeout=10
                            )
                            if resp.status_code == 200:
                                resp_data = resp.json()
                                if resp_data.get('success'):
                                    db_update_success = True
                                    print(f"✅ Successfully updated User_print_jobs and Vendor_print_jobs tables synchronously for {filename} (user_email: {user_email if user_email else 'not provided'})")
                                else:
                                    error_msg = resp_data.get('error', 'Unknown error')
                                    print(f"❌ Update endpoint returned success=false: {error_msg}")
                                    return JsonResponse({
                                        'success': False,
                                        'error': f'Failed to update database tables: {error_msg}'
                                    }, status=500)
                            else:
                                error_text = resp.text[:200] if resp.text else 'No response text'
                                print(f"❌ Failed to update job status: {resp.status_code} - {error_text}")
                                return JsonResponse({
                                    'success': False,
                                    'error': f'Failed to update database tables: HTTP {resp.status_code}'
                                }, status=500)
                        except Exception as e:
                            print(f"❌ Error updating job status in D1: {e}")
                            return JsonResponse({
                                'success': False,
                                'error': f'Database update error: {str(e)}'
                            }, status=500)
                    else:
                        print("❌ Worker API not configured - cannot update database")
                        return JsonResponse({
                            'success': False,
                            'error': 'Worker API not configured'
                        }, status=500)
                except Exception as e:
                    print(f"❌ Error updating D1 database: {e}")
                    return JsonResponse({
                        'success': False,
                        'error': f'Database update error: {str(e)}'
                    }, status=500)
                
                # Only proceed with other operations if database update was successful
                if db_update_success:
                    # Get pricing details from job metadata for transaction report
                    total_price = 0.0
                    platform_profit = 0.0
                    
                    # First, try to get from D1 database Vendor_print_jobs table
                    try:
                        api_url = getattr(settings, 'WORKER_API_URL', '')
                        api_key = getattr(settings, 'WORKER_API_KEY', '')
                        
                        if api_url and api_key:
                            # Get job from Vendor_print_jobs table
                            worker_endpoint = api_url.rstrip('/') + '/get-vendor-print-jobs'
                            job_response = requests.post(
                                worker_endpoint,
                                json={'vendor_id': vendor_id, 'filename': filename},
                                headers={
                                    'Content-Type': 'application/json',
                                    'x-api-key': api_key
                                },
                                timeout=10
                            )
                            
                            if job_response.status_code == 200:
                                job_data = job_response.json()
                                if job_data.get('success') and job_data.get('data'):
                                    jobs = job_data.get('data', [])
                                    # Filter by filename to get exact match
                                    matching_jobs = [j for j in jobs if j.get('filename') == filename]
                                    if matching_jobs:
                                        job = matching_jobs[0]  # Get first matching job
                                        total_price = float(job.get('total_price', 0.0))
                                        platform_profit = float(job.get('platform_profit', 0.0))
                                        print(f"✅ Retrieved pricing from D1: total_price={total_price}, platform_profit={platform_profit}")
                    except Exception as e:
                        print(f"⚠️ Error getting pricing from D1: {e}")
                    
                    # Fallback to R2 metadata if D1 doesn't have pricing
                    if total_price == 0.0 or platform_profit == 0.0:
                        try:
                            s3 = boto3.client('s3',
                                              aws_access_key_id=settings.R2_ACCESS_KEY,
                                              aws_secret_access_key=settings.R2_SECRET_KEY,
                                              endpoint_url=settings.R2_ENDPOINT,
                                              region_name='auto')
                            
                            vendor_key = f'vendor_print_jobs/{vendor_id}/{filename}'
                            try:
                                result = s3.head_object(Bucket=settings.R2_BUCKET, Key=vendor_key)
                                metadata = result.get('Metadata', {})
                                
                                # Extract pricing_details
                                pricing_details_str = metadata.get('pricing_details')
                                if pricing_details_str:
                                    try:
                                        pricing_obj = json.loads(pricing_details_str)
                                        if total_price == 0.0:
                                            total_price = float(pricing_obj.get('total', pricing_obj.get('total_price', 0.0)))
                                        if platform_profit == 0.0:
                                            platform_profit = float(pricing_obj.get('platform_profit', 0.0))
                                    except:
                                        pass
                                
                                # Fallback to metadata if pricing_details not available
                                if total_price == 0.0:
                                    total_price = float(metadata.get('total_price', 0.0))
                                if platform_profit == 0.0:
                                    platform_profit = float(metadata.get('platform_profit', 0.0))
                            except Exception as e:
                                print(f"⚠️ Error getting pricing from R2 metadata: {e}")
                        except Exception as e:
                            print(f"⚠️ Error accessing S3 for pricing: {e}")
                    
                    # Create or update vendor transaction report for 2-day period
                    completion_date = datetime.datetime.now().date()
                    try:
                        create_or_update_vendor_transaction(
                            vendor_email, vendor_id, vendor_name, 
                            completion_date, filename, total_price, platform_profit
                        )
                    except Exception as e:
                        print(f"⚠️ Error creating transaction report: {e}")
                    
                    # Free the token if it was found
                    token_freed = False
                    if job_token:
                        try:
                            token_freed = free_token_in_vendor_pool(vendor_email, job_token)
                            if token_freed:
                                print(f"✅ Token {job_token} freed for vendor {vendor_email}")
                            else:
                                print(f"❌ Failed to free token {job_token} for vendor {vendor_email}")
                        except Exception as e:
                            print(f"⚠️ Error freeing token {job_token}: {str(e)}")
                    
                    # Send notification to user and store vendor notification
                    if user_email:
                        try:
                            print(f"📧 Preparing to send notification to user: {user_email} for job: {filename}")
                            send_job_completion_notification(user_email, filename, vendor_id, 'completed', completion_time, job_token)
                            print(f"✅ Completion notification process completed for user: {user_email}")
                            
                            # Also store vendor notification directly with vendor email from session
                            try:
                                store_vendor_notification_direct(vendor_email, filename, vendor_id, user_email, completion_time, job_token)
                                print(f"✅ Stored vendor notification for {vendor_email}")
                            except Exception as e:
                                print(f"⚠️ Error storing vendor notification: {e}")
                                import traceback
                                traceback.print_exc()
                        except Exception as e:
                            print(f"❌ Error in notification process: {e}")
                            import traceback
                            traceback.print_exc()
                    else:
                        # Store vendor notification even if no user email (fallback)
                        try:
                            store_vendor_notification_direct(vendor_email, filename, vendor_id, 'unknown@user.com', completion_time, job_token)
                            print(f"✅ Stored vendor notification for {vendor_email} (no user email)")
                        except Exception as e:
                            print(f"⚠️ Error storing vendor notification: {e}")
                    
                    return JsonResponse({
                        'success': True, 
                        'message': f'Job "{filename}" marked as completed successfully!',
                        'vendor_name': vendor_name,
                        'user_notified': user_email is not None,
                        'token_freed': token_freed,
                        'token_number': job_token if token_freed else None
                    })
            except Exception as e:
                print(f"Error in mark_job_completed: {str(e)}")
                import traceback
                traceback.print_exc()
                return JsonResponse({'success': False, 'error': str(e)})
            finally:
                # Always release the lock, even if there's an error
                file_lock.release()
                # Clean up lock if no longer needed (optional, to prevent memory leak)
                with _job_completion_lock:
                    if filename in _job_completion_locks:
                        # Only remove if lock is not locked (i.e., no one is waiting)
                        if not _job_completion_locks[filename].locked():
                            del _job_completion_locks[filename]
                
        except Exception as e:
            print(f"Error in mark_job_completed: {str(e)}")
            import traceback
            traceback.print_exc()
            return JsonResponse({'success': False, 'error': str(e)})
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'})

@csrf_exempt
def retry_failed_job(request):
    """Retry a failed print job"""
    try:
        data = json.loads(request.body)
        filename = data.get('filename')
        vendor_id = data.get('vendor_id')
        
        if not filename or not vendor_id:
            return JsonResponse({'error': 'Filename and vendor_id required'}, status=400)
        
        # Set job status to accepted so the vendor client will poll and print
        success = update_job_vendor_status(filename, vendor_id, 'accepted')
        
        if success:
            return JsonResponse({
                'success': True,
                'message': f'Job {filename} queued for retry'
            })
        else:
            return JsonResponse({
                'success': False,
                'error': 'Failed to retry job'
            }, status=500)
            
    except Exception as e:
        print(f"Error retrying job: {e}")
        return JsonResponse({'error': 'Internal server error'}, status=500)

@csrf_exempt
def cancel_failed_job(request):
    """Cancel a failed print job and compensate user"""
    try:
        data = json.loads(request.body)
        filename = data.get('filename')
        vendor_id = data.get('vendor_id')
        user_email = data.get('user_email')
        job_price = data.get('job_price', 0.0)
        
        if not filename:
            return JsonResponse({'error': 'Filename required'}, status=400)
        
        # Get vendor_id from session if not provided
        if not vendor_id:
            vendor_id = request.session.get('vendor_id')
            if not vendor_id:
                return JsonResponse({'error': 'Vendor ID required'}, status=400)
        
        # Mark job as cancelled in R2
        success = update_job_vendor_status(filename, vendor_id, 'cancelled')
        
        if success:
            # Also update job_failed status to YES to mark it as failed
            update_job_failed_status_in_r2(filename, vendor_id, 'YES')
            
            # Compensate user with points
            if user_email and job_price > 0:
                points = float(job_price)  # Preserve decimal values (1 rupee = 1 point)
                # Add points to user account
                add_user_points(user_email, points, f"Compensation for cancelled job: {filename}")
                print(f"💰 Compensated user {user_email} with {points} points for cancelled job: {filename}")
            
            # Send notification to user
            if user_email:
                # Get the token associated with this job
                job_token = get_token_from_file_metadata(filename, vendor_id)
                if job_token:
                    print(f"🔍 Found token {job_token} for cancelled job {filename}")
                else:
                    print(f"⚠️ No token found for cancelled job {filename}")
                
                send_job_completion_notification(user_email, filename, vendor_id, 'cancelled', datetime.datetime.now().isoformat(), job_token)
            
            return JsonResponse({
                'success': True,
                'message': f'Job {filename} cancelled and user compensated with {int(job_price)} points'
            })
        else:
            return JsonResponse({
                'success': False,
                'error': 'Failed to cancel job'
            }, status=500)
            
    except Exception as e:
        print(f"Error cancelling job: {e}")
        traceback.print_exc()
        return JsonResponse({'error': 'Internal server error'}, status=500)

@csrf_exempt
def cancel_print_job(request):
    """Cancel a print job and refund user - updates D1 database only, not R2"""
    try:
        data = json.loads(request.body)
        filename = data.get('filename')
        user_email = data.get('user_email')
        
        # Handle job_price - can be None, string, or number
        job_price_raw = data.get('job_price')
        if job_price_raw is None or job_price_raw == '':
            job_price = 0.0
        else:
            try:
                job_price = float(job_price_raw)
            except (ValueError, TypeError):
                job_price = 0.0
        
        if not filename:
            return JsonResponse({
                'success': False,
                'error': 'Missing filename parameter'
            }, status=400)
        
        # Get vendor_id from session
        vendor_id = request.session.get('vendor_id')
        if not vendor_id:
            return JsonResponse({
                'success': False,
                'error': 'Vendor not authenticated'
            }, status=401)
        
        # Get vendor_email from session
        vendor_email = request.session.get('vendor_email')
        if not vendor_email:
            return JsonResponse({
                'success': False,
                'error': 'Vendor not authenticated'
            }, status=401)
        
        print(f"🔄 Cancelling job {filename} for vendor {vendor_id}")
        print(f"📧 User email (provided): {user_email}")
        print(f"💰 Job price: {job_price}")
        
        # Get user_email from D1 database if not provided
        if not user_email:
            try:
                api_url = getattr(settings, 'WORKER_API_URL', '')
                api_key = getattr(settings, 'WORKER_API_KEY', '')
                
                if api_url and api_key:
                    # Construct the Worker API endpoint to get job details
                    if '/add-contact' in api_url:
                        worker_endpoint = api_url.replace('/add-contact', '/get-vendor-print-jobs')
                    elif '/add-vendor-register' in api_url:
                        worker_endpoint = api_url.replace('/add-vendor-register', '/get-vendor-print-jobs')
                    else:
                        worker_endpoint = api_url.rstrip('/') + '/get-vendor-print-jobs'
                    
                    # Get job details from D1 database
                    payload = {
                        'filename': filename,
                        'vendor_id': vendor_id,
                        'vendor_email': vendor_email
                    }
                    
                    resp = requests.post(
                        worker_endpoint,
                        json=payload,
                        headers={
                            'Content-Type': 'application/json',
                            'x-api-key': api_key
                        },
                        timeout=10
                    )
                    
                    if resp.status_code == 200:
                        job_data = resp.json()
                        if job_data.get('success') and job_data.get('data'):
                            jobs = job_data.get('data', [])
                            if jobs and len(jobs) > 0:
                                user_email = jobs[0].get('user_email', '')
                                if user_email:
                                    print(f"📧 Retrieved user email from D1 database: {user_email}")
                                    # Also get job_price from D1 if not provided
                                    if job_price == 0.0:
                                        job_price = float(jobs[0].get('total_price', jobs[0].get('final_amount', 0.0)) or 0.0)
                                        print(f"💰 Retrieved job price from D1 database: {job_price}")
            except Exception as e:
                print(f"⚠️ Error getting user email from D1 database: {str(e)}")
                # Continue without user_email - will skip points refund
        
        # Update Vendor_print_jobs and User_print_jobs tables in D1 database only
        completion_time = datetime.datetime.now().isoformat()
        db_update_success = False
        
        try:
            # Update job_completed status in Vendor_print_jobs and User_print_jobs tables via Worker API
            api_url = getattr(settings, 'WORKER_API_URL', '')
            api_key = getattr(settings, 'WORKER_API_KEY', '')
            
            if api_url and api_key:
                # Construct the Worker API endpoint
                if '/add-contact' in api_url:
                    worker_endpoint = api_url.replace('/add-contact', '/update-job-completed')
                elif '/add-vendor-register' in api_url:
                    worker_endpoint = api_url.replace('/add-vendor-register', '/update-job-completed')
                else:
                    worker_endpoint = api_url.rstrip('/') + '/update-job-completed'
                
                # Set job_completed to 'CANCELLED' (uppercase) to mark as cancelled (worker API expects uppercase)
                payload = {
                    'filename': filename,
                    'job_completed': 'CANCELLED',  # Worker API expects 'CANCELLED' (uppercase)
                    'vendor_email': vendor_email,
                    'vendor_id': vendor_id,
                    'completion_time': completion_time,
                    'user_email': user_email if user_email else ''  # Include user_email to update User_print_jobs table
                }
                
                resp = requests.post(
                    worker_endpoint,
                    json=payload,
                    headers={
                        'Content-Type': 'application/json',
                        'x-api-key': api_key
                    },
                    timeout=10
                )
                
                if resp.status_code == 200:
                    db_update_success = True
                    print(f"✅ Updated Vendor_print_jobs and User_print_jobs tables in D1: {filename} -> job_completed='CANCELLED'")
                else:
                    print(f"⚠️ Failed to update D1 database tables: {resp.status_code} - {resp.text}")
            else:
                print(f"⚠️ Worker API not configured")
        except Exception as e:
            print(f"⚠️ Error updating D1 database tables: {str(e)}")
            traceback.print_exc()
        
        if not db_update_success:
            return JsonResponse({
                'success': False,
                'error': 'Failed to update database - job may not exist in D1 database'
            }, status=500)
        
        # Add points to user's account in D1 database - preserve decimal values
        if user_email and job_price > 0:
            points = float(job_price)  # Preserve decimal values (2.4 -> 2.4, 2.5 -> 2.5)
            points_success = add_user_points(user_email, points, f"Refund for cancelled job: {filename}")
            if points_success:
                print(f"💰 Refunded user {user_email} with {points} points (₹{job_price}) for cancelled job: {filename}")
            else:
                print(f"⚠️ Failed to add points for user {user_email}")
        
        # Clear vendor cache to ensure fresh data on next load
        clear_vendor_cache(vendor_email, vendor_id)
        
        return JsonResponse({
            'success': True,
            'message': f'Job {filename} cancelled and user refunded with {round(job_price)} points'
        })
            
    except Exception as e:
        print(f"Error cancelling print job: {e}")
        traceback.print_exc()
        return JsonResponse({'error': 'Internal server error'}, status=500)

def simple_update_vendor_status(filename, vendor_id, status):
    """Simple function to update vendor status directly in R2"""
    try:
        print(f"🔄 Attempting to update vendor status for {filename} to {status}")
        
        s3 = boto3.client('s3',
            aws_access_key_id=settings.R2_ACCESS_KEY,
            aws_secret_access_key=settings.R2_SECRET_KEY,
            endpoint_url=settings.R2_ENDPOINT,
            region_name='auto')
        
        # Try to update the job in vendor_print_jobs folder
        vendor_key = f'vendor_print_jobs/{vendor_id}/{filename}'
        print(f"🔍 Looking for job at: {vendor_key}")
        
        try:
            # Get the existing job data
            response = s3.get_object(Bucket=settings.R2_BUCKET, Key=vendor_key)
            metadata = json.loads(response['Body'].read().decode('utf-8'))
            print(f"📄 Found job metadata for {filename}")
            
            # Update vendor status
            metadata['vendor_status'] = status
            metadata['updated_at'] = datetime.datetime.now().isoformat()
            
            # Save back to R2
            s3.put_object(
                Bucket=settings.R2_BUCKET,
                Key=vendor_key,
                Body=json.dumps(metadata, indent=2),
                ContentType='application/json'
            )
            
            print(f"✅ Updated vendor status for {filename}: {status}")
            return True
            
        except Exception as e:
            print(f"❌ Could not update vendor status for {filename}: {e}")
            print(f"🔍 Vendor key attempted: {vendor_key}")
            return False
            
    except Exception as e:
        print(f"❌ Error in simple_update_vendor_status: {e}")
        import traceback
        traceback.print_exc()
        return False
def vendor_dashboard_notification(request):
    """Handle vendor dashboard notifications including compensation requests"""
    try:
        data = json.loads(request.body)
        notification_type = data.get('notification_type')
        vendor_id = data.get('vendor_id')
        action = data.get('action')
        message = data.get('message', '')
        level = data.get('level', 'info')
        
        # Get vendor email from vendor_id to append notification
        # Vendor notification functionality removed
        
        if action == 'compensate_failed_job':
            # Handle compensation request from vendor client
            job_filename = data.get('job_filename')
            user_email = data.get('user_email')
            job_price = float(data.get('job_price', 0))
            reason = data.get('reason', 'Print job failed')
            
            if not job_filename or not user_email or job_price <= 0:
                return JsonResponse({
                    'success': False,
                    'error': 'Missing required compensation data'
                }, status=400)
            
            # Compensate user with points (1 rupee = 1 point) - preserve decimal values
            points = float(job_price)
            success = add_user_points(user_email, points, reason)
            
            if success:
                print(f"💰 Compensated user {user_email} with {points} points for failed job: {job_filename}")
                return JsonResponse({
                    'success': True,
                    'message': f'User compensated with {points} points for failed job: {job_filename}'
                })
            else:
                return JsonResponse({
                    'success': False,
                    'error': 'Failed to compensate user'
                }, status=500)
        
        elif notification_type == 'printer_issue':
            # Handle printer issue notification
            printer_name = data.get('printer_name')
            issue_type = data.get('issue_type')
            error_message = data.get('error_message')
            
            print(f"🖨️ Printer issue notification: {printer_name} - {issue_type}: {error_message}")
            return JsonResponse({
                'success': True,
                'message': f'Printer issue notification received for {printer_name}'
            })
        
        elif notification_type == 'job_waiting':
            # Handle job waiting notification
            job_filename = data.get('job_filename')
            printer_name = data.get('printer_name')
            issue_type = data.get('issue_type')
            error_message = data.get('error_message')
            
            print(f"⏳ Job waiting notification: {job_filename} - {printer_name}: {error_message}")
            return JsonResponse({
                'success': True,
                'message': f'Job waiting notification received for {job_filename}'
            })
        
        else:
            return JsonResponse({
                'success': False,
                'error': 'Unknown notification type or action'
            }, status=400)
            
    except Exception as e:
        print(f"Error handling vendor dashboard notification: {e}")
        traceback.print_exc()
        return JsonResponse({'error': 'Internal server error'}, status=500)

def update_job_vendor_status(filename, vendor_id, status):
    """Update vendor status for a specific job"""
    try:
        s3 = boto3.client('s3',
            aws_access_key_id=settings.R2_ACCESS_KEY,
            aws_secret_access_key=settings.R2_SECRET_KEY,
            endpoint_url=settings.R2_ENDPOINT,
            region_name='auto')
        
        # First, try to find the job in print_requests folder
        source_key = f'print_requests/{filename}'
        vendor_key = f'vendor_print_jobs/{vendor_id}/{filename}'
        
        try:
            # Get the job from print_requests
            response = s3.get_object(Bucket=settings.R2_BUCKET, Key=source_key)
            metadata = json.loads(response['Body'].read().decode('utf-8'))
            
            # Update metadata with vendor status
            metadata['vendor_status'] = status
            metadata['vendor_id'] = vendor_id
            metadata['updated_at'] = datetime.datetime.now().isoformat()
            
            # Copy to vendor_print_jobs folder
            s3.put_object(
                Bucket=settings.R2_BUCKET,
                Key=vendor_key,
                Body=json.dumps(metadata, indent=2),
                ContentType='application/json'
            )
            
            print(f"✅ Updated vendor status for {filename}: {status}")
            return True
            
        except Exception as e:
            # If not found in print_requests, try to update existing vendor job
            try:
                response = s3.get_object(Bucket=settings.R2_BUCKET, Key=vendor_key)
                metadata = json.loads(response['Body'].read().decode('utf-8'))
                metadata['vendor_status'] = status
                metadata['updated_at'] = datetime.datetime.now().isoformat()
                
                s3.put_object(
                    Bucket=settings.R2_BUCKET,
                    Key=vendor_key,
                    Body=json.dumps(metadata, indent=2),
                    ContentType='application/json'
                )
                
                print(f"✅ Updated existing vendor job status for {filename}: {status}")
                return True
                
            except Exception as e2:
                print(f"❌ Job not found in either location: {filename}")
                return False
            
    except Exception as e:
        print(f"❌ Error updating vendor status: {e}")
        return False

def update_job_failed_status_in_r2(filename, vendor_id, failed_status):
    """Update job_failed status in R2 storage"""
    try:
        s3 = boto3.client('s3',
                         aws_access_key_id=settings.R2_ACCESS_KEY,
                         aws_secret_access_key=settings.R2_SECRET_KEY,
                         endpoint_url=settings.R2_ENDPOINT,
                         region_name='auto')
        
        # Update in vendor_print_jobs folder
        vendor_key = f'vendor_print_jobs/{vendor_id}/{filename}'
        
        try:
            result = s3.get_object(Bucket=settings.R2_BUCKET, Key=vendor_key)
            metadata = json.loads(result['Body'].read().decode('utf-8'))
            metadata['job_failed'] = failed_status
            metadata['updated_at'] = datetime.datetime.now().isoformat()
            
            s3.put_object(
                Bucket=settings.R2_BUCKET,
                Key=vendor_key,
                Body=json.dumps(metadata, indent=2),
                ContentType='application/json'
            )
            
            print(f"✅ Updated job_failed status for {filename}: {failed_status}")
            return True
        except Exception as e:
            print(f"❌ Job not found in R2: {filename} - {e}")
            return False
            
    except Exception as e:
        print(f"❌ Error updating job_failed status: {e}")
        traceback.print_exc()
        return False

def add_user_points(user_email, points, reason):
    """Add points to user account - stores in D1 database (User_points table)"""
    try:
        # Store points data using IST timestamps
        now_ist = timezone.localtime(timezone.now())
        date_str = now_ist.strftime('%Y-%m-%d')
        time_str = now_ist.strftime('%H:%M:%S')
        timestamp = now_ist.isoformat()
        
        # Store in D1 database via Worker API
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            print(f"⚠️ Worker API not configured, skipping database storage")
            return False

        # Construct the Worker API endpoint
        if '/add-contact' in api_url:
            worker_endpoint = api_url.replace('/add-contact', '/add-user-points')
        elif '/add-vendor-register' in api_url:
            worker_endpoint = api_url.replace('/add-vendor-register', '/add-user-points')
        else:
            worker_endpoint = api_url.rstrip('/') + '/add-user-points'
        
        payload = {
            'user_email': user_email,
            'points': points,
            'date': date_str,
            'time': time_str,
            'reason': reason or '',
            'transaction_timestamp': timestamp
        }
        
        resp = requests.post(
            worker_endpoint,
            json=payload,
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        if resp.status_code == 200:
            print(f"✅ Stored {points} points for {user_email} in database on {date_str} at {time_str}")
            return True
        else:
            error_text = resp.text[:200] if resp.text else 'Unknown error'
            print(f"⚠️ Failed to store points in database: {resp.status_code} - {error_text}")
            return False
            
    except Exception as e:
        print(f"❌ Error adding points: {e}")
        return False

@csrf_exempt
def get_user_notifications(request):
    """Get user notifications from D1 database - only daily completed print job notifications"""
    try:
        if not request.user.is_authenticated:
            return JsonResponse({'success': False, 'error': 'Not authenticated'}, status=401)
        
        user_email = request.user.email
        notifications = []
        
        # Get notifications from D1 database via Worker API
        try:
            api_url = getattr(settings, 'WORKER_API_URL', '')
            api_key = getattr(settings, 'WORKER_API_KEY', '')
            
            if not api_url or not api_key:
                print("⚠️ Worker API not configured for notifications")
                return JsonResponse({
                    'success': True,
                    'notifications': [],
                    'unread_count': 0
                })
            
            # Construct Worker API endpoint
            worker_endpoint = api_url.rstrip('/') + '/get-user-notifications'
            
            # Get today's date for filtering (YYYY-MM-DD format)
            today = datetime.datetime.now().date()
            today_str = today.strftime('%Y-%m-%d')
            
            payload = {
                'user_email': user_email,
                'date': today_str  # Filter for today's notifications
            }
            
            response = requests.post(
                worker_endpoint,
                json=payload,
                headers={
                    'x-api-key': api_key,
                    'Content-Type': 'application/json'
                },
                timeout=10
            )
            
            if response.status_code == 200:
                data = response.json()
                if data.get('success') and data.get('notifications'):
                    notifications = data.get('notifications', [])
                    
                    # Filter for completed print job notifications only and today's date
                    filtered_notifications = []
                    for notif in notifications:
                        # Check if it's a completed job notification
                        if (notif.get('type') == 'job_completed' and 
                            notif.get('status') == 'completed'):
                            
                            # Check if notification is from today
                            notification_date = None
                            if 'completion_time' in notif and notif['completion_time']:
                                try:
                                    notification_date = datetime.datetime.fromisoformat(
                                        notif['completion_time'].replace('Z', '+00:00')
                                    ).date()
                                except:
                                    try:
                                        notification_date = datetime.datetime.fromisoformat(
                                            notif['completion_time'].split('T')[0]
                                        ).date()
                                    except:
                                        pass
                            
                            if not notification_date and 'created_at' in notif and notif['created_at']:
                                try:
                                    notification_date = datetime.datetime.fromisoformat(
                                        notif['created_at'].replace('Z', '+00:00')
                                    ).date()
                                except:
                                    try:
                                        notification_date = datetime.datetime.fromisoformat(
                                            notif['created_at'].split('T')[0]
                                        ).date()
                                    except:
                                        pass
                            
                            # Only include notifications from today
                            if notification_date == today:
                                # Ensure required fields are present
                                if 'filename' not in notif:
                                    notif['filename'] = 'Document'
                                if 'token' not in notif or not notif['token']:
                                    notif['token'] = 'Unknown'
                                
                                # Convert read from int to bool if needed
                                if isinstance(notif.get('read'), int):
                                    notif['read'] = bool(notif['read'])
                                
                                filtered_notifications.append(notif)
                    
                    notifications = filtered_notifications
                    
                    # Sort by completion_time or created_at (newest first)
                    notifications.sort(
                        key=lambda x: x.get('completion_time') or x.get('created_at') or '', 
                        reverse=True
                    )
                    
                    # Count unread notifications
                    unread_count = sum(1 for n in notifications if not n.get('read', False))
                    
                    print(f"✅ Retrieved {len(notifications)} notifications for {user_email} from D1 database")
                    return JsonResponse({
                        'success': True,
                        'notifications': notifications,
                        'unread_count': unread_count
                    })
                else:
                    print(f"⚠️ No notifications found in D1 database for {user_email}")
                    return JsonResponse({
                        'success': True,
                        'notifications': [],
                        'unread_count': 0
                    })
            else:
                print(f"⚠️ Error getting notifications from D1: {response.status_code} - {response.text[:200]}")
                return JsonResponse({
                    'success': True,
                    'notifications': [],
                    'unread_count': 0
                })
            
        except Exception as e:
            print(f"❌ Error getting notifications from D1 database: {e}")
            import traceback
            traceback.print_exc()
            return JsonResponse({
                'success': True,
                'notifications': [],
                'unread_count': 0
            })
            
    except Exception as e:
        print(f"❌ Error in get_user_notifications: {e}")
        import traceback
        traceback.print_exc()
        return JsonResponse({'success': False, 'error': 'Internal server error'}, status=500)

@csrf_exempt
def get_user_points(request):
    """Get user total points"""
    try:
        if not request.user.is_authenticated:
            return JsonResponse({'success': False, 'error': 'Not authenticated'}, status=401)
        
        user_email = request.user.email
        total_points = get_total_user_points(user_email)
        
        return JsonResponse({
            'success': True,
            'total_points': total_points
        })
        
    except Exception as e:
        print(f"❌ Error getting user points: {e}")
        return JsonResponse({
            'success': False,
            'error': str(e),
            'total_points': 0
        })

@csrf_exempt
def mark_notification_read(request):
    """Mark a notification as read"""
    try:
        if not request.user.is_authenticated:
            return JsonResponse({'success': False, 'error': 'Not authenticated'}, status=401)
        
        data = json.loads(request.body)
        notification_id = data.get('notification_id')
        
        if not notification_id:
            return JsonResponse({'success': False, 'error': 'Notification ID required'}, status=400)
        
        user_email = request.user.email
        notification_key = f'user_notifications/{sanitize_email(user_email)}/{notification_id}.json'
        
        try:
            # Initialize R2 client
            s3 = boto3.client('s3',
                              aws_access_key_id=settings.R2_ACCESS_KEY,
                              aws_secret_access_key=settings.R2_SECRET_KEY,
                              endpoint_url=settings.R2_ENDPOINT,
                              region_name='auto')
            
            # Get current notification
            result = s3.get_object(Bucket=settings.R2_BUCKET, Key=notification_key)
            notification_data = json.loads(result['Body'].read().decode('utf-8'))
            
            # Mark as read
            notification_data['read'] = True
            notification_data['read_at'] = datetime.datetime.now().isoformat()
            
            # Update notification
            s3.put_object(
                Bucket=settings.R2_BUCKET,
                Key=notification_key,
                Body=json.dumps(notification_data, indent=2),
                ContentType='application/json'
            )
            
            return JsonResponse({'success': True, 'message': 'Notification marked as read'})
            
        except Exception as e:
            print(f"Error marking notification as read: {e}")
            return JsonResponse({'success': False, 'error': 'Notification not found'}, status=404)
            
    except Exception as e:
        print(f"Error in mark_notification_read: {e}")
        return JsonResponse({'success': False, 'error': 'Internal server error'}, status=500)

def get_total_user_points(user_email: str) -> float:
    """Get total points for the user from D1 database. Returns 0 if none."""
    try:
        # Get from D1 database via Worker API
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            print(f"⚠️ Worker API not configured, returning 0 points")
            return 0

        # Construct the Worker API endpoint
        if '/add-contact' in api_url:
            worker_endpoint = api_url.replace('/add-contact', '/get-user-total-points')
        elif '/add-vendor-register' in api_url:
            worker_endpoint = api_url.replace('/add-vendor-register', '/get-user-total-points')
        else:
            worker_endpoint = api_url.rstrip('/') + '/get-user-total-points'
        
        resp = requests.post(
            worker_endpoint,
            json={'user_email': user_email},
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        if resp.status_code == 200:
            response_data = resp.json()
            if response_data.get('success'):
                total_points = response_data.get('total_points', 0)
                # Preserve decimal precision - round to 1 decimal place
                total_points = round(float(total_points), 1)
                print(f"💰 Total points for {user_email}: {total_points}")
                return total_points
        
        print(f"⚠️ Failed to get points from database: {resp.status_code}")
        return 0
    except Exception as e:
        print(f"❌ Error getting user points: {e}")
        return 0

def deduct_user_points(user_email: str, points: int, reason: str) -> bool:
    """Deduct points from a user's account - stores in D1 database"""
    try:
        if points <= 0:
            return True  # No points to deduct
        
        # Create negative points record (use add_user_points with negative value)
        return add_user_points(user_email, -points, reason)
    except Exception as e:
        print(f"❌ Error deducting points: {e}")
        return False

@csrf_exempt
def return_user_points(request):
    """Return points to user - used when payment fails after points were deducted"""
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Invalid method'}, status=405)
    
    try:
        body = json.loads(request.body.decode('utf-8')) if request.body else {}
        user_email = body.get('user_email')
        points = float(body.get('points', 0))  # Preserve decimal values
        reason = body.get('reason', 'Points returned due to payment failure')
        
        if not user_email:
            return JsonResponse({'success': False, 'error': 'User email required'}, status=400)
        
        if points <= 0:
            return JsonResponse({'success': False, 'error': 'Invalid points amount'}, status=400)
        
        # Return points by adding them back
        success = add_user_points(user_email, points, reason)
        
        if success:
            print(f"✅ Returned {points} points to {user_email}: {reason}")
            return JsonResponse({'success': True, 'message': f'Returned {points} points'})
        else:
            return JsonResponse({'success': False, 'error': 'Failed to return points'}, status=500)
            
    except Exception as e:
        print(f"❌ Error returning points: {str(e)}")
        return JsonResponse({'success': False, 'error': str(e)}, status=500)

def store_vendor_print_job_in_db(vendor_id, vendor_email, user_email, filename, storage_folder, r2_path, metadata, pricing_details=None, user_id=None, shop_id=None):
    """Store vendor print job in D1 database via Worker API"""
    try:
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            print(f"⚠️ Worker API not configured, skipping database storage")
            return False

        # Construct the Worker API endpoint
        if '/add-contact' in api_url:
            worker_endpoint = api_url.replace('/add-contact', '/add-vendor-print-job')
        elif '/add-vendor-register' in api_url:
            worker_endpoint = api_url.replace('/add-vendor-register', '/add-vendor-print-job')
        else:
            worker_endpoint = api_url.rstrip('/') + '/add-vendor-print-job'
        
        # Extract pricing information from pricing_details if available
        total_price = None
        platform_profit = None
        price_per_page = None
        page_count = None
        num_copies = None
        
        if pricing_details:
            if isinstance(pricing_details, dict):
                total_price = pricing_details.get('total_price', pricing_details.get('total'))
                platform_profit = pricing_details.get('platform_profit')
                if isinstance(pricing_details.get('pricing_breakdown'), dict):
                    breakdown = pricing_details.get('pricing_breakdown', {})
                    price_per_page = breakdown.get('price_per_page')
                    page_count = breakdown.get('page_count')
                    num_copies = breakdown.get('num_copies')
            elif isinstance(pricing_details, str):
                try:
                    pricing_obj = json.loads(pricing_details)
                    total_price = pricing_obj.get('total_price', pricing_obj.get('total'))
                    platform_profit = pricing_obj.get('platform_profit')
                    price_per_page = pricing_obj.get('price_per_page', pricing_obj.get('per_page'))
                    page_count = pricing_obj.get('page_count', pricing_obj.get('pages'))
                    num_copies = pricing_obj.get('num_copies', pricing_obj.get('copies'))
                except:
                    pass
        
        # Also check metadata for pricing info
        if total_price is None:
            total_price = metadata.get('total_price')
        if price_per_page is None:
            price_per_page = metadata.get('price_per_page')
        if page_count is None:
            page_count = metadata.get('page_count')
        if num_copies is None:
            num_copies = metadata.get('num_copies')
        pages_value = metadata.get('page_count') or metadata.get('pages')

        # Normalize color and page range for Vendor_print_jobs (Mixed/Both support)
        raw_color = metadata.get('color_mode', metadata.get('color', 'Black and White'))
        bw_range_value_db = (metadata.get('bwPageRangeValue') or '').strip()
        color_range_value_db = (metadata.get('colorPageRangeValue') or '').strip()

        # If both BW and Color ranges are present, treat this as a mixed (Both) job
        # even if the upstream code forgot to set color="Mix"
        if bw_range_value_db and color_range_value_db and raw_color in ['Black and White', 'Color', 'bw', 'color']:
            normalized_color = 'Mix'
        else:
            # Preserve explicit Mix/Mixed values, otherwise use whatever came from metadata
            if str(raw_color).lower() in ['mix', 'mixed', 'both']:
                normalized_color = 'Mix'
            else:
                normalized_color = raw_color

        # Ensure a human-readable combined pageRange string for mixed jobs
        page_range_combined = metadata.get('pageRange', '') or ''
        if normalized_color == 'Mix':
            bw_label = bw_range_value_db if bw_range_value_db else 'all'
            color_label = color_range_value_db if color_range_value_db else 'all'
            page_range_combined = f"BW: {bw_label} | Color: {color_label}"
        
        # Convert string values to appropriate types
        if total_price:
            try:
                total_price = float(total_price)
            except:
                total_price = None
        if platform_profit:
            try:
                platform_profit = float(platform_profit)
            except:
                platform_profit = None
        if price_per_page:
            try:
                price_per_page = float(price_per_page)
            except:
                price_per_page = None
        if page_count:
            try:
                page_count = int(page_count)
            except:
                page_count = None
        if num_copies:
            try:
                num_copies = int(num_copies)
            except:
                num_copies = None
        if pages_value:
            try:
                pages_value = int(pages_value)
            except:
                pages_value = None
        
        # Store only base_price in pricing_details (not full JSON)
        pricing_details_str = None
        if pricing_details:
            if isinstance(pricing_details, dict):
                # Extract base_price from pricing_breakdown or directly from pricing_details
                breakdown = pricing_details.get('pricing_breakdown', {})
                base_price = breakdown.get('base_price', 0) if isinstance(breakdown, dict) else 0
                if base_price == 0:
                    base_price = pricing_details.get('base_price', 0)
                pricing_details_str = str(base_price) if base_price else None
            elif isinstance(pricing_details, str):
                # If it's already a string, try to parse and extract base_price
                try:
                    parsed = json.loads(pricing_details)
                    if isinstance(parsed, dict):
                        breakdown = parsed.get('pricing_breakdown', {})
                        base_price = breakdown.get('base_price', 0) if isinstance(breakdown, dict) else 0
                        if base_price == 0:
                            base_price = parsed.get('base_price', 0)
                        pricing_details_str = str(base_price) if base_price else None
                    else:
                        pricing_details_str = str(parsed) if parsed else None
                except:
                    # If parsing fails, assume it's already just a base_price value
                    pricing_details_str = pricing_details
        
        # Get shop address and shop name from metadata or vendor data
        shop_address = metadata.get('shop_address', '')
        shop_name = metadata.get('shop_name', '')
        if (not shop_address or not shop_name) and vendor_email:
            try:
                vendor_data = get_vendor_coordinates_from_email(vendor_email)
                if vendor_data:
                    if not shop_address:
                        shop_address = vendor_data.get('shop_address', '')
                    if not shop_name:
                        shop_name = vendor_data.get('vendor_name', vendor_data.get('shop_name', ''))
            except:
                pass
        
        payload = {
            'vendor_id': vendor_id,
            'vendor_email': vendor_email or '',
            'user_email': user_email,
            'user_id': user_id or metadata.get('user_id', ''),
            'filename': filename,
            'storage_folder': storage_folder,
            'shop_id': shop_id or metadata.get('shop_id', vendor_id),
            'r2_path': r2_path,
            'service_type': metadata.get('service_type', ''),
            'status': metadata.get('status', 'pending'),
            'job_completed': metadata.get('job_completed', 'NO'),
            'vendor_status': metadata.get('vendor_status', 'not sended'),
            'token': metadata.get('token', ''),
            'job_id': metadata.get('job_id', ''),
            'copies': metadata.get('copies', '1'),
            'color': normalized_color,
            'print_type': metadata.get('print_type', metadata.get('layout_type', 'single_side')),
            'orientation': metadata.get('orientation', ''),
            'pageSize': metadata.get('pageSize', ''),
            'pageRange': page_range_combined,
            'specificPages': metadata.get('specificPages', ''),
            # Mixed (Both) color support: store separate page ranges too
            'bwPageRange': metadata.get('bwPageRange', ''),
            'bwPageRangeValue': metadata.get('bwPageRangeValue', ''),
            'colorPageRange': metadata.get('colorPageRange', ''),
            'colorPageRangeValue': metadata.get('colorPageRangeValue', ''),
            'spiralBinding': metadata.get('spiralBinding', 'No'),
            'lamination': metadata.get('lamination', 'No'),
            'service_name': metadata.get('service_name', ''),
            'feedback': metadata.get('feedback', ''),
            'quality': metadata.get('quality', ''),
            'thickness': metadata.get('thickness', ''),
            'points_applied': metadata.get('points_applied', 'false'),
            'points_used': metadata.get('points_used', '0'),
            'timestamp': metadata.get('timestamp', datetime.datetime.now().isoformat()),
            'completion_time': metadata.get('completion_time', ''),
            'rendered_status': metadata.get('rendered_status', 'NO'),
            'trash': metadata.get('trash', 'NO'),
            'total_price': total_price,
            'platform_profit': platform_profit,
            'price_per_page': price_per_page,
            'final_amount': metadata.get('final_amount'),
            'page_count': page_count,
            'pages': pages_value,
            'num_copies': num_copies,
            'paper_type': metadata.get('paper_type', ''),
            'color_mode': normalized_color,
            'layout_type': metadata.get('layout_type', ''),
            'pricing_details': pricing_details_str,
            'shop_address': shop_address,
            'shop_name': shop_name
        }
        
        # Convert final_amount to float if present
        if payload['final_amount']:
            try:
                payload['final_amount'] = float(payload['final_amount'])
            except:
                payload['final_amount'] = None
        
        print(f"🔍 Attempting to store vendor print job: {filename} at endpoint: {worker_endpoint}")
        print(f"🔍 Payload keys: {list(payload.keys())}")
        
        resp = requests.post(
            worker_endpoint,
            json=payload,
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        print(f"🔍 Response status: {resp.status_code}")
        print(f"🔍 Response text: {resp.text[:500]}")
        
        if resp.status_code == 200:
            try:
                response_data = resp.json()
                if response_data.get('success'):
                    print(f"✅ Stored vendor print job in database: {filename} in {storage_folder}")
                    return True
                else:
                    error_msg = response_data.get('error', 'Unknown error')
                    print(f"⚠️ Worker API returned success=false: {error_msg}")
                    return False
            except json.JSONDecodeError:
                print(f"⚠️ Invalid JSON response from worker: {resp.text[:200]}")
                return False
        else:
            print(f"⚠️ Failed to store vendor print job in database: {resp.status_code} - {resp.text[:200]}")
            return False
            
    except requests.exceptions.RequestException as e:
        print(f"❌ Network error storing vendor print job in database: {e}")
        return False
    except Exception as e:
        print(f"❌ Error storing vendor print job in database: {e}")
        import traceback
        traceback.print_exc()
        return False

def store_user_print_job_in_db(vendor_id, vendor_email, user_email, filename, storage_folder, r2_path, metadata, pricing_details=None, user_id=None, shop_id=None):
    """Store user print job in D1 database via Worker API"""
    try:
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')

        if not api_url or not api_key:
            print(f"⚠️ Worker API not configured, skipping user print job storage")
            return False

        if '/add-contact' in api_url:
            worker_endpoint = api_url.replace('/add-contact', '/add-user-print-job')
        elif '/add-vendor-register' in api_url:
            worker_endpoint = api_url.replace('/add-vendor-register', '/add-user-print-job')
        else:
            worker_endpoint = api_url.rstrip('/') + '/add-user-print-job'

        total_price = None
        platform_profit = None
        price_per_page = None
        page_count = None
        num_copies = None

        if pricing_details:
            if isinstance(pricing_details, dict):
                total_price = pricing_details.get('total_price', pricing_details.get('total'))
                platform_profit = pricing_details.get('platform_profit')
                if isinstance(pricing_details.get('pricing_breakdown'), dict):
                    breakdown = pricing_details.get('pricing_breakdown', {})
                    price_per_page = breakdown.get('price_per_page')
                    page_count = breakdown.get('page_count')
                    num_copies = breakdown.get('num_copies')
            elif isinstance(pricing_details, str):
                try:
                    pricing_obj = json.loads(pricing_details)
                    total_price = pricing_obj.get('total_price', pricing_obj.get('total'))
                    platform_profit = pricing_obj.get('platform_profit')
                    price_per_page = pricing_obj.get('price_per_page', pricing_obj.get('per_page'))
                    page_count = pricing_obj.get('page_count', pricing_obj.get('pages'))
                    num_copies = pricing_obj.get('num_copies', pricing_obj.get('copies'))
                except:
                    pass

        if total_price is None:
            total_price = metadata.get('total_price')
        if price_per_page is None:
            price_per_page = metadata.get('price_per_page')
        if page_count is None:
            page_count = metadata.get('page_count') or metadata.get('pages')
        if num_copies is None:
            num_copies = metadata.get('num_copies') or metadata.get('copies')
        pages_value = metadata.get('page_count') or metadata.get('pages')

        def to_float(value):
            if value is None:
                return None
            try:
                return float(value)
            except (TypeError, ValueError):
                return None

        def to_int(value):
            if value is None:
                return None
            try:
                return int(value)
            except (TypeError, ValueError):
                return None

        total_price = to_float(total_price)
        platform_profit = to_float(platform_profit)
        price_per_page = to_float(price_per_page)
        page_count = to_int(page_count)
        num_copies = to_int(num_copies)
        pages_value = to_int(pages_value)

        # Normalize color and page range for User_print_jobs as well
        raw_color = metadata.get('color_mode', metadata.get('color', 'Black and White'))
        bw_range_value_db = (metadata.get('bwPageRangeValue') or '').strip()
        color_range_value_db = (metadata.get('colorPageRangeValue') or '').strip()

        if bw_range_value_db and color_range_value_db and raw_color in ['Black and White', 'Color', 'bw', 'color']:
            normalized_color = 'Mix'
        else:
            if str(raw_color).lower() in ['mix', 'mixed', 'both']:
                normalized_color = 'Mix'
            else:
                normalized_color = raw_color

        page_range_combined = metadata.get('pageRange', '') or ''
        if normalized_color == 'Mix':
            bw_label = bw_range_value_db if bw_range_value_db else 'all'
            color_label = color_range_value_db if color_range_value_db else 'all'
            page_range_combined = f"BW: {bw_label} | Color: {color_label}"

        # Store only base_price in pricing_details (not full JSON)
        pricing_details_str = None
        if pricing_details:
            if isinstance(pricing_details, dict):
                # Extract base_price from pricing_breakdown or directly from pricing_details
                breakdown = pricing_details.get('pricing_breakdown', {})
                base_price = breakdown.get('base_price', 0) if isinstance(breakdown, dict) else 0
                if base_price == 0:
                    base_price = pricing_details.get('base_price', 0)
                pricing_details_str = str(base_price) if base_price else None
            elif isinstance(pricing_details, str):
                # If it's already a string, try to parse and extract base_price
                try:
                    parsed = json.loads(pricing_details)
                    if isinstance(parsed, dict):
                        breakdown = parsed.get('pricing_breakdown', {})
                        base_price = breakdown.get('base_price', 0) if isinstance(breakdown, dict) else 0
                        if base_price == 0:
                            base_price = parsed.get('base_price', 0)
                        pricing_details_str = str(base_price) if base_price else None
                    else:
                        pricing_details_str = str(parsed) if parsed else None
                except:
                    # If parsing fails, assume it's already just a base_price value
                    pricing_details_str = pricing_details

        # Get shop address and shop name from metadata or vendor data
        shop_address = metadata.get('shop_address', '')
        shop_name = metadata.get('shop_name', '')
        if (not shop_address or not shop_name) and vendor_email:
            try:
                vendor_data = get_vendor_coordinates_from_email(vendor_email)
                if vendor_data:
                    if not shop_address:
                        shop_address = vendor_data.get('shop_address', '')
                    if not shop_name:
                        shop_name = vendor_data.get('vendor_name', vendor_data.get('shop_name', ''))
            except:
                pass
        
        payload = {
            'vendor_id': vendor_id,
            'vendor_email': vendor_email or '',
            'user_email': user_email,
            'user_id': user_id or metadata.get('user_id', ''),
            'filename': filename,
            'storage_folder': storage_folder or 'users',
            'shop_id': shop_id or metadata.get('shop_id', vendor_id),
            'r2_path': r2_path,
            'service_type': metadata.get('service_type', ''),
            'status': metadata.get('status', 'pending'),
            'job_completed': metadata.get('job_completed', 'NO'),
            'vendor_status': metadata.get('vendor_status', 'not sended'),
            'token': metadata.get('token', ''),
            'job_id': metadata.get('job_id', ''),
            'copies': metadata.get('copies', '1'),
            'color': normalized_color,
            'print_type': metadata.get('print_type', metadata.get('layout_type', 'single_side')),
            'orientation': metadata.get('orientation', ''),
            'pageSize': metadata.get('pageSize', ''),
            'pageRange': page_range_combined,
            'specificPages': metadata.get('specificPages', ''),
            # Mixed (Both) color support: store separate page ranges too
            'bwPageRange': metadata.get('bwPageRange', ''),
            'bwPageRangeValue': metadata.get('bwPageRangeValue', ''),
            'colorPageRange': metadata.get('colorPageRange', ''),
            'colorPageRangeValue': metadata.get('colorPageRangeValue', ''),
            'spiralBinding': metadata.get('spiralBinding', 'No'),
            'lamination': metadata.get('lamination', 'No'),
            'service_name': metadata.get('service_name', ''),
            'feedback': metadata.get('feedback', ''),
            'quality': metadata.get('quality', ''),
            'thickness': metadata.get('thickness', ''),
            'points_applied': metadata.get('points_applied', 'false'),
            'points_used': metadata.get('points_used', '0'),
            'timestamp': metadata.get('timestamp', datetime.datetime.now().isoformat()),
            'completion_time': metadata.get('completion_time', ''),
            'rendered_status': metadata.get('rendered_status', 'NO'),
            'trash': metadata.get('trash', 'NO'),
            'total_price': total_price,
            'platform_profit': platform_profit,
            'price_per_page': price_per_page,
            'final_amount': metadata.get('final_amount'),
            'page_count': page_count,
            'pages': pages_value,
            'num_copies': num_copies,
            'paper_type': metadata.get('paper_type', ''),
            'color_mode': metadata.get('color_mode', metadata.get('color', '')),
            'layout_type': metadata.get('layout_type', ''),
            'pricing_details': pricing_details_str,
            'shop_address': shop_address,
            'shop_name': shop_name
        }

        if payload['final_amount']:
            payload['final_amount'] = to_float(payload['final_amount'])

        print(f"🔍 Attempting to store user print job: {filename} at endpoint: {worker_endpoint}")
        print(f"🔍 Payload keys: {list(payload.keys())}")
        
        resp = requests.post(
            worker_endpoint,
            json=payload,
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        print(f"🔍 Response status: {resp.status_code}")
        print(f"🔍 Response text: {resp.text[:500]}")
        
        if resp.status_code == 200:
            try:
                response_data = resp.json()
                if response_data.get('success'):
                    print(f"✅ Stored user print job in database: {filename}")
                    return True
                else:
                    error_msg = response_data.get('error', 'Unknown error')
                    print(f"⚠️ Worker API returned success=false: {error_msg}")
                    return False
            except json.JSONDecodeError:
                print(f"⚠️ Invalid JSON response from worker: {resp.text[:200]}")
                return False
        else:
            print(f"⚠️ Failed to store user print job in database: {resp.status_code} - {resp.text[:200]}")
            return False

    except requests.exceptions.RequestException as e:
        print(f"❌ Network error storing user print job in database: {e}")
        return False
    except Exception as e:
        print(f"❌ Error storing user print job in database: {e}")
        import traceback
        traceback.print_exc()
        return False

def store_user_notification_in_db(notification_data):
    """Store user notification in D1 database via Worker API"""
    try:
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            print(f"⚠️ Worker API not configured, skipping database storage")
            return False

        # Construct the Worker API endpoint
        if '/add-contact' in api_url:
            worker_endpoint = api_url.replace('/add-contact', '/add-user-notification')
        elif '/add-vendor-register' in api_url:
            worker_endpoint = api_url.replace('/add-vendor-register', '/add-user-notification')
        else:
            worker_endpoint = api_url.rstrip('/') + '/add-user-notification'
        
        platform_profit = _to_decimal(notification_data.get('platform_profit'))
        total_price = _to_decimal(notification_data.get('total_price'))

        payload = {
            'notification_id': notification_data.get('notification_id', ''),
            'user_email': notification_data.get('user_email', ''),
            'filename': notification_data.get('filename', ''),
            'vendor_id': notification_data.get('vendor_id', ''),
            'status': notification_data.get('status', ''),
            'completion_time': notification_data.get('completion_time', ''),
            'created_at': notification_data.get('created_at', get_ist_timestamp()),
            'read': notification_data.get('read', False),
            'type': notification_data.get('type', ''),
            'token': notification_data.get('token', ''),
            'service_type': notification_data.get('service_type', ''),
            'platform_profit': platform_profit,
            'total_price': total_price
        }
        
        resp = requests.post(
            worker_endpoint,
            json=payload,
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        if resp.status_code == 200:
            print(f"✅ Stored user notification in database: {notification_data.get('notification_id')}")
            return True
        else:
            print(f"⚠️ Failed to store user notification in database: {resp.status_code} - {resp.text[:200]}")
            return False
            
    except Exception as e:
        print(f"❌ Error storing user notification in database: {e}")
        return False

def store_vendor_notification_in_db(notification_data):
    """Store vendor notification in D1 database via Worker API"""
    try:
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            print(f"⚠️ Worker API not configured, skipping database storage")
            return False

        # Construct the Worker API endpoint
        if '/add-contact' in api_url:
            worker_endpoint = api_url.replace('/add-contact', '/add-vendor-notification')
        elif '/add-vendor-register' in api_url:
            worker_endpoint = api_url.replace('/add-vendor-register', '/add-vendor-notification')
        else:
            worker_endpoint = api_url.rstrip('/') + '/add-vendor-notification'
        
        platform_profit = _to_decimal(notification_data.get('platform_profit'))
        total_price = _to_decimal(notification_data.get('total_price'))

        payload = {
            'notification_id': notification_data.get('notification_id', ''),
            'vendor_id': notification_data.get('vendor_id', ''),
            'vendor_email': notification_data.get('vendor_email', ''),
            'user_email': notification_data.get('user_email', ''),
            'filename': notification_data.get('filename', ''),
            'service_type': notification_data.get('service_type', ''),
            'platform_profit': platform_profit,
            'total_price': total_price,
            'completion_time': notification_data.get('completion_time', ''),
            'timestamp': notification_data.get('timestamp', get_ist_timestamp()),
            'token': notification_data.get('token', ''),
            'read': notification_data.get('read', False)
        }
        
        resp = requests.post(
            worker_endpoint,
            json=payload,
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        if resp.status_code == 200:
            print(f"✅ Stored vendor notification in database: {notification_data.get('notification_id')}")
            return True
        else:
            print(f"⚠️ Failed to store vendor notification in database: {resp.status_code} - {resp.text[:200]}")
            return False
            
    except Exception as e:
        print(f"❌ Error storing vendor notification in database: {e}")
        return False

def get_token_from_vendor_print_jobs(filename, vendor_id, vendor_email):
    """Get token number from Vendor_print_jobs table in D1"""
    try:
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            return None
        
        # Get from Vendor_print_jobs table
        base_url = api_url.rstrip('/')
        if '/add-contact' in base_url:
            worker_endpoint = base_url.replace('/add-contact', '/get-vendor-print-jobs')
        elif '/add-vendor-register' in base_url:
            worker_endpoint = base_url.replace('/add-vendor-register', '/get-vendor-print-jobs')
        else:
            worker_endpoint = base_url + '/get-vendor-print-jobs'
        
        job_response = requests.post(
            worker_endpoint,
            json={
                'filename': filename,
                'vendor_id': vendor_id,
                'vendor_email': vendor_email
            },
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        if job_response.status_code == 200:
            job_data = job_response.json()
            if job_data.get('success') and job_data.get('data'):
                jobs = job_data.get('data', [])
                # Filter by filename to get exact match
                matching_jobs = [j for j in jobs if j.get('filename') == filename or filename in j.get('filename', '')]
                if matching_jobs:
                    job = matching_jobs[0]  # Get first matching job
                    token = job.get('token')
                    if token:
                        print(f"✅ Retrieved token from Vendor_print_jobs: {token} for {filename}")
                        return str(token).strip()
        
        return None
    except Exception as e:
        print(f"⚠️ Error getting token from Vendor_print_jobs: {e}")
        return None


def get_pricing_from_user_print_jobs(user_email, filename):
    """Get pricing (total_price and platform_profit) from User_print_jobs table in D1"""
    try:
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            return None, None, None
        
        # Get from User_print_jobs table
        base_url = api_url.rstrip('/')
        if '/add-contact' in base_url:
            worker_endpoint = base_url.replace('/add-contact', '/get-user-print-jobs')
        elif '/add-vendor-register' in base_url:
            worker_endpoint = base_url.replace('/add-vendor-register', '/get-user-print-jobs')
        else:
            worker_endpoint = base_url + '/get-user-print-jobs'
        
        job_response = requests.post(
            worker_endpoint,
            json={'user_email': user_email, 'filename': filename},
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        total_price = None
        platform_profit = None
        service_type = None
        
        if job_response.status_code == 200:
            job_data = job_response.json()
            if job_data.get('success') and job_data.get('data'):
                jobs = job_data.get('data', [])
                # Filter by filename to get exact match
                matching_jobs = [j for j in jobs if j.get('filename') == filename or filename in j.get('filename', '')]
                if matching_jobs:
                    job = matching_jobs[0]  # Get first matching job
                    # Get service_type
                    service_type = job.get('service_type', 'Print Job')
                    
                    # Get total_price
                    total_price_val = job.get('total_price') or job.get('final_amount')
                    if total_price_val is not None:
                        try:
                            total_price = float(total_price_val)
                        except (ValueError, TypeError):
                            total_price = 0.0
                    else:
                        total_price = 0.0
                    
                    # Get platform_profit - check both platform_profit and platform_pr
                    platform_profit_val = job.get('platform_profit') or job.get('platform_pr')
                    if platform_profit_val is not None:
                        try:
                            platform_profit = float(platform_profit_val)
                        except (ValueError, TypeError):
                            platform_profit = 0.0
                    else:
                        platform_profit = 0.0
                    
                    if total_price > 0 or platform_profit > 0:
                        print(f"✅ Retrieved pricing from User_print_jobs: total_price={total_price}, platform_profit={platform_profit}, service_type={service_type}")
                        return total_price, platform_profit, service_type
        
        return None, None, None
    except Exception as e:
        print(f"⚠️ Error getting pricing from User_print_jobs: {e}")
        return None, None, None


def send_job_completion_notification(user_email, filename, vendor_id, status, completion_time, token=None):
    """Send job completion notification to user and vendor"""
    try:
        # Get token from Vendor_print_jobs table if not provided
        if not token:
            # Try to get vendor_email from vendor_id
            vendor_email = get_vendor_email_by_id(vendor_id) if vendor_id else None
            
            if vendor_id and vendor_email:
                token = get_token_from_vendor_print_jobs(filename, vendor_id, vendor_email)
            
            # Fallback to filename if token still not found
            if not token:
                token = os.path.splitext(filename)[0]
                print(f"⚠️ Using filename as token fallback for {filename}")
        
        # FIRST: Try to get pricing from User_print_jobs table (most accurate source)
        service_type = "Print Job"
        platform_profit = 0.0
        total_price = 0.0
        
        if user_email and filename:
            d1_total_price, d1_platform_profit, d1_service_type = get_pricing_from_user_print_jobs(user_email, filename)
            if d1_total_price is not None and d1_total_price > 0:
                total_price = d1_total_price
                print(f"📊 Using total_price from User_print_jobs: {total_price}")
            if d1_platform_profit is not None and d1_platform_profit > 0:
                platform_profit = d1_platform_profit
                print(f"📊 Using platform_profit from User_print_jobs: {platform_profit}")
            if d1_service_type:
                service_type = d1_service_type
                print(f"📊 Using service_type from User_print_jobs: {service_type}")
        
        # FALLBACK: Try to get service type and pricing from R2 metadata if D1 didn't have pricing
        if total_price == 0.0 or platform_profit == 0.0:
            try:
                s3 = boto3.client('s3',
                                  aws_access_key_id=settings.R2_ACCESS_KEY,
                                  aws_secret_access_key=settings.R2_SECRET_KEY,
                                  endpoint_url=settings.R2_ENDPOINT,
                                  region_name='auto')
                
                # Try to get service type and platform_profit from vendor_print_jobs
                vendor_key = f'vendor_print_jobs/{vendor_id}/{filename}'
                try:
                    result = s3.get_object(Bucket=settings.R2_BUCKET, Key=vendor_key)
                    job_data = json.loads(result['Body'].read().decode('utf-8'))
                    if not service_type or service_type == "Print Job":
                        service_type = job_data.get('service_type', 'Print Job')
                    
                    # Only use R2 metadata if D1 didn't provide pricing
                    if total_price == 0.0:
                        # Extract platform_profit and total_price from pricing_details if available
                        pricing_details = job_data.get('pricing_details')
                        if pricing_details:
                            try:
                                if isinstance(pricing_details, str):
                                    pricing_obj = json.loads(pricing_details)
                                else:
                                    pricing_obj = pricing_details
                                
                                # Try multiple possible keys for total_price
                                total_price = float(pricing_obj.get('total', pricing_obj.get('total_price', 0.0)))
                                
                                print(f"📊 Extracted total_price from R2 pricing_details: {total_price}")
                            except Exception as e:
                                print(f"⚠️ Error parsing pricing_details: {e}")
                        
                        # If total_price still not found, try metadata
                        if total_price == 0.0:
                            total_price = float(job_data.get('total_price', 0.0))
                            print(f"📊 Using total_price from R2 metadata: {total_price}")
                    
                    # Only use R2 metadata if D1 didn't provide platform_profit
                    if platform_profit == 0.0:
                        pricing_details = job_data.get('pricing_details')
                        if pricing_details:
                            try:
                                if isinstance(pricing_details, str):
                                    pricing_obj = json.loads(pricing_details)
                                else:
                                    pricing_obj = pricing_details
                                
                                # Try multiple possible keys for platform_profit
                                platform_profit = float(pricing_obj.get('platform_profit', 0.0))
                                
                                print(f"📊 Extracted platform_profit from R2 pricing_details: {platform_profit}")
                            except Exception as e:
                                print(f"⚠️ Error parsing pricing_details: {e}")
                        
                        # If platform_profit still not found, try metadata
                        if platform_profit == 0.0:
                            platform_profit = float(job_data.get('platform_profit', 0.0))
                            print(f"📊 Using platform_profit from R2 metadata: {platform_profit}")
                            
                except:
                    pass
            except:
                pass
        
        # Create notification data in format expected by user dashboard
        notification_id = f"{filename}_{int(time.time())}"
        
        # Format service type for better display
        formatted_service_type = service_type.replace('_', ' ').title()
        if formatted_service_type == 'Print Job':
            formatted_service_type = 'Document Printing'
        
        # Extract document name from filename - show PDF name with extension
        # Get the base filename (without path)
        base_filename = os.path.basename(filename)
        file_ext = os.path.splitext(base_filename)[1]  # Get extension (.pdf, .jpg, etc.)
        document_name_without_ext = os.path.splitext(base_filename)[0]
        
        # Extract the actual document name (remove token if present)
        # Format is usually: DocumentName_Token.pdf
        if '_' in document_name_without_ext:
            parts = document_name_without_ext.split('_')
            if len(parts) > 1:
                # Take all parts except the last one (which is usually the token)
                document_name_base = '_'.join(parts[:-1])
            else:
                document_name_base = document_name_without_ext
        else:
            document_name_base = document_name_without_ext
        
        # Create display name with extension (e.g., "Azfar...pdf")
        if len(document_name_base) > 20:
            display_name = document_name_base[:17] + '...' + file_ext
        else:
            display_name = document_name_base + file_ext
        
        # Format token without # symbol
        token_display = str(token) if token else 'Unknown'
        
        notification_data = {
            'notification_id': notification_id,
            'user_email': user_email,
            'filename': filename,
            'vendor_id': vendor_id,
            'status': status,
            'completion_time': completion_time,
            'timestamp': get_ist_timestamp(),
            'created_at': get_ist_timestamp(),
            'read': False,
            'type': 'job_completed',
            'title': 'Print Job Completed',
            'message': f'Your document "{display_name}" is ready for pickup. Token: {token_display}',
            'detailed_message': f'Document: {display_name}\nService Type: {formatted_service_type}\nStatus: Completed ✅\nToken: {token_display}\nCompleted at: {timezone.localtime(timezone.now()).strftime("%B %d, %Y at %I:%M %p")}',
            'token': token,
            'document_name': document_name_base,
            'service_type': formatted_service_type,
            'platform_profit': platform_profit,
            'total_price': total_price
        }
        
        # Store user notification in D1 database only
        try:
            store_user_notification_in_db(notification_data)
            print(f"📧 Stored {status} notification for {user_email} for job {filename} in database")
        except Exception as e:
            print(f"❌ Error storing user notification in database: {e}")
        
        # Store vendor notification in D1 database
        try:
            store_vendor_notification(vendor_id, notification_data, completion_time)
        except Exception as e:
            print(f"❌ Error storing vendor notification: {e}")
        
        # Send FCM push notification - CRITICAL: Must send even if other operations fail
        fcm_sent = False
        try:
            print(f"📱 Attempting to send FCM notification to {user_email} for job {filename}")
            fcm_sent = send_fcm_notification(user_email, notification_data)
            if fcm_sent:
                print(f"✅ FCM notification sent successfully to {user_email}")
            else:
                print(f"⚠️ FCM notification failed to send to {user_email} - no active tokens or send failed")
        except Exception as e:
            print(f"❌ Error sending FCM notification: {e}")
            import traceback
            traceback.print_exc()
        
        # Send WebSocket notification for instant real-time updates
        try:
            channel_layer = get_channel_layer()
            if channel_layer:
                # Create room group name for user (same format as UserConsumer)
                room_group_name = f'user_{user_email.replace("@", "_").replace(".", "_")}'
                
                # Send notification via WebSocket
                async_to_sync(channel_layer.group_send)(
                    room_group_name,
                    {
                        'type': 'job_completion_notification',
                        'filename': filename,
                        'vendor_id': vendor_id,
                        'completion_time': completion_time,
                        'message': notification_data.get('message', ''),
                        'notification_data': notification_data
                    }
                )
                print(f"📡 Sent WebSocket notification to {user_email} for job {filename}")
            else:
                print(f"⚠️ Channel layer not configured, skipping WebSocket notification")
        except Exception as e:
            print(f"⚠️ Error sending WebSocket notification: {e}")
        
        # Here you could also send email, push notification, etc.
        print(f"📧 Sent {status} notification to {user_email} for job {filename}")
        return True
    except Exception as e:
        print(f"❌ Error sending notification: {e}")
        return False

# ─────────────────────────────────────────────────────────────
# FCM (Firebase Cloud Messaging) Functions
# ─────────────────────────────────────────────────────────────

def save_fcm_token(request):
    """Save FCM registration token for a user in D1 database"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            token = data.get('token')
            user_email = data.get('user_email')
            device_type = data.get('device_type', 'web')
            user_agent = request.META.get('HTTP_USER_AGENT', '')
            
            if not token or not user_email:
                return JsonResponse({'success': False, 'error': 'Token and user_email are required'}, status=400)
            
            # Store token in D1 database via Worker API
            try:
                api_url = getattr(settings, 'WORKER_API_URL', '')
                api_key = getattr(settings, 'WORKER_API_KEY', '')
                
                if not api_url or not api_key:
                    return JsonResponse({'success': False, 'error': 'Worker API not configured'}, status=500)
                
                # Construct the Worker API endpoint
                base_url = api_url.rstrip('/')
                for endpoint in ['/add-fcm-token', '/get-fcm-tokens', '/delete-fcm-token']:
                    if base_url.endswith(endpoint):
                        base_url = base_url[:-len(endpoint)]
                worker_endpoint = base_url.rstrip('/') + '/add-fcm-token'
                
                payload = {
                    'token': token,
                    'user_email': user_email,
                    'device_type': device_type,
                    'user_agent': user_agent,
                    'is_active': True
                }
                
                response = requests.post(
                    worker_endpoint,
                    json=payload,
                    headers={
                        'Content-Type': 'application/json',
                        'x-api-key': api_key
                    },
                    timeout=10
                )
                
                if response.status_code == 200:
                    print(f"✅ FCM token saved in D1 database for {user_email}")
                    return JsonResponse({'success': True, 'message': 'FCM token saved successfully'})
                else:
                    error_msg = response.text[:200] if response.text else 'Unknown error'
                    print(f"❌ Error saving FCM token to D1: {response.status_code} - {error_msg}")
                    return JsonResponse({'success': False, 'error': f'Failed to save token: {error_msg}'}, status=response.status_code)
                
            except Exception as e:
                print(f"❌ Error saving FCM token: {e}")
                return JsonResponse({'success': False, 'error': str(e)}, status=500)
                
        except json.JSONDecodeError:
            return JsonResponse({'success': False, 'error': 'Invalid JSON'}, status=400)
        except Exception as e:
            print(f"❌ Error in save_fcm_token: {e}")
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)


def get_user_fcm_tokens(user_email):
    """Get all active FCM tokens for a user from D1 database"""
    try:
        api_url = getattr(settings, 'WORKER_API_URL', '')
        api_key = getattr(settings, 'WORKER_API_KEY', '')
        
        if not api_url or not api_key:
            print("⚠️ Worker API not configured for FCM tokens")
            return []
        
        # Construct the Worker API endpoint
        base_url = api_url.rstrip('/')
        for endpoint in ['/add-fcm-token', '/get-fcm-tokens', '/delete-fcm-token']:
            if base_url.endswith(endpoint):
                base_url = base_url[:-len(endpoint)]
        worker_endpoint = base_url.rstrip('/') + '/get-fcm-tokens'
        
        payload = {
            'user_email': user_email
        }
        
        response = requests.post(
            worker_endpoint,
            json=payload,
            headers={
                'Content-Type': 'application/json',
                'x-api-key': api_key
            },
            timeout=10
        )
        
        if response.status_code == 200:
            data = response.json()
            if data.get('success') and data.get('tokens'):
                # Extract just the token strings
                tokens = [token_info.get('token') for token_info in data.get('tokens', []) if token_info.get('token')]
                return tokens
            return []
        else:
            print(f"⚠️ Error getting FCM tokens from D1: {response.status_code} - {response.text[:200]}")
            return []
            
    except Exception as e:
        print(f"❌ Error getting FCM tokens: {e}")
        return []


def send_fcm_notification(user_email, notification_data):
    """Send FCM push notification to user"""
    try:
        from firebase_admin import messaging
        if messaging is None:
            print("⚠️ Firebase messaging not available")
            return False
        # Get user's FCM tokens
        tokens = get_user_fcm_tokens(user_email)
        
        if not tokens:
            print(f"⚠️ No FCM tokens found for {user_email}")
            return False
        
        # Prepare notification message
        title = notification_data.get('title', 'PrintMax Notification')
        message = notification_data.get('message', 'You have a new notification')
        
        # Construct absolute HTTPS URL for FCM link
        # Use the site domain from settings, or construct from request if available
        site_domain = getattr(settings, 'SITE_DOMAIN', '')
        if not site_domain:
            # Try to get from ALLOWED_HOSTS or construct default
            allowed_hosts = getattr(settings, 'ALLOWED_HOSTS', [])
            if allowed_hosts and allowed_hosts[0] != '*':
                site_domain = allowed_hosts[0]
            else:
                # Default to production domain
                site_domain = 'printmax.onrender.com'
        
        # Ensure HTTPS protocol (FCM requires HTTPS)
        if not site_domain.startswith('http'):
            # Always use HTTPS for FCM (even for localhost, use a tunnel or production domain)
            if 'localhost' in site_domain or '127.0.0.1' in site_domain:
                # For localhost development, use production domain or skip webpush link
                # FCM requires HTTPS, so we'll use production domain as fallback
                fcm_link = 'https://printmax.onrender.com/userdashboard/'
            else:
                fcm_link = f'https://{site_domain}/userdashboard/'
        else:
            if site_domain.startswith('http://'):
                # Convert HTTP to HTTPS for FCM
                fcm_link = site_domain.replace('http://', 'https://') + '/userdashboard/'
            else:
                fcm_link = f'{site_domain}/userdashboard/'
        
        # Ensure absolute URLs for icons (required for stable delivery on some devices)
        if site_domain and not site_domain.startswith('http'):
            # Basic protocol addition if missing (assumes https for prod)
            full_domain = f"https://{site_domain}" if 'localhost' not in site_domain else f"http://{site_domain}"
        else:
             full_domain = site_domain or 'https://printmax.onrender.com'
             
        # Normalize domain to not have trailing slash
        full_domain = full_domain.rstrip('/')
        
        # Use PrintMax day-time logo for mobile notifications (visible in system notification shade)
        icon_url = f"{full_domain}/static/images/printmaxdaylogo.png"
        # Reuse the same logo for badge so branding is consistent
        badge_url = f"{full_domain}/static/images/printmaxdaylogo.png"
        fallback_icon_url = f"{full_domain}/static/images/printmax-color-512.png"

        # Create FCM message
        # Build webpush config conditionally to avoid link issues
        webpush_notification = messaging.WebpushNotification(
            title=title,
            body=message,
            icon=icon_url,
            badge=badge_url,
            require_interaction=True
        )
        
        # Always use HTTPS link for FCM (required by FCM)
        # Ensure fcm_link is HTTPS
        if not fcm_link.startswith('https://') and 'localhost' not in fcm_link and '127.0.0.1' not in fcm_link:
            fcm_link = 'https://printmax.onrender.com/userdashboard/'
        
        # Create FCM message configuration (reusuable parts)
        # Create webpush config with FCM options (always HTTPS)
        webpush_config = messaging.WebpushConfig(
            notification=webpush_notification,
            fcm_options=messaging.WebpushFCMOptions(
                link=fcm_link
            )
        )

        sent_count = 0
        failed_count = 0
        
        # Send to tokens one by one, STOPPING after the first success to avoid duplicates
        # We reverse the list to try the most recent tokens first (assuming they are appended)
        print(f"📡 Attempting to send FCM notification to user {user_email} (has {len(tokens)} tokens)")
        
        for token in reversed(tokens):
            try:
                # NOTE: firebase_admin.messaging.Notification does NOT accept an `icon` argument.
                # Icon is provided via Webpush/Web app options instead (icon_url above and data payload).
                fcm_message = messaging.Message(
                    notification=messaging.Notification(
                        title=title,
                        body=message,
                        image=notification_data.get('icon', None)
                    ),
                    data={
                        'notification_id': notification_data.get('notification_id', ''),
                        'type': notification_data.get('type', 'job_completed'),
                        'filename': notification_data.get('filename', ''),
                        'token': notification_data.get('token', ''),
                        'status': notification_data.get('status', 'completed'),
                        'click_action': fcm_link,
                        'icon': icon_url,  # Include icon in data payload for service worker
                        'badge': badge_url  # Include badge in data payload for service worker
                    },
                    webpush=webpush_config,
                    token=token
                )
                
                # Send message
                response = messaging.send(fcm_message)
                print(f"✅ FCM notification sent to token {token[:10]}...: {response}")
                sent_count += 1
                
                # Allow up to 3 successful deliveries to cover multiple active devices (e.g. Mobile + Desktop)
                # But stop there to avoid excessive duplicates if they have many old active tokens
                if sent_count >= 3:
                     print("🛑 Reached 3 successful notifications, stopping to prevent spam.")
                     break
                
            except Exception as e:
                print(f"⚠️ Failed to send to token {token[:10]}...: {e}")
                failed_count += 1
        
        if sent_count == 0:
            print(f"❌ Failed to send notification to any of {len(tokens)} tokens")
            
        return sent_count > 0
        
    except Exception as e:
        print(f"❌ Error sending FCM notification: {e}")
        import traceback
        traceback.print_exc()
        return False


def store_vendor_notification_direct(vendor_email, filename, vendor_id, user_email, completion_time, token=None):
    """Store vendor notification directly using vendor email from session"""
    try:
        # Get token from Vendor_print_jobs table if not provided
        if not token:
            token = get_token_from_vendor_print_jobs(filename, vendor_id, vendor_email)
            if not token:
                token = os.path.splitext(filename)[0]
                print(f"⚠️ Using filename as token fallback for {filename}")
        
        # Parse completion time to get date
        completion_date = datetime.datetime.fromisoformat(completion_time.replace('Z', '+00:00')).date()
        
        # Get the appropriate 2-day date folder
        date_folder = get_vendor_notification_date_folder(completion_date)
        
        # FIRST: Try to get pricing from User_print_jobs table (most accurate source)
        service_type = "Print Job"
        platform_profit = 0.0
        total_price = 0.0
        
        if user_email and filename:
            d1_total_price, d1_platform_profit, d1_service_type = get_pricing_from_user_print_jobs(user_email, filename)
            if d1_total_price is not None and d1_total_price > 0:
                total_price = d1_total_price
                print(f"📊 Using total_price from User_print_jobs for vendor notification: {total_price}")
            if d1_platform_profit is not None and d1_platform_profit > 0:
                platform_profit = d1_platform_profit
                print(f"📊 Using platform_profit from User_print_jobs for vendor notification: {platform_profit}")
            if d1_service_type:
                service_type = d1_service_type
                print(f"📊 Using service_type from User_print_jobs for vendor notification: {service_type}")
        
        # FALLBACK: Try to get service type and pricing from R2 metadata if D1 didn't have pricing
        if total_price == 0.0 or platform_profit == 0.0:
            try:
                s3 = boto3.client('s3',
                                  aws_access_key_id=settings.R2_ACCESS_KEY,
                                  aws_secret_access_key=settings.R2_SECRET_KEY,
                                  endpoint_url=settings.R2_ENDPOINT,
                                  region_name='auto')
                
                # Try to get service type and platform_profit from vendor_print_jobs metadata
                vendor_key = f'vendor_print_jobs/{vendor_id}/{filename}'
                try:
                    # Get object metadata instead of trying to read JSON from PDF
                    result = s3.head_object(Bucket=settings.R2_BUCKET, Key=vendor_key)
                    metadata = result.get('Metadata', {})
                    
                    if not service_type or service_type == "Print Job":
                        service_type = metadata.get('service_type', 'Print Job')
                    
                    # Only use R2 metadata if D1 didn't provide pricing
                    if total_price == 0.0:
                        # Extract platform_profit and total_price from pricing_details if available
                        pricing_details_str = metadata.get('pricing_details')
                        if pricing_details_str:
                            try:
                                pricing_obj = json.loads(pricing_details_str)
                                
                                # Try multiple possible keys for total_price
                                total_price = float(pricing_obj.get('total', pricing_obj.get('total_price', 0.0)))
                                
                                print(f"📊 Extracted total_price from R2 pricing_details: {total_price}")
                            except Exception as e:
                                print(f"⚠️ Error parsing pricing_details: {e}")
                        
                        # If total_price still not found, try metadata
                        if total_price == 0.0:
                            total_price = float(metadata.get('total_price', 0.0))
                            print(f"📊 Using total_price from R2 metadata: {total_price}")
                    
                    # Only use R2 metadata if D1 didn't provide platform_profit
                    if platform_profit == 0.0:
                        pricing_details_str = metadata.get('pricing_details')
                        if pricing_details_str:
                            try:
                                pricing_obj = json.loads(pricing_details_str)
                                
                                # Try multiple possible keys for platform_profit
                                platform_profit = float(pricing_obj.get('platform_profit', 0.0))
                                
                                print(f"📊 Extracted platform_profit from R2 pricing_details: {platform_profit}")
                            except Exception as e:
                                print(f"⚠️ Error parsing pricing_details: {e}")
                        
                        # If platform_profit still not found, try metadata
                        if platform_profit == 0.0:
                            platform_profit = float(metadata.get('platform_profit', 0.0))
                            print(f"📊 Using platform_profit from R2 metadata: {platform_profit}")
                            
                except Exception as e:
                    print(f"⚠️ Error reading metadata for {vendor_key}: {e}")
                    pass
            except Exception as e:
                print(f"⚠️ Error accessing S3: {e}")
                pass
        
        # Create notification ID
        notification_id = f"{filename}_{int(time.time())}"
        
        # Format service type for better display
        formatted_service_type = service_type.replace('_', ' ').title()
        if formatted_service_type == 'Print Job':
            formatted_service_type = 'Document Printing'
        
        # Extract document name from filename (remove extension and token)
        document_name = os.path.splitext(filename)[0]
        if '_' in document_name:
            # Remove token part if present
            parts = document_name.split('_')
            if len(parts) > 1:
                document_name = '_'.join(parts[:-1])
        
        # Create vendor notification data
        vendor_notification_data = {
            'notification_id': notification_id,
            'vendor_id': vendor_id,
            'vendor_email': vendor_email,
            'user_email': user_email,
            'filename': filename,
            'service_type': formatted_service_type,
            'platform_profit': platform_profit,
            'total_price': total_price,
            'completion_time': completion_time,
            'timestamp': get_ist_timestamp(),
            'token': token,
            'document_name': document_name
        }
        
        # Store vendor notification in D1 database only
        try:
            store_vendor_notification_in_db(vendor_notification_data)
            print(f"📧 Stored vendor notification for {vendor_email} in database")
            print(f"📊 Notification data: {json.dumps(vendor_notification_data, indent=2)}")
            return True
        except Exception as e:
            print(f"❌ Error storing vendor notification in database: {e}")
            return False
        
    except Exception as e:
        print(f"❌ Error storing vendor notification: {e}")
        return False

def store_vendor_notification(vendor_id, notification_data, completion_time):
    """Store vendor notification in D1 database only"""
    try:
        # Get vendor email from vendor_id
        vendor_email = get_vendor_email_by_id(vendor_id)
        if not vendor_email:
            print(f"❌ Could not find vendor email for vendor_id: {vendor_id}")
            return False
        
        # Parse completion time to get date
        completion_date = datetime.datetime.fromisoformat(completion_time.replace('Z', '+00:00')).date()
        
        # Get the appropriate 2-day date folder
        date_folder = get_vendor_notification_date_folder(completion_date)
        
        # Create vendor notification data (optimized for vendor dashboard)
        vendor_notification_data = {
            'notification_id': notification_data['notification_id'],
            'vendor_id': vendor_id,
            'vendor_email': vendor_email,
            'user_email': notification_data['user_email'],
            'filename': notification_data['filename'],
            'service_type': notification_data['service_type'],
            'platform_profit': notification_data['platform_profit'],
            'total_price': notification_data['total_price'],
            'completion_time': completion_time,
            'timestamp': notification_data['timestamp'],
            'token': notification_data['token'],
            'document_name': notification_data.get('document_name', '')
        }
        
        # Store vendor notification in D1 database only
        try:
            store_vendor_notification_in_db(vendor_notification_data)
            print(f"📧 Stored vendor notification for {vendor_email} in database")
        except Exception as e:
            print(f"❌ Error storing vendor notification in database: {e}")
        
        return True
        
    except Exception as e:
        print(f"❌ Error storing vendor notification: {e}")
        return False
def get_vendor_notification_date_folder(completion_date):
    """Get the appropriate 2-day date folder for vendor notifications"""
    try:
        # Convert to date object if it's a string
        if isinstance(completion_date, str):
            completion_date = datetime.datetime.fromisoformat(completion_date.replace('Z', '+00:00')).date()
        
        # Calculate the 2-day folder range
        # Each folder covers 2 consecutive days starting from odd days
        # Pattern: 1-2, 3-4, 5-6, 7-8, 9-10, etc.
        
        # Get the day of year
        day_of_year = completion_date.timetuple().tm_yday
        
        # Calculate which 2-day period this falls into
        # Each period starts on odd days (1-2, 3-4, 5-6, etc.)
        period_start_day = ((day_of_year - 1) // 2) * 2 + 1
        
        # Handle year boundaries
        year = completion_date.year
        if period_start_day > 365:
            # If we're near the end of the year, adjust
            if completion_date.month == 12 and completion_date.day >= 30:
                # Last few days of December - create a special folder
                period_start_day = 365
            else:
                # Move to next year
                year += 1
                period_start_day = 1
        
        # Create start date
        start_date = datetime.date(year, 1, 1) + datetime.timedelta(days=period_start_day - 1)
        
        # Create end date (2 days later)
        end_date = start_date + datetime.timedelta(days=1)
        
        # Handle year boundaries for end date
        if end_date.year != start_date.year:
            end_date = datetime.date(start_date.year, 12, 31)
        
        # Format as folder name: YYYY-MM-DD_to_YYYY-MM-DD
        folder_name = f"{start_date.strftime('%Y-%m-%d')}_to_{end_date.strftime('%Y-%m-%d')}"
        
        print(f"📅 Created 2-day folder for {completion_date}: {folder_name}")
        return folder_name
        
    except Exception as e:
        print(f"❌ Error calculating date folder: {e}")
        # Fallback to current date
        today = datetime.date.today()
        return f"{today.strftime('%Y-%m-%d')}_to_{(today + datetime.timedelta(days=1)).strftime('%Y-%m-%d')}"

def get_vendor_email_by_id(vendor_id):
    """Get vendor email by vendor ID"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        # List all vendor registration details
        reg_prefix = "vendor_register_details/"
        reg_objects = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=reg_prefix)
        
        if 'Contents' in reg_objects:
            for obj in reg_objects['Contents']:
                key = obj["Key"]
                if key.endswith('/registration.json'):
                    try:
                        result = s3.get_object(Bucket=settings.R2_BUCKET, Key=key)
                        vendor_data = json.loads(result['Body'].read().decode('utf-8'))
                        
                        if vendor_data.get('vendor_id') == vendor_id:
                            return vendor_data.get('vendor_email')
                    except Exception as e:
                        print(f"Error reading vendor data from {key}: {e}")
                        continue
        
        return None
        
    except Exception as e:
        print(f"Error getting vendor email by ID: {e}")
        return None

@csrf_exempt
def accept_print_job(request):
    """Accept a print job by updating vendor_status to 'accepted'"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            filename = data.get('filename')
            
            if not filename:
                return JsonResponse({'success': False, 'error': 'Filename required'})
            
            # Get vendor email from session
            vendor_email = request.session.get('vendor_email')
            if not vendor_email:
                return JsonResponse({'success': False, 'error': 'Vendor not authenticated'})
            
            # Get vendor details
            vendor_details = get_vendor_details_by_email(vendor_email)
            if not vendor_details:
                return JsonResponse({'success': False, 'error': 'Vendor details not found'})
            
            vendor_id = vendor_details.get('vendor_id', 'vendor1')
            
            # Update vendor_status in R2 metadata
            success = update_vendor_status_in_r2(filename, 'accepted', vendor_id)
            
            if success:
                # Ensure the accepted job exists under vendor_print_jobs/<vendor_id>/ so vendor client can poll it
                try:
                    ensured = ensure_job_in_vendor_folder(filename, vendor_id)
                    if not ensured:
                        print(f"⚠️ Could not ensure vendor folder copy for {filename}")
                except Exception as _e:
                    print(f"⚠️ ensure_job_in_vendor_folder error: {_e}")
                # Reset failed status for retry
                update_job_failed_status_in_r2(filename, vendor_id, 'NO')
                
                # Clear vendor cache to ensure fresh data on next load
                clear_vendor_cache(vendor_email, vendor_id)
                return JsonResponse({
                    'success': True,
                    'message': f'Job "{filename}" accepted successfully',
                    'vendor_status': 'accepted'
                })
            else:
                return JsonResponse({
                    'success': False,
                    'error': 'Failed to update vendor status - file not found'
                })
                
        except json.JSONDecodeError:
            return JsonResponse({'success': False, 'error': 'Invalid JSON'}, status=400)
        except Exception as e:
            print(f"Error accepting print job: {str(e)}")
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)

def update_vendor_status_in_r2(filename, status, vendor_id):
    """Update vendor_status in R2 storage metadata - searches in correct folder structure"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        # Define all possible folder structures to search
        possible_paths = [
            f'vendor_print_jobs/{vendor_id}/{filename}',
            f'vendor_manual_print_jobs/{vendor_id}/{filename}',
            f'vendor_register_details/{vendor_id}/firozshop/{filename}',
            f'vendor_register_details/{sanitize_email(vendor_id)}/firozshop/{filename}',
            f'users/*/{filename}',  # This will be handled separately
        ]
        
        print(f"🔍 Searching for {filename} in vendor {vendor_id}...")
        
        # First, try the most common paths
        for path in possible_paths[:4]:  # Skip the users path for now
            try:
                print(f"   🔍 Trying path: {path}")
                head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=path)
                current_metadata = head_response.get('Metadata', {})
                
                # Update vendor_status and appropriate timestamps
                current_metadata['vendor_status'] = status
                current_metadata['updated_at'] = datetime.datetime.now().isoformat()
                
                if status == 'accepted':
                    current_metadata['accepted_time'] = datetime.datetime.now().isoformat()
                    current_metadata['print_round'] = '1'  # Set print_round to 1 when accepted
                elif status == 'cancelled':
                    current_metadata['cancelled_time'] = datetime.datetime.now().isoformat()
                    current_metadata['job_failed'] = 'YES'  # Mark as failed for cancelled jobs
                
                # Avoid storing metadata for vendor/user print jobs; rely on DB tables
                if path.startswith('vendor_print_jobs/') or path.startswith('users/'):
                    print(f"ℹ️ Skipping R2 metadata write for {path}; DB holds metadata.")
                    print(f"✅ Treated vendor_status as '{status}' for {filename} at {path} (DB-driven).")
                    return True
                else:
                    # Copy object with updated metadata for non print-job assets
                    copy_source = {'Bucket': settings.R2_BUCKET, 'Key': path}
                    s3.copy_object(
                        CopySource=copy_source,
                        Bucket=settings.R2_BUCKET,
                        Key=path,
                        Metadata=current_metadata,
                        MetadataDirective='REPLACE'
                    )
                    
                    print(f"✅ Updated vendor_status to '{status}' for {filename} at {path}")
                    return True
                
            except Exception as e:
                print(f"   ⚠️ Not found at {path}: {str(e)}")
                continue
        
        # If not found in vendor folders, search in users folder
        print(f"   🔍 Searching in users folder...")
        try:
            # List all users and search for the file
            users_prefix = 'users/'
            users_response = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=users_prefix)
            
            for obj in users_response.get('Contents', []):
                key = obj['Key']
                if key.endswith(f'/{filename}'):
                    print(f"   🎯 Found file at: {key}")
                    
                    # Get current metadata
                    head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                    current_metadata = head_response.get('Metadata', {})
                    
                    # Update vendor_status and appropriate timestamps
                    current_metadata['vendor_status'] = status
                    current_metadata['updated_at'] = datetime.datetime.now().isoformat()
                    
                    if status == 'accepted':
                        current_metadata['accepted_time'] = datetime.datetime.now().isoformat()
                        current_metadata['print_round'] = '1'  # Set print_round to 1 when accepted
                    elif status == 'cancelled':
                        current_metadata['cancelled_time'] = datetime.datetime.now().isoformat()
                        current_metadata['job_failed'] = 'YES'  # Mark as failed for cancelled jobs
                    
                    # Copy object with updated metadata
                    copy_source = {'Bucket': settings.R2_BUCKET, 'Key': key}
                    s3.copy_object(
                        CopySource=copy_source,
                        Bucket=settings.R2_BUCKET,
                        Key=key,
                        Metadata=current_metadata,
                        MetadataDirective='REPLACE'
                    )
                    
                    print(f"✅ Updated vendor_status to '{status}' for {filename} at {key}")
                    return True
                    
        except Exception as e:
            print(f"   ⚠️ Error searching users folder: {str(e)}")
        
        print(f"❌ File {filename} not found in any expected location")
        return False
                
    except Exception as e:
        print(f"❌ Error updating vendor status: {str(e)}")
        return False

def ensure_job_in_vendor_folder(filename: str, vendor_id: str) -> bool:
    """Copy the job object to vendor_print_jobs/<vendor_id>/ if it is not already there.
    Preserve existing metadata and force vendor_status='accepted', job_completed='NO', rendered_status default.
    """
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        vendor_key = f'vendor_print_jobs/{vendor_id}/{filename}'
        # If already exists in vendor folder, nothing to do
        try:
            s3.head_object(Bucket=settings.R2_BUCKET, Key=vendor_key)
            return True
        except Exception:
            pass

        # Find source key from likely locations
        search_prefixes = [
            f'vendor_manual_print_jobs/{vendor_id}/',
            f'users/',
        ]
        source_key = None
        src_metadata = {}
        # Try manual jobs first
        try:
            mk = f'vendor_manual_print_jobs/{vendor_id}/{filename}'
            head = s3.head_object(Bucket=settings.R2_BUCKET, Key=mk)
            source_key = mk
            src_metadata = head.get('Metadata', {})
        except Exception:
            # Search users tree
            resp = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix='users/')
            for obj in resp.get('Contents', []):
                if obj['Key'].endswith(f'/{filename}'):
                    source_key = obj['Key']
                    try:
                        head = s3.head_object(Bucket=settings.R2_BUCKET, Key=source_key)
                        src_metadata = head.get('Metadata', {})
                    except Exception:
                        src_metadata = {}
                    break

        if not source_key:
            print(f"❌ ensure_job_in_vendor_folder: source not found for {filename}")
            return False

        # Prepare metadata for vendor folder
        md = dict(src_metadata or {})
        md['vendor_status'] = 'accepted'
        md['job_completed'] = md.get('job_completed', 'NO') or 'NO'
        if md.get('rendered_status') is None:
            md['rendered_status'] = 'NO'
        md['vendor'] = vendor_id
        md['updated_at'] = datetime.datetime.now().isoformat()

        # Perform server-side copy to vendor folder with metadata replace
        s3.copy_object(
            CopySource={'Bucket': settings.R2_BUCKET, 'Key': source_key},
            Bucket=settings.R2_BUCKET,
            Key=vendor_key,
            Metadata=md,
            MetadataDirective='REPLACE'
        )
        print(f"✅ ensure_job_in_vendor_folder: copied {source_key} -> {vendor_key}")
        return True
    except Exception as e:
        print(f"❌ ensure_job_in_vendor_folder error: {e}")
        return False

def update_job_completed_status_in_r2(filename, status, vendor_id):
    """Update job_completed_status in R2 storage metadata"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        # Define all possible folder structures to search
        possible_paths = [
            f'vendor_print_jobs/{vendor_id}/{filename}',
            f'vendor_manual_print_jobs/{vendor_id}/{filename}',
            f'vendor_register_details/{vendor_id}/firozshop/{filename}',
            f'vendor_register_details/{sanitize_email(vendor_id)}/firozshop/{filename}',
        ]
        
        print(f"🔍 Searching for {filename} to mark as completed...")
        
        # First, try the most common paths
        for path in possible_paths:
            try:
                print(f"   🔍 Trying path: {path}")
                head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=path)
                current_metadata = head_response.get('Metadata', {})
                
                # Update job_completed (not job_completed_status)
                current_metadata['job_completed'] = status
                current_metadata['completion_time'] = datetime.datetime.now().isoformat()
                
                # Copy object with updated metadata
                copy_source = {'Bucket': settings.R2_BUCKET, 'Key': path}
                s3.copy_object(
                    CopySource=copy_source,
                    Bucket=settings.R2_BUCKET,
                    Key=path,
                    Metadata=current_metadata,
                    MetadataDirective='REPLACE'
                )
                
                print(f"✅ Updated job_completed to '{status}' for {filename} at {path}")
                return True
                
            except Exception as e:
                print(f"   ⚠️ Not found at {path}: {str(e)}")
                continue
        
        # If not found in vendor folders, search in users folder
        print(f"   🔍 Searching in users folder...")
        try:
            # List all users and search for the file
            users_prefix = 'users/'
            users_response = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=users_prefix)
            
            for obj in users_response.get('Contents', []):
                key = obj['Key']
                if key.endswith(f'/{filename}'):
                    print(f"   🎯 Found file at: {key}")
                    
                    # Get current metadata
                    head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                    current_metadata = head_response.get('Metadata', {})
                    
                    # Update job_completed (not job_completed_status)
                    current_metadata['job_completed'] = status
                    current_metadata['completion_time'] = datetime.datetime.now().isoformat()
                    
                    # Copy object with updated metadata
                    copy_source = {'Bucket': settings.R2_BUCKET, 'Key': key}
                    s3.copy_object(
                        CopySource=copy_source,
                        Bucket=settings.R2_BUCKET,
                        Key=key,
                        Metadata=current_metadata,
                        MetadataDirective='REPLACE'
                    )
                    
                    print(f"✅ Updated job_completed to '{status}' for {filename} at {key}")
                    return True
                    
        except Exception as e:
            print(f"   ⚠️ Error searching users folder: {str(e)}")
        
        print(f"❌ File {filename} not found in any expected location")
        return False
                
    except Exception as e:
        print(f"❌ Error updating job completion status: {str(e)}")
        return False

@csrf_exempt
def get_user_completed_jobs(request):
    """Get completed print jobs for a specific user from users folder in R2"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            user_email = data.get('user_email')
            
            if not user_email:
                return JsonResponse({'success': False, 'error': 'User email required'})
            
            s3 = boto3.client('s3',
                              aws_access_key_id=settings.R2_ACCESS_KEY,
                              aws_secret_access_key=settings.R2_SECRET_KEY,
                              endpoint_url=settings.R2_ENDPOINT,
                              region_name='auto')
            
            completed_jobs = []
            
            # Search in users folder for the specific user
            user_folder = f'users/{user_email}/'
            
            try:
                response = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=user_folder)
                
                if 'Contents' in response:
                    for obj in response['Contents']:
                        key = obj['Key']
                        filename = key.split('/')[-1]
                        
                        # Skip folder itself and non-document files
                        if not filename or filename.endswith('.json') or filename.endswith('/'):
                            continue
                        
                        try:
                            # Get object metadata
                            head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                            metadata = head_response.get('Metadata', {})
                            
                            # Check if job is completed
                            job_completed = metadata.get('job_completed', 'NO').upper()
                            if job_completed != 'YES':
                                continue
                            
                            # Generate download URL
                            download_url = s3.generate_presigned_url(
                                ClientMethod='get_object',
                                Params={
                                    'Bucket': settings.R2_BUCKET,
                                    'Key': key
                                },
                                ExpiresIn=3600
                            )
                            
                            # Build job info
                            job_info = {
                                'filename': filename,
                                'download_url': download_url,
                                'preview_url': download_url,  # For compatibility
                                'file_size': obj.get('Size', 0),
                                'file_extension': filename.split('.')[-1].lower() if '.' in filename else '',
                                'service_type': metadata.get('service_type', 'Unknown'),
                                'pages': metadata.get('pages', 'Unknown'),
                                'copies': metadata.get('copies', '1'),
                                'color': metadata.get('color', 'Color'),
                                'pageSize': metadata.get('pageSize', 'A4'),
                                'user_email': user_email,
                                'uploaded_at': metadata.get('uploaded_at', 'Unknown'),
                                'completion_time': metadata.get('completion_time', ''),
                                'token': metadata.get('token', ''),
                                'price': metadata.get('price', '0'),
                                'vendor': metadata.get('vendor', 'Unknown'),
                                'vendor_lat': metadata.get('vendor_lat', ''),
                                'vendor_lng': metadata.get('vendor_lng', ''),
                                'job_completed': 'YES',
                                'status': 'completed'
                            }

                            # Prefer full pricing_details from metadata if available
                            pricing_details_full = None
                            pd_raw = metadata.get('pricing_details')
                            if pd_raw:
                                try:
                                    pricing_details_full = json.loads(pd_raw) if isinstance(pd_raw, str) else pd_raw
                                except Exception:
                                    pricing_details_full = None
                            if isinstance(pricing_details_full, dict):
                                job_info['pricing_details'] = pricing_details_full
                            else:
                                # Fallback minimal structure
                                job_info['pricing_details'] = {
                                    'total_price': metadata.get('price', '0'),
                                    'pricing_breakdown': {
                                        'page_count': metadata.get('pages', '1'),
                                        'num_copies': metadata.get('copies', '1'),
                                        'pages': metadata.get('pages', '1'),
                                        'copies': metadata.get('copies', '1'),
                                        'total_pages': metadata.get('pages', '1')
                                    }
                                }
                            
                            completed_jobs.append(job_info)
                            
                        except Exception as e:
                            print(f"Error processing job {filename}: {str(e)}")
                            continue
                            
            except Exception as e:
                print(f"Error listing objects in {user_folder}: {str(e)}")
                return JsonResponse({'success': False, 'error': f'Error accessing user folder: {str(e)}'})
            
            # Sort by completion time (newest first)
            completed_jobs.sort(key=lambda x: x.get('completion_time', ''), reverse=True)
            
            return JsonResponse({
                'success': True,
                'jobs': completed_jobs,
                'total_count': len(completed_jobs)
            })
            
        except Exception as e:
            print(f"Error in get_user_completed_jobs: {str(e)}")
            return JsonResponse({'success': False, 'error': str(e)})
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'})

def get_user_email_from_file_metadata(filename, vendor_id):
    """Extract user email from file metadata"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        # Search in users folder for the file
        users_prefix = 'users/'
        users_response = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=users_prefix)
        
        for obj in users_response.get('Contents', []):
            key = obj['Key']
            if key.endswith(f'/{filename}'):
                # Get metadata to extract user email
                head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                metadata = head_response.get('Metadata', {})
                return metadata.get('user_email', '')
        
        return None
    except Exception as e:
        print(f"Error getting user email from metadata: {str(e)}")
        return None

@csrf_exempt
def hide_completed_job(request):
    """Clear a completed job by updating vendor status to 'clear'"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            filename = data.get('filename')
            token = data.get('token', '')
            action = data.get('action', 'hide')  # 'hide' or 'clear'
            
            if not filename:
                return JsonResponse({'success': False, 'error': 'Filename required'})
            
            # Get vendor email from session
            vendor_email = request.session.get('vendor_email')
            if not vendor_email:
                return JsonResponse({'success': False, 'error': 'Vendor not authenticated'})
            
            # Get vendor details
            vendor_details = get_vendor_details_by_email(vendor_email)
            if not vendor_details:
                return JsonResponse({'success': False, 'error': 'Vendor details not found'})
            
            vendor_id = vendor_details.get('vendor_id', 'vendor1')
            
            if action == 'clear':
                # Update vendor status to 'clear' instead of hiding
                success = update_job_vendor_status_in_r2(filename, 'clear', vendor_id)
                message = f'Job "{filename}" cleared from completed section'
                
                # CRITICAL: Also update database to set rendered_status='YES' and job_completed='YES' so job disappears from completed section
                try:
                    # Update D1 database via Worker API
                    api_url = getattr(settings, 'WORKER_API_URL', '')
                    api_key = getattr(settings, 'WORKER_API_KEY', '')
                    
                    if api_url and api_key:
                        # Get user_email from the job metadata
                        user_email = None
                        try:
                            s3 = boto3.client('s3',
                                              aws_access_key_id=settings.R2_ACCESS_KEY,
                                              aws_secret_access_key=settings.R2_SECRET_KEY,
                                              endpoint_url=settings.R2_ENDPOINT,
                                              region_name='auto')
                            
                            # Try to get user_email from R2 metadata
                            possible_paths = [
                                f'vendor_print_jobs/{vendor_id}/{filename}',
                                f'vendor_manual_print_jobs/{vendor_id}/{filename}',
                            ]
                            
                            for path in possible_paths:
                                try:
                                    head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=path)
                                    metadata = head_response.get('Metadata', {})
                                    user_email = metadata.get('user_email', '')
                                    if user_email:
                                        break
                                except:
                                    continue
                        except Exception as e:
                            print(f"⚠️ Error getting user_email from R2: {e}")
                        
                        if user_email:
                            # Construct worker endpoint
                            if '/add-contact' in api_url:
                                worker_endpoint = api_url.replace('/add-contact', '/update-job-completed')
                            elif '/add-vendor-register' in api_url:
                                worker_endpoint = api_url.replace('/add-vendor-register', '/update-job-completed')
                            else:
                                worker_endpoint = api_url.rstrip('/') + '/update-job-completed'
                            
                            # Update both Vendor_print_jobs and User_print_jobs to set rendered_status='YES' and job_completed='YES'
                            worker_payload = {
                                'filename': filename,
                                'vendor_id': vendor_id,
                                'vendor_email': vendor_email,
                                'user_email': user_email,
                                'job_completed': 'YES',
                                'rendered_status': 'YES',
                                'completion_time': datetime.datetime.now().isoformat()
                            }
                            
                            try:
                                resp = requests.post(
                                    worker_endpoint,
                                    json=worker_payload,
                                    headers={
                                        'x-api-key': api_key,
                                        'Content-Type': 'application/json'
                                    },
                                    timeout=10
                                )
                                if resp.status_code == 200:
                                    resp_data = resp.json()
                                    if resp_data.get('success'):
                                        print(f"✅ Updated database: rendered_status='YES' and job_completed='YES' for {filename}")
                                    else:
                                        print(f"⚠️ Database update returned success=false: {resp_data.get('error')}")
                                else:
                                    print(f"⚠️ Failed to update database: {resp.status_code} - {resp.text[:200]}")
                            except Exception as e:
                                print(f"⚠️ Error updating database via Worker API: {e}")
                        else:
                            print(f"⚠️ Could not find user_email for {filename}, skipping database update")
                    else:
                        print("⚠️ Worker API not configured - skipping database update")
                except Exception as db_error:
                    print(f"⚠️ Error updating database for clear completed job: {db_error}")
                    # Don't fail the request if database update fails
            else:
                # Legacy behavior - update is_hidden status
                success = update_job_hidden_status_in_r2(filename, 'true', vendor_id)
                message = f'Job "{filename}" hidden from completed section'
            
            if success:
                # Clear vendor cache to ensure fresh data on next load
                clear_vendor_cache(vendor_email, vendor_id)
                return JsonResponse({
                    'success': True,
                    'message': message
                })
            else:
                return JsonResponse({
                    'success': False,
                    'error': 'Failed to update job - file not found'
                })
                
        except json.JSONDecodeError:
            return JsonResponse({'success': False, 'error': 'Invalid JSON'}, status=400)
        except Exception as e:
            print(f"Error updating completed job: {str(e)}")
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)

def update_job_vendor_status_in_r2(filename, vendor_status, vendor_id):
    """Update vendor_status in R2 storage metadata"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        # Search for the file in vendor folders
        possible_paths = [
            f'vendor_print_jobs/{vendor_id}/{filename}',
            f'vendor_manual_print_jobs/{vendor_id}/{filename}',
            f'vendor_register_details/{vendor_id}/firozshop/{filename}',
        ]
        
        for path in possible_paths:
            try:
                # Get current metadata
                head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=path)
                current_metadata = head_response.get('Metadata', {})
                
                # Update vendor_status
                current_metadata['vendor_status'] = vendor_status
                
                # Copy object with updated metadata
                s3.copy_object(
                    Bucket=settings.R2_BUCKET,
                    CopySource={'Bucket': settings.R2_BUCKET, 'Key': path},
                    Key=path,
                    Metadata=current_metadata,
                    MetadataDirective='REPLACE'
                )
                
                print(f"✅ Updated vendor_status to '{vendor_status}' for {filename}")
                return True
                
            except Exception as e:
                continue
        
        print(f"❌ File {filename} not found in vendor folders")
        return False
        
    except Exception as e:
        print(f"❌ Error updating vendor_status for {filename}: {str(e)}")
        return False

def update_job_hidden_status_in_r2(filename, hidden_status, vendor_id):
    """Update is_hidden status in R2 storage metadata"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        # Define all possible folder structures to search
        possible_paths = [
            f'vendor_print_jobs/{vendor_id}/{filename}',
            f'vendor_manual_print_jobs/{vendor_id}/{filename}',
            f'vendor_register_details/{vendor_id}/firozshop/{filename}',
            f'vendor_register_details/{sanitize_email(vendor_id)}/firozshop/{filename}',
        ]
        
        print(f"🔍 Searching for {filename} to hide...")
        
        # First, try the most common paths
        for path in possible_paths:
            try:
                print(f"   🔍 Trying path: {path}")
                head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=path)
                current_metadata = head_response.get('Metadata', {})
                
                # Update is_hidden status
                current_metadata['is_hidden'] = hidden_status
                
                # Copy object with updated metadata
                copy_source = {'Bucket': settings.R2_BUCKET, 'Key': path}
                s3.copy_object(
                    CopySource=copy_source,
                    Bucket=settings.R2_BUCKET,
                    Key=path,
                    Metadata=current_metadata,
                    MetadataDirective='REPLACE'
                )
                
                print(f"✅ Updated is_hidden to '{hidden_status}' for {filename} at {path}")
                return True
                
            except Exception as e:
                print(f"   ⚠️ Not found at {path}: {str(e)}")
                continue
        
        # If not found in vendor folders, search in users folder
        print(f"   🔍 Searching in users folder...")
        try:
            # List all users and search for the file
            users_prefix = 'users/'
            users_response = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=users_prefix)
            
            for obj in users_response.get('Contents', []):
                key = obj['Key']
                if key.endswith(f'/{filename}'):
                    print(f"   🎯 Found file at: {key}")
                    
                    # Get current metadata
                    head_response = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                    current_metadata = head_response.get('Metadata', {})
                    
                    # Update is_hidden status
                    current_metadata['is_hidden'] = hidden_status
                    
                    # Copy object with updated metadata
                    copy_source = {'Bucket': settings.R2_BUCKET, 'Key': key}
                    s3.copy_object(
                        CopySource=copy_source,
                        Bucket=settings.R2_BUCKET,
                        Key=key,
                        Metadata=current_metadata,
                        MetadataDirective='REPLACE'
                    )
                    
                    print(f"✅ Updated is_hidden to '{hidden_status}' for {filename} at {key}")
                    return True
                    
        except Exception as e:
            print(f"   ⚠️ Error searching users folder: {str(e)}")
        
        print(f"❌ File {filename} not found in any expected location")
        return False
                
    except Exception as e:
        print(f"❌ Error updating hidden status: {str(e)}")
        return False
def debug_file_locations(request):
    """Debug endpoint to find where files are stored"""
    if request.method == 'GET':
        try:
            filename = request.GET.get('filename')
            vendor_email = request.session.get('vendor_email')
            
            if not filename:
                return JsonResponse({'success': False, 'error': 'Filename parameter required'})
            
            if not vendor_email:
                return JsonResponse({'success': False, 'error': 'Vendor not authenticated'})
            
            # Get vendor details
            vendor_details = get_vendor_details_by_email(vendor_email)
            if not vendor_details:
                return JsonResponse({'success': False, 'error': 'Vendor details not found'})
            
            vendor_id = vendor_details.get('vendor_id', 'vendor1')
            
            s3 = boto3.client('s3',
                              aws_access_key_id=settings.R2_ACCESS_KEY,
                              aws_secret_access_key=settings.R2_SECRET_KEY,
                              endpoint_url=settings.R2_ENDPOINT,
                              region_name='auto')
            
            # Search in all possible locations
            search_locations = [
                f'vendor_print_jobs/{vendor_id}/',
                f'vendor_manual_print_jobs/{vendor_id}/',
                f'vendor_register_details/{vendor_id}/firozshop/',
                f'vendor_register_details/{sanitize_email(vendor_id)}/firozshop/',
                'users/'
            ]
            
            found_locations = []
            
            for location in search_locations:
                try:
                    if location == 'users/':
                        # Search all user folders
                        response = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=location)
                        for obj in response.get('Contents', []):
                            if obj['Key'].endswith(f'/{filename}'):
                                found_locations.append({
                                    'path': obj['Key'],
                                    'size': obj.get('Size', 0),
                                    'last_modified': obj.get('LastModified', '').isoformat() if obj.get('LastModified') else '',
                                    'location_type': 'users'
                                })
                    else:
                        # Search specific vendor folder
                        response = s3.list_objects_v2(Bucket=settings.R2_BUCKET, Prefix=location)
                        for obj in response.get('Contents', []):
                            if obj['Key'].endswith(f'/{filename}'):
                                found_locations.append({
                                    'path': obj['Key'],
                                    'size': obj.get('Size', 0),
                                    'last_modified': obj.get('LastModified', '').isoformat() if obj.get('LastModified') else '',
                                    'location_type': 'vendor'
                                })
                except Exception as e:
                    print(f"Error searching {location}: {str(e)}")
            
            return JsonResponse({
                'success': True,
                'filename': filename,
                'vendor_id': vendor_id,
                'vendor_email': vendor_email,
                'found_locations': found_locations,
                'search_locations': search_locations
            })
            
        except Exception as e:
            return JsonResponse({'success': False, 'error': str(e)})
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'})

def logout(request):
    """
    Logout view that clears session data and redirects to home
    """
    # Use Django's logout function to properly clear authentication
    from django.contrib.auth import logout as django_logout
    django_logout(request)
    
    # Add success message (before session is cleared if possible)
    try:
        messages.success(request, 'You have been successfully logged out!')
    except:
        pass  # Session might already be cleared
    
    # Clear any remaining session data
    request.session.flush()
    
    # Clear authentication cookies
    response = redirect('home')
    response.delete_cookie('sessionid')
    response.delete_cookie('csrftoken')
    
    return response


# Connection monitoring system for vendor client
import threading
import time
from datetime import timedelta

# Global dictionary to track vendor client connections
vendor_connections = {}
connection_lock = threading.Lock()

# Global dictionary to track vendor authentication timestamps
vendor_auth_timestamps = {}
auth_lock = threading.Lock()

def update_vendor_connection_status(vendor_id, status='connected'):
    """Update vendor client connection status"""
    with connection_lock:
        vendor_connections[vendor_id] = {
            'status': status,
            'last_seen': datetime.datetime.now(),
            'is_connected': status == 'connected'
        }

def update_vendor_auth_timestamp(vendor_id):
    """Update vendor authentication timestamp"""
    with auth_lock:
        vendor_auth_timestamps[vendor_id] = datetime.datetime.now()
        print(f"🔐 Updated auth timestamp for vendor {vendor_id}")

def get_vendor_auth_status(vendor_id):
    """Check if vendor authentication is still valid (within 8 hours)"""
    with auth_lock:
        if vendor_id not in vendor_auth_timestamps:
            return False
        
        auth_time = vendor_auth_timestamps[vendor_id]
        current_time = datetime.datetime.now()
        time_diff = current_time - auth_time
        
        # Check if more than 8 hours have passed
        is_valid = time_diff < timedelta(hours=8)
        
        if not is_valid:
            print(f"🔐 Vendor {vendor_id} authentication expired (8+ hours old)")
        
        return is_valid

def get_vendor_connection_info(vendor_id):
    """Get vendor client connection status"""
    with connection_lock:
        return vendor_connections.get(vendor_id, {
            'status': 'disconnected',
            'last_seen': None,
            'is_connected': False
        })

def check_vendor_connection_timeout():
    """Check if vendor connections have timed out.

    Use a generous inactivity window to avoid flapping on slow/unstable networks.
    """
    with connection_lock:
        current_time = datetime.datetime.now()
        # Flip to disconnected if no heartbeat/poll within 5 seconds
        timeout_threshold = timedelta(seconds=5)
        
        for vendor_id, connection_info in vendor_connections.items():
            if connection_info['is_connected']:
                time_since_last_seen = current_time - connection_info['last_seen']
                if time_since_last_seen > timeout_threshold:
                    # Mark as disconnected
                    vendor_connections[vendor_id]['status'] = 'disconnected'
                    vendor_connections[vendor_id]['is_connected'] = False
                    
                    # Auto-disable services for this vendor
                    auto_disable_vendor_services(vendor_id)
                    print(f"🔌 Vendor {vendor_id} connection timed out - services disabled")

def auto_disable_vendor_services(vendor_id):
    """Automatically disable vendor services when connection is lost"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        # Get vendor email from vendor_id
        vendor_email = get_vendor_email_by_vendor_id(vendor_id)
        if not vendor_email:
            print(f"❌ Could not find vendor email for vendor_id: {vendor_id}")
            return False
        
        # Get current vendor details
        vendor_details = get_vendor_details_by_email(vendor_email)
        if not vendor_details:
            print(f"❌ Could not find vendor details for email: {vendor_email}")
            return False
        
        # Update service availability to disable specific services
        vendor_folder = vendor_email_folder(vendor_email)
        availability_key = f'{vendor_folder}/service.json'
        
        try:
            # Get current availability data
            response = s3.get_object(Bucket=settings.R2_BUCKET, Key=availability_key)
            availability_data = json.loads(response['Body'].read().decode('utf-8'))
        except:
            # Create default availability data if not found
            availability_data = {
                'vendor_shop_avaliability': 'online',
                'passport_photo_available': True,
                'digital_print_available': True,
                'regular_print_available': True,
                'last_updated': datetime.datetime.now().isoformat()
            }
        
        # Disable specific services due to connection loss
        availability_data.update({
            'passport_photo_available': False,
            'digital_print_available': False,
            'regular_print_available': False,
            'connection_lost_at': datetime.datetime.now().isoformat(),
            'connection_status': 'disconnected',
            'last_updated': datetime.datetime.now().isoformat()
        })
        
        # Save updated availability
        s3.put_object(
            Bucket=settings.R2_BUCKET,
            Key=availability_key,
            Body=json.dumps(availability_data, indent=2),
            ContentType='application/json'
        )
        
        # Reset printer counts to zero when vendor goes offline
        try:
            vendor_folder = vendor_email_folder(vendor_email)
            pricing_key = f'{vendor_folder}/pricing.json'
            response = s3.get_object(Bucket=settings.R2_BUCKET, Key=pricing_key)
            pricing_data = json.loads(response['Body'].read().decode('utf-8'))
            # Reset all printer counts to zero
            if isinstance(pricing_data, dict) and 'printer_counts' in pricing_data:
                for count_key in list(pricing_data['printer_counts'].keys()):
                    pricing_data['printer_counts'][count_key] = 0
                s3.put_object(
                    Bucket=settings.R2_BUCKET,
                    Key=pricing_key,
                    Body=json.dumps(pricing_data, indent=4),
                    ContentType='application/json'
                )
                print(f"🔄 Reset printer counts to zero for vendor {vendor_id}")
        except Exception as e:
            print(f"⚠️ Could not reset printer counts for vendor {vendor_id}: {str(e)}")

        print(f"🔌 Auto-disabled services for vendor {vendor_id} due to connection loss")
        return True
        
    except Exception as e:
        print(f"❌ Error auto-disabling vendor services: {str(e)}")
        return False

def auto_enable_vendor_services(vendor_id):
    """Automatically re-enable vendor services when connection is restored"""
    try:
        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')
        
        # Get vendor email from vendor_id
        vendor_email = get_vendor_email_by_vendor_id(vendor_id)
        if not vendor_email:
            print(f"❌ Could not find vendor email for vendor_id: {vendor_id}")
            return False
        
        # Update service availability to re-enable specific services
        vendor_folder = vendor_email_folder(vendor_email)
        availability_key = f'{vendor_folder}/service.json'
        
        try:
            # Get current availability data
            response = s3.get_object(Bucket=settings.R2_BUCKET, Key=availability_key)
            availability_data = json.loads(response['Body'].read().decode('utf-8'))
        except:
            # Create default availability data if not found
            availability_data = {
                'vendor_shop_avaliability': 'online',
                'passport_photo_available': True,
                'digital_print_available': True,
                'regular_print_available': True,
                'last_updated': datetime.datetime.now().isoformat()
            }
        
        # Re-enable specific services due to connection restoration
        availability_data.update({
            'passport_photo_available': True,
            'digital_print_available': True,
            'regular_print_available': True,
            'connection_restored_at': datetime.datetime.now().isoformat(),
            'connection_status': 'connected',
            'last_updated': datetime.datetime.now().isoformat()
        })
        
        # Save updated availability
        s3.put_object(
            Bucket=settings.R2_BUCKET,
            Key=availability_key,
            Body=json.dumps(availability_data, indent=2),
            ContentType='application/json'
        )
        
        print(f"🔌 Auto-enabled services for vendor {vendor_id} due to connection restoration")
        return True
        
    except Exception as e:
        print(f"❌ Error auto-enabling vendor services: {str(e)}")
        return False

@csrf_exempt
def vendor_connection_heartbeat(request):
    """Endpoint for vendor client to report connection status"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            vendor_id = data.get('vendor_id')
            status = data.get('status', 'connected')  # Get status from client
            timestamp = data.get('timestamp', time.time())
            
            if not vendor_id:
                return JsonResponse({'success': False, 'error': 'Missing vendor_id'})
            
            # Check if authentication is still valid (8 hours)
            is_authenticated = get_vendor_auth_status(vendor_id)
            if not is_authenticated:
                return JsonResponse({
                    'success': False, 
                    'error': 'Authentication expired',
                    'auth_expired': True,
                    'message': 'Please re-authenticate - 8 hours have passed'
                }, status=401)
            
            # Update connection status with the status from client
            update_vendor_connection_status(vendor_id, status)
            
            # Re-enable services if they were disabled due to connection loss and status is connected
            if status == 'connected':
                auto_enable_vendor_services(vendor_id)
            
            # Check for timed out connections
            check_vendor_connection_timeout()
            
            return JsonResponse({
                'success': True, 
                'message': f'Connection status updated to {status}',
                'status': status
            })
            
        except json.JSONDecodeError:
            return JsonResponse({'success': False, 'error': 'Invalid JSON'}, status=400)
        except Exception as e:
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)


@csrf_exempt
@require_POST
def mark_job_rendered(request):
    """Mark given filenames' rendered_status to YES in vendor folders for current vendor session."""
    try:
        data = json.loads(request.body)
        filenames = data.get('filenames', [])
        if not filenames:
            return JsonResponse({'success': False, 'error': 'No filenames provided'}, status=400)

        vendor_email = request.session.get('vendor_email')
        vendor_id = request.session.get('vendor_id')
        if not vendor_id and vendor_email:
            vd = get_vendor_details_by_email(vendor_email)
            vendor_id = vd.get('vendor_id') if vd else None
        if not vendor_id:
            return JsonResponse({'success': False, 'error': 'Vendor not authenticated'}, status=401)

        s3 = boto3.client('s3',
                          aws_access_key_id=settings.R2_ACCESS_KEY,
                          aws_secret_access_key=settings.R2_SECRET_KEY,
                          endpoint_url=settings.R2_ENDPOINT,
                          region_name='auto')

        updated = 0
        for fn in filenames:
            for base in [f'vendor_print_jobs/{vendor_id}/', f'vendor_manual_print_jobs/{vendor_id}/']:
                key = f"{base}{fn}"
                try:
                    head = s3.head_object(Bucket=settings.R2_BUCKET, Key=key)
                    md = head.get('Metadata', {})
                    if md.get('rendered_status', 'NO').upper() != 'YES':
                        md['rendered_status'] = 'YES'
                        s3.copy_object(CopySource={'Bucket': settings.R2_BUCKET, 'Key': key},
                                       Bucket=settings.R2_BUCKET,
                                       Key=key,
                                       Metadata=md,
                                       MetadataDirective='REPLACE')
                        updated += 1
                    break
                except Exception:
                    continue

        return JsonResponse({'success': True, 'updated': updated})
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)

@csrf_exempt
def get_vendor_connection_status(request):
    """Get vendor client connection status for dashboard"""
    if request.method == 'GET':
        try:
            vendor_email = request.session.get('vendor_email')
            if not vendor_email:
                return JsonResponse({'success': False, 'error': 'Vendor not authenticated'})
            
            # Get vendor details
            vendor_details = get_vendor_details_by_email(vendor_email)
            if not vendor_details:
                return JsonResponse({'success': False, 'error': 'Vendor details not found'})
            
            vendor_id = vendor_details.get('vendor_id')
            if not vendor_id:
                return JsonResponse({'success': False, 'error': 'Vendor ID not found'})
            
            # Get connection status
            connection_info = get_vendor_connection_info(vendor_id)
            
            # Check for timeout
            check_vendor_connection_timeout()
            
            # Get updated status after timeout check
            connection_info = get_vendor_connection_info(vendor_id)
            
            # Check authentication status (8-hour validity)
            is_authenticated = get_vendor_auth_status(vendor_id)
            
            # Determine overall status
            if not is_authenticated:
                overall_status = 'auth_expired'
                is_connected = False
            else:
                overall_status = connection_info['status']
                is_connected = connection_info['is_connected']
            
            return JsonResponse({
                'success': True,
                'connection_status': overall_status,
                'is_connected': is_connected,
                'is_authenticated': is_authenticated,
                'last_seen': connection_info['last_seen'].isoformat() if connection_info['last_seen'] else None
            })
            
        except Exception as e:
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    
    return JsonResponse({'success': False, 'error': 'Invalid request method'}, status=405)


def start_vendor_pending_jobs_scheduler():
    """Start automatic vendor pending jobs snapshot scheduler"""
    # Note: Scheduled snapshots are disabled because we need active vendor session
    # Snapshots will be triggered on-demand when vendors access their dashboard
    print("✅ Vendor pending jobs snapshot scheduler initialized (on-demand mode)")
    print("📋 Snapshots will be stored when vendors access their dashboard")


# Initialize automatic vendor pending jobs snapshot when the module is imported
start_vendor_pending_jobs_scheduler()
